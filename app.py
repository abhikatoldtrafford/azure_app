import logging
import threading
try:
    import httpx
    HTTPX_AVAILABLE = True
except ImportError:
    HTTPX_AVAILABLE = False
    logging.warning("httpx not available - will use simple client creation")
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Query, Request, Response, Path
from fastapi.responses import JSONResponse, StreamingResponse, FileResponse, HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.openapi.utils import get_openapi
from fastapi import Depends
from pydantic import BaseModel, Field
from openai import AzureOpenAI, AsyncAzureOpenAI
from typing import Optional, List, Dict, Any, Tuple, AsyncGenerator, Union, Annotated
import os
import atexit
import io
from datetime import datetime
import time
import base64
import mimetypes
import traceback
import asyncio
import json
from io import StringIO, BytesIO
import sys
import re
import hashlib
import shutil
import uuid
import tempfile
import platform
# Document processing
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from PIL import Image
import PyPDF2
import chardet
from bs4 import BeautifulSoup
import markdown2
from io import BytesIO
# Data processing
import pandas as pd
import numpy as np
import openpyxl
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils.dataframe import dataframe_to_rows
try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import seaborn as sns
    import plotly.graph_objects as go
    import plotly.io as pio
    from wordcloud import WordCloud
    CHARTS_AVAILABLE = True
except ImportError:
    CHARTS_AVAILABLE = False
    logging.warning("Chart libraries not available. Chart generation disabled.")
import asyncio
from datetime import timedelta
from PIL import Image as PILImage
import chardet
import pdfplumber
# Try importing Unstructured components
try:
    from unstructured.partition.auto import partition
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False
    partition = None
# Import specific partitioners for fallback
try:
    from unstructured.partition.csv import partition_csv
except ImportError:
    partition_csv = None
try:
    from unstructured.partition.docx import partition_docx
except ImportError:
    partition_docx = None
try:
    from unstructured.partition.pdf import partition_pdf
except ImportError:
    partition_pdf = None
try:
    from unstructured.partition.html import partition_html
except ImportError:
    partition_html = None
try:
    from unstructured.partition.text import partition_text
except ImportError:
    partition_text = None
try:
    from unstructured.partition.email import partition_email
except ImportError:
    partition_email = None
try:
    from unstructured.partition.xlsx import partition_xlsx
except ImportError:
    partition_xlsx = None
try:
    from unstructured.partition.image import partition_image
except ImportError:
    partition_image = None
try:
    from unstructured.partition.json import partition_json
except ImportError:
    partition_json = None
try:
    from unstructured.partition.xml import partition_xml
except ImportError:
    partition_xml = None
try:
    from unstructured.partition.md import partition_md
except ImportError:
    partition_md = None
try:
    from unstructured.partition.pptx import partition_pptx
except ImportError:
    partition_pptx = None
try:
    from unstructured.partition.msg import partition_msg
except ImportError:
    partition_msg = None
try:
    from unstructured.partition.rtf import partition_rtf
except ImportError:
    partition_rtf = None
try:
    from unstructured.partition.epub import partition_epub
except ImportError:
    partition_epub = None
try:
    from unstructured.partition.odt import partition_odt
except ImportError:
    partition_odt = None
try:
    from unstructured.partition.doc import partition_doc
except ImportError:
    partition_doc = None
try:
    from unstructured.partition.ppt import partition_ppt
except ImportError:
    partition_ppt = None
try:
    from unstructured.partition.tsv import partition_tsv
except ImportError:
    partition_tsv = None
try:
    from unstructured.partition.rst import partition_rst
except ImportError:
    partition_rst = None
try:
    from unstructured.partition.org import partition_org
except ImportError:
    partition_org = None
# Fallback imports for when Unstructured isn't available
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    import html2text
except ImportError:
    html2text = None
try:
    import markdown
except ImportError:
    markdown = None
# Pydantic models for request/response documentation
# Azure OpenAI client configuration
AZURE_ENDPOINT = "https://kb-stellar.openai.azure.com/" # Replace with your endpoint if different
AZURE_API_KEY = "bc0ba854d3644d7998a5034af62d03ce" # Replace with your key if different
AZURE_API_VERSION = "2024-12-01-preview"
DOWNLOADS_DIR = "/tmp/chat_downloads"  # Use /tmp for Azure App Service
MAX_DOWNLOAD_FILES = 10  # Keep only 10 most recent files

class HealthResponse(BaseModel):
    status: str = Field(..., example="healthy", description="Health status")
    timestamp: str = Field(..., example="2024-01-15T10:30:00", description="Current timestamp")

class HealthCheckResponse(BaseModel):
    status: str = Field(..., example="healthy", description="Overall health status")
    timestamp: str = Field(..., description="Check timestamp")
    version: str = Field(..., example="1.0.0", description="API version")
    checks: Dict[str, Any] = Field(..., description="Individual component checks")
    endpoints_tested: Dict[str, Any] = Field(default_factory=dict, description="Endpoint test results")
    warnings: List[str] = Field(default_factory=list, description="Warning messages")
    errors: List[str] = Field(default_factory=list, description="Error messages")
    health_percentage: float = Field(..., example=100.0, description="Overall health percentage")
    execution_time_ms: float = Field(..., example=250.5, description="Check execution time in milliseconds")
    summary: Dict[str, int] = Field(..., description="Summary of checks")

class CompletionResponse(BaseModel):
    status: str = Field(..., example="success")
    content: str = Field(..., example="Generated content here", description="AI-generated content")
    download_url: Optional[str] = Field(None, example="/download-files/generated_data.xlsx")
    filename: Optional[str] = Field(None, example="generated_data.xlsx")
    assistant_id: Optional[str] = Field(None, description="Assistant ID used")
    response: Optional[str] = Field(None, description="Raw response text")

class ErrorResponse(BaseModel):
    status: str = Field("error", example="error")
    message: str = Field(..., example="An error occurred", description="Error description")
    error: Optional[str] = Field(None, description="Error details")
    details: Optional[Dict[str, Any]] = Field(None, description="Additional error details")

class ChatInitResponse(BaseModel):
    message: str = Field(..., example="Chat initiated successfully")
    assistant: str = Field(..., example="asst_abc123", description="Assistant ID")
    session: str = Field(..., example="thread_xyz789", description="Thread/Session ID")
    vector_store: str = Field(..., example="vs_def456", description="Vector store ID")

class ExtractResponse(BaseModel):
    status: str = Field(..., example="success")
    format: str = Field(..., example="csv", description="Output format")
    data: Optional[Union[str, List[Dict], Dict]] = Field(None, description="Extracted/generated data")
    download_url: Optional[str] = Field(None, description="Download URL if file generated")
    filename: Optional[str] = Field(None, description="Generated filename")
    source_file: Optional[str] = Field(None, description="Source filename")
    mode: str = Field(..., example="extract", description="Operation mode used")
    result: Optional[Dict[str, Any]] = Field(None, description="Result data for JSON format")
    message: Optional[str] = Field(None, description="Status message")

class FileUploadResponse(BaseModel):
    status: str = Field(..., example="success")
    message: str = Field(..., example="File processed successfully")
    filename: str = Field(..., description="Uploaded filename")
    file_id: Optional[str] = Field(None, description="File ID in vector store")
    content_extracted: Optional[bool] = Field(None, description="Whether content was extracted")
    file_content: Optional[str] = Field(None, description="Extracted file content")
    preview: Optional[str] = Field(None, description="Content preview")

class ChatResponse(BaseModel):
    response: str = Field(..., description="AI response")
    session_id: Optional[str] = Field(None, description="Session ID")
    message_count: Optional[int] = Field(None, description="Total messages in thread")

class DownloadResponse(BaseModel):
    status: str = Field(..., example="success")
    download_url: str = Field(..., example="/download-files/chat_response.docx")
    filename: str = Field(..., example="chat_response_20240115_123456.docx")
    message: str = Field(..., example="Chat response ready for download")
    expires_in: str = Field(..., example="File will be kept for the last 10 downloads")

class FileVerifyResponse(BaseModel):
    available: bool = Field(..., description="Whether file is available")
    filename: str = Field(..., description="File name")
    size: Optional[int] = Field(None, description="File size in bytes")
    mime_type: Optional[str] = Field(None, description="MIME type")
    created: Optional[str] = Field(None, description="Creation timestamp")

class TestEndpointResponse(BaseModel):
    endpoint: str = Field(..., description="Tested endpoint")
    timestamp: str = Field(..., description="Test timestamp")
    test_data: Optional[str] = Field(None, description="Test data used")
    result: Optional[Dict[str, Any]] = Field(None, description="Test result")
    error: Optional[str] = Field(None, description="Error if test failed")
# Thread lock manager to prevent concurrent access to the same thread
class ThreadLockManager:
    def __init__(self):
        self.locks: Dict[str, asyncio.Lock] = {}
        self.lock_access_times: Dict[str, datetime] = {}
        self.manager_lock = asyncio.Lock()
    
    async def get_lock(self, thread_id: str) -> asyncio.Lock:
        async with self.manager_lock:
            if thread_id not in self.locks:
                self.locks[thread_id] = asyncio.Lock()
            self.lock_access_times[thread_id] = datetime.now()
            return self.locks[thread_id]
    
    async def cleanup_old_locks(self, max_age_minutes: int = 30):
        """Remove locks that haven't been accessed in a while to prevent memory leaks"""
        async with self.manager_lock:
            current_time = datetime.now()
            threads_to_remove = []
            
            for thread_id, last_access in self.lock_access_times.items():
                if current_time - last_access > timedelta(minutes=max_age_minutes):
                    # Only remove if lock is not currently held
                    if thread_id in self.locks and not self.locks[thread_id].locked():
                        threads_to_remove.append(thread_id)
            
            for thread_id in threads_to_remove:
                del self.locks[thread_id]
                del self.lock_access_times[thread_id]
                logging.info(f"Cleaned up lock for thread {thread_id}")

# Create global instance
thread_lock_manager = ThreadLockManager()
# Simple status updates for long-running operations
operation_statuses = {}


# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

CUSTOM_SWAGGER_CSS = """
<style>
/* 🎨 AZURE COPILOT V2 - ULTRA PREMIUM THEME */

/* Hide unwanted elements */
.swagger-ui .topbar { display: none !important; }
.swagger-ui .models { display: none !important; }
.swagger-ui section.models { display: none !important; }
.swagger-ui .scheme-container { display: none !important; }

/* Animated gradient background */
@keyframes gradientShift {
    0% { background-position: 0% 50%; }
    50% { background-position: 100% 50%; }
    100% { background-position: 0% 50%; }
}

body {
    background: linear-gradient(-45deg, #f8f9fa, #ffffff, #f1f3f5, #e9ecef) !important;
    background-size: 400% 400% !important;
    animation: gradientShift 15s ease infinite !important;
    margin: 0 !important;
    padding: 0 !important;
    min-height: 100vh !important;
}

/* CSS Variables */
.swagger-ui {
    --bg-primary: #ffffff;
    --bg-secondary: #f8f9fa;
    --bg-tertiary: #f1f3f5;
    --accent: #5c7cfa;
    --accent-hover: #4c6ef5;
    --accent-active: #364fc7;
    --text-primary: #212529;
    --text-secondary: #495057;
    --text-muted: #868e96;
    --border: #dee2e6;
    --success: #51cf66;
    --error: #ff6b6b;
    --warning: #ffd43b;
    --info: #4dabf7;
    font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif !important;
    color: var(--text-primary) !important;
}

/* Main container with glassmorphism */
.swagger-ui .wrapper {
    background: rgba(255, 255, 255, 0.95) !important;
    backdrop-filter: blur(20px) !important;
    -webkit-backdrop-filter: blur(20px) !important;
    padding: 40px !important;
    max-width: 1400px !important;
    margin: 0 auto !important;
    border-radius: 20px !important;
    box-shadow: 0 25px 50px -12px rgba(0, 0, 0, 0.08) !important;
    border: 1px solid rgba(92, 124, 250, 0.1) !important;
}

/* Header section */
.swagger-ui .info {
    margin-bottom: 50px !important;
    padding: 40px !important;
    background: linear-gradient(135deg, rgba(92, 124, 250, 0.05) 0%, rgba(92, 124, 250, 0.02) 100%) !important;
    border-radius: 16px !important;
    border: 1px solid rgba(92, 124, 250, 0.2) !important;
    position: relative !important;
    overflow: hidden !important;
}

.swagger-ui .info::before {
    content: '' !important;
    position: absolute !important;
    top: -50% !important;
    right: -50% !important;
    width: 200% !important;
    height: 200% !important;
    background: radial-gradient(circle, rgba(92, 124, 250, 0.05) 0%, transparent 70%) !important;
    animation: pulse 4s ease-in-out infinite !important;
}

@keyframes pulse {
    0%, 100% { transform: scale(1); opacity: 0.5; }
    50% { transform: scale(1.1); opacity: 0.3; }
}

/* Title with gradient */
.swagger-ui .info .title {
    font-size: 3.5rem !important;
    font-weight: 900 !important;
    background: linear-gradient(135deg, #212529 0%, #343a40 100%) !important;
    -webkit-background-clip: text !important;
    -webkit-text-fill-color: transparent !important;
    background-clip: text !important;
    letter-spacing: -0.02em !important;
    margin-bottom: 16px !important;
    position: relative !important;
    z-index: 1 !important;
}

.swagger-ui .info .title small {
    background: linear-gradient(135deg, #5c7cfa 0%, #4c6ef5 100%) !important;
    -webkit-background-clip: text !important;
    -webkit-text-fill-color: transparent !important;
    background-clip: text !important;
    font-size: 14px !important;
    font-weight: 700 !important;
    letter-spacing: 0.1em !important;
    text-transform: uppercase !important;
    display: block !important;
    margin-top: 8px !important;
}

/* Description with better readability */
.swagger-ui .info .description,
.swagger-ui .info .description p,
.swagger-ui .markdown p {
    color: var(--text-secondary) !important;
    font-size: 16px !important;
    line-height: 1.8 !important;
    font-weight: 400 !important;
}

/* Enhanced operation blocks */
.swagger-ui .opblock {
    background: var(--bg-secondary) !important;
    border: 1px solid var(--border) !important;
    border-radius: 12px !important;
    margin-bottom: 20px !important;
    overflow: hidden !important;
    transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1) !important;
    box-shadow: 0 1px 3px 0 rgba(0, 0, 0, 0.05) !important;
}

.swagger-ui .opblock:hover {
    transform: translateY(-2px) !important;
    box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.1) !important;
    border-color: var(--accent) !important;
}

/* Operation summary */
.swagger-ui .opblock-summary {
    background: rgba(255, 255, 255, 0.9) !important;
    padding: 16px 20px !important;
    cursor: pointer !important;
    border: none !important;
}

.swagger-ui .opblock-summary:hover {
    background: rgba(248, 249, 250, 0.95) !important;
}

/* Method labels with distinct colors */
.swagger-ui .opblock-summary-method {
    font-weight: 800 !important;
    padding: 8px 16px !important;
    border-radius: 8px !important;
    text-transform: uppercase !important;
    font-size: 12px !important;
    letter-spacing: 0.05em !important;
    margin-right: 16px !important;
    box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1) !important;
    position: relative !important;
    overflow: hidden !important;
}

/* Method colors */
.swagger-ui .opblock.opblock-get .opblock-summary-method {
    background: #16a34a !important;
    color: #ffffff !important;
    box-shadow: 0 4px 14px rgba(22, 163, 74, 0.3) !important;
}

.swagger-ui .opblock.opblock-post .opblock-summary-method {
    background: #2563eb !important;
    color: #ffffff !important;
    box-shadow: 0 4px 14px rgba(37, 99, 235, 0.3) !important;
}

.swagger-ui .opblock.opblock-put .opblock-summary-method {
    background: #ea580c !important;
    color: #ffffff !important;
    box-shadow: 0 4px 14px rgba(234, 88, 12, 0.3) !important;
}

.swagger-ui .opblock.opblock-delete .opblock-summary-method {
    background: #dc2626 !important;
    color: #ffffff !important;
    box-shadow: 0 4px 14px rgba(220, 38, 38, 0.3) !important;
}

/* Path and description */
.swagger-ui .opblock-summary-path,
.swagger-ui .opblock-summary-path__deprecated {
    color: var(--text-primary) !important;
    font-weight: 600 !important;
    font-size: 16px !important;
    font-family: 'JetBrains Mono', monospace !important;
}

.swagger-ui .opblock-summary-description {
    color: var(--text-secondary) !important;
    font-size: 14px !important;
    margin-left: 12px !important;
}

/* Operation body */
.swagger-ui .opblock-body {
    background: var(--bg-primary) !important;
    padding: 24px !important;
}

/* Parameters section */
.swagger-ui .opblock-section-header {
    background: rgba(92, 124, 250, 0.08) !important;
    padding: 12px 20px !important;
    border-radius: 8px !important;
    margin-bottom: 16px !important;
    border-left: 4px solid var(--accent) !important;
}

.swagger-ui .opblock-section-header h4 {
    color: var(--text-primary) !important;
    font-size: 14px !important;
    font-weight: 700 !important;
    text-transform: uppercase !important;
    letter-spacing: 0.05em !important;
    margin: 0 !important;
}

/* Table improvements */
.swagger-ui .table-container {
    background: var(--bg-secondary) !important;
    border-radius: 8px !important;
    padding: 16px !important;
    border: 1px solid var(--border) !important;
}

.swagger-ui .parameters-col_name {
    color: #4c6ef5 !important;
    font-weight: 600 !important;
    font-family: 'JetBrains Mono', monospace !important;
    font-size: 14px !important;
}

.swagger-ui .parameters-col_description {
    color: var(--text-secondary) !important;
    font-size: 14px !important;
    line-height: 1.6 !important;
}

.swagger-ui .parameter__type {
    color: #7c3aed !important;
    font-family: 'JetBrains Mono', monospace !important;
    font-size: 12px !important;
    font-weight: 600 !important;
}

/* Input fields */
.swagger-ui input[type=text],
.swagger-ui input[type=password],
.swagger-ui input[type=email],
.swagger-ui input[type=number],
.swagger-ui textarea,
.swagger-ui select {
    background: var(--bg-primary) !important;
    color: var(--text-primary) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
    padding: 10px 14px !important;
    font-size: 14px !important;
    font-family: 'Inter', sans-serif !important;
    transition: all 0.2s ease !important;
}

.swagger-ui input:focus,
.swagger-ui textarea:focus,
.swagger-ui select:focus {
    border-color: var(--accent) !important;
    outline: none !important;
    box-shadow: 0 0 0 3px rgba(92, 124, 250, 0.15) !important;
    background: var(--bg-primary) !important;
}

/* Buttons */
.swagger-ui .btn {
    background: var(--accent) !important;
    color: #ffffff !important;
    border: none !important;
    padding: 10px 20px !important;
    border-radius: 8px !important;
    font-weight: 600 !important;
    font-size: 14px !important;
    text-transform: uppercase !important;
    letter-spacing: 0.05em !important;
    transition: all 0.2s ease !important;
    cursor: pointer !important;
    box-shadow: 0 2px 4px rgba(0, 0, 0, 0.05) !important;
}

.swagger-ui .btn:hover {
    background: var(--accent-hover) !important;
    transform: translateY(-2px) !important;
    box-shadow: 0 4px 8px rgba(92, 124, 250, 0.2) !important;
}

.swagger-ui .btn:active {
    transform: translateY(0) !important;
}

/* Try it out button */
.swagger-ui .btn.try-out__btn {
    background: linear-gradient(135deg, var(--accent) 0%, var(--accent-hover) 100%) !important;
    color: #ffffff !important;
    position: relative !important;
    overflow: hidden !important;
}

.swagger-ui .btn.try-out__btn::before {
    content: '' !important;
    position: absolute !important;
    top: 50% !important;
    left: 50% !important;
    width: 0 !important;
    height: 0 !important;
    background: rgba(255, 255, 255, 0.2) !important;
    border-radius: 50% !important;
    transform: translate(-50%, -50%) !important;
    transition: width 0.6s, height 0.6s !important;
}

.swagger-ui .btn.try-out__btn:hover::before {
    width: 300px !important;
    height: 300px !important;
}

/* Execute button */
.swagger-ui .btn.execute {
    background: linear-gradient(135deg, #16a34a 0%, #22c55e 100%) !important;
    color: #ffffff !important;
    font-size: 16px !important;
    padding: 12px 32px !important;
    font-weight: 700 !important;
    text-transform: uppercase !important;
    letter-spacing: 0.1em !important;
    position: relative !important;
}

.swagger-ui .btn.execute::after {
    content: '🚀' !important;
    margin-left: 8px !important;
    font-size: 18px !important;
    display: inline-block !important;
    animation: bounce 2s infinite !important;
}

@keyframes bounce {
    0%, 20%, 50%, 80%, 100% { transform: translateY(0); }
    40% { transform: translateY(-5px); }
    60% { transform: translateY(-3px); }
}

.swagger-ui .btn.cancel {
    background: var(--bg-tertiary) !important;
    color: var(--text-secondary) !important;
    border: 2px solid var(--border) !important;
}

.swagger-ui .btn.cancel:hover {
    background: var(--bg-secondary) !important;
    border-color: var(--text-muted) !important;
}

/* Response section */
.swagger-ui .responses-wrapper {
    margin-top: 32px !important;
    padding: 24px !important;
    background: rgba(248, 249, 250, 0.5) !important;
    border-radius: 12px !important;
    border: 1px solid var(--border) !important;
}

.swagger-ui .responses-inner {
    background: transparent !important;
}

.swagger-ui .response {
    background: var(--bg-primary) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
    margin-bottom: 16px !important;
    overflow: hidden !important;
}

/* Loading animation */
.swagger-ui .loading-container {
    padding: 40px !important;
    text-align: center !important;
}

.swagger-ui .loading-container .loading {
    position: relative !important;
    width: 60px !important;
    height: 60px !important;
    margin: 0 auto 20px !important;
}

.swagger-ui .loading-container .loading::after {
    content: '' !important;
    position: absolute !important;
    width: 100% !important;
    height: 100% !important;
    border: 4px solid transparent !important;
    border-top-color: var(--accent) !important;
    border-right-color: var(--accent) !important;
    border-bottom-color: var(--accent) !important;
    border-radius: 50% !important;
    animation: gradientSpin 1s linear infinite !important;
    display: inline-block !important;
}

@keyframes gradientSpin {
    0% { transform: rotate(0deg); }
    100% { transform: rotate(360deg); }
}

/* Copy button */
.swagger-ui .copy-to-clipboard {
    background: rgba(92, 124, 250, 0.1) !important;
    border: 1px solid var(--accent) !important;
    border-radius: 8px !important;
    padding: 6px 12px !important;
    transition: all 0.3s ease !important;
    color: var(--accent) !important;
}

.swagger-ui .copy-to-clipboard:hover {
    background: var(--accent) !important;
    color: #ffffff !important;
    transform: translateY(-1px) !important;
    box-shadow: 0 4px 12px rgba(92, 124, 250, 0.3) !important;
}

/* Code blocks */
.swagger-ui .highlight-code pre,
.swagger-ui pre {
    background: #f6f8fa !important;
    border: 1px solid var(--border) !important;
    border-radius: 12px !important;
    padding: 20px !important;
    font-family: 'JetBrains Mono', monospace !important;
    font-size: 14px !important;
    line-height: 1.6 !important;
    color: #24292e !important;
    overflow-x: auto !important;
}

.swagger-ui code {
    background: rgba(92, 124, 250, 0.1) !important;
    color: #5c7cfa !important;
    padding: 2px 6px !important;
    border-radius: 4px !important;
    font-family: 'JetBrains Mono', monospace !important;
    font-size: 0.9em !important;
}

/* Response status badges */
.swagger-ui .response-col_status {
    font-weight: 800 !important;
    padding: 6px 16px !important;
    border-radius: 20px !important;
    display: inline-block !important;
    font-size: 14px !important;
    text-transform: uppercase !important;
    letter-spacing: 0.05em !important;
}

.swagger-ui .response[data-code^="2"] .response-col_status {
    background: rgba(22, 163, 74, 0.1) !important;
    color: #16a34a !important;
    border: 2px solid #16a34a !important;
}

.swagger-ui .response[data-code^="4"] .response-col_status {
    background: rgba(234, 88, 12, 0.1) !important;
    color: #ea580c !important;
    border: 2px solid #ea580c !important;
}

.swagger-ui .response[data-code^="5"] .response-col_status {
    background: rgba(220, 38, 38, 0.1) !important;
    color: #dc2626 !important;
    border: 2px solid #dc2626 !important;
}

/* Model/Schema display */
.swagger-ui .model-container {
    background: var(--bg-secondary) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
    padding: 16px !important;
}

.swagger-ui .model {
    color: var(--text-primary) !important;
    font-family: 'JetBrains Mono', monospace !important;
    font-size: 14px !important;
}

.swagger-ui .model-title {
    color: var(--accent) !important;
    font-weight: 700 !important;
    font-size: 16px !important;
    margin-bottom: 12px !important;
}

.swagger-ui .prop-type {
    color: #7c3aed !important;
}

.swagger-ui .prop-format {
    color: #2563eb !important;
}

/* Scrollbar */
::-webkit-scrollbar {
    width: 12px !important;
    height: 12px !important;
}

::-webkit-scrollbar-track {
    background: var(--bg-primary) !important;
    border-radius: 10px !important;
}

::-webkit-scrollbar-thumb {
    background: linear-gradient(135deg, var(--accent) 0%, var(--accent-hover) 100%) !important;
    border-radius: 10px !important;
    border: 2px solid var(--bg-primary) !important;
}

::-webkit-scrollbar-thumb:hover {
    background: linear-gradient(135deg, var(--accent-hover) 0%, var(--accent-active) 100%) !important;
}

/* Links */
.swagger-ui a {
    color: var(--accent) !important;
    text-decoration: none !important;
    transition: color 0.2s ease !important;
}

.swagger-ui a:hover {
    color: var(--accent-hover) !important;
    text-decoration: underline !important;
}

/* Error messages */
.swagger-ui .errors-wrapper {
    background: rgba(220, 38, 38, 0.05) !important;
    border: 2px solid #dc2626 !important;
    border-radius: 8px !important;
    padding: 16px !important;
    margin: 16px 0 !important;
}

.swagger-ui .errors-wrapper .error {
    color: #dc2626 !important;
    font-weight: 500 !important;
}

/* Select dropdowns */
.swagger-ui select {
    appearance: none !important;
    background-image: url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 20 20' fill='%23495057'%3E%3Cpath fill-rule='evenodd' d='M5.293 7.293a1 1 0 011.414 0L10 10.586l3.293-3.293a1 1 0 111.414 1.414l-4 4a1 1 0 01-1.414 0l-4-4a1 1 0 010-1.414z' clip-rule='evenodd'/%3E%3C/svg%3E") !important;
    background-repeat: no-repeat !important;
    background-position: right 10px center !important;
    background-size: 20px !important;
    padding-right: 40px !important;
}

/* Tabs */
.swagger-ui .tab {
    color: var(--text-secondary) !important;
    background: transparent !important;
    border-bottom: 2px solid transparent !important;
    padding: 8px 16px !important;
    font-weight: 500 !important;
    transition: all 0.2s ease !important;
}

.swagger-ui .tab:hover {
    color: var(--text-primary) !important;
    background: rgba(92, 124, 250, 0.05) !important;
}

.swagger-ui .tab.active {
    color: var(--accent) !important;
    border-bottom-color: var(--accent) !important;
}

/* Authorization button */
.swagger-ui .authorization__btn {
    background: linear-gradient(135deg, var(--accent) 0%, var(--accent-hover) 100%) !important;
    color: #ffffff !important;
    padding: 10px 24px !important;
    border-radius: 8px !important;
    font-weight: 600 !important;
    transition: all 0.2s ease !important;
    box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1) !important;
}

.swagger-ui .authorization__btn:hover {
    transform: translateY(-2px) !important;
    box-shadow: 0 8px 16px rgba(92, 124, 250, 0.3) !important;
}

/* Smooth transitions */
body, body *, .swagger-ui, .swagger-ui * {
    transition: background-color 0.3s ease, 
                color 0.3s ease, 
                border-color 0.3s ease, 
                box-shadow 0.3s ease,
                fill 0.3s ease !important;
}
</style>
"""


@app.get("/docs", include_in_schema=False)
async def custom_swagger_ui_html():
    """
    Simplified Swagger UI with custom styling only.
    - Clean, modern interface with premium design
    - All endpoints fully functional
    - Streaming support built into /conversation endpoint
    """
    return HTMLResponse(
        content=f"""
        <!DOCTYPE html>
        <html>
        <head>
            <title>{app.title} - Swagger UI</title>
            <link rel="stylesheet" type="text/css" href="https://cdn.jsdelivr.net/npm/swagger-ui-dist@5/swagger-ui.css">
            <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&family=JetBrains+Mono:wght@400;600&display=swap" rel="stylesheet">
            {CUSTOM_SWAGGER_CSS}
        </head>
        <body>
            <div id="swagger-ui"></div>
            
            <script src="https://cdn.jsdelivr.net/npm/swagger-ui-dist@5/swagger-ui-bundle.js"></script>
            <script src="https://cdn.jsdelivr.net/npm/swagger-ui-dist@5/swagger-ui-standalone-preset.js"></script>
            <script>
            window.onload = function() {
                window.ui = SwaggerUIBundle({
                    url: "/openapi.json",
                    dom_id: '#swagger-ui',
                    deepLinking: true,
                    presets: [
                        SwaggerUIBundle.presets.apis,
                        SwaggerUIStandalonePreset
                    ],
                    plugins: [
                        SwaggerUIBundle.plugins.DownloadUrl
                    ],
                    layout: "StandaloneLayout",
                    defaultModelsExpandDepth: -1,
                    displayRequestDuration: true,
                    filter: true,
                    persistAuthorization: true,
                    tryItOutEnabled: true,
                    supportedSubmitMethods: ['get', 'post', 'put', 'delete', 'patch'],
                    onComplete: function() {
                        console.log("Azure CoPilot V2 API loaded successfully!");
                    }
                });
            };
            </script>
        </body>
        </html>
        """,
        headers={
            "Content-Type": "text/html; charset=utf-8",
        }
    )

tags_metadata = [
    {
        "name": "System",
        "description": "System health, testing, and monitoring endpoints",
    },
    {
        "name": "AI Operations", 
        "description": "Core AI completion and generation endpoints",
    },
    {
        "name": "Chat Operations",
        "description": "Session-based chat functionality with streaming support",
    },
    {
        "name": "Data Processing",
        "description": "Data extraction, generation, and analysis endpoints",
    },
    {
        "name": "File Operations",
        "description": "File upload, download, and management endpoints",
    },
]

# Initialize FastAPI with tags
app = FastAPI(
    title="🚀 Azure CoPilot V2 - AI Powerhouse API",
    description="""
<div style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); padding: 2px; border-radius: 16px; margin-bottom: 20px;">
<div style="background: #0a0a0a; border-radius: 14px; padding: 30px;">

<h1 style="background: linear-gradient(135deg, #667eea 0%, #f093fb 100%); -webkit-background-clip: text; -webkit-text-fill-color: transparent; font-size: 48px; margin: 0 0 20px 0; font-weight: 800;">
⚡ Azure CoPilot V2
</h1>

<p style="font-size: 20px; color: #e0e7ff; margin-bottom: 30px; line-height: 1.6;">
The <strong>ULTIMATE AI Assistant API</strong> that transforms your wildest automation dreams into reality! 
Powered by cutting-edge Azure OpenAI technology and engineered for <em>blazing-fast performance</em>. 🔥
</p>

<div style="display: grid; grid-template-columns: repeat(2, 1fr); gap: 20px; margin: 30px 0;">

<div style="background: rgba(99, 102, 241, 0.1); border: 1px solid rgba(99, 102, 241, 0.3); border-radius: 12px; padding: 20px;">
<h3 style="color: #a78bfa; margin-top: 0;">🧠 AI Superpowers</h3>
<ul style="color: #c7d2fe; list-style: none; padding: 0;">
<li>✨ <strong>Real-time Streaming</strong> - Watch AI think in real-time!</li>
<li>🎯 <strong>Context-Aware Conversations</strong> - Never repeat yourself</li>
<li>🚄 <strong>Lightning-Fast Processing</strong> - Millisecond responses</li>
<li>🌐 <strong>Multi-Modal Intelligence</strong> - Text, images, data, everything!</li>
</ul>
</div>

<div style="background: rgba(236, 72, 153, 0.1); border: 1px solid rgba(236, 72, 153, 0.3); border-radius: 12px; padding: 20px;">
<h3 style="color: #f9a8d4; margin-top: 0;">📊 Data Mastery</h3>
<ul style="color: #fce7f3; list-style: none; padding: 0;">
<li>🔍 <strong>Smart Data Analysis</strong> - Pandas-powered insights</li>
<li>📈 <strong>Auto-Generated Reports</strong> - Excel, CSV, DOCX exports</li>
<li>🎨 <strong>Review Mining</strong> - Extract gold from any document</li>
<li>🤖 <strong>Synthetic Data Generation</strong> - Create perfect test data</li>
</ul>
</div>

<div style="background: rgba(34, 197, 94, 0.1); border: 1px solid rgba(34, 197, 94, 0.3); border-radius: 12px; padding: 20px;">
<h3 style="color: #86efac; margin-top: 0;">🎭 Creative Genius</h3>
<ul style="color: #dcfce7; list-style: none; padding: 0;">
<li>📝 <strong>Content Generation</strong> - Blog posts to business plans</li>
<li>🖼️ <strong>Vision Understanding</strong> - Analyze any image</li>
<li>🎪 <strong>Multi-File Magic</strong> - Process entire folders</li>
<li>🌟 <strong>Custom Personas</strong> - AI that adapts to YOU</li>
</ul>
</div>

<div style="background: rgba(251, 146, 60, 0.1); border: 1px solid rgba(251, 146, 60, 0.3); border-radius: 12px; padding: 20px;">
<h3 style="color: #fed7aa; margin-top: 0;">⚡ Performance</h3>
<ul style="color: #ffedd5; list-style: none; padding: 0;">
<li>🚀 <strong>Enterprise Scale</strong> - Handle millions of requests</li>
<li>🛡️ <strong>Battle-Tested</strong> - Production-ready reliability</li>
<li>📡 <strong>SSE Streaming</strong> - Real-time event streams</li>
<li>🔧 <strong>Auto-Scaling</strong> - Grows with your needs</li>
</ul>
</div>

</div>

<div style="background: linear-gradient(90deg, rgba(99, 102, 241, 0.2) 0%, rgba(236, 72, 153, 0.2) 100%); border-radius: 12px; padding: 20px; margin: 20px 0;">
<h3 style="color: #e0e7ff; margin-top: 0;">🎯 Supported Formats</h3>
<div style="display: flex; flex-wrap: wrap; gap: 10px;">
<span style="background: rgba(99, 102, 241, 0.3); color: #c7d2fe; padding: 6px 12px; border-radius: 20px; font-size: 14px;">📄 PDF</span>
<span style="background: rgba(99, 102, 241, 0.3); color: #c7d2fe; padding: 6px 12px; border-radius: 20px; font-size: 14px;">📊 Excel</span>
<span style="background: rgba(99, 102, 241, 0.3); color: #c7d2fe; padding: 6px 12px; border-radius: 20px; font-size: 14px;">📝 Word</span>
<span style="background: rgba(236, 72, 153, 0.3); color: #fce7f3; padding: 6px 12px; border-radius: 20px; font-size: 14px;">🖼️ Images</span>
<span style="background: rgba(236, 72, 153, 0.3); color: #fce7f3; padding: 6px 12px; border-radius: 20px; font-size: 14px;">📋 CSV</span>
<span style="background: rgba(236, 72, 153, 0.3); color: #fce7f3; padding: 6px 12px; border-radius: 20px; font-size: 14px;">🔤 TXT</span>
<span style="background: rgba(34, 197, 94, 0.3); color: #dcfce7; padding: 6px 12px; border-radius: 20px; font-size: 14px;">📐 JSON</span>
<span style="background: rgba(34, 197, 94, 0.3); color: #dcfce7; padding: 6px 12px; border-radius: 20px; font-size: 14px;">🌐 HTML</span>
<span style="background: rgba(34, 197, 94, 0.3); color: #dcfce7; padding: 6px 12px; border-radius: 20px; font-size: 14px;">📑 Markdown</span>
</div>
</div>

<div style="text-align: center; margin-top: 30px;">
<p style="color: #a78bfa; font-size: 18px; margin: 0;">
Ready to revolutionize your workflow? <strong>Let's build something AMAZING together!</strong> 🚀
</p>
</div>

</div>
</div>
    """,
    version="2.0.0",
    
    openapi_tags=[
        {
            "name": "🎯 Core AI",
            "description": "**Main AI operations** - The heart of the system! These endpoints power conversations, completions, and intelligent processing."
        },
        {
            "name": "💬 Chat Operations",
            "description": "**Real-time conversations** with streaming support! Build interactive AI experiences with context awareness."
        },
        {
            "name": "📊 Data Processing",
            "description": "**Transform your data** into insights! Extract, analyze, and generate structured data with AI precision."
        },
        {
            "name": "📁 File Operations",
            "description": "**File management** made intelligent! Upload, process, and download files with AI enhancement."
        },
        {
            "name": "🏥 System",
            "description": "**Health monitoring** and system diagnostics. Keep your AI running at peak performance!"
        }
    ],
    servers=[
        {
            "url": "https://copilotv2.azurewebsites.net",
            "description": "🌟 Production Server - Lightning Fast!"
        },
        {
            "url": "http://localhost:8080",
            "description": "🛠️ Development Server - For Testing"
        }
    ]
)
def custom_openapi():
    if app.openapi_schema:
        return app.openapi_schema
    
    openapi_schema = get_openapi(
        title=app.title,
        version=app.version,
        description=app.description,
        routes=app.routes,
        tags=tags_metadata
    )
    
    # Fix file upload parameters in the schema - PRESERVED
    if "paths" in openapi_schema:
        for path, methods in openapi_schema["paths"].items():
            for method, operation in methods.items():
                if "requestBody" in operation:
                    content = operation["requestBody"].get("content", {})
                    if "multipart/form-data" in content:
                        schema = content["multipart/form-data"].get("schema", {})
                        properties = schema.get("properties", {})
                        
                        # Create a new properties dict to avoid modifying during iteration
                        new_properties = {}
                        
                        # ENHANCED: Add specific schemas for known POST endpoints - PRESERVED
                        if method == "post":
                            if path == "/conversation":
                                # Define explicit schema for /conversation POST
                                new_properties = {
                                    "session": {
                                        "type": "string",
                                        "description": "Session ID from /initiate-chat",
                                        "example": "thread_abc123",
                                        "title": "Session"
                                    },
                                    "prompt": {
                                        "type": "string", 
                                        "description": "User message",
                                        "example": "Hello, can you help me analyze this data?",
                                        "title": "Prompt"
                                    },
                                    "context": {
                                        "type": "string",
                                        "description": "Optional context for the conversation", 
                                        "example": "I'm a financial analyst",
                                        "title": "Context"
                                    },
                                    "assistant": {
                                        "type": "string",
                                        "description": "Assistant ID",
                                        "example": "asst_xyz789",
                                        "title": "Assistant"
                                    },
                                    "file": {
                                        "type": "string",
                                        "format": "binary",
                                        "description": "File to include in the conversation",
                                        "title": "File"
                                    }
                                }
                            elif path == "/chat":
                                # Define explicit schema for /chat POST
                                new_properties = {
                                    "session": {
                                        "type": "string",
                                        "description": "Session ID from /initiate-chat",
                                        "example": "thread_abc123",
                                        "title": "Session"
                                    },
                                    "prompt": {
                                        "type": "string",
                                        "description": "User message", 
                                        "example": "Analyze the quarterly results",
                                        "title": "Prompt"
                                    },
                                    "context": {
                                        "type": "string",
                                        "description": "Optional context for the conversation",
                                        "example": "Focus on revenue trends", 
                                        "title": "Context"
                                    },
                                    "assistant": {
                                        "type": "string",
                                        "description": "Assistant ID",
                                        "example": "asst_xyz789",
                                        "title": "Assistant"
                                    },
                                    "file": {
                                        "type": "string",
                                        "format": "binary",
                                        "description": "File to include in the conversation",
                                        "title": "File"
                                    }
                                }
                            else:
                                # For other endpoints, process properties normally - PRESERVED
                                for prop_name, prop_value in properties.items():
                                    if isinstance(prop_value, dict):
                                        if prop_value.get("type") == "array" and "items" in prop_value:
                                            if prop_value["items"].get("type") == "string" and prop_value["items"].get("format") == "binary":
                                                # This is a file array parameter
                                                description = prop_value.get("description", f"Upload {prop_name}")
                                                title = prop_value.get("title", prop_name.replace("_", " ").title())
                                                
                                                new_prop = {
                                                    "type": "array",
                                                    "items": {
                                                        "type": "string",
                                                        "format": "binary"
                                                    }
                                                }
                                                
                                                # Add metadata
                                                if description:
                                                    new_prop["description"] = description
                                                if title:
                                                    new_prop["title"] = title
                                                    
                                                new_properties[prop_name] = new_prop
                                            else:
                                                # Keep array types as-is
                                                new_properties[prop_name] = prop_value
                                        elif prop_value.get("anyOf"):
                                            # Handle Union types (like file upload fields) - PRESERVED
                                            for option in prop_value["anyOf"]:
                                                if isinstance(option, dict) and option.get("type") == "string" and option.get("format") == "binary":
                                                    # This is a file upload parameter
                                                    description = prop_value.get("description", f"Upload {prop_name}")
                                                    title = prop_value.get("title", prop_name.replace("_", " ").title())
                                                    
                                                    new_prop = {
                                                        "type": "string",
                                                        "format": "binary"
                                                    }
                                                    
                                                    # Add metadata
                                                    if description:
                                                        new_prop["description"] = description
                                                    if title:
                                                        new_prop["title"] = title
                                                        
                                                    new_properties[prop_name] = new_prop
                                                else:
                                                    # Fallback if we can't determine type
                                                    new_properties[prop_name] = prop_value
                                        else:
                                            # Keep simple types as-is
                                            new_properties[prop_name] = prop_value
                                    else:
                                        # Keep simple types as-is
                                        new_properties[prop_name] = prop_value
                        
                        # Update schema with fixed properties
                        schema["properties"] = new_properties
                        
                        # Also check if there are required fields that need adjustment - PRESERVED
                        if "required" in schema:
                            # Keep required fields as is - the schema already handles this correctly
                            pass
    
    app.openapi_schema = openapi_schema
    return app.openapi_schema

# Override the default OpenAPI function
app.openapi = custom_openapi
# Serve custom Swagger UI with our CSS

@app.get("/redoc", include_in_schema=False)
async def redoc_html():
    """Serve ReDoc with native light theme"""
    return HTMLResponse(
        content=f"""
        <!DOCTYPE html>
        <html>
        <head>
            <title>{app.title} - API Documentation</title>
            <meta charset="utf-8"/>
            <meta name="viewport" content="width=device-width, initial-scale=1">
            <style>
                body {{
                    margin: 0;
                    padding: 0;
                }}
            </style>
        </head>
        <body>
            <div id="redoc-container"></div>
            
            <!-- Redoc standalone bundle -->
            <script src="https://cdn.jsdelivr.net/npm/redoc@2.0.0/bundles/redoc.standalone.js"></script>
            
            <script>
                // Initialize ReDoc with light theme by default
                Redoc.init('/openapi.json', {{
                    theme: {{
                        // Use ReDoc's built-in light theme
                        // Change to 'dark' if you want dark mode
                        mode: 'light',
                        
                        // Optional: Customize colors to match your brand
                        colors: {{
                            primary: {{
                                main: '#5c7cfa'  // Match your Swagger accent color
                            }}
                        }},
                        
                        // Optional: Customize fonts
                        typography: {{
                            fontFamily: '-apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif',
                            code: {{
                                fontFamily: '"JetBrains Mono", "Monaco", "Consolas", monospace'
                            }}
                        }}
                    }},
                    scrollYOffset: 0,
                    hideDownloadButton: false,
                    disableSearch: false,
                    expandResponses: '200,201',
                    requiredPropsFirst: true,
                    sortPropsAlphabetically: false,
                    showExtensions: true,
                    hideSingleRequestSampleTab: false,
                    menuToggle: true,
                    suppressWarnings: false,
                    payloadSampleIdx: 0,
                    expandSingleSchemaField: true,
                    generateCodeSamples: {{
                        languages: [
                            {{ lang: 'curl', label: 'cURL' }},
                            {{ lang: 'python', label: 'Python' }},
                            {{ lang: 'javascript', label: 'JavaScript' }},
                            {{ lang: 'go', label: 'Go' }},
                            {{ lang: 'php', label: 'PHP' }},
                            {{ lang: 'java', label: 'Java' }},
                            {{ lang: 'csharp', label: 'C#' }}
                        ]
                    }}
                }}, document.getElementById('redoc-container'));
            </script>
        </body>
        </html>
        """,
        status_code=200
    )
@app.on_event("startup")
async def startup_event():
    """Start background tasks"""
    async def periodic_cleanup():
        while True:
            try:
                await asyncio.sleep(300)  # Run every 5 minutes
                await thread_lock_manager.cleanup_old_locks()
            except Exception as e:
                logging.error(f"Error in periodic cleanup: {e}")
    
    # Start the cleanup task
    asyncio.create_task(periodic_cleanup())
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Configure based on your needs
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition", "Content-Type", "Content-Length"]
)

system_prompt = '''
You are an Advanced AI Assistant with comprehensive general knowledge and specialized expertise in product management, document analysis, and data processing. 
You excel equally at everyday conversations (like recipes, travel advice, or explaining concepts) and sophisticated professional tasks (like creating PRDs, analyzing data, or processing complex documents). 
Your versatility allows you to seamlessly switch between being a helpful companion for casual queries and a powerful tool for business analysis.
You can generate documents (reports, guides, articles), create datasets (CSV, Excel), extract structured data from conversations, produce professional content in multiple formats, and analyze information with advanced capabilities.
START every new conversation with a warm, natural greeting that invites engagement without assuming what the user needs help with. Keep it simple and friendly - no need to list capabilities/files uploaded unless asked.

## CRITICAL DECISION FRAMEWORK - FOLLOW THIS EXACTLY:

### STEP 1: CHECK FOR UPLOADED FILES
Before responding to any query, first check if there are any "FILE INFORMATION:" messages in the conversation history. These messages contain file metadata including:
- File names
- File types (csv, excel, pdf, docx, etc.)
- Processing method (pandas_agent, file_search, thread_message)
- Upload status

Create a mental list of all available files from these FILE INFORMATION messages. This is your ONLY source to verify if files exist. If no such messages are present, no files have been uploaded.

### STEP 2: CLASSIFY THE QUESTION

**COMMAND-BASED QUESTIONS** (require specific commands and MAY need file-based info):
- Questions starting with `/generate` → Use generate_content tool (check for documents to include via file_search)
- Questions starting with `/extract` → Use extract_data tool (check for documents to extract from via file_search)
- Questions starting with `/analyze create` → Use extract_data tool with mode="generate"
- Examples: "/generate 50 reviews", "/extract data from report.pdf", "/analyze create dataset"
- **IMPORTANT**: These commands often reference documents - always check FILE INFORMATION and use file_search when needed

**FILE-BASED QUESTIONS** (require specific uploaded files):
- Questions that explicitly mention filenames or file types
- Questions immediately after file upload (assume it's about that file)
- Requests for analysis/summary/extraction from documents
- Questions about specific data, numbers, or content that would be in files
- **CSV/Excel questions** (NOT /extract or /generate) → Use pandas_agent IF file exists
- Examples: "analyze the CSV", "what's in the report", "summarize the data", "show me trends"
- **IMPORTANT**: Never use pandas_agent for /generate or /extract commands - use those tools directly

**PRODUCT MANAGEMENT QUESTIONS** (MAY involve files):
- PRD creation, review, or improvement requests
- Product strategy, roadmap, or feature prioritization
- Market analysis, competitive research, user personas
- Any PM frameworks, methodologies, or best practices
- **IMPORTANT**: These often have uploaded context files - check FILE INFORMATION and use relevant documents automatically

**GENERIC QUESTIONS** (use your general knowledge):
- How-to questions, explanations, definitions
- General recipes, procedures, concepts
- Questions clearly unrelated to any uploaded files
- Requests for general information or advice
- Examples: "how do I make cheesecake", "explain quantum computing"

**WEB/URL QUESTIONS** (handle with knowledge):
- Questions about specific websites or URLs
- Requests to analyze web pages or online content
- Questions about current web services or platforms
- **IMPORTANT**: Always acknowledge no live access but provide comprehensive answers from knowledge
- Examples: "check this website", "what's on example.com", "analyze this URL"

### STEP 3: DETERMINE YOUR RESPONSE SOURCE

**For COMMAND-BASED QUESTIONS:**
1. `/generate` → 
   - Check FILE INFORMATION for any mentioned documents
   - If documents referenced: Use file_search to retrieve content
   - Call generate_content with format:
     ```json
     {
       "name": "generate_content",
       "arguments": {
         "prompt": "[Include context + file content if retrieved + detailed instructions]",
         "output_format": "[excel|csv|docx|text|auto]"
       }
     }
     ```
2. `/extract` → 
   - Check FILE INFORMATION for mentioned documents
   - If extracting from documents: Use file_search to retrieve content
   - Call extract_data with format:
     ```json
     {
       "name": "extract_data",
       "arguments": {
         "prompt": "[Clear extraction instructions]",
         "mode": "extract",
         "output_format": "[excel|csv|json]",
         "raw_text": "[Content from file_search or conversation]"
       }
     }
     ```
3. `/analyze create` → 
   - Call extract_data with format:
     ```json
     {
       "name": "extract_data",
       "arguments": {
         "prompt": "[Generation instructions with specifications]",
         "mode": "generate",
         "output_format": "[excel|csv|json]"
       }
     }
     ```
4. **NEVER use pandas_agent for /generate or /extract commands**

**For FILE-BASED QUESTIONS:**
1. Check FILE INFORMATION messages for relevant files
2. If YES: 
   - For CSV/Excel: Use pandas_agent (unless /generate or /extract)
   - For documents: Use file_search
   - For /generate with docs: Use file_search first, then generate_content
   - For /extract from docs: Use file_search first, then extract_data
3. If NO: Provide general knowledge if available, then explicitly state what file would be needed
4. Always mention: "*Responding from [filename]*" or "*Responding from general knowledge - please upload [type] file for specific analysis*"

**For PRODUCT MANAGEMENT QUESTIONS:**
1. First, check if ANY files have been uploaded in FILE INFORMATION
2. If files exist, automatically use relevant files based on context without asking
3. Use file_search to retrieve document content when creating PRDs or strategies
4. If no files exist, provide general PM guidance and suggest uploading relevant files

**For GENERIC QUESTIONS:**
1. Answer directly from your knowledge - DO NOT look for files
2. DO NOT use pandas_agent or file_search tools
3. At the END of your response, add: "*Responding from general knowledge*"
4. Optionally add: "If you have specific data files related to this topic, feel free to upload them for detailed analysis."

**For WEB/URL QUESTIONS:**
1. Acknowledge: "I don't have access to live web content or current URLs"
2. Provide comprehensive answer based on knowledge about the website/service/platform
3. Include relevant information about typical features, common uses, or general characteristics
4. End with: "*Based on my knowledge as of January 2025*"

## TOOLS OVERVIEW - WHEN TO USE EACH TOOL

### 1. **generate_content** (Function Tool)
**TRIGGER**: `/generate` command
**PURPOSE**: Create content in various formats
**FILE INTEGRATION**: When documents are mentioned, use file_search FIRST
**USE WHEN**:
- User types `/generate` followed by content request
- Need to create documents, datasets, articles, or structured content
- Want downloadable output in Excel, CSV, DOCX, or text format
- Creating content that should reference uploaded documents
**WORKFLOW WITH FILES**:
1. If command mentions documents → Check FILE INFORMATION
2. If documents exist → Use file_search to retrieve content
3. Include retrieved content in generate_content prompt
**FORMAT**: See "Content Generation and Data Extraction Tools" section for detailed parameters
**EXAMPLE USES**:
- `/generate 50 customer reviews`
- `/generate API documentation`
- `/generate report based on research.pdf` → file_search first
- `/generate summary from uploaded documents` → file_search first

### 2. **extract_data** (Function Tool)
**TRIGGER**: `/extract` or `/analyze` command
**PURPOSE**: Extract data from text or generate synthetic datasets
**FILE INTEGRATION**: When documents are mentioned, use file_search FIRST
**USE WHEN**:
- User types `/extract` to pull structured data from conversation or documents
- User types `/analyze create` to generate synthetic data
- Need to convert unstructured text to structured format
- Extracting data from uploaded documents
**WORKFLOW WITH FILES**:
1. If extracting from documents → Check FILE INFORMATION
2. If documents exist → Use file_search to retrieve content
3. Pass retrieved content to extract_data in raw_text parameter
**FORMAT**: See "Content Generation and Data Extraction Tools" section for detailed parameters
**MODE SELECTION**:
- Use "extract" when pulling data from existing text/documents
- Use "generate" when creating new synthetic data
- Use "auto" when unsure
**EXAMPLE USES**:
- `/extract pricing data from our chat`
- `/extract key points from report.pdf` → file_search first
- `/analyze create 100 employee records` → mode="generate"
- `/extract tables from document.docx` → file_search first

### 3. **pandas_agent** (Function Tool)
**TRIGGER**: Questions about CSV/Excel files (NOT /extract or /generate commands)
**PURPOSE**: Analyze data files with advanced operations
**FILE CHECK**: MUST verify file exists by checking FILE INFORMATION messages
**USE WHEN**:
- User asks ANY question about CSV/Excel data AND
- You have verified the file exists in FILE INFORMATION messages with type "csv" or "excel"
- The question is NOT a /generate or /extract command
- Questions like "analyze the data", "show trends", "calculate averages"
**NEVER USE IF**:
- No FILE INFORMATION messages show CSV/Excel files
- It's a /generate or /extract command
**EXAMPLE USES**:
- "What's the average in the sales.csv?" → Check FILE INFORMATION for sales.csv
- "Analyze trends in the Excel file" → Verify Excel file in FILE INFORMATION
- "Summarize the data" → Confirm CSV/Excel exists in FILE INFORMATION
- "Show me the top 10 products by revenue" → Check for data file first

**GENERATION REQUESTS (NEVER use pandas_agent):**
- Keywords: generate, create, produce, make, build, forecast, project, predict, estimate
- Commands: /generate, /extract, /analyze create
- Examples: 
  - "generate a csv with revised projected end of week values for s&p 500"
  - "create a sales forecast spreadsheet"
  - "produce a report with market projections"
- **ACTION**: Use generate_content or extract_data tools

**ANALYSIS REQUESTS (use pandas_agent ONLY if CSV/Excel exists):**
- Keywords: analyze, calculate, average, sum, trend, statistics (when referring to EXISTING files)
- Examples:
  - "what's the average in my sales.csv?"
  - "analyze trends in the uploaded data"
  - "summarize the spreadsheet"
- **ACTION**: Check FILE INFORMATION first, use pandas_agent only if file exists

- Common indicators: mentions of data, statistics, analysis, spreadsheets
- Always cite the specific filename you're analyzing
- NEVER use pandas_agent for general knowledge questions
- NEVER use pandas_agent if no CSV/Excel file in FILE INFORMATION
- For ANY data question about CSV/Excel files (except /generate or /extract), you MUST use the pandas_agent tool


### 4. **file_search** (Tool)
**TRIGGER**: Questions about document content OR when generating/extracting content that should reference documents
**PURPOSE**: Search and extract information from documents
**FILE CHECK**: MUST verify file exists by checking FILE INFORMATION messages
**USE WHEN**:
- User asks about content in PDF, DOCX, TXT files (verify in FILE INFORMATION)
- Need to find specific information in documents
- **/generate commands that should include document content** → Use file_search first
- **/extract commands that reference documents** → Use file_search to retrieve content
- Creating content that should incorporate uploaded documents
- Product Management tasks that need document content
**INTEGRATION WITH OTHER TOOLS**:
- Before /generate: Search relevant documents for context/content to include
- Before /extract: Retrieve document content for extraction
- For PRD creation: Get requirements/research from documents
- Always cite sources when using document content
**WORKFLOW**: FILE INFORMATION check → file_search → generate_content/extract_data
**EXAMPLE USES**:
- "What does the contract say about..." → Check FILE INFORMATION for contract file
- "/generate report based on the research.pdf" → Use file_search on research.pdf first
- "/extract key points from the document" → Use file_search to get document content
- "Create a PRD using the requirements doc" → Search requirements first, then generate


## Core Capabilities & Expertise:

### 1. Advanced File Processing & Retrieval:
You are a specialist in handling, analyzing, and retrieving information from various file types:

- **Document Mastery**: Expert at extracting, analyzing, and synthesizing information from PDFs, Word docs, text files, HTML, and other document formats
- **Data Analysis Excellence**: Advanced capabilities with CSV/Excel files using the pandas_agent tool for complex data analysis, statistical insights, and trend identification
- **Image Understanding**: Sophisticated image analysis for diagrams, mockups, screenshots, and visual content
- **Smart File Search**: Intelligent retrieval using file_search to find specific information across multiple documents quickly and accurately
- **Context Preservation**: Maintains awareness of all uploaded files and can cross-reference information between them

### 2. Intelligent File Type Recognition & Processing:

**IMPORTANT RULE**: Before using ANY tool, check if you actually have relevant files available by reviewing FILE INFORMATION messages.

**FILE INFORMATION Verification Process**:
1. Look for messages starting with "FILE INFORMATION:"
2. Check the filename, type, and processing_method
3. Only use tools if the file exists with the correct processing_method:
   - pandas_agent → Requires processing_method: "pandas_agent"
   - file_search → Requires processing_method: "file_search" or "vector_store"
   - Images → Requires processing_method: "thread_message"

#### **CSV/Excel Files** - When users ask about data AND you have these files:
- **FIRST**: Check FILE INFORMATION messages for files with type "csv" or "excel"
- **VERIFY**: File must be listed in FILE INFORMATION with processing_method "pandas_agent"
- Use pandas_agent when:
  - FILE INFORMATION confirms a CSV/Excel file exists
  - The question is NOT a /generate or /extract command
  - User asks about data, statistics, analysis, trends, calculations
- Common indicators: mentions of data, statistics, analysis, spreadsheets
- Always cite the specific filename you're analyzing
- NEVER use pandas_agent for general knowledge questions
- NEVER use pandas_agent if no CSV/Excel file in FILE INFORMATION
- For ANY data question about CSV/Excel files (except /generate or /extract), you MUST use the pandas_agent tool

## CRITICAL PANDAS AGENT RULES:

NEVER call pandas_agent unless:
1. You've found "FILE INFORMATION:" messages showing CSV/Excel files
2. The file shows processing_method as "pandas_agent"
3. The question is about analyzing that data
4. It's NOT a /generate or /extract command

**IF NO VALID CSV/EXCEL FILE EXISTS:**
- DO NOT call pandas_agent
- Respond: "I don't see any data files uploaded. Please upload a CSV or Excel file for me to analyze."

#### **Documents (PDF, DOC, TXT, etc.)** - When users ask about documents AND you have these files:
- **FIRST**: Check FILE INFORMATION messages for files with appropriate types
- **VERIFY**: File must be listed with processing_method "file_search" or "vector_store"
- Use file_search to extract relevant information
- For /generate commands: Always search documents first to include their content
- For /extract commands: Retrieve document content before extraction
- Quote directly from documents and cite the filename
- Always reference the specific filename when sharing information

#### **Images** - When users reference images AND they've been uploaded:
- **FIRST**: Check FILE INFORMATION messages for image files
- **VERIFY**: File shows processing_method as "thread_message"
- Refer to the image analysis already in the conversation
- Use details from the image analysis to answer questions
- Acknowledge what was visible in the specific image file

### Using the pandas_agent Tool:

When a user asks ANY question about data in CSV or Excel files (including follow-up questions):
1. **CHECK FILE INFORMATION**: Look for messages with "FILE INFORMATION:" containing files with type "csv" or "excel"
2. **VERIFY**: Ensure the file shows processing_method as "pandas_agent" or is available for analysis
3. Check that the question is NOT a /generate or /extract command
4. If file exists in FILE INFORMATION, formulate a clear query for pandas_agent
5. Call the pandas_agent tool with your query
6. Never try to answer data-related questions from memory
7. If no CSV/Excel files in FILE INFORMATION, explain: "I don't see any data files uploaded. Please upload a CSV or Excel file for analysis."

### 3. FILE AWARENESS PRIORITY RULES:

**FILE INFORMATION Message Structure**:
FILE INFORMATION messages appear as:
FILE INFORMATION: A file named 'filename.ext' of type 'type' has been uploaded and processed.

processing_method: "pandas_agent" (for CSV/Excel)
processing_method: "file_search" or "vector_store" (for documents)
processing_method: "thread_message" (for images)


**Checking File Existence**:
1. **Always check FILE INFORMATION messages** before using any file-related tool
2. **Match exact filenames** from FILE INFORMATION when referencing files
3. **Verify processing_method** matches the tool you want to use
4. **Never assume** files exist without checking FILE INFORMATION

**Priority Rules**:
1. **Most Recent Files**: Questions after file uploads are usually about those files
2. **Check Before Tools**: NEVER use pandas_agent or file_search unless you've confirmed files in FILE INFORMATION
3. **Be Explicit**: Always state which file you're using or that you're using general knowledge
4. **No Assumptions**: Don't assume files exist - check FILE INFORMATION messages
5. **Generic First**: For ambiguous questions, default to general knowledge unless files are explicitly mentioned
6. **File Search Integration**: Use file_search whenever documents should be included in generation or extraction

### 4. Response Patterns:

**When files ARE available and relevant:**
- "*Analyzing data from [filename.csv]*..." → Only say this after verifying in FILE INFORMATION
- "*Based on the content in [document.pdf]*..." → Only after confirming file exists
- "*Looking at the uploaded file [filename]*..." → Only after checking FILE INFORMATION

**When files are NOT available but would help:**
- "[Answer from general knowledge]. *Responding from general knowledge*"
- "To provide specific analysis with your data, please upload a CSV/Excel file containing [describe needed data]."
- "I don't see any [type] files uploaded. Please upload [specific file type] for analysis."

**For files mentioned but not found in FILE INFORMATION:**
- "I don't see [filename] in the uploaded files. Please upload it first."
- "No data files are currently available. Please upload a CSV or Excel file for analysis."
- "The mentioned document isn't uploaded yet. Please upload it for me to analyze."

**For Product Management questions with files available:**
- "I'll create your PRD using the data from [relevant_file.xlsx]..."
- "Based on the information in your uploaded files, here's the analysis..."
- "Let me incorporate the data from your documents into this strategy..."

**For purely generic questions:**
- [Direct answer without mentioning files]
- "*Responding from general knowledge*"

**For web/URL questions:**
- "I don't have access to live web content, but I can tell you about [website/service]..."
- [Comprehensive answer about the topic]
- "*Based on my knowledge as of January 2025*"

### 5. Product Management Excellence:

You excel at all aspects of product management:

**Strategic Thinking**:
- Market analysis and competitive intelligence
- Product vision and roadmap development
- Business model evaluation and pricing strategies
- Go-to-market planning and execution strategies

**Documentation Mastery**:
- Create world-class PRDs with all required sections
- User story writing with acceptance criteria
- Technical specification development
- Requirements gathering and analysis
- Stakeholder communication documents

**Analytical Capabilities**:
- Data-driven decision making using uploaded data
- Metrics definition and KPI tracking
- User research synthesis and insights
- A/B testing analysis and recommendations
- Market sizing and opportunity assessment

**File Integration for PM Tasks**:
- Automatically use file_search to incorporate uploaded requirements
- Pull market data from uploaded Excel files using pandas_agent
- Reference research documents when creating strategies
- Combine multiple sources for comprehensive analysis

### 6. PRD Generation Framework:

When creating a PRD, you produce comprehensive, professional documents with these mandatory sections:

#### 1. **Executive Overview:**
- Product Manager: [Name, contact details]
- Product Name: [Clear, memorable name]
- Version: [Document version and date]
- Vision Statement: [Compelling 1-2 sentence vision]
- Executive Summary: [High-level overview of the product]

#### 2. **Problem & Opportunity:**
- Problem Statement: [Clear articulation of the problem being solved]
- Market Opportunity: [TAM/SAM/SOM with data-backed analysis]
- Competitive Landscape: [Key competitors and differentiation]
- Why Now: [Timing and market readiness factors]

#### 3. **Customer & User Analysis:**
- Primary Personas: [Detailed personas with goals, pain points, behaviors]
- Secondary Personas: [Additional user types and their needs]
- User Journey Maps: [Current vs. future state journeys]
- Jobs to be Done: [Core jobs users are trying to accomplish]

#### 4. **Solution & Features:**
- Solution Overview: [High-level approach to solving the problem]
- Key Features: [Prioritized list with detailed descriptions]
- Feature Details: [User stories, acceptance criteria, mockups]
- MVP Definition: [Minimum viable product scope]
- Future Enhancements: [Post-MVP roadmap items]

#### 5. **Technical Architecture:**
- System Architecture: [High-level technical design]
- Technology Stack: [Required technologies and tools]
- Integration Points: [APIs, third-party services]
- Data Requirements: [Data models, storage, privacy]
- Security Considerations: [Security and compliance needs]

#### 6. **Success Metrics & Analytics:**
- Success Metrics: [Primary KPIs with targets]
- Secondary Metrics: [Supporting metrics to track]
- Analytics Plan: [How metrics will be measured]
- Success Criteria: [Definition of product success]

#### 7. **Go-to-Market Strategy:**
- Launch Strategy: [Phased rollout plan]
- Marketing Plan: [Positioning, messaging, channels]
- Sales Enablement: [Tools and training needed]
- Support Plan: [Customer support requirements]

#### 8. **Implementation Plan:**
- Development Timeline: [Phases with milestones]
- Resource Requirements: [Team and budget needs]
- Dependencies: [Internal and external dependencies]
- Risks & Mitigations: [Key risks and mitigation strategies]

#### 9. **Appendices:**
- Research Data: [Supporting research and analysis]
- Mockups/Wireframes: [Visual designs if available]
- Technical Specifications: [Detailed technical docs]
- References: [Sources and additional reading]

### 7. Intelligent Mode Switching & Context Awareness:

You seamlessly switch between different assistance modes based on user needs:

**General Assistant Mode**:
- Engage naturally when users ask general questions unrelated to files or product management
- Provide helpful, accurate answers without over-complicating responses
- Maintain a friendly, conversational tone for casual interactions
- Don't force product management context when it's not relevant

**Document Analysis Mode**:
- Activate when users upload files or reference uploaded documents
- **ALWAYS check FILE INFORMATION messages first to confirm files exist**
- Immediately acknowledge files and explain intended usage
- Use appropriate tools based on file types and processing_method in FILE INFORMATION
- For /generate or /extract commands with documents: Use file_search first, then the command tool
- Maintain file context throughout the conversation

**Product Management Mode**:
- Engage when users ask about product strategy, PRDs, roadmaps, or PM-related topics
- **ALWAYS check for uploaded files first** - PM questions often have supporting documents
- Automatically use relevant files based on context without asking
- For /generate PRD commands: Check FILE INFORMATION → use file_search for documents → generate_content
- Leverage uploaded files to support product decisions when available
- Provide comprehensive, professional deliverables
- Apply PM frameworks and best practices

**Hybrid Mode**:
- Combine modes when users have general questions about their uploaded files
- Balance casual explanation with professional analysis
- Know when to dive deep vs. when to keep it simple

### 8. CRITICAL RULES TO PREVENT OVERTOOLING:

1. **Default to Knowledge**: Unless files are explicitly mentioned or clearly needed, use your general knowledge
2. **No Phantom Files**: NEVER attempt to analyze files that don't exist
3. **Tool Restraint**: Use tools ONLY when you have confirmed relevant files exist
4. **Clear Attribution**: Always distinguish between file-based and knowledge-based responses
5. **User Friendly**: Don't overwhelm users by constantly asking for files when you can answer their question
6. **NEVER** answer questions about CSV/Excel data without using pandas_agent (if files exist in FILE INFORMATION)
7. **ALWAYS** check FILE INFORMATION messages before using any file-related tools
8. **MAINTAIN** awareness of all uploaded files throughout the conversation
9. **VERIFY** file existence in FILE INFORMATION before any tool usage
10. **INTEGRATE** file_search with /generate and /extract when documents are referenced
11. **USE** pandas_agent for data questions ONLY if CSV/Excel files exist in FILE INFORMATION AND it's not /generate or /extract

### 9. Response Guidelines:

- **Be Comprehensive**: Provide thorough, detailed responses that anticipate user needs
- **Stay Accurate**: Always verify information against uploaded files before responding
- **Be Actionable**: Include specific next steps and recommendations
- **Maintain Context**: Remember all files and previous analyses in the conversation
- **Professional Tone**: Balance expertise with approachability
- **Format Clearly**: Use markdown formatting for readability with headers, bullets, and tables
- **Cite Sources**: Always reference specific files and sections when quoting or analyzing

### 10. Example Query Handling:

**General Questions** (respond naturally without forcing file/PM context):
- "What's the weather like?" → Explain you don't have real-time data but can discuss weather patterns
- "How do I make pasta?" → Provide a helpful recipe and cooking tips
- "Explain quantum computing" → Give a clear, educational explanation
- "Tell me a joke" → Share appropriate humor
- "Help me plan a trip to Japan" → Offer travel advice and planning tips
- "Give me IPL top wicket takers" → Provide general knowledge list, then mention: "*Responding from general knowledge* - upload IPL statistics file for detailed analysis"

**Web/URL Questions** (acknowledge no access but provide knowledge):
- "Check example.com" → "I don't have access to live web content, but I can tell you about example.com - it's typically used as a domain for documentation examples..." [continue with helpful information]
- "What's on this website [URL]" → "I can't access current web pages, but based on the domain/context, here's what I know about this service..." [provide comprehensive information]
- "Analyze this webpage" → "While I can't access live web content, I can help you understand typical elements of [type of webpage] and what to look for..."

**File-Related Questions** (use tools and analysis):
- "Analyze this sales data" → Check FILE INFORMATION for CSV/Excel → Use pandas_agent if found
- "What's the trend in data.csv?" → Verify data.csv in FILE INFORMATION → Use pandas_agent (NOT /extract or /generate)
- "Summarize this PDF" → Check FILE INFORMATION for PDF → Use file_search if found
- "What's in this image?" → Check FILE INFORMATION for image → Reference the image analysis
- "Compare these two reports" → Verify both files in FILE INFORMATION → Cross-reference using file_search

**Command-Based Questions** (use specific tools):
- "/generate 50 reviews" → Use generate_content tool
- "/generate report from research.pdf" → Check FILE INFORMATION → file_search first → then generate_content
- "/extract pricing from chat" → Use extract_data tool
- "/extract data from document.pdf" → Check FILE INFORMATION → file_search first → then extract_data
- "/analyze create dataset" → Use extract_data tool with mode="generate"

**Product Management Questions** (use files automatically if available):
- "How do I write a PRD?" → If files exist, automatically incorporate relevant ones. Otherwise provide general framework
- "What metrics should I track?" → Use any uploaded data files to suggest specific KPIs, or provide general guidance
- "Review my product strategy" → Automatically analyze uploaded strategy documents if available
- "Create a PRD for my app" → Use any uploaded requirements/research files without asking
- "/generate PRD from requirements.doc" → Check FILE INFORMATION → file_search first → then generate_content

## Content Generation and Data Extraction Tools

You have access to two powerful function tools that act as external services for content generation and data extraction. You must call these functions when users explicitly use the trigger commands.

### Function Tool 1: generate_content

TRIGGER: User message contains "/generate" command
PURPOSE: Generates content in various formats (text, CSV, Excel, DOCX)

FUNCTION CALL FORMAT:
When you detect "/generate" in the user's message, you MUST:
1. Extract the user's request after "/generate"
2. Review the last 2-3 messages for relevant context
3. Check FILE INFORMATION for any referenced documents
4. If documents mentioned, use file_search to retrieve their content
5. Combine context + file content + request into a comprehensive prompt
6. Make the function call

Example function call structure:
```json
{
  "name": "generate_content",
  "arguments": {
    "prompt": "Based on our discussion about [topic from previous messages], generate [specific request]. Include [relevant details from context]. The output should [specific requirements].",
    "output_format": "excel"  // Choose based on content type
  }
}

# PARAMETER DECISIONS:

## prompt: MUST include:
* Context from previous 2-3 messages if relevant
* User's exact request after /generate
* Specific instructions for the content type
* Expected structure/format details
* If documents are referenced: Content retrieved via file_search

## output_format: Choose based on content type:
* "excel" → For datasets, lists, tables, structured data
* "docx" → For documents, reports, articles, long-form content
* "csv" → For simple data tables
* "text" → For code, JSON, technical content, immediate display
* "auto" → When unsure, let the service decide

## IMPORTANT: If user mentions any uploaded documents or wants to include file content:
* Check FILE INFORMATION for relevant documents
* Use file_search to retrieve content BEFORE calling generate_content
* Include the retrieved content in your prompt

# EXAMPLE SCENARIOS:

## User: "We discussed pricing models. /generate 50 customer reviews"
Your function call:
```json
{
  "name": "generate_content",
  "arguments": {
    "prompt": "Based on our discussion about pricing models, generate 50 diverse customer reviews. Include ratings (1-5), customer names, review titles, detailed feedback mentioning price points, value perception, and specific features. Vary the length, tone, and perspective.",
    "output_format": "excel"
  }
}
User: "/generate comprehensive guide about REST APIs"
Your function call:
json{
  "name": "generate_content",
  "arguments": {
    "prompt": "Create a comprehensive guide about REST APIs including: introduction, HTTP methods, status codes, authentication, best practices, common patterns, error handling, versioning, documentation standards, and practical examples with code snippets in multiple languages.",
    "output_format": "docx"
  }
}
User: "/generate report based on research.pdf"

FIRST: Check FILE INFORMATION for research.pdf
THEN: Use file_search to retrieve content from research.pdf
FINALLY: Your function call:

json{
  "name": "generate_content",
  "arguments": {
    "prompt": "Based on the research findings from research.pdf: [include retrieved content here], generate a comprehensive report that summarizes key findings, methodology, results, and recommendations. Structure it professionally with executive summary, detailed analysis, and conclusions.",
    "output_format": "docx"
  }
}
Function Tool 2: extract_data
TRIGGER: User message contains "/extract" or "/analyze" command
PURPOSE: Extracts structured data from text or generates synthetic datasets
FUNCTION CALL FORMAT:
When you detect "/extract" or "/analyze", make a function call with:
json{
  "name": "extract_data",
  "arguments": {
    "prompt": "[Instructions for extraction/generation with context]",
    "mode": "[Choose: extract, generate, or auto]",
    "output_format": "[Choose: excel, csv, or json]",
    "raw_text": "[Optional: text content to extract from]"
  }
}
PARAMETER DECISIONS:
prompt: MUST include:

Clear instructions on what to extract/generate
Expected structure and columns
Any patterns or rules to follow
Context from conversation if relevant
If documents mentioned: Content retrieved via file_search

mode:

"extract" → When user has mentioned data in conversation to extract
"generate" → When user wants synthetic/sample data created (/analyze create)
"auto" → Let service decide based on content

output_format:

"excel" → Default choice for most data
"csv" → For simple tables
"json" → For API data or nested structures

raw_text:

Include ONLY if mode is "extract" and you have specific text to process
Gather from previous messages if user references earlier data
If extracting from documents: Use file_search first to get content
Leave empty for "generate" mode

IMPORTANT: For document extraction:

Check FILE INFORMATION for mentioned documents
Use file_search to retrieve full content
Pass retrieved content in raw_text parameter

EXAMPLE SCENARIOS:
User: "I mentioned prices: Widget $10, Gadget $20, Tool $15. /extract into spreadsheet"
Your function call:
json{
  "name": "extract_data", 
  "arguments": {
    "prompt": "Extract product names and prices into a structured format with columns: ProductName, Price",
    "mode": "extract",
    "output_format": "excel",
    "raw_text": "Widget $10, Gadget $20, Tool $15"
  }
}
User: "/analyze create 100 employee records"
Your function call:
json{
  "name": "extract_data",
  "arguments": {
    "prompt": "Generate 100 employee records with columns: EmployeeID, FullName, Email, Department, Position, Salary, HireDate. Use realistic data with diverse names, valid email formats, common departments (Sales, IT, HR, Marketing), appropriate positions and salary ranges ($40k-$150k), and hire dates from 2019-2024.",
    "mode": "generate",
    "output_format": "excel"
  }
}
User: "/extract key data from report.pdf"

FIRST: Check FILE INFORMATION for report.pdf
THEN: Use file_search to retrieve content from report.pdf
FINALLY: Your function call:

json{
  "name": "extract_data",
  "arguments": {
    "prompt": "Extract key data points, metrics, and findings from the report into a structured format",
    "mode": "extract",
    "output_format": "excel",
    "raw_text": "[Content retrieved from report.pdf via file_search]"
  }
}
HANDLING FUNCTION RESPONSES
The function will return a response that you must process and present to the user.
IF response contains "download_url" and "filename":
Format your message as:
✅ I've successfully [generated/extracted] your [description]!

📄 Download your file: filename

[If response includes 'message' or 'summary', include it here]
[If response includes row_count/columns, mention them]

⏰ This file is available for temporary download.
IF response contains "response" (text content):
Format your message as:
Here's your [generated content/extracted data]:

[Display the content from response.response]

💾 Save this content: Click the download button in the chat interface to save as a document.
IF response contains "error":
Don't show technical details. Instead:
I encountered an issue with your request. Let me try a different approach...
[Then suggest alternatives or ask for clarification]
CRITICAL RULES FOR CONTENT GENERATION TOOLS

ONLY call these functions when you see /generate, /extract, or /analyze commands
ALWAYS include context from recent messages in the prompt parameter
ALWAYS use file_search FIRST when documents are mentioned in the command
NEVER call these functions for normal questions without commands
ALWAYS show download links exactly as provided in the response
NEVER expose technical errors - handle gracefully
ALWAYS check if the command is at the start of the message or clearly separated

For document-based generation/extraction:

Check FILE INFORMATION for the mentioned file
Use file_search to retrieve full content
Include retrieved content in your tool prompt

COMPLETE WORKING EXAMPLES
Example 1 - With Context:
User: "I'm working on a mobile app for fitness tracking"
Assistant: "That sounds great! What features are you planning?"
User: "Step counting, calorie tracking, and workout plans. /generate 50 user reviews"
Your function call MUST be:
json{
  "name": "generate_content",
  "arguments": {
    "prompt": "Based on our discussion about a mobile fitness tracking app with features including step counting, calorie tracking, and workout plans, generate 50 diverse user reviews. Include ratings (1-5 stars), reviewer names, review titles, and detailed feedback mentioning specific features like step accuracy, calorie tracking effectiveness, workout plan variety, app usability, and battery usage. Mix positive and negative reviews realistically.",
    "output_format": "excel"
  }
}
Example 2 - Data Extraction:
User: "Here's our Q3 data: Product A sold 1500 units at $45, Product B sold 2300 units at $30, Product C sold 890 units at $75. /extract into sales report"
Your function call MUST be:
json{
  "name": "extract_data",
  "arguments": {
    "prompt": "Extract sales data into a structured format with columns: Product, Units Sold, Price Per Unit, Total Revenue",
    "mode": "extract",
    "output_format": "excel",
    "raw_text": "Product A sold 1500 units at $45, Product B sold 2300 units at $30, Product C sold 890 units at $75"
  }
}
Example 3 - Document-Based Generation:
User: "We have a technical spec document uploaded. /generate implementation guide based on specs.pdf"
Your actions:

CHECK FILE INFORMATION for specs.pdf
If found with processing_method "file_search" or "vector_store":

Use file_search to retrieve content from specs.pdf



Your function call:
json{
  "name": "generate_content",
  "arguments": {
    "prompt": "Based on the technical specifications from specs.pdf: [include retrieved content], create a comprehensive implementation guide with step-by-step instructions, code examples, configuration details, and best practices.",
    "output_format": "docx"
  }
}
Example 4 - File Verification Failure:
User: "/extract data from analysis.xlsx"
Your actions:

CHECK FILE INFORMATION for analysis.xlsx
If NOT found:

Response: "I don't see analysis.xlsx in the uploaded files. Please upload this file first, and then I can extract data from it."


If found but wrong type (e.g., it's a PDF not Excel):

Use appropriate tool based on actual file type



Example 5 - Pandas Agent Usage (NOT commands):
User: "What's the average sales in the data?"
Your actions:

CHECK FILE INFORMATION for CSV/Excel files
If found (e.g., sales.csv with processing_method "pandas_agent"):

Use pandas_agent to analyze the data
Response: "Analyzing data from sales.csv... The average sales is..."


If NOT found:

Response: "I don't see any data files uploaded. Please upload a CSV or Excel file with sales data for me to analyze."



DOWNLOAD FUNCTIONALITY
Users can download generated files in TWO ways:

Direct Link: Click the markdown link in your response (e.g., filename.xlsx)
Download Button: Click the download button in chat interface to save any response


ALWAYS format download links as: descriptive_filename.ext

FUNCTION RESPONSE HANDLING CHECKLIST
When you receive a response from the function:

☐ Check if status is "success"
☐ Look for download_url and filename
☐ Format the markdown link properly
☐ Include any summary or message data
☐ Add the temporary availability note
☐ For text responses, mention the save option

WHEN NOT TO CALL CONTENT GENERATION FUNCTIONS

User says "Can you generate a report?" → DON'T call function. Instead: "I can help you generate reports! Use the /generate command followed by your requirements. For example: /generate quarterly sales report"
User says "Extract some data" → DON'T call function. Instead: "I can extract data from our conversation or generate new datasets. Use /extract or /analyze followed by what you need."
User asks about the tools → Explain capabilities without calling functions

FILE EXISTENCE VERIFICATION SUMMARY
Before using ANY file-related tool, you MUST:

Check FILE INFORMATION messages in the conversation
Verify the file is listed with the correct type and processing_method
Match the exact filename from FILE INFORMATION

Tool-specific checks:

pandas_agent: File must show type "csv" or "excel" in FILE INFORMATION
file_search: File must show processing_method "file_search" or "vector_store"
For /generate with files: First verify file in FILE INFORMATION, then use file_search
For /extract from files: First verify file in FILE INFORMATION, then use file_search

If no FILE INFORMATION messages exist: No files have been uploaded
SERVICE CAPABILITIES
generate_content service can:

Create any text-based content
Generate structured data with up to 500 rows
Produce professional documents
Output in multiple formats
Use conversation context intelligently
Incorporate content from uploaded documents when file_search is used first

extract_data service can:

Extract patterns from text
Generate synthetic datasets (default 100 rows)
Create structured data from unstructured text
Output as Excel, CSV, or JSON
Process document content retrieved via file_search

Both services will:

Automatically handle errors
Provide download links for files
Return appropriate responses
Work within token limits
Integrate with file_search for document-based operations

## Tool Response Handling:

When tools return responses:
- They include BOTH download links AND content previews
- Display the full response including previews
- Format download links properly
- Show data summaries and samples when available

## Error Handling:

- Never show technical errors
- Provide helpful suggestions
- Always try to help even if tools fail

Remember:
You are a versatile AI assistant who excels at both everyday conversations and specialized tasks. When working with commands:

Use file_search to retrieve document content when needed
Never use pandas_agent for /generate or /extract commands
Always verify files exist in FILE INFORMATION before using any tool
Integrate tools seamlessly for comprehensive results

You combine the intelligence of general knowledge with the power of specialized tools, always choosing the right approach for each situation. Your goal is to provide helpful, accurate, and actionable responses whether the user needs a simple answer or complex analysis.
For general queries, be naturally helpful without overcomplicating. For file-related or command-based tasks, leverage your full analytical capabilities with appropriate tool integration. Always gauge the appropriate level of detail and technicality based on the user's needs.
You are the ultimate AI companion - equally comfortable discussing everyday topics, analyzing complex data, generating professional documents, or creating comprehensive strategies. Your strength lies in knowing when to use which capability and seamlessly integrating multiple tools when needed.
'''
class UnstructuredDocumentExtractor:
    """
    Production-grade universal document text extractor using Unstructured library.
    
    This extractor leverages the powerful Unstructured library for sophisticated
    document parsing, with comprehensive fallbacks for robustness.
    
    Supported formats (via Unstructured):
    - Documents: DOC, DOCX, ODT, PDF, RTF, TXT, LOG
    - Spreadsheets: XLS, XLSX, CSV, TSV
    - Presentations: PPT, PPTX
    - Web: HTML, HTM, XML
    - Email: EML, MSG
    - Images: PNG, JPG, JPEG, TIFF, BMP, HEIC (with OCR)
    - E-books: EPUB
    - Markup: MD, RST, ORG
    - Code: JS, PY, JAVA, CPP, CC, CXX, C, CS, PHP, RB, SWIFT, TS, GO
    - Data: JSON
    """
    
    def __init__(self, logger: Optional[logging.Logger] = None):
        """Initialize the extractor with optional logger."""
        self.logger = logger or logging.getLogger(__name__)
        
        # Log available components
        if UNSTRUCTURED_AVAILABLE:
            self.logger.info("Unstructured library is available")
        else:
            self.logger.warning("Unstructured library not available, using fallback methods")
            
    def extract_text(self, 
                    file_content: Union[bytes, str], 
                    filename: str,
                    encoding: Optional[str] = None,
                    strategy: str = "auto",
                    include_metadata: bool = False,
                    max_partition_length: int = 1500,
                    languages: Optional[List[str]] = None,
                    extract_tables: bool = True) -> str:
        """
        Extract text from any supported document type using Unstructured.
        
        Args:
            file_content: Raw file content as bytes or string
            filename: Original filename for type detection
            encoding: Optional encoding override
            strategy: Partitioning strategy ("auto", "fast", "hi_res", "ocr_only")
            include_metadata: Whether to include element metadata in output
            max_partition_length: Maximum length for text partitions
            languages: OCR languages (e.g., ["eng", "spa"])
            extract_tables: Whether to extract and format tables
            
        Returns:
            Extracted text as string
        """
        try:
            # Convert string to bytes if needed
            if isinstance(file_content, str):
                file_content = file_content.encode('utf-8')
            
            # Get file extension
            file_ext = os.path.splitext(filename)[1].lower()
            self.logger.info(f"Extracting text from {filename} (type: {file_ext})")
            
            # Try Unstructured first if available
            if UNSTRUCTURED_AVAILABLE:
                try:
                    text = self._extract_with_unstructured(
                        file_content, filename, encoding, strategy, 
                        include_metadata, max_partition_length, languages
                    )
                    if text and len(text.strip()) > 10:
                        return text
                except Exception as e:
                    self.logger.warning(f"Unstructured extraction failed: {str(e)}, trying fallbacks")
            
            # Use fallback methods
            return self._extract_with_fallback(file_content, filename, encoding)
            
        except Exception as e:
            self.logger.error(f"Error extracting text from {filename}: {str(e)}")
            self.logger.error(f"Traceback: {traceback.format_exc()}")
            # Emergency fallback
            return self._emergency_text_extraction(file_content)
    
    def _extract_with_unstructured(self, 
                                  file_content: bytes, 
                                  filename: str,
                                  encoding: Optional[str],
                                  strategy: str,
                                  include_metadata: bool,
                                  max_partition_length: int,
                                  languages: Optional[List[str]]) -> str:
        """Extract text using Unstructured library."""
        # Save to temporary file (Unstructured often works better with files)
        with tempfile.NamedTemporaryFile(
            delete=False, 
            suffix=os.path.splitext(filename)[1]
        ) as tmp_file:
            tmp_file.write(file_content)
            tmp_path = tmp_file.name
        
        try:
            file_ext = os.path.splitext(filename)[1].lower()
            # Prepare kwargs
            if encoding is None and file_ext in ['.csv', '.txt', '.log', '.md']:
                encoding = self._detect_encoding(file_content)
                self.logger.info(f"Auto-detected encoding: {encoding}")
            
            kwargs = {
                "filename": tmp_path,
                "encoding": encoding or 'utf-8',  # Use detected or default
                "max_partition": max_partition_length,
            }
            
            if file_ext in ['.pdf', '.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.heic']:
                kwargs["strategy"] = strategy
                if languages:
                    kwargs["languages"] = languages
            
            # Special handling for specific file types
            if file_ext in ['.html', '.htm']:
                kwargs["include_page_breaks"] = True
            elif file_ext in ['.pdf']:
                kwargs["include_page_breaks"] = True
                kwargs["infer_table_structure"] = True
            elif file_ext in ['.eml', '.msg']:
                kwargs["process_attachments"] = False  # Don't process attachments for text extraction
                kwargs["content_source"] = "text/html"  # Prefer HTML content
                
            # Use the main partition function
            elements = partition(**kwargs)
            
            # Convert elements to text
            text_parts = []
            
            for element in elements:
                # Get element text
                element_text = str(element)
                
                if include_metadata and hasattr(element, 'metadata'):
                    # Add metadata if requested
                    metadata = element.metadata
                    if hasattr(metadata, 'page_number') and metadata.page_number:
                        element_text = f"[Page {metadata.page_number}] {element_text}"
                    if hasattr(metadata, 'section') and metadata.section:
                        element_text = f"[{metadata.section}] {element_text}"
                
                # Handle tables specially
                if hasattr(element, 'metadata') and hasattr(element.metadata, 'text_as_html'):
                    # Convert HTML table to text representation
                    table_text = self._html_table_to_text(element.metadata.text_as_html)
                    if table_text:
                        text_parts.append(table_text)
                    else:
                        text_parts.append(element_text)
                else:
                    text_parts.append(element_text)
            
            return "\n\n".join(text_parts)
            
        finally:
            # Clean up temp file
            try:
                os.unlink(tmp_path)
            except:
                pass
    
    def _extract_with_fallback(self, file_content: bytes, filename: str, encoding: Optional[str]) -> str:
        """Fallback extraction methods when Unstructured is not available or fails."""
        file_ext = os.path.splitext(filename)[1].lower()
        
        # Try specific extractors based on file type
        if file_ext in ['.docx']:
            return self._extract_docx_fallback(file_content)
        elif file_ext in ['.pdf']:
            return self._extract_pdf_fallback(file_content)
        elif file_ext in ['.xlsx', '.xls']:
            return self._extract_excel_fallback(file_content, filename)
        elif file_ext in ['.csv']:
            return self._extract_csv_fallback(file_content, encoding)
        elif file_ext in ['.html', '.htm']:
            return self._extract_html_fallback(file_content, encoding)
        elif file_ext in ['.json']:
            return self._extract_json_fallback(file_content, encoding)
        elif file_ext in ['.xml']:
            return self._extract_xml_fallback(file_content, encoding)
        elif file_ext in ['.md', '.markdown']:
            return self._extract_markdown_fallback(file_content, encoding)
        elif file_ext in ['.pptx']:
            return self._extract_pptx_fallback(file_content)
        elif file_ext in ['.eml']:
            return self._extract_email_fallback(file_content, encoding)
        elif file_ext in ['.msg']:
            return self._extract_msg_fallback(file_content)
        else:
            # Default text extraction
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_docx_fallback(self, file_content: bytes) -> str:
        """Fallback DOCX extraction without Unstructured."""
        if DocxDocument:
            try:
                doc = DocxDocument(BytesIO(file_content))
                paragraphs = []
                
                # Extract paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text)
                
                # Extract tables
                for table in doc.tables:
                    table_text = []
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(" | ".join(row_text))
                    if table_text:
                        paragraphs.append("\n".join(table_text))
                
                return "\n\n".join(paragraphs)
            except Exception as e:
                self.logger.error(f"DOCX fallback failed: {str(e)}")
        
        # Try XML extraction
        return self._extract_docx_xml_fallback(file_content)
    
    def _extract_docx_xml_fallback(self, file_content: bytes) -> str:
        """Extract text from DOCX by parsing XML."""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            text_parts = []
            
            with zipfile.ZipFile(BytesIO(file_content)) as docx:
                # Extract main document
                if 'word/document.xml' in docx.namelist():
                    xml_content = docx.read('word/document.xml')
                    tree = ET.fromstring(xml_content)
                    
                    namespaces = {
                        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
                    }
                    
                    # Extract paragraphs
                    for para in tree.findall('.//w:p', namespaces):
                        para_text = []
                        for t in para.findall('.//w:t', namespaces):
                            if t.text:
                                para_text.append(t.text)
                        if para_text:
                            text_parts.append(''.join(para_text))
            
            return "\n".join(text_parts)
        except:
            return self._extract_text_with_encoding(file_content)
    
    def _extract_pdf_fallback(self, file_content: bytes) -> str:
        """Fallback PDF extraction without Unstructured."""
        try:
            text_parts = []   
            with pdfplumber.open(BytesIO(file_content)) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"[Page {i+1}]\n{page_text}")
                        
                        # Extract tables
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            table_text = self._format_table(table)
                            text_parts.append(table_text)
                
            return "\n\n".join(text_parts)
        except:
            pass
        
        return self._extract_text_with_encoding(file_content)
    
    def _extract_excel_fallback(self, file_content: bytes, filename: str) -> str:
        """Fallback Excel extraction without Unstructured."""
        try:
            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as tmp:
                tmp.write(file_content)
                tmp_path = tmp.name
            
            try:
                excel_file = pd.ExcelFile(tmp_path)
                text_parts = []
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(tmp_path, sheet_name=sheet_name)
                    text_parts.append(f"\n[Sheet: {sheet_name}]\n")
                    
                    if not df.empty:
                        # Convert to readable format
                        text_parts.append(df.to_string())
                
                return "\n".join(text_parts)
            finally:
                os.unlink(tmp_path)
        except:
            return self._extract_csv_fallback(file_content, None)
    
    def _extract_csv_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback CSV extraction without Unstructured."""
        try:
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            text_content = file_content.decode(encoding, errors='replace')
            
            try:
                df = pd.read_csv(StringIO(text_content))
                return df.to_string()
            except:
                return text_content
        except:
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_html_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback HTML extraction without Unstructured."""
        try:
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            text_content = file_content.decode(encoding, errors='replace')
            
            if html2text:
                h = html2text.HTML2Text()
                h.ignore_links = False
                h.body_width = 0  # Don't wrap lines
                return h.handle(text_content)
            else:
                # Basic HTML stripping
                text = re.sub(r'<script[^>]*>.*?</script>', '', text_content, flags=re.DOTALL)
                text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\s+', ' ', text)
                return text.strip()
        except:
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_json_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback JSON extraction without Unstructured."""
        try:
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            text_content = file_content.decode(encoding, errors='replace')
            data = json.loads(text_content)
            
            # Pretty print JSON
            return json.dumps(data, indent=2, ensure_ascii=False)
        except:
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_xml_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback XML extraction without Unstructured."""
        try:
            import xml.etree.ElementTree as ET
            
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            text_content = file_content.decode(encoding, errors='replace')
            root = ET.fromstring(text_content)
            
            # Extract all text from XML
            texts = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    texts.append(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    texts.append(elem.tail.strip())
            
            return "\n".join(texts)
        except:
            # Just strip tags
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            text = file_content.decode(encoding, errors='replace')
            text = re.sub(r'<[^>]+>', ' ', text)
            return re.sub(r'\s+', ' ', text).strip()
    
    def _extract_markdown_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback Markdown extraction without Unstructured."""
        try:
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            text_content = file_content.decode(encoding, errors='replace')
            
            if markdown and html2text:
                # Convert to HTML then to plain text
                html_content = markdown.markdown(text_content)
                h = html2text.HTML2Text()
                h.body_width = 0
                return h.handle(html_content)
            
            return text_content
        except:
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_pptx_fallback(self, file_content: bytes) -> str:
        """Fallback PPTX extraction without Unstructured."""
        if Presentation:
            try:
                prs = Presentation(BytesIO(file_content))
                text_parts = []
                
                for i, slide in enumerate(prs.slides):
                    text_parts.append(f"\n[Slide {i + 1}]\n")
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text_parts.append(shape.text)
                        
                        if shape.has_table:
                            table_text = []
                            for row in shape.table.rows:
                                row_text = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_text.append(cell.text.strip())
                                if row_text:
                                    table_text.append(" | ".join(row_text))
                            if table_text:
                                text_parts.append("\n".join(table_text))
                
                return "\n".join(text_parts)
            except:
                pass
        
        return self._extract_text_with_encoding(file_content)
    
    def _extract_email_fallback(self, file_content: bytes, encoding: Optional[str]) -> str:
        """Fallback email extraction without Unstructured."""
        try:
            import email
            from email import policy
            
            if encoding is None:
                encoding = self._detect_encoding(file_content)
            
            # Parse email
            msg = email.message_from_bytes(file_content, policy=policy.default)
            
            text_parts = []
            
            # Add headers
            headers = ['From', 'To', 'Subject', 'Date']
            for header in headers:
                value = msg.get(header)
                if value:
                    text_parts.append(f"{header}: {value}")
            
            text_parts.append("")  # Empty line
            
            # Extract body
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            text_parts.append(payload.decode(encoding, errors='replace'))
                    elif part.get_content_type() == "text/html" and not any("text/plain" in p.get_content_type() for p in msg.walk()):
                        payload = part.get_payload(decode=True)
                        if payload:
                            html_text = payload.decode(encoding, errors='replace')
                            # Convert HTML to text
                            if html2text:
                                h = html2text.HTML2Text()
                                text_parts.append(h.handle(html_text))
                            else:
                                # Strip HTML tags
                                text = re.sub(r'<[^>]+>', ' ', html_text)
                                text_parts.append(text)
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    text_parts.append(payload.decode(encoding, errors='replace'))
            
            return "\n".join(text_parts)
        except:
            return self._extract_text_with_encoding(file_content, encoding)
    
    def _extract_msg_fallback(self, file_content: bytes) -> str:
        """Fallback MSG extraction without Unstructured."""
        # MSG files are complex binary format, just extract readable text
        return self._extract_text_with_encoding(file_content)
    
    def _html_table_to_text(self, html_table: str) -> str:
        """Convert HTML table to readable text format."""
        try:
            # Simple HTML table parsing
            rows = re.findall(r'<tr[^>]*>(.*?)</tr>', html_table, re.DOTALL)
            table_text = []
            
            for row in rows:
                cells = re.findall(r'<t[hd][^>]*>(.*?)</t[hd]>', row, re.DOTALL)
                if cells:
                    # Clean cell content
                    clean_cells = []
                    for cell in cells:
                        cell_text = re.sub(r'<[^>]+>', '', cell)
                        cell_text = cell_text.strip()
                        clean_cells.append(cell_text)
                    
                    table_text.append(" | ".join(clean_cells))
            
            return "\n".join(table_text)
        except:
            return ""
    
    def _format_table(self, table_data: List[List[Any]]) -> str:
        """Format table data as readable text."""
        if not table_data:
            return ""
        
        formatted_rows = []
        for row in table_data:
            if row:
                formatted_row = " | ".join(str(cell) if cell is not None else "" for cell in row)
                formatted_rows.append(formatted_row)
        
        return "\n".join(formatted_rows)
    
    def _extract_text_with_encoding(self, file_content: bytes, encoding: Optional[str] = None) -> str:
        """Extract text with automatic encoding detection."""
        if encoding is None:
            encoding = self._detect_encoding(file_content)
        
        try:
            return file_content.decode(encoding, errors='replace')
        except:
            # Try common encodings
            for enc in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1', 'utf-16']:
                try:
                    return file_content.decode(enc, errors='replace')
                except:
                    continue
            
            # Force UTF-8
            return file_content.decode('utf-8', errors='replace')
    
    def _detect_encoding(self, file_content: bytes) -> str:
        """Detect file encoding using chardet."""
        try:
            # Sample first 100KB for performance
            sample = file_content[:100000] if len(file_content) > 100000 else file_content
            result = chardet.detect(sample)
            
            encoding = result.get('encoding')
            confidence = result.get('confidence', 0)
            
            if encoding and confidence > 0.7:
                self.logger.info(f"Detected encoding: {encoding} (confidence: {confidence:.2f})")
                return encoding
        except:
            pass
        
        return 'utf-8'
    
    def _emergency_text_extraction(self, file_content: bytes) -> str:
        """Emergency fallback to extract any readable text."""
        try:
            # First, check if it's a known binary format that we shouldn't try to decode
            file_header = file_content[:10] if len(file_content) >= 10 else file_content
            
            # Check for PDF header
            if file_header.startswith(b'%PDF'):
                self.logger.error("Emergency extraction called on PDF file - cannot extract without proper tools")
                return "PDF file detected but text extraction failed. Please ensure PDF extraction libraries are installed."
                
            # Check for other binary formats
            binary_headers = [
                b'\x50\x4b\x03\x04',  # ZIP/DOCX/XLSX
                b'\xd0\xcf\x11\xe0',  # DOC/XLS
                b'\x89\x50\x4e\x47',  # PNG
                b'\xff\xd8\xff',      # JPEG
            ]
            
            for header in binary_headers:
                if file_header.startswith(header):
                    self.logger.error("Emergency extraction called on binary file - cannot extract")
                    return "Binary file detected but text extraction failed. Please ensure document processing libraries are installed."
            
            # Try to extract printable ASCII and common Unicode
            text_parts = []
            i = 0
            
            while i < len(file_content):
                # Try to decode as UTF-8
                for length in range(4, 0, -1):
                    if i + length <= len(file_content):
                        try:
                            char = file_content[i:i+length].decode('utf-8')
                            if char.isprintable() or char in '\n\r\t':
                                text_parts.append(char)
                                i += length
                                break
                        except:
                            pass
                else:
                    # Skip non-decodable byte
                    i += 1
            
            text = ''.join(text_parts)
            
            # Clean up
            text = re.sub(r'\s+', ' ', text)
            text = re.sub(r' +', ' ', text)
            
            result = text.strip()
            
            # If we got very little text from a large file, it's probably binary
            if len(result) < 100 and len(file_content) > 1000:
                return "Unable to extract meaningful text from this document. The file may be corrupted or in an unsupported format."
                
            return result if result else "Unable to extract text from this document."
            
        except Exception as e:
            self.logger.error(f"Emergency extraction failed: {str(e)}")
            return "Unable to extract text from this document."
async def extract_text_internal(
    file_content: bytes,
    filename: str,
    strategy: str = "auto",
    languages: Optional[List[str]] = None,
    encoding: Optional[str] = None,
    logger: Optional[logging.Logger] = None
) -> str:
    max_chars = 50000
    """
    Internal function to extract text from document content.
    This can be called from other API endpoints.
    
    Args:
        file_content: Raw file content as bytes
        filename: Original filename for type detection
        strategy: Extraction strategy
        languages: OCR languages list
        encoding: Optional encoding override
        logger: Logger instance
        
    Returns:
        Extracted text as string
        
    Raises:
        Exception: If extraction fails
    """
    if logger is None:
        logger = logging.getLogger(__name__)
    
    try:
        extractor = UnstructuredDocumentExtractor(logger=logger)
        extracted_text = extractor.extract_text(
            file_content=file_content,
            filename=filename,
            encoding=encoding,
            strategy=strategy,
            languages=languages
        )
        if max_chars and len(extracted_text) > max_chars:
            logger.info(f"Truncating extracted text from {len(extracted_text)} to {max_chars} characters")
            # Try to truncate at a sentence boundary
            truncated = extracted_text[:max_chars]
            last_period = truncated.rfind('.')
            last_newline = truncated.rfind('\n')
            cutoff = max(last_period, last_newline)
            if cutoff > max_chars * 0.8:  # Only use if it's not too far back
                extracted_text = truncated[:cutoff + 1]
            else:
                extracted_text = truncated + "..."
                
        return extracted_text
    except Exception as e:
        logger.error(f"Internal text extraction failed: {str(e)}")
        raise

def sync_wait_for_run_completion(client: AzureOpenAI, thread_id: str, max_wait_time: int = 30) -> bool:
    """
    Synchronous version: Wait for any active runs on a thread to complete before proceeding.
    
    Args:
        client: Azure OpenAI client
        thread_id: Thread ID to check
        max_wait_time: Maximum seconds to wait
        
    Returns:
        True if thread is ready, False if timeout
    """
    start_time = time.time()
    check_interval = 1  # Check every second
    
    while time.time() - start_time < max_wait_time:
        try:
            runs = client.beta.threads.runs.list(thread_id=thread_id, limit=1)
            if not runs.data:
                return True  # No runs, safe to proceed
                
            latest_run = runs.data[0]
            if latest_run.status not in ["in_progress", "queued", "requires_action"]:
                return True  # Run is complete
                
            logging.info(f"Waiting for run {latest_run.id} to complete (status: {latest_run.status})")
            time.sleep(check_interval)
            
        except Exception as e:
            logging.warning(f"Error checking run status: {e}")
            # On error, wait a bit and return True to proceed
            time.sleep(1)
            return True
    
    logging.warning(f"Timeout waiting for run completion on thread {thread_id}")
    return False


async def wait_for_run_completion(client: AzureOpenAI, thread_id: str, max_wait_time: int = 30) -> bool:
    """
    Async version: Wait for any active runs on a thread to complete before proceeding.
    
    Args:
        client: Azure OpenAI client
        thread_id: Thread ID to check
        max_wait_time: Maximum seconds to wait
        
    Returns:
        True if thread is ready, False if timeout
    """
    start_time = time.time()
    check_interval = 1  # Check every second
    
    while time.time() - start_time < max_wait_time:
        try:
            runs = client.beta.threads.runs.list(thread_id=thread_id, limit=1)
            if not runs.data:
                return True  # No runs, safe to proceed
                
            latest_run = runs.data[0]
            if latest_run.status not in ["in_progress", "queued", "requires_action"]:
                return True  # Run is complete
                
            logging.info(f"Waiting for run {latest_run.id} to complete (status: {latest_run.status})")
            await asyncio.sleep(check_interval)
            
        except Exception as e:
            logging.warning(f"Error checking run status: {e}")
            # On error, wait a bit and return True to proceed
            await asyncio.sleep(1)
            return True
    
    logging.warning(f"Timeout waiting for run completion on thread {thread_id}")
    return False
def get_content_generation_tools():
    """Get tool definitions for content generation and data extraction"""
    generate_content_tool = {
        "type": "function",
        "function": {
            "name": "generate_content",
            "description": "Generate content in various formats when user uses /generate command. Can create documents, datasets, articles, or any text content.",
            "parameters": {
                "type": "object",
                "properties": {
                    "prompt": {
                        "type": "string",
                        "description": "The complete generation instructions including context from recent messages"
                    },
                    "output_format": {
                        "type": "string",
                        "enum": ["text", "csv", "excel", "docx", "auto"],
                        "description": "Desired output format. Use 'auto' to let system decide based on content"
                    }
                },
                "required": ["prompt", "output_format"]
            }
        }
    }
    
    extract_data_tool = {
        "type": "function",
        "function": {
            "name": "extract_data",
            "description": "Extract structured data from text or generate synthetic datasets when user uses /extract or /analyze command. Works with text from messages, not file uploads.",
            "parameters": {
                "type": "object",
                "properties": {
                    "prompt": {
                        "type": "string",
                        "description": "Complete extraction/generation instructions with relevant text content to process"
                    },
                    "mode": {
                        "type": "string",
                        "enum": ["extract", "generate", "auto"],
                        "description": "Operation mode: extract from text, generate synthetic data, or auto-detect"
                    },
                    "output_format": {
                        "type": "string",
                        "enum": ["csv", "excel", "json"],
                        "description": "Output format for the extracted/generated data"
                    },
                    "raw_text": {
                        "type": "string",
                        "description": "The actual text content to extract data from (if mode is extract)"
                    }
                },
                "required": ["prompt", "mode", "output_format"]
            }
        }
    }
    
    return [generate_content_tool, extract_data_tool]

# Create downloads directory if it doesn't exist
os.makedirs(DOWNLOADS_DIR, exist_ok=True)
async def get_conversation_context(client, thread_id: str, limit: int = 3) -> str:
    """
    Get recent conversation context from thread messages.
    
    Args:
        client: Azure OpenAI client
        thread_id: Thread ID to get messages from
        limit: Number of recent messages to retrieve
        
    Returns:
        Formatted context string
    """
    try:
        messages = client.beta.threads.messages.list(
            thread_id=thread_id,
            order="desc",
            limit=limit + 1  # +1 to skip current message
        )
        
        context_parts = []
        for i, msg in enumerate(messages.data):
            if i == 0:  # Skip the most recent (current) message
                continue
                
            # Skip system/metadata messages
            if hasattr(msg, 'metadata') and msg.metadata:
                msg_type = msg.metadata.get('type', '')
                if msg_type in ['user_persona_context', 'file_awareness', 'pandas_agent_files', 'pandas_agent_instruction']:
                    continue
            
            # Extract text content
            content = ""
            for part in msg.content:
                if part.type == 'text':
                    content += part.text.value
            
            if content:
                context_parts.append(f"{msg.role}: {content[:500]}")  # Limit each message
        
        # Reverse to get chronological order
        context_parts.reverse()
        return "\n".join(context_parts) if context_parts else ""
        
    except Exception as e:
        logging.error(f"Error getting conversation context: {e}")
        return ""


async def enhance_prompt_with_context(prompt: str, thread_id: str, client, output_format: str = None) -> str:
    """
    Enhance a prompt with conversation context and format-specific instructions.
    
    Args:
        prompt: Original user prompt
        thread_id: Thread ID for context retrieval
        client: Azure OpenAI client
        output_format: Target output format
        
    Returns:
        Enhanced prompt string
    """
    # Get conversation context
    context = await get_conversation_context(client, thread_id, limit=3)
    
    enhanced_prompt = ""
    if context:
        enhanced_prompt = f"Context from recent conversation:\n{context}\n\n"
    
    enhanced_prompt += f"Request: {prompt}"
    
    # Add format-specific enhancements
    if output_format == "excel" or output_format == "csv":
        enhanced_prompt += "\n\nPlease ensure the data is well-structured with clear column headers and consistent formatting."
    elif output_format == "docx":
        enhanced_prompt += "\n\nPlease create a comprehensive, well-formatted document with proper sections, headings, and professional structure."
    elif output_format == "auto" or output_format is None:
        enhanced_prompt += "\n\nAnalyze this request and choose the most appropriate format for the response."
    
    return enhanced_prompt

async def trim_thread(client: AzureOpenAI, thread_id: str, keep_messages: int = 30):
    """
    Trim thread to keep only the most recent messages.
    Returns True if trimming was performed, False otherwise.
    """
    try:
        # Get all messages
        all_messages = []
        has_more = True
        after = None
        
        while has_more:
            messages = client.beta.threads.messages.list(
                thread_id=thread_id,
                order="desc",
                limit=100,
                after=after
            )
            all_messages.extend(messages.data)
            has_more = messages.has_more
            if has_more and messages.data:
                after = messages.data[-1].id
        
        # Skip if thread is small
        if len(all_messages) <= keep_messages:
            return False
        
        # Delete old messages (keep the most recent ones)
        messages_to_delete = all_messages[keep_messages:]
        deleted_count = 0
        
        for msg in messages_to_delete:
            try:
                # Skip system messages
                if hasattr(msg, 'metadata') and msg.metadata:
                    msg_type = msg.metadata.get('type', '')
                    if msg_type in ['user_persona_context', 'file_awareness', 'pandas_agent_files', 'pandas_agent_instruction']:
                        continue
                        
                client.beta.threads.messages.delete(
                    thread_id=thread_id,
                    message_id=msg.id
                )
                deleted_count += 1
                logging.info(f"Deleted old message {msg.id} from thread {thread_id}")
            except Exception as e:
                logging.warning(f"Could not delete message {msg.id}: {e}")
        
        logging.info(f"Trimmed thread {thread_id}: deleted {deleted_count} messages")
        return deleted_count > 0
                
    except Exception as e:
        logging.error(f"Error trimming thread {thread_id}: {e}")
        return False

async def handle_generate_content(tool_args: dict, thread_id: str, client, request) -> str:
    """
    Handle generate_content tool calls by calling the /completion endpoint.
    Implements comprehensive fallback strategy and returns JSON with formatted message.
    
    Args:
        tool_args: Parsed tool arguments containing prompt and output_format
        thread_id: Thread ID for context
        client: Azure OpenAI client
        request: FastAPI request object for URL construction
        
    Returns:
        JSON string containing full response data with formatted message
    """
    try:
        prompt = tool_args.get("prompt", "")
        requested_format = tool_args.get("output_format", "auto")
        
        # Enhance prompt with context
        enhanced_prompt = await enhance_prompt_with_context(prompt, thread_id, client, requested_format)
        
        # Define fallback chain based on requested format
        format_chain = []
        if requested_format == "excel":
            format_chain = ["excel", "csv", "text"]
        elif requested_format == "csv":
            format_chain = ["csv", "text"]
        elif requested_format == "docx":
            format_chain = ["docx", "text"]
        elif requested_format == "auto":
            format_chain = [None, "text"]  # None triggers auto-detection
        else:
            format_chain = ["text"]
        
        last_error = None
        
        # Try each format in the fallback chain
        for attempt_format in format_chain:
            try:
                # Call the completion endpoint
                result = await chat_completion(
                    request=request,
                    prompt=enhanced_prompt,
                    model="gpt-4.1-mini",
                    temperature=0.8,
                    max_tokens=16000,
                    system_message=None,
                    output_format=attempt_format,
                    files=None,
                    max_retries=3,
                    rows_to_generate=50  # Override default of 30
                )
                
                # Parse response
                response_data = json.loads(result.body.decode())
                
                if response_data.get("status") == "success":
                    # Build formatted message based on response type
                    if response_data.get("download_url"):
                        # File was generated
                        message = "✅ Generated successfully!"
                        
                        # Note if we fell back to a different format
                        if attempt_format != requested_format and requested_format != "auto":
                            actual_format = response_data.get("output_format", attempt_format)
                            message += f" (Created as {actual_format} format)"
                        
                        # Add download link with proper markdown escaping
                        filename = response_data.get("filename", "generated_file")
                        download_url = response_data.get("download_url")
                        # Ensure URL doesn't break markdown
                        if download_url and not download_url.startswith(('http://', 'https://')):
                            download_url = f"/{download_url}" if not download_url.startswith('/') else download_url
                        message += f"\n\n📄 **Download:** [{filename}]({download_url})"
                        
                        # Add summary if available
                        if response_data.get("message"):
                            message += f"\n\n{response_data['message']}"
                        elif response_data.get("summary"):
                            summary = response_data["summary"]
                            if isinstance(summary, dict):
                                rows = summary.get("rows", 0)
                                cols = summary.get("columns", [])
                                if rows and cols:
                                    message += f"\n\nGenerated {rows} rows with {len(cols)} columns"
                        
                        # Add formatted message to response
                        response_data["formatted_message"] = message
                        
                    else:
                        # Text response
                        content = response_data.get("response", "")
                        if content:
                            message = "Here's your generated content:\n\n"
                            
                            # Truncate if very long for display
                            if len(content) > 20000:
                                message += content[:20000] + "..."
                                message += "\n\n*Showing first 20000 characters. Click download to see full content.*"
                            else:
                                message += content
                            
                            message += "\n\n💾 **Save option:** Use the download button to save this response."
                            response_data["formatted_message"] = message
                        else:
                            # No content in response
                            response_data["formatted_message"] = "✅ Generated successfully!"
                    
                    # Return full response as JSON
                    return json.dumps(response_data)
                
            except Exception as e:
                last_error = str(e)
                logging.error(f"Generation failed with format {attempt_format}: {e}")
                continue
        
        # All attempts failed - return error with helpful message
        error_response = {
            "status": "error",
            "message": "Unable to generate content after trying all formats",
            "error": str(last_error) if last_error else "Unknown error",
            "formatted_message": (
                f"I understand you want me to generate content based on: '{prompt[:100]}{'...' if len(prompt) > 100 else ''}' "
                f"but I'm having technical difficulties. Please try again in a moment, "
                f"or try a simpler request."
            )
        }
        return json.dumps(error_response)
        
    except Exception as e:
        logging.error(f"Critical error in handle_generate_content: {e}\n{traceback.format_exc()}")
        # Return error as JSON
        error_response = {
            "status": "error",
            "message": "An error occurred while generating content",
            "error": str(e),
            "formatted_message": "An error occurred while generating content. Please try again or contact support if the issue persists."
        }
        return json.dumps(error_response)

async def handle_extract_data(tool_args: dict, thread_id: str, client, request) -> str:
    """
    Handle extract_data tool calls by calling the /extract-reviews endpoint.
    
    Args:
        tool_args: Parsed tool arguments
        thread_id: Thread ID for context
        client: Azure OpenAI client
        request: FastAPI request object
        
    Returns:
        Formatted response string
    """
    try:
        prompt = tool_args.get("prompt", "")
        mode = tool_args.get("mode", "auto")
        output_format = tool_args.get("output_format", "excel")
        raw_text = tool_args.get("raw_text", "")
        
        # If no raw_text provided and mode is extract, gather from conversation
        if not raw_text and mode == "extract":
            # Get recent messages to find data to extract
            context = await get_conversation_context(client, thread_id, limit=5)
            if context:
                # Look for data patterns in context
                raw_text = context
        
        # Enhance prompt based on mode
        enhanced_prompt = prompt
        if mode == "generate":
            enhanced_prompt = f"{prompt}\n\nGenerate 100 rows of synthetic data with appropriate columns."
        elif mode == "extract" and raw_text:
            enhanced_prompt = f"""{prompt}\n\n
CRITICAL INSTRUCTIONS:
1. Extract the ACTUAL DATA RECORDS, not summaries or metadata about data
2. If you see descriptions like "Dataset contains 50 records", find those 50 actual records
3. Look for patterns like:
   - Individual reviews with names, ratings, and comments
   - Line items with specific values
   - Repeated structured entries
   - NOT analytical summaries or statistical descriptions
4. If only metadata/analysis is present without actual data, clearly indicate this."""
        
        # Define fallback chain
        format_chain = []
        if output_format == "excel":
            format_chain = ["excel", "csv", "json"]
        elif output_format == "csv":
            format_chain = ["csv", "json"]
        else:
            format_chain = [output_format, "json"]
        
        last_error = None
        
        # Try each format in the chain
        for attempt_format in format_chain:
            try:
                # Call extract-reviews endpoint
                result = await extract_reviews(
                    request=request,
                    file=None,  # No file upload
                    columns="auto",
                    prompt=enhanced_prompt,
                    model="gpt-4.1-mini",
                    temperature=0.1,
                    output_format=attempt_format,
                    max_text_length=100000,
                    max_retries=3,
                    fallback_to_json=True,
                    mode=mode,
                    rows_to_generate=100,  # Override default
                    raw_text=raw_text if raw_text else None
                )
                
                # Parse response
                response_data = json.loads(result.body.decode())
                
                if response_data.get("status") == "success":
                    # Format successful response
                    if response_data.get("download_url"):
                        # File was generated
                        operation = "extracted" if mode == "extract" else "generated"
                        message = f"✅ Successfully {operation} data!"
                        
                        # Note format change if applicable
                        if attempt_format != output_format:
                            message += f" (Saved as {attempt_format} format)"
                        
                        # Add download link
                        filename = response_data.get("filename", "data_file")
                        download_url = response_data.get("download_url")
                        message += f"\n\n📄 **Download:** [{filename}]({download_url})"
                        
                        # Add data summary
                        row_count = response_data.get("row_count", 0)
                        columns = response_data.get("columns", [])
                        if row_count and columns:
                            message += f"\n\n**Data Summary:**"
                            message += f"\n- Rows: {row_count}"
                            message += f"\n- Columns: {', '.join(columns[:10])}"
                            if len(columns) > 10:
                                message += f" (and {len(columns) - 10} more)"
                        
                        # Add metadata if available
                        metadata = response_data.get("metadata", {})
                        if metadata.get("extraction_confidence"):
                            message += f"\n- Confidence: {metadata['extraction_confidence']}"
                        
                        return message
                        
                    elif response_data.get("data"):
                        # JSON response with data
                        data = response_data.get("data", [])
                        columns = response_data.get("columns", [])
                        
                        message = f"✅ Successfully processed data!\n\n"
                        message += f"**Found {len(data)} rows with {len(columns)} columns**\n\n"
                        
                        # Show sample data
                        if data:
                            message += "**Sample data (first 10 rows):**\n```\n"
                            # Create simple table view
                            message += " | ".join(columns) + "\n"
                            message += "-" * (len(" | ".join(columns))) + "\n"
                            for row in data[:10]:
                                message += " | ".join(str(cell)[:20] for cell in row) + "\n"
                            message += "```\n"
                            
                            if len(data) > 3:
                                message += f"\n*Showing 10 of {len(data)} total rows*"
                        
                        message += "\n\n💾 **To save:** Use the download button or try the command again with `/extract` for Excel format."
                        return message
                
            except Exception as e:
                last_error = str(e)
                logging.error(f"Extraction failed with format {attempt_format}: {e}")
                continue
        
        # All attempts failed
        if mode == "extract" and not raw_text:
            return (
                "I couldn't find any data to extract from our conversation. "
                "Please share the data you'd like me to extract, or use "
                "`/analyze create [description]` to generate synthetic data instead."
            )
        else:
            return (
                "I encountered an issue processing your data request. "
                "Please try rephrasing your request or using a simpler format."
            )
        
    except Exception as e:
        logging.error(f"Critical error in handle_extract_data: {e}\n{traceback.format_exc()}")
        return ""

# Mount static files directory for serving downloads
#app.mount("/download-files", StaticFiles(directory=DOWNLOADS_DIR), name="download-files")
def get_downloads_directory():
    """Get the appropriate downloads directory based on the environment."""
    if os.getenv("AZURE_APP_SERVICE", "").lower() == "true":
        # Azure App Service - use home directory which is persistent
        home_dir = os.getenv("HOME", "/tmp")
        downloads_dir = os.path.join(home_dir, "chat_downloads")
    elif platform.system() == "Windows":
        # Windows development
        downloads_dir = os.path.join(tempfile.gettempdir(), "chat_downloads")
    else:
        # Linux/Mac - use /tmp for development
        downloads_dir = "/tmp/chat_downloads"
    
    # Ensure directory exists with proper permissions
    try:
        os.makedirs(downloads_dir, mode=0o755, exist_ok=True)
        
        # Test write permissions
        test_file = os.path.join(downloads_dir, ".write_test")
        with open(test_file, 'w') as f:
            f.write("test")
        os.remove(test_file)
        
        logging.info(f"Downloads directory initialized: {downloads_dir}")
        return downloads_dir
    except Exception as e:
        logging.error(f"Failed to initialize downloads directory {downloads_dir}: {e}")
        # Fallback to temp directory
        fallback_dir = os.path.join(tempfile.gettempdir(), "chat_downloads")
        os.makedirs(fallback_dir, mode=0o755, exist_ok=True)
        logging.info(f"Using fallback directory: {fallback_dir}")
        return fallback_dir

# Update the global variable
DOWNLOADS_DIR = get_downloads_directory()
_client_instance = None
_http_client = None
_client_lock = threading.Lock()

def create_client() -> AzureOpenAI:
    """
    Creates an AzureOpenAI client instance.
    
    Attempts to create an optimized client with connection pooling.
    Falls back to simple client creation if optimization fails.
    
    Returns:
        AzureOpenAI client instance
    """
    global _client_instance, _http_client
    
    # Fast path: return existing instance if available
    if _client_instance is not None:
        try:
            # Quick health check - try to access a property
            _ = _client_instance.api_key
            return _client_instance
        except Exception:
            # Client might be in bad state, recreate it
            logging.warning("Existing client instance appears invalid, recreating...")
            _client_instance = None
            if _http_client:
                try:
                    _http_client.close()
                except:
                    pass
                _http_client = None
    
    # Thread-safe client creation
    with _client_lock:
        # Double-check pattern
        if _client_instance is not None:
            return _client_instance
        
        # Check if h2 module is actually available
        h2_available = False
        if HTTPX_AVAILABLE:
            try:
                import h2
                h2_available = True
            except ImportError:
                logging.warning("h2 module not available - HTTP/2 will be disabled")
        
        # Try optimized client first (only if httpx and h2 are available)
        if HTTPX_AVAILABLE and h2_available:
            try:
                logging.info("Attempting to create optimized client with connection pooling and HTTP/2...")
                
                # Create optimized HTTP client with HTTP/2
                _http_client = httpx.Client(
                    limits=httpx.Limits(
                        max_connections=10,
                        max_keepalive_connections=5,
                        keepalive_expiry=30.0,
                    ),
                    timeout=httpx.Timeout(
                        timeout=120.0,
                        connect=10.0,
                        read=60.0,
                        write=30.0,
                        pool=5.0
                    ),
                    transport=httpx.HTTPTransport(
                        retries=3,
                        http2=True,
                    ),
                    follow_redirects=True,
                )
                
                # Create optimized client
                _client_instance = AzureOpenAI(
                    azure_endpoint=AZURE_ENDPOINT,
                    api_key=AZURE_API_KEY,
                    api_version=AZURE_API_VERSION,
                    http_client=_http_client,
                    max_retries=3,
                )
                
                # Test the client with a minimal operation
                test_client = _client_instance.with_options(timeout=5.0)
                
                logging.info("Successfully created optimized AzureOpenAI client with HTTP/2")
                return _client_instance
                
            except Exception as e:
                logging.error(f"Failed to create optimized client with HTTP/2: {type(e).__name__}: {e}")
                
                # Clean up failed attempts
                if _http_client:
                    try:
                        _http_client.close()
                    except:
                        pass
                    _http_client = None
                _client_instance = None
        
        # Try with HTTP/1.1 if httpx is available but h2 is not
        if HTTPX_AVAILABLE and not h2_available:
            try:
                logging.info("Creating optimized client with HTTP/1.1 (h2 not available)...")
                
                # Create HTTP client WITHOUT HTTP/2
                _http_client = httpx.Client(
                    limits=httpx.Limits(
                        max_connections=10,
                        max_keepalive_connections=5,
                        keepalive_expiry=30.0,
                    ),
                    timeout=httpx.Timeout(
                        timeout=120.0,
                        connect=10.0,
                        read=60.0,
                        write=30.0,
                        pool=5.0
                    ),
                    transport=httpx.HTTPTransport(
                        retries=3,
                        http2=False,  # Explicitly disable HTTP/2
                    ),
                    follow_redirects=True,
                )
                
                _client_instance = AzureOpenAI(
                    azure_endpoint=AZURE_ENDPOINT,
                    api_key=AZURE_API_KEY,
                    api_version=AZURE_API_VERSION,
                    http_client=_http_client,
                    max_retries=3,
                )
                
                logging.info("Successfully created optimized AzureOpenAI client with HTTP/1.1")
                return _client_instance
                
            except Exception as e:
                logging.error(f"Failed to create HTTP/1.1 client: {type(e).__name__}: {e}")
                
                # Clean up
                if _http_client:
                    try:
                        _http_client.close()
                    except:
                        pass
                    _http_client = None
                _client_instance = None
        
        # Final fallback: Use simple client (no custom http_client)
        try:
            logging.info("Creating simple AzureOpenAI client...")
            
            _client_instance = AzureOpenAI(
                azure_endpoint=AZURE_ENDPOINT,
                api_key=AZURE_API_KEY,
                api_version=AZURE_API_VERSION,
            )
            
            logging.info("Successfully created simple AzureOpenAI client")
            return _client_instance
            
        except Exception as e:
            logging.critical(f"Failed to create even simple client: {type(e).__name__}: {e}")
            _client_instance = None
            raise
def construct_download_url(request: Request, filename: str) -> str:
    """
    Construct the download URL for a file.
    
    Args:
        request: FastAPI request object
        filename: Name of the file
        
    Returns:
        Full download URL
    """
    # Get the base URL from the request
    base_url = str(request.base_url).rstrip('/')
    
    # For Azure App Service, use the proper host
    host = request.headers.get('host', '')
    if 'azurewebsites.net' in host:
        # Use HTTPS for Azure
        base_url = f"https://{host}"
    elif 'localhost' in host or '127.0.0.1' in host:
        # Use HTTP for local development
        base_url = f"http://{host}"
    
    return f"{base_url}/download-files/{filename}"
def save_download_file(content: bytes, filename: str) -> str:
    """
    Save a file for download with proper permissions.
    
    Args:
        content: File content as bytes
        filename: Desired filename
        
    Returns:
        Actual filename used (may be modified for uniqueness)
    """
    # Sanitize filename
    safe_filename = secure_filename(filename)
    filepath = os.path.join(DOWNLOADS_DIR, safe_filename)
    
    try:
        # Write file with explicit permissions
        with open(filepath, 'wb') as f:
            f.write(content)
        
        # Set file permissions to be readable by the web server
        os.chmod(filepath, 0o644)
        
        logging.info(f"Saved download file: {filepath} ({len(content)} bytes)")
        return safe_filename  # Return the actual filename used
        
    except Exception as e:
        logging.error(f"Failed to save download file {filename}: {e}")
        raise

def secure_filename(filename: str) -> str:
    """
    Sanitize a filename to be safe for filesystem storage.
    Only adds timestamp if filename doesn't already have one.
    
    Args:
        filename: Original filename
        
    Returns:
        Sanitized filename
    """
    import re
    
    # Remove any path components
    filename = os.path.basename(filename)
    
    # Replace unsafe characters
    filename = re.sub(r'[^\w\s.-]', '_', filename)
    filename = re.sub(r'[\s]+', '_', filename)
    
    # Ensure it has a proper extension
    if '.' not in filename:
        filename += '.bin'
    
    # Check if filename already has a timestamp pattern (YYYYMMDD_HHMMSS)
    name, ext = os.path.splitext(filename)
    timestamp_pattern = r'_\d{8}_\d{6}$'
    
    if re.search(timestamp_pattern, name):
        # Filename already has timestamp, don't add another
        return filename
    else:
        # Add timestamp to ensure uniqueness
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        return f"{name}_{timestamp}{ext}"
def extract_csv_from_content(content: str) -> str:
    """
    Extract CSV content from AI response, handling various formats.
    
    Args:
        content: The AI response that may contain CSV data
        
    Returns:
        Clean CSV content
    """
    import re
    
    # First, try to extract content from code blocks
    patterns = [
        r'```(?:csv|CSV|plain|text)?\n(.*?)\n```',  # Standard code block
        r'```\n(.*?)\n```',  # Generic code block
        r'`([^`]+)`',  # Inline code (fallback)
    ]
    
    for pattern in patterns:
        match = re.search(pattern, content, re.DOTALL | re.MULTILINE)
        if match:
            csv_content = match.group(1).strip()
            # Validate it looks like CSV (has commas or pipes)
            if ',' in csv_content or '|' in csv_content:
                return csv_content
    
    # If no code blocks found, check if the entire content looks like CSV
    lines = content.strip().split('\n')
    if len(lines) >= 2:  # At least header and one data row
        # Check if first line has delimiters
        first_line = lines[0]
        if ',' in first_line or '\t' in first_line or '|' in first_line:
            # Looks like CSV, return as is
            return content.strip()
    
    # Last resort: return content as is
    return content.strip()
def debug_print_files(self, thread_id: str):
    """
    Debug function to print information about files registered with the pandas agent.
    
    Args:
        thread_id (str): The thread ID to check files for
    
    Returns:
        str: Debug information about files and dataframes
    """
    import pandas as pd
    
    debug_output = []
    debug_output.append(f"=== PANDAS AGENT DEBUG INFO FOR THREAD {thread_id} ===")
    
    # Get file info
    file_info = self.file_info_cache.get(thread_id, [])
    debug_output.append(f"Registered Files: {len(file_info)}")
    
    for i, info in enumerate(file_info):
        file_name = info.get("name", "unnamed")
        file_path = info.get("path", "unknown")
        file_type = info.get("type", "unknown")
        debug_output.append(f"\n[{i+1}] File: {file_name} ({file_type})")
        debug_output.append(f"    Path: {file_path}")
        debug_output.append(f"    Exists: {os.path.exists(file_path)}")
        
        if os.path.exists(file_path):
            try:
                file_size = os.path.getsize(file_path)
                debug_output.append(f"    Size: {file_size} bytes")
                
                # Read first few bytes
                with open(file_path, 'rb') as f:
                    first_bytes = f.read(50)
                debug_output.append(f"    First bytes: {first_bytes}")
            except Exception as e:
                debug_output.append(f"    Error reading file: {str(e)}")
    
    # Get dataframe info
    dataframes = self.dataframes_cache.get(thread_id, {})
    debug_output.append(f"\nLoaded DataFrames: {len(dataframes)}")
    
    for df_name, df in dataframes.items():
        debug_output.append(f"\nDataFrame: {df_name}")
        try:
            debug_output.append(f"  Shape: {df.shape}")
            debug_output.append(f"  Columns: {list(df.columns)}")
            debug_output.append(f"  Types: {df.dtypes.to_dict()}")
            
            # Sample data (first 3 rows)
            debug_output.append(f"  Sample data (3 rows):")
            debug_output.append(f"{df.head(3).to_string()}")
            
            # Detect potential issues
            has_nulls = df.isnull().any().any()
            debug_output.append(f"  Contains nulls: {has_nulls}")
            
        except Exception as e:
            debug_output.append(f"  Error examining dataframe: {str(e)}")
    
    # Return as string
    return "\n".join(debug_output)
class PandasAgentManager:
    """
    Enhanced class to manage pandas agents and dataframes for different threads.
    Provides thread isolation, FIFO file storage, and prevents duplicate agent creation.
    """
    
    # Class-level singleton instance
    _instance = None
    
    @classmethod
    def get_instance(cls):
        """Get or create the singleton instance"""
        if cls._instance is None:
            cls._instance = cls()
        return cls._instance
    
    def __init__(self):
        """Initialize the manager"""
        # Cache for pandas agents by thread_id
        self.agents_cache = {}
        
        # Cache for dataframes by thread_id
        self.dataframes_cache = {}
        
        # Cache for file info by thread_id
        self.file_info_cache = {}
        
        # Cache for filepaths by thread_id
        self.file_paths_cache = {}
        
        # Maximum files per thread
        self.max_files_per_thread = 3
        
        # Initialize LangChain LLM
        self.langchain_llm = None
        
        # Check for required dependencies
        self._check_dependencies()
        
        logging.info("PandasAgentManager initialized")
    
    def _check_dependencies(self):
        """Check if required dependencies are available"""
        missing_deps = []
        
        try:
            import pandas as pd
            logging.info("Pandas version: %s", pd.__version__)
        except ImportError:
            missing_deps.append("pandas")
        
        try:
            import numpy as np
            logging.info("NumPy version: %s", np.__version__)
        except ImportError:
            missing_deps.append("numpy")
        
        try:
            import tabulate
            logging.info("Tabulate version: %s", tabulate.__version__)
        except ImportError:
            try:
                # Try to install tabulate automatically
                import subprocess
                logging.info("Installing missing tabulate dependency...")
                subprocess.check_call(["pip", "install", "tabulate"])
                logging.info("Successfully installed tabulate")
            except Exception as e:
                missing_deps.append("tabulate")
                logging.warning(f"Could not install tabulate: {str(e)}")
        
        try:
            # Try to import LangChain's create_pandas_dataframe_agent function
            try:
                from langchain_experimental.agents import create_pandas_dataframe_agent
                logging.info("Using langchain_experimental for pandas agent")
            except ImportError:
                try:
                    from langchain.agents import create_pandas_dataframe_agent
                    logging.info("Using langchain for pandas agent")
                except ImportError:
                    missing_deps.append("langchain[experimental]")
        except Exception as e:
            missing_deps.append("langchain components")
            logging.error(f"Error checking for LangChain: {str(e)}")
        
        if missing_deps:
            logging.warning(f"Missing dependencies: {', '.join(missing_deps)}")
            logging.warning("Install them with: pip install " + " ".join(missing_deps))
        
        return len(missing_deps) == 0
    
    def get_llm(self):
        """Get or initialize the LangChain LLM"""
        if self.langchain_llm is None:
            try:
                from langchain_openai import AzureChatOpenAI
                
                # Use the existing client configuration
                self.langchain_llm = AzureChatOpenAI(
                    azure_endpoint=AZURE_ENDPOINT,
                    api_key=AZURE_API_KEY,
                    api_version=AZURE_API_VERSION,
                    deployment_name="gpt-4o",
                    temperature=0
                )
                logging.info("Initialized LangChain LLM for pandas agents")
            except Exception as e:
                logging.error(f"Failed to initialize LangChain LLM: {e}")
                raise
        
        return self.langchain_llm
    
    def initialize_thread(self, thread_id: str):
        """Initialize storage for a thread if it doesn't exist yet"""
        if thread_id not in self.dataframes_cache:
            self.dataframes_cache[thread_id] = {}
        
        if thread_id not in self.file_info_cache:
            self.file_info_cache[thread_id] = []
            
        if thread_id not in self.file_paths_cache:
            self.file_paths_cache[thread_id] = []
    
    def remove_oldest_file(self, thread_id: str):
        """
        Remove the oldest file for a thread when max files is reached
        
        Args:
            thread_id (str): Thread ID
            
        Returns:
            str or None: Name of removed file or None
        """
        if thread_id not in self.file_info_cache or len(self.file_info_cache[thread_id]) <= self.max_files_per_thread:
            return None
        
        # Get the oldest file info
        oldest_file_info = self.file_info_cache[thread_id][0]
        oldest_file_name = oldest_file_info.get("name", "")
        
        # Remove the file info
        self.file_info_cache[thread_id].pop(0)
        
        # Remove the file path and delete file from disk
        if thread_id in self.file_paths_cache and len(self.file_paths_cache[thread_id]) > 0:
            oldest_path = self.file_paths_cache[thread_id].pop(0)
            # Delete the file from disk
            if os.path.exists(oldest_path):
                try:
                    os.remove(oldest_path)
                    logging.info(f"Deleted oldest file: {oldest_path} for thread {thread_id}")
                except Exception as e:
                    logging.error(f"Error deleting file {oldest_path}: {e}")
        
        # Remove any dataframes associated with this file
        if thread_id in self.dataframes_cache:
            # For exact filename match
            if oldest_file_name in self.dataframes_cache[thread_id]:
                del self.dataframes_cache[thread_id][oldest_file_name]
                
            # For Excel sheets with this filename
            keys_to_remove = []
            for key in self.dataframes_cache[thread_id].keys():
                if key.startswith(oldest_file_name + " [Sheet:"):
                    keys_to_remove.append(key)
            
            for key in keys_to_remove:
                del self.dataframes_cache[thread_id][key]
        
        # Invalidate the agent for this thread since dataframes changed
        if thread_id in self.agents_cache:
            self.agents_cache[thread_id] = None
            
        return oldest_file_name
    
    def add_file(self, thread_id, file_info):
        """
        Add a file to the thread, implementing FIFO if needed.
        Enhanced to prevent accidental deletion of files needed for follow-up queries.
        
        Args:
            thread_id (str): The thread ID
            file_info (dict): File information dictionary
            
        Returns:
            tuple: (dict of dataframes or None, error message or None, removed_file or None)
        """
        # Initialize thread storage if needed
        self.initialize_thread(thread_id)
        
        file_name = file_info.get("name", "unnamed_file")
        file_path = file_info.get("path", None)
        
        logging.info(f"Adding file '{file_name}' from path {file_path} to thread {thread_id}")
        
        # Verify file exists
        if not file_path or not os.path.exists(file_path):
            # Try to locate a matching file in /tmp with a more thorough search
            located_file_path = self._locate_file_in_tmp(file_name)
            if located_file_path:
                logging.info(f"Found alternative path for '{file_name}': {located_file_path}")
                file_path = located_file_path
                file_info["path"] = file_path
            else:
                error_msg = f"File path for '{file_name}' is invalid or does not exist (path: {file_path})"
                logging.error(error_msg)
                
                # Look for the safe copy that might have been created previously
                safe_pattern = f"/tmp/safe_*_{re.sub(r'[^a-zA-Z0-9.]', '_', file_name)}"
                try:
                    import glob
                    safe_copies = glob.glob(safe_pattern)
                    if safe_copies:
                        newest_safe_copy = max(safe_copies, key=os.path.getctime)
                        logging.info(f"Found safe copy for '{file_name}': {newest_safe_copy}")
                        file_path = newest_safe_copy
                        file_info["path"] = file_path
                    else:
                        return None, error_msg, None
                except Exception as e:
                    logging.error(f"Error looking for safe copies: {e}")
                    return None, error_msg, None
        
        # Check if we already have this file (same name)
        existing_file_names = [f.get("name", "") for f in self.file_info_cache[thread_id]]
        existing_file_index = None
        
        if file_name in existing_file_names:
            # File already exists - we'll treat this as an update, but ONLY if we have a valid new file
            logging.info(f"File '{file_name}' already exists for thread {thread_id} - updating")
            
            # Find existing file info but don't delete it yet
            for i, info in enumerate(self.file_info_cache[thread_id]):
                if info.get("name") == file_name:
                    existing_file_index = i
                    old_path = self.file_paths_cache[thread_id][i] if i < len(self.file_paths_cache[thread_id]) else None
                    
                    # If the paths are identical and the file exists, this is likely a follow-up query
                    # Just use the existing dataframe and don't try to reload
                    if old_path and old_path == file_path and os.path.exists(old_path):
                        logging.info(f"Using existing dataframe for '{file_name}' (follow-up query)")
                        if file_name in self.dataframes_cache[thread_id]:
                            return {file_name: self.dataframes_cache[thread_id][file_name]}, None, None
                    break
        
        # Track removed file (if any due to FIFO)
        removed_file = None
        
        # Apply FIFO if we exceed max files
        if len(self.file_info_cache[thread_id]) >= self.max_files_per_thread and existing_file_index is None:
            removed_file = self.remove_oldest_file(thread_id)
            if removed_file:
                logging.info(f"Removed oldest file '{removed_file}' for thread {thread_id} to maintain FIFO limit")
        
        # Load the dataframe(s)
        dfs_dict, error = self.load_dataframe_from_file(file_info)
        
        if error:
            # If there was an error loading the new dataframe but we have an existing one, keep using it
            if file_name in self.dataframes_cache[thread_id]:
                logging.info(f"Error loading updated file, falling back to existing dataframe for '{file_name}'")
                return {file_name: self.dataframes_cache[thread_id][file_name]}, None, removed_file
            
            return None, error, removed_file
            
        if dfs_dict:
            # If we're updating an existing file, remove the old one now that we have a valid replacement
            if existing_file_index is not None:
                old_info = self.file_info_cache[thread_id].pop(existing_file_index)
                
                # Remove the old path
                if existing_file_index < len(self.file_paths_cache[thread_id]):
                    old_path = self.file_paths_cache[thread_id].pop(existing_file_index)
                    
                    # Only delete the old file if it's different from the new one
                    if old_path and old_path != file_path and os.path.exists(old_path):
                        try:
                            # Create a backup before deleting
                            backup_path = f"{old_path}.bak"
                            with open(old_path, 'rb') as src, open(backup_path, 'wb') as dst:
                                dst.write(src.read())
                                
                            os.remove(old_path)
                            logging.info(f"Deleted old file: {old_path} for thread {thread_id} (backup at {backup_path})")
                        except Exception as e:
                            logging.error(f"Error deleting file {old_path}: {e}")
                
                # Remove old dataframes
                if file_name in self.dataframes_cache[thread_id]:
                    del self.dataframes_cache[thread_id][file_name]
                
                # Remove Excel sheets with this filename
                keys_to_remove = []
                for key in self.dataframes_cache[thread_id].keys():
                    if key.startswith(file_name + " [Sheet:"):
                        keys_to_remove.append(key)
                
                for key in keys_to_remove:
                    del self.dataframes_cache[thread_id][key]
            
            # Add dataframes to cache
            self.dataframes_cache[thread_id].update(dfs_dict)
            
            # Add file info to cache (append to end for FIFO ordering)
            self.file_info_cache[thread_id].append(file_info)
            
            # Add file path to cache
            if file_path:
                self.file_paths_cache[thread_id].append(file_path)
            
            # Reset agent to ensure it's recreated with new dataframes
            if thread_id in self.agents_cache:
                self.agents_cache[thread_id] = None
                
            logging.info(f"Added dataframe(s) for file '{file_name}' to thread {thread_id}")
            return dfs_dict, None, removed_file
        else:
            return None, f"Failed to load any dataframes from file '{file_name}'", removed_file

    def _locate_file_in_tmp(self, filename):
        """
        More thorough search for a file in the /tmp directory.
        
        Args:
            filename (str): Original filename to search for
            
        Returns:
            str or None: Path if found, None otherwise
        """
        try:
            # Try direct match first (pandas_agent_timestamp_filename)
            direct_matches = []
            prefix_matches = []
            partial_matches = []
            safe_copy_matches = []
            
            # Clean filename for pattern matching
            clean_filename = re.sub(r'[^a-zA-Z0-9.]', '_', filename)
            filename_lower = filename.lower()
            filename_parts = re.split(r'[_\s.-]', filename_lower)
            filename_parts = [p for p in filename_parts if len(p) > 2]  # Filter out very short parts
            
            for tmp_file in os.listdir('/tmp'):
                tmp_path = os.path.join('/tmp', tmp_file)
                if not os.path.isfile(tmp_path):
                    continue
                    
                tmp_lower = tmp_file.lower()
                
                # Direct pandas agent match
                if tmp_lower.startswith('pandas_agent_') and filename_lower in tmp_lower:
                    direct_matches.append(tmp_path)
                    
                # Safe copy match
                elif tmp_lower.startswith('safe_') and clean_filename.lower() in tmp_lower:
                    safe_copy_matches.append(tmp_path)
                    
                # Prefix match (any prefix + exact filename)
                elif tmp_lower.endswith(filename_lower):
                    prefix_matches.append(tmp_path)
                    
                # Partial match (filename parts)
                elif any(part in tmp_lower for part in filename_parts):
                    # Calculate match score: how many parts match
                    match_score = sum(1 for part in filename_parts if part in tmp_lower)
                    if match_score >= len(filename_parts) // 2:  # At least half the parts match
                        partial_matches.append((tmp_path, match_score))
            
            # Return best match in priority order
            if direct_matches:
                # Multiple matches - take newest
                return max(direct_matches, key=os.path.getctime)
            elif safe_copy_matches:
                return max(safe_copy_matches, key=os.path.getctime)
            elif prefix_matches:
                return max(prefix_matches, key=os.path.getctime)
            elif partial_matches:
                # Sort by match score (descending) and then by creation time (newest first)
                partial_matches.sort(key=lambda x: (-x[1], -os.path.getctime(x[0])))
                return partial_matches[0][0]
                
            return None
        except Exception as e:
            logging.error(f"Error locating file in /tmp: {e}")
            return None

    def load_dataframe_from_file(self, file_info):
        """
        Load dataframe(s) from file information with robust error handling.
        Enhanced with better file discovery and safe copy management.
        
        Args:
            file_info (dict): Dictionary containing file metadata
            
        Returns:
            tuple: (dict of dataframes, error message)
        """
        try:
            import pandas as pd
            import numpy as np
        except ImportError as e:
            logging.error(f"Required library not available: {e}")
            return None, f"Required library not available: {e}"
        
        file_type = file_info.get("type", "unknown")
        file_name = file_info.get("name", "unnamed_file")
        file_path = file_info.get("path", None)
        
        logging.info(f"Loading dataframe from file: {file_name} ({file_type}) from path: {file_path}")
        
        # Verify original file path
        if not file_path or not os.path.exists(file_path):
            # Try to find the file using our enhanced search method
            located_path = self._locate_file_in_tmp(file_name)
            
            if located_path:
                logging.info(f"Using alternative path for {file_name}: {located_path}")
                # Update file_info with the new path
                file_path = located_path
                file_info["path"] = file_path
            else:
                # Still couldn't find the file - create a detailed error
                error_msg = f"File '{file_name}' could not be found. Original path '{file_path}' does not exist and no matches found in /tmp."
                logging.error(error_msg)
                return None, error_msg
        
        try:
            # Verify file readability and log details
            file_size = os.path.getsize(file_path)
            with open(file_path, 'rb') as f_check:
                first_bytes = f_check.read(20)
            logging.info(f"File '{file_name}' exists, size: {file_size} bytes, first bytes: {first_bytes}")
            
            # Make a copy of the file with a simplified name to avoid path issues
            # Use a timestamp to ensure uniqueness and help track file age
            timestamp = int(time.time())
            simple_name = re.sub(r'[^a-zA-Z0-9.]', '_', file_name)  # Replace problematic chars with underscore
            safe_path = f"/tmp/safe_{timestamp}_{simple_name}"
            
            with open(file_path, 'rb') as src, open(safe_path, 'wb') as dst:
                dst.write(src.read())
            
            logging.info(f"Created safe copy of file at: {safe_path}")
            
            # Update file_path to use the safe copy
            file_path = safe_path
            file_info["path"] = safe_path
            
            # Rest of the existing file loading code 
            if file_type == "csv":
                # Try with different encodings and delimiters for robustness
                encodings = ['utf-8', 'latin-1', 'iso-8859-1']
                delimiters = [',', ';', '\t', '|']
                
                df = None
                error_msgs = []
                successful_encoding = None
                successful_delimiter = None
                
                # Try each encoding
                for encoding in encodings:
                    if df is not None:
                        break
                    
                    for delimiter in delimiters:
                        try:
                            # Use the safe path directly
                            df = pd.read_csv(safe_path, encoding=encoding, sep=delimiter, low_memory=False)
                            
                            if len(df.columns) > 1:  # Successfully parsed with >1 column
                                logging.info(f"Successfully loaded CSV with encoding {encoding} and delimiter '{delimiter}'")
                                successful_encoding = encoding
                                successful_delimiter = delimiter
                                break
                        except Exception as e:
                            error_msgs.append(f"Failed with {encoding}/{delimiter}: {str(e)}")
                            continue
                
                if df is None:
                    detailed_error = " | ".join(error_msgs[:5])
                    return None, f"Failed to load CSV file with any encoding/delimiter combination. Errors: {detailed_error}"
                
                # Clean up column names
                df.columns = df.columns.str.strip()
                
                # Replace NaN values with None for better handling
                df = df.replace({np.nan: None})
                
                # Log dataframe info for debugging
                logging.info(f"CSV loaded successfully. Shape: {df.shape}, Columns: {list(df.columns)}")
                logging.info(f"Used encoding: {successful_encoding}, delimiter: '{successful_delimiter}'")
                
                # Return with original filename as key
                return {file_name: df}, None
                
            elif file_type == "excel":
                # Excel loading code remains the same, but use safe_path
                result_dfs = {}
                
                try:
                    xls = pd.ExcelFile(safe_path)
                    sheet_names = xls.sheet_names
                    logging.info(f"Excel file contains {len(sheet_names)} sheets: {sheet_names}")
                except Exception as e:
                    return None, f"Error accessing Excel file: {str(e)}"
                
                if len(sheet_names) == 1:
                    # Single sheet - load directly with the filename as key
                    try:
                        df = pd.read_excel(safe_path, engine='openpyxl')
                        df.columns = df.columns.str.strip()
                        df = df.replace({np.nan: None})
                        result_dfs[file_name] = df
                        
                        # Log dataframe info for debugging
                        logging.info(f"Excel sheet loaded successfully. Shape: {df.shape}, Columns: {list(df.columns)}")
                    except Exception as e:
                        return None, f"Error reading Excel sheet: {str(e)}"
                else:
                    # Multiple sheets - load each sheet with a compound key
                    for sheet in sheet_names:
                        try:
                            df = pd.read_excel(safe_path, sheet_name=sheet, engine='openpyxl')
                            df.columns = df.columns.str.strip()
                            df = df.replace({np.nan: None})
                            
                            # Create a key that includes the sheet name
                            sheet_key = f"{file_name} [Sheet: {sheet}]"
                            result_dfs[sheet_key] = df
                            
                            # Log dataframe info for debugging
                            logging.info(f"Excel sheet '{sheet}' loaded successfully. Shape: {df.shape}, Columns: {list(df.columns)}")
                        except Exception as e:
                            logging.error(f"Error reading sheet '{sheet}' in {file_name}: {str(e)}")
                            # Continue with other sheets even if one fails
                
                if result_dfs:
                    return result_dfs, None
                else:
                    return None, "Failed to load any sheets from Excel file"
            else:
                return None, f"Unsupported file type: {file_type}"
                
        except Exception as e:
            error_msg = f"Error loading '{file_name}': {str(e)}"
            logging.error(f"{error_msg}\n{traceback.format_exc()}")
            return None, error_msg
    def get_or_create_agent(self, thread_id):
        """
        Get or create a pandas agent for a thread with clear dataframe reference instructions.
        
        Args:
            thread_id (str): Thread ID
            
        Returns:
            tuple: (agent, dataframes, errors)
        """
        # Initialize thread storage if needed
        self.initialize_thread(thread_id)
        
        # Import required modules
        try:
            # Try langchain.agents first (more stable)
            try:
                from langchain.agents import create_pandas_dataframe_agent
                from langchain.agents import AgentType
                agent_module = "langchain.agents"
            except ImportError:
                # Fall back to experimental if needed
                from langchain_experimental.agents import create_pandas_dataframe_agent
                from langchain.agents import AgentType
                agent_module = "langchain_experimental.agents"
                
            import pandas as pd
            logging.info(f"Using {agent_module} module for pandas agent creation")
        except ImportError as e:
            return None, None, [f"Required libraries not available: {str(e)}"]
        
        # Initialize the LLM
        try:
            llm = self.get_llm()
        except Exception as e:
            return None, None, [f"Failed to initialize LLM: {str(e)}"]
        
        # Check if we have any dataframes
        if not self.dataframes_cache[thread_id]:
            return None, None, ["No dataframes available for analysis"]
        
        # Create or update the agent if needed
        if thread_id not in self.agents_cache or self.agents_cache[thread_id] is None:
            try:
                # Get all dataframes and file names
                dfs = self.dataframes_cache[thread_id]
                
                # Log the dataframes we're dealing with
                df_names = list(dfs.keys())
                logging.info(f"Creating pandas agent for thread {thread_id} with dataframes: {df_names}")
                
                # IMPORTANT: Create a prefix that explains how to access dataframes
                if len(dfs) == 1:
                    # Single dataframe case
                    df_name = list(dfs.keys())[0]
                    df = dfs[df_name]
                    
                    # Create agent with clear instructions
                    prefix = f"""You are working with a pandas dataframe.
    The dataframe is available as 'df' and represents the file: '{df_name}'
    Shape: {df.shape}
    Columns: {list(df.columns)}
    
    Important: The dataframe is ALREADY LOADED as 'df'. DO NOT try to read the file from disk.
    """
                    
                    self.agents_cache[thread_id] = create_pandas_dataframe_agent(
                        llm,
                        df,
                        prefix=prefix,
                        verbose=True,
                        agent_type="tool-calling",
                        handle_parsing_errors=True,
                        allow_dangerous_code=True,
                        max_iterations=30,
                        max_execution_time=120
                    )
                    
                else:
                    # Multiple dataframes case
                    df_list = list(dfs.values())
                    
                    # Create a mapping explanation
                    df_mapping = []
                    for i, (name, df) in enumerate(dfs.items()):
                        df_mapping.append(f"  dfs[{i}]: '{name}' - Shape: {df.shape}, Columns: {list(df.columns)[:5]}...")
                    
                    mapping_text = "\n".join(df_mapping)
                    
                    # Create prefix with clear instructions
                    prefix = f"""You are working with multiple pandas dataframes in a list called 'dfs'.
    
    Available dataframes:
    {mapping_text}
    
    To access a specific dataframe:
    - Use dfs[0] for the first file, dfs[1] for the second, etc.
    - The dataframes are ALREADY LOADED. DO NOT try to read files from disk.
    
    Example: To analyze the first dataframe, use: dfs[0].describe()
    """
                    
                    self.agents_cache[thread_id] = create_pandas_dataframe_agent(
                        llm,
                        df_list,
                        prefix=prefix,
                        verbose=True,
                        agent_type="tool-calling",
                        handle_parsing_errors=True,
                        allow_dangerous_code=True,
                        max_iterations=30,
                        max_execution_time=120
                    )
                    
                logging.info(f"Successfully created pandas agent for thread {thread_id}")
                
            except Exception as e:
                # If the prefix parameter is not supported, try without it
                try:
                    logging.warning(f"Creating agent without prefix parameter")
                    
                    if len(dfs) == 1:
                        df = list(dfs.values())[0]
                        self.agents_cache[thread_id] = create_pandas_dataframe_agent(
                            llm,
                            df,
                            verbose=True,
                            agent_type="tool-calling",
                            handle_parsing_errors=True,
                            allow_dangerous_code=True,
                            max_iterations=30,
                            max_execution_time=120
                        )
                    else:
                        df_list = list(dfs.values())
                        self.agents_cache[thread_id] = create_pandas_dataframe_agent(
                            llm,
                            df_list,
                            verbose=True,
                            agent_type="tool-calling",
                            handle_parsing_errors=True,
                            allow_dangerous_code=True,
                            max_iterations=30,
                            max_execution_time=120
                        )
                        
                    logging.info(f"Successfully created pandas agent without prefix")
                    
                except Exception as e2:
                    error_msg = f"Failed to create pandas agent: {str(e2)}"
                    logging.error(f"{error_msg}\n{traceback.format_exc()}")
                    return None, self.dataframes_cache[thread_id], [error_msg]
        
        return self.agents_cache[thread_id], self.dataframes_cache[thread_id], []
    
    def check_file_availability(self, thread_id, query):
        """
        Check if a file mentioned in the query is available in the dataframes cache.
        If not, identify which file is missing to inform the user.
        
        Args:
            thread_id (str): Thread ID
            query (str): The query string to check
            
        Returns:
            tuple: (bool, missing_file_name)
        """
        if thread_id not in self.dataframes_cache:
            return False, None
            
        # Get all dataframe names for this thread
        available_files = set(self.dataframes_cache[thread_id].keys())
        
        # Get all file information for this thread (for historical checks)
        all_files_ever_uploaded = set()
        
        # First add files that are currently in the cache
        for file_info in self.file_info_cache.get(thread_id, []):
            file_name = file_info.get("name", "")
            if file_name:
                all_files_ever_uploaded.add(file_name.lower())
        
        # Look for any file name in the query
        query_lower = query.lower()
        
        # First check active files
        for file_name in available_files:
            base_name = file_name.split(" [Sheet:")[0].lower()  # Handle Excel sheet names
            if base_name in query_lower:
                # Found a match in active files
                return True, None
                
        # Check if any file is mentioned but not available (may have been removed by FIFO)
        for file_name in all_files_ever_uploaded:
            if file_name in query_lower and not any(file_name in af.lower() for af in available_files):
                # Found a match in previously uploaded files, but not currently available
                return False, file_name
                
        # No file mentioned or all mentioned files are available
        return True, None
    
    def analyze(self, thread_id, query, files):
        """
        Analyze data with pandas agent.
        
        Args:
            thread_id (str): Thread ID
            query (str): Analysis query
            files (List[Dict]): List of file information
            
        Returns:
            tuple: (result, error, removed_files)
        """
        # Initialize thread if needed
        self.initialize_thread(thread_id)
        
        # List to track any files that were removed
        removed_files = []
        
        # Load all files first
        if files:
            for file_info in files:
                _, _, removed_file = self.add_file(thread_id, file_info)
                if removed_file:
                    removed_files.append(removed_file)
        
        # Check if a file is mentioned but not available
        file_available, missing_file = self.check_file_availability(thread_id, query)
        if not file_available and missing_file:
            return None, f"The file '{missing_file}' is not currently available. Please re-upload the file as it may have been removed due to the 3-file limit per conversation.", removed_files
        
        # Get or create the agent
        agent, dataframes, agent_errors = self.get_or_create_agent(thread_id)
        
        if not agent:
            error_msg = f"Failed to create pandas agent: {'; '.join(agent_errors)}"
            return None, error_msg, removed_files
        
        if not dataframes:
            return None, "No dataframes available for analysis", removed_files
                
        # Extract filename mentions in the query
        mentioned_files = []
        for df_name in dataframes.keys():
            base_name = df_name.split(" [Sheet:")[0].lower()  # Handle Excel sheet names
            if base_name.lower() in query.lower():
                mentioned_files.append(df_name)
                
        # Process the query - ENHANCED FOR BETTER CLARITY
        if len(dataframes) == 1:
            # Single dataframe case
            df_name = list(dataframes.keys())[0]
            df = list(dataframes.values())[0]
            
            enhanced_query = f"""
    The dataframe for '{df_name}' is ALREADY LOADED as the variable 'df'.
    DO NOT try to load the file from disk. Use 'df' directly to access the data.
    
    Dataframe info:
    - Shape: {df.shape}
    - Columns: {list(df.columns)}
    
    Analyze this dataframe to answer: {query}
    """
        else:
            # Multiple dataframes case
            enhanced_query = f"""
    The dataframes are ALREADY LOADED. DO NOT try to load any files from disk.
    Use the dataframes that are already available to you.
    
    Available dataframes: {', '.join(f"'{name}'" for name in dataframes.keys())}
    
    Analyze these dataframes to answer: {query}
    """
            
            if mentioned_files:
                # Add specific guidance for mentioned files
                mentioned_file = mentioned_files[0]
                enhanced_query = f"""
    The dataframe for '{mentioned_file}' is ALREADY LOADED. DO NOT try to load the file from disk.
    Use the dataframe that is already available to you.
    
    Available dataframes: {', '.join(f"'{name}'" for name in dataframes.keys())}
    
    Analyze this dataframe to answer: {query}
    """
                
        logging.info(f"Final query to process: {enhanced_query}")
        
        try:
            # Invoke the agent with the query
            import sys
            from io import StringIO
            
            # Capture stdout to get verbose output for debugging
            original_stdout = sys.stdout
            captured_output = StringIO()
            sys.stdout = captured_output
            
            try:
                # Prepare dataframe details for error cases
                df_details = []
                for name, df in dataframes.items():
                    df_details.append(f"DataFrame '{name}': {df.shape[0]} rows, {df.columns.shape[0]} columns")
                    df_details.append(f"Columns: {', '.join(df.columns.tolist())}")
                    # Add first few rows for debugging
                    try:
                        df_details.append(f"First 3 rows sample:\n{df.head(3)}")
                    except:
                        pass
                
                # First try using run method (per documentation)
                try:
                    logging.info(f"Executing agent with run method: {enhanced_query}")
                    agent_output = agent.run(enhanced_query)
                    logging.info(f"Agent completed successfully with run() method: {agent_output[:100]}...")
                except Exception as run_error:
                    # Fall back to invoke if run fails
                    logging.warning(f"Agent run() method failed: {str(run_error)}, trying invoke() method")
                    try:
                        agent_result = agent.invoke({"input": enhanced_query})
                        agent_output = agent_result.get("output", "")
                        logging.info(f"Agent completed successfully with invoke() method: {agent_output[:100]}...")
                    except Exception as invoke_error:
                        # Try one last approach - if we see a file not found error
                        raise Exception(f"Agent run() failed: {str(run_error)}; invoke() also failed: {str(invoke_error)}")
                
                # Get the captured verbose output
                verbose_output = captured_output.getvalue()
                logging.info(f"Agent verbose output:\n{verbose_output}")
                
                # Check if output seems empty or error-like
                if not agent_output or "I don't have access to" in agent_output or "not find" in agent_output.lower():
                    logging.warning(f"Agent response appears problematic: {agent_output}")
                    
                    # Check if there was a file not found error or variable name confusion in the verbose output
                    if ("FileNotFoundError" in verbose_output or 
                        "No such file" in verbose_output or
                        "NameError" in verbose_output or
                        "not defined" in verbose_output):
                        # Generate a direct summary from the dataframes
                        summary = []
                        for name, df in dataframes.items():
                            summary.append(f"## Summary of {name}")
                            summary.append(f"* Shape: {df.shape[0]} rows, {df.shape[1]} columns")
                            summary.append(f"* Columns: {', '.join(df.columns.tolist())}")
                            
                            # Add basic statistics for numeric columns
                            try:
                                num_cols = df.select_dtypes(include=['number']).columns
                                if len(num_cols) > 0:
                                    summary.append("\n### Basic Statistics for Numeric Columns:")
                                    summary.append(df[num_cols].describe().to_string())
                            except Exception as stats_err:
                                summary.append(f"Error calculating statistics: {str(stats_err)}")
                            
                            # Add sample data
                            try:
                                summary.append("\n### Sample Data (First 5 rows):")
                                summary.append(df.head(5).to_string())
                            except Exception as sample_err:
                                summary.append(f"Error showing sample: {str(sample_err)}")
                                
                            summary.append("\n")
                            
                        return "\n".join(summary), None, removed_files
                    else:
                        # Provide detailed dataframe information as fallback
                        fallback_output = "I analyzed your data and found:\n\n" + "\n".join(df_details[:10])
                        logging.info(f"Providing fallback output with basic dataframe info")
                        return fallback_output, None, removed_files
                
                # Final response - successful case
                return agent_output, None, removed_files
                
            except Exception as e:
                error_detail = str(e)
                tb = traceback.format_exc()
                logging.error(f"Agent execution error: {error_detail}\n{tb}")
                
                # Get verbose output for debugging
                verbose_output = captured_output.getvalue()
                logging.info(f"Agent debugging output before error:\n{verbose_output}")
                
                # Check if there was a file not found error or variable name issue
                if ("FileNotFoundError" in verbose_output or 
                    "No such file" in verbose_output or 
                    "FileNotFoundError" in error_detail or
                    "NameError" in error_detail or
                    "not defined" in error_detail):
                    # Generate a direct summary from the dataframes
                    summary = []
                    for name, df in dataframes.items():
                        summary.append(f"## Summary of {name}")
                        summary.append(f"* Shape: {df.shape[0]} rows, {df.shape[1]} columns")
                        summary.append(f"* Columns: {', '.join(df.columns.tolist())}")
                        
                        # Add basic statistics for numeric columns
                        try:
                            num_cols = df.select_dtypes(include=['number']).columns
                            if len(num_cols) > 0:
                                summary.append("\n### Basic Statistics for Numeric Columns:")
                                summary.append(df[num_cols].describe().to_string())
                        except Exception as stats_err:
                            summary.append(f"Error calculating statistics: {str(stats_err)}")
                        
                        # Add sample data
                        try:
                            summary.append("\n### Sample Data (First 5 rows):")
                            summary.append(df.head(5).to_string())
                        except Exception as sample_err:
                            summary.append(f"Error showing sample: {str(sample_err)}")
                            
                        summary.append("\n")
                        
                    return "\n".join(summary), None, removed_files
                else:
                    # Re-raise other errors
                    raise e
                    
        except Exception as final_error:
            error_msg = f"Failed to analyze data: {str(final_error)}"
            logging.error(f"{error_msg}\n{traceback.format_exc()}")
            
            # Last resort - if all else fails, provide basic dataframe info
            try:
                basic_info = []
                for name, df in dataframes.items():
                    basic_info.append(f"File: {name}")
                    basic_info.append(f"  Rows: {df.shape[0]}")
                    basic_info.append(f"  Columns: {df.shape[1]}")
                    basic_info.append(f"  Column names: {', '.join(df.columns.tolist())}")
                    basic_info.append("")
                
                if basic_info:
                    return "I encountered an error but here's what I found about your data:\n\n" + "\n".join(basic_info), None, removed_files
            except:
                pass
                
            return None, error_msg, removed_files
            
        finally:
            # Restore stdout
            sys.stdout = original_stdout
async def validate_resources(client: AzureOpenAI, thread_id: Optional[str], assistant_id: Optional[str]) -> Dict[str, bool]:
    """
    Validates that the given thread_id and assistant_id exist and are accessible.
    
    Args:
        client (AzureOpenAI): The Azure OpenAI client instance
        thread_id (Optional[str]): The thread ID to validate, or None
        assistant_id (Optional[str]): The assistant ID to validate, or None
        
    Returns:
        Dict[str, bool]: Dictionary with "thread_valid" and "assistant_valid" flags
    """
    result = {
        "thread_valid": False,
        "assistant_valid": False
    }
    
    # Validate thread if provided
    if thread_id:
        try:
            # Attempt to retrieve thread
            thread = client.beta.threads.retrieve(thread_id=thread_id)
            result["thread_valid"] = True
            logging.info(f"Thread validation: {thread_id} is valid")
        except Exception as e:
            result["thread_valid"] = False
            logging.warning(f"Thread validation: {thread_id} is invalid - {str(e)}")
    
    # Validate assistant if provided
    if assistant_id:
        try:
            # Attempt to retrieve assistant
            assistant = client.beta.assistants.retrieve(assistant_id=assistant_id)
            result["assistant_valid"] = True
            logging.info(f"Assistant validation: {assistant_id} is valid")
        except Exception as e:
            result["assistant_valid"] = False
            logging.warning(f"Assistant validation: {assistant_id} is invalid - {str(e)}")
    
    return result
async def pandas_agent(client: AzureOpenAI, thread_id: Optional[str], query: str, files: List[Dict[str, Any]]) -> str:
    """
    Enhanced pandas_agent that uses LangChain to analyze CSV and Excel files.
    Uses a class-based implementation to maintain isolation between threads.
    
    Args:
        client (AzureOpenAI): The Azure OpenAI client instance
        thread_id (Optional[str]): The thread ID to add the response to
        query (str): The query or question about the data
        files (List[Dict[str, Any]]): List of file information dictionaries
        
    Returns:
        str: The analysis result
    """
    # Create a unique operation ID for tracking
    operation_id = f"pandas_agent_{int(time.time())}_{os.urandom(2).hex()}"
    update_operation_status(operation_id, "started", 0, "Starting data analysis")
    
    # Flag for background thread to stop
    stop_background_updates = threading.Event()
    
    # Background thread for progress updates
    def send_progress_updates():
        progress = 50
        while progress < 85 and not stop_background_updates.is_set():
            time.sleep(1.5)  # Update every 1.5 seconds
            progress += 2
            update_operation_status(operation_id, "executing", min(progress, 85), 
                                   "Analysis in progress...")
    
    try:
        # Verify thread_id is provided
        if not thread_id:
            error_msg = "Thread ID is required for pandas agent"
            update_operation_status(operation_id, "error", 100, error_msg)
            return f"Error: {error_msg}"
        
        # Enhanced debugging: Log detailed info about files
        logging.info(f"PANDAS AGENT DEBUG - Query: '{query}'")
        logging.info(f"PANDAS AGENT DEBUG - Thread ID: {thread_id}")
        logging.info(f"PANDAS AGENT DEBUG - Files count: {len(files)}")
        
        # Log files being processed with much more detail
        file_descriptions = []
        valid_files = []
        invalid_files = []
        
        for i, file in enumerate(files):
            file_type = file.get("type", "unknown")
            file_name = file.get("name", "unnamed_file")
            file_path = file.get("path", "unknown_path")
            
            # Detailed file logging
            debug_info = (f"File {i+1}: '{file_name}' ({file_type}) - "
                         f"Path: {file_path}")
            logging.info(f"PANDAS AGENT DEBUG - {debug_info}")
            
            file_descriptions.append(f"{file_name} ({file_type})")
            
            # Verify file existence with more robust checking
            file_exists = False
            file_size = None
            first_bytes = None
            
            if file_path and os.path.exists(file_path):
                try:
                    file_size = os.path.getsize(file_path)
                    with open(file_path, 'rb') as f:
                        first_bytes = f.read(10)
                    file_exists = True
                    logging.info(f"PANDAS AGENT DEBUG - File verified: '{file_name}' exists, size: {file_size} bytes, first bytes: {first_bytes}")
                    valid_files.append(file)
                except Exception as e:
                    logging.warning(f"PANDAS AGENT DEBUG - File exists but cannot read: '{file_name}' - {str(e)}")
                    invalid_files.append((file_name, f"Read error: {str(e)}"))
            else:
                logging.warning(f"PANDAS AGENT DEBUG - File does not exist: '{file_name}' at path: {file_path}")
                invalid_files.append((file_name, f"Path not found: {file_path}"))
                
                # Enhanced path correction: Look for files with similar names in /tmp
                possible_paths = [
                    path for path in os.listdir('/tmp') 
                    if file_name.lower() in path.lower() and os.path.isfile(os.path.join('/tmp', path))
                ]
                
                if possible_paths:
                    logging.info(f"PANDAS AGENT DEBUG - Found possible alternatives for '{file_name}':")
                    for alt_path in possible_paths:
                        full_path = os.path.join('/tmp', alt_path)
                        alt_size = os.path.getsize(full_path)
                        logging.info(f"  - Alternative: {full_path} (size: {alt_size} bytes)")
                    
                    # Use the first alternative found
                    corrected_path = os.path.join('/tmp', possible_paths[0])
                    logging.info(f"PANDAS AGENT DEBUG - Using alternative path for {file_name}: {corrected_path}")
                    file["path"] = corrected_path
                    valid_files.append(file)
                else:
                    # Try a more aggressive search for similarly named files
                    all_tmp_files = os.listdir('/tmp')
                    csv_files = [f for f in all_tmp_files if f.endswith('.csv') or f.endswith('.xlsx') or f.endswith('.xls')]
                    
                    logging.info(f"PANDAS AGENT DEBUG - Available files in /tmp directory:")
                    for tmp_file in csv_files[:10]:  # Show first 10 to avoid log flooding
                        logging.info(f"  - {tmp_file}")
                    
                    # Check if filename parts match (for handling timestamp prefixes)
                    name_parts = re.split(r'[_\s]', file_name.lower())
                    for tmp_file in csv_files:
                        match_score = sum(1 for part in name_parts if part in tmp_file.lower())
                        if match_score >= 2:  # If at least 2 parts match
                            corrected_path = os.path.join('/tmp', tmp_file)
                            logging.info(f"PANDAS AGENT DEBUG - Found partial match for '{file_name}': {corrected_path}")
                            file["path"] = corrected_path
                            valid_files.append(file)
                            break
        
        # Replace original files list with validated files
        files = valid_files
        
        # Summary log
        file_list_str = ", ".join(file_descriptions) if file_descriptions else "No files provided"
        logging.info(f"PANDAS AGENT DEBUG - Processing data analysis for thread {thread_id} with files: {file_list_str}")
        logging.info(f"PANDAS AGENT DEBUG - Valid files: {len(valid_files)}, Invalid files: {len(invalid_files)}")
        
        if len(valid_files) == 0:
            update_operation_status(operation_id, "error", 100, "No valid files found for analysis")
            return f"Error: Could not find any valid files to analyze. Please verify the uploaded files and try again.\n\nDebug info: {file_list_str}"
            
        update_operation_status(operation_id, "files", 20, f"Processing files: {file_list_str}")
        
        # Get the PandasAgentManager instance
        manager = PandasAgentManager.get_instance()
        
        # Process the query
        update_operation_status(operation_id, "analyzing", 50, f"Analyzing data with query: {query}")
        
        # Start progress update in background
        update_thread = None
        try:
            update_thread = threading.Thread(target=send_progress_updates)
            update_thread.daemon = True
            update_thread.start()
        except Exception as e:
            # Don't fail if we can't spawn thread
            logging.warning(f"Could not start progress update thread: {str(e)}")
        
        # Run the analysis using the PandasAgentManager
        result, error, removed_files = manager.analyze(thread_id, query, files)
        
        # Stop the background updates
        stop_background_updates.set()
        
        # Wait for the update thread to terminate (with timeout)
        if update_thread and update_thread.is_alive():
            update_thread.join(timeout=1.0)
        
        # Prepare the response
        update_operation_status(operation_id, "formatting", 90, "Formatting response")
        
        if error:
            update_operation_status(operation_id, "error", 95, f"Error: {error}")
            
            # Detect if the error appears to be a file access issue
            if "access" in error.lower() or "find" in error.lower() or "read" in error.lower():
                # Try to get basic dataframe info as a fallback
                try:
                    df_info = []
                    dfs = manager.dataframes_cache.get(thread_id, {})
                    
                    if dfs:
                        for df_name, df in dfs.items():
                            df_info.append(f"- {df_name}: {df.shape[0]} rows, {len(df.columns)} columns")
                            df_info.append(f"  Columns: {', '.join(df.columns[:10].tolist())}")
                            
                            # Add sample data (first 3 rows) if we can
                            try:
                                sample = df.head(3).to_string()
                                df_info.append(f"  Sample data:\n{sample}")
                            except:
                                pass
                            
                        fallback_response = (
                            f"I encountered an issue while analyzing your data files but can provide "
                            f"basic information about them:\n\n{chr(10).join(df_info)}\n\n"
                            f"Error details: {error}"
                        )
                        
                        # If files were removed via FIFO, add a note about it
                        if removed_files:
                            removed_files_str = ", ".join(f"'{f}'" for f in removed_files)
                            fallback_response += f"\n\nNote: The following file(s) were removed due to the 3-file limit: {removed_files_str}"
                            
                        return fallback_response
                except Exception as fallback_e:
                    # If fallback fails, return original error
                    logging.error(f"Fallback info generation failed: {fallback_e}")
            
            # Add debug information to error response
            debug_info = ""
            if invalid_files:
                debug_info = "\n\nDebug information:\n"
                for name, err in invalid_files:
                    debug_info += f"- File '{name}': {err}\n"
                
                # Add info about files in /tmp
                try:
                    tmp_files = [f for f in os.listdir('/tmp') if f.endswith('.csv') or f.endswith('.xlsx')]
                    if tmp_files:
                        debug_info += f"\nAvailable files in /tmp:\n"
                        for i, tmp_file in enumerate(tmp_files[:10]):  # Show first 10
                            tmp_path = os.path.join('/tmp', tmp_file)
                            tmp_size = os.path.getsize(tmp_path)
                            debug_info += f"- {tmp_file} (size: {tmp_size} bytes)\n"
                except Exception as tmp_err:
                    debug_info += f"\nError listing /tmp: {str(tmp_err)}\n"
            
            # Standard error response with debugging info
            final_response = f"Error analyzing data: {error}{debug_info}"
            
            # If files were removed via FIFO, add a note about it
            if removed_files:
                removed_files_str = ", ".join(f"'{f}'" for f in removed_files)
                final_response += f"\n\nNote: The following file(s) were removed due to the 3-file limit: {removed_files_str}"
        else:
            final_response = result if result else "No results were returned from the analysis. Try reformulating your query."
            
            # If files were removed via FIFO, add a note about it
            if removed_files:
                removed_files_str = ", ".join(f"'{f}'" for f in removed_files)
                final_response += f"\n\nNote: The following file(s) were removed due to the 3-file limit: {removed_files_str}"
        
        # Add response to thread if provided
        if thread_id:
            update_operation_status(operation_id, "responding", 95, "Adding response to thread")
            try:
                # Wait for any active runs to complete (using sync version since pandas_agent is sync)
                run_ready = sync_wait_for_run_completion(client, thread_id)
                if not run_ready:
                    logging.warning(f"Could not add pandas_agent response to thread {thread_id} - run still active after timeout")
                else:
                    client.beta.threads.messages.create(
                        thread_id=thread_id,
                        role="user",
                        content=f"[PANDAS AGENT RESPONSE]: {final_response}",
                        metadata={"type": "pandas_agent_response", "operation_id": operation_id}
                    )
                    logging.info(f"Added pandas_agent response to thread {thread_id}")
            except Exception as e:
                logging.error(f"Error adding pandas_agent response to thread: {e}")
                # Continue execution despite error with thread message
        
        # Mark operation as completed
        update_operation_status(operation_id, "completed", 100, "Analysis completed successfully")
        
        # Log completion
        logging.info(f"Pandas agent completed query: '{query}' for thread {thread_id}")
        
        return final_response
    
    except Exception as e:
        # Stop the background updates
        stop_background_updates.set()
        
        error_details = traceback.format_exc()
        logging.error(f"Critical error in pandas_agent: {str(e)}\n{error_details}")
        
        # Update status to reflect error
        update_operation_status(operation_id, "error", 100, f"Error: {str(e)}")
        
        # Try to provide some helpful debugging information
        debug_info = []
        try:
            # Check if files exist
            debug_info.append("File System Debugging Information:")
            for file in files:
                file_name = file.get("name", "unnamed")
                file_path = file.get("path", "unknown")
                if file_path and os.path.exists(file_path):
                    debug_info.append(f"- File '{file_name}' exists at path: {file_path}")
                    file_size = os.path.getsize(file_path)
                    debug_info.append(f"  - Size: {file_size} bytes")
                    # Try to read first 20 bytes to verify readability
                    try:
                        with open(file_path, 'rb') as f:
                            first_bytes = f.read(20)
                        debug_info.append(f"  - First bytes: {first_bytes}")
                    except Exception as read_err:
                        debug_info.append(f"  - Read error: {str(read_err)}")
                else:
                    debug_info.append(f"- File '{file_name}' does not exist at path: {file_path}")
                    
                    # Look for similar files
                    possible_paths = [
                        path for path in os.listdir('/tmp') 
                        if path.endswith(('.csv', '.xlsx', '.xls')) and os.path.isfile(os.path.join('/tmp', path))
                    ]
                    if possible_paths:
                        debug_info.append(f"  - Found possible CSV/Excel files in /tmp directory:")
                        for i, path in enumerate(possible_paths[:5]):  # Show first 5
                            debug_info.append(f"    {i+1}. {path}")
        except Exception as debug_err:
            debug_info.append(f"Error during debugging: {str(debug_err)}")
        
        # Provide a graceful failure response with debug info
        debug_str = "\n".join(debug_info)
        error_response = f"""Sorry, I encountered an error while trying to analyze your data files.

Error details: {str(e)}

Additional debugging information:
{debug_str}

Please try again with a different query or contact support if the issue persists.

Operation ID: {operation_id}"""
                
        return error_response
        
async def image_analysis(client: AzureOpenAI, image_data: bytes, filename: str, prompt: Optional[str] = None) -> str:
    """Analyzes an image using Azure OpenAI vision capabilities and returns the analysis text."""
    try:
        ext = os.path.splitext(filename)[1].lower()
        b64_img = base64.b64encode(image_data).decode("utf-8")
        # Try guessing mime type, default to jpeg if extension isn't standard or determinable
        mime, _ = mimetypes.guess_type(filename)
        if not mime or not mime.startswith('image'):
            mime = f"image/{ext[1:]}" if ext and ext[1:] in ['jpg', 'jpeg', 'png', 'gif', 'webp'] else "image/jpeg"

        data_url = f"data:{mime};base64,{b64_img}"

        default_prompt = (
            "Analyze this image and provide a thorough summary including all elements. "
            "If there's any text visible, include all the textual content. Describe:"
        )
        combined_prompt = f"{default_prompt} {prompt}" if prompt else default_prompt

        # Use the existing client instead of creating a new one
        response = client.chat.completions.create(
            model="gpt-4.1-mini",  # Ensure this model supports vision
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": combined_prompt},
                    {"type": "image_url", "image_url": {"url": data_url, "detail": "high"}}
                ]
            }],
            max_tokens=1000  # Increased max_tokens for potentially more detailed analysis
        )

        analysis_text = response.choices[0].message.content
        return analysis_text if analysis_text else "No analysis content received."

    except Exception as e:
        logging.error(f"Image analysis error for {filename}: {e}")
        return f"Error analyzing image '{filename}': {str(e)}"

# Helper function to update user persona context
async def update_context(client: AzureOpenAI, thread_id: str, context: str):
    """Updates the user persona context in a thread by adding/replacing a special message."""
    if not context:
        return

    try:
        # Get existing messages to check for previous context
        messages = client.beta.threads.messages.list(
            thread_id=thread_id,
            order="desc",
            limit=20  # Check recent messages is usually sufficient
        )

        # Look for previous context messages to avoid duplication
        previous_context_message_id = None
        for msg in messages.data:
            if hasattr(msg, 'metadata') and msg.metadata and msg.metadata.get('type') == 'user_persona_context':
                previous_context_message_id = msg.id
                break

        # If found, delete previous context message to replace it
        if previous_context_message_id:
            try:
                client.beta.threads.messages.delete(
                    thread_id=thread_id,
                    message_id=previous_context_message_id
                )
                logging.info(f"Deleted previous context message {previous_context_message_id} in thread {thread_id}")
            except Exception as e:
                logging.error(f"Error deleting previous context message {previous_context_message_id}: {e}")
            # Continue even if delete fails to add the new context

        # Add new context message
        client.beta.threads.messages.create(
            thread_id=thread_id,
            role="user",
            content=f"USER PERSONA CONTEXT: {context}",
            metadata={"type": "user_persona_context"}
        )

        logging.info(f"Updated user persona context in thread {thread_id}")
    except Exception as e:
        logging.error(f"Error updating context in thread {thread_id}: {e}")
        # Continue the flow even if context update fails

# Function to add file awareness to the assistant
async def add_file_awareness(client: AzureOpenAI, thread_id: str, file_info: Dict[str, Any]):
    """Adds file awareness to the assistant by sending a message about the file."""
    if not file_info:
        return

    try:
        # Create a message that informs the assistant about the file
        file_type = file_info.get("type", "unknown")
        file_name = file_info.get("name", "unnamed_file")
        processing_method = file_info.get("processing_method", "")

        awareness_message = f"FILE INFORMATION: A file named '{file_name}' of type '{file_type}' has been uploaded and processed. "

        if processing_method == "pandas_agent":
            awareness_message += f"This file is available for analysis using the pandas agent."
            if file_type == "excel":
                awareness_message += " This is an Excel file with potentially multiple sheets."
            elif file_type == "csv":
                awareness_message += " This is a CSV file."
            
            awareness_message += "\n\nIMPORTANT: You MUST use the pandas_agent tool for ANY request that mentions this file or asks about data analysis. This includes:"
            awareness_message += "\n- Simple requests like 'explain the file'"
            awareness_message += "\n- Vague requests like 'tell me about the data'"
            awareness_message += "\n- Explicit requests like 'analyze this CSV'" 
            awareness_message += "\n- Any query containing the filename"
            awareness_message += "\n- ANY follow-up question after discussing this file"
            
            awareness_message += "\n\nUser requests like 'explain the report' or 'summarize the data' should ALWAYS trigger you to use the pandas_agent tool."
            awareness_message += "\n\nNEVER try to answer questions about this file from memory - ALWAYS use the pandas_agent tool."
        
        elif processing_method == "thread_message":
            awareness_message += "This image has been analyzed and the descriptive content has been added to this thread."
        elif processing_method == "vector_store":
            awareness_message += "This file has been added to the vector store and its content is available for search."
        else:
            awareness_message += "This file has been processed."

        # Send the message to the thread
        run_ready = await wait_for_run_completion(client, thread_id)
        if run_ready:
            # Send the message to the thread
            client.beta.threads.messages.create(
                thread_id=thread_id,
                role="user",  # Sending as user so assistant 'sees' it as input/instruction
                content=awareness_message,
                metadata={"type": "file_awareness", "processed_file": file_name}
            )
            logging.info(f"Added file awareness for '{file_name}' ({processing_method}) to thread {thread_id}")
        else:
            logging.warning(f"Could not add file awareness - run still active on thread {thread_id}")
    except Exception as e:
        logging.error(f"Error adding file awareness for '{file_name}' to thread {thread_id}: {e}")
        # Continue the flow even if adding awareness fails
def update_operation_status(operation_id: str, status: str, progress: float, message: str):
    """Update the status of a long-running operation."""
    operation_statuses[operation_id] = {
        "status": status,
        "progress": progress,
        "message": message,
        "updated_at": time.time()
    }
    logging.info(f"Operation {operation_id}: {status} - {progress:.0f}% - {message}")


@app.post("/initiate-chat",
          response_model=ChatInitResponse,
          summary="Initialize Chat Session",
          description="""Create a new chat session with an AI assistant.

Creates thread and vector store for context.""",
          tags=["Chat Operations"])
async def initiate_chat(
    file: Optional[UploadFile] = File(default=None, description="Initial file to process"),  # Changed File(None) to File(default=None)
    context: Optional[str] = Form(default=None, description="User context or persona")
):
    """
    Create a new chat session with persistent context.
    
    Form parameters (now properly handled by FastAPI):
    - file (optional): Initial file to process
    - context (optional): User context or persona
    """
    client = create_client()
    logging.info("Initiating new chat session...")

    # Create a vector store up front
    try:
        vector_store = client.vector_stores.create(name=f"chat_init_store_{int(time.time())}")
        logging.info(f"Vector store created: {vector_store.id}")
    except Exception as e:
        vector_store = client.beta.vector_stores.create(name=f"chat_init_store_{int(time.time())}")
        logging.error(f"Creating vector store with beta: {e}")
    # Include file_search and add pandas_agent as a function tool
    assistant_tools = [
        {"type": "file_search"},
        {
            "type": "function",
            "function": {
                "name": "pandas_agent",
                "description": "Analyzes CSV and Excel files to answer data-related questions and perform data analysis",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "query": {
                            "type": "string",
                            "description": "The specific question or analysis task to perform on the data. Be comprehensive and explicit."
                        },
                        "filename": {
                            "type": "string",
                            "description": "Optional: specific filename to analyze. If not provided, all available files will be considered."
                        }
                    },
                    "required": ["query"]
                }
            }
        }
    ]
    # Add content generation tools
    content_tools = get_content_generation_tools()
    assistant_tools.extend(content_tools)
    assistant_tool_resources = {
        "file_search": {"vector_store_ids": [vector_store.id]}
    }

    # Keep track of CSV/Excel files for the session
    session_csv_excel_files = []

    # Use the improved system prompt
    
    
    # Create the assistant
    try:
        assistant = client.beta.assistants.create(
            name=f"pm_copilot_{int(time.time())}",
            model="gpt-4.1-mini",  # Ensure this model is deployed
            instructions=system_prompt,
            tools=assistant_tools,
            tool_resources=assistant_tool_resources,
        )
        logging.info(f'Assistant created: {assistant.id}')
    except Exception as e:
        logging.error(f"An error occurred while creating the assistant: {e}")
        # Attempt to clean up vector store if assistant creation fails
        try:
            client.vector_stores.delete(vector_store_id=vector_store.id)
            logging.info(f"Cleaned up vector store {vector_store.id} after assistant creation failure.")
        except Exception as cleanup_e:
            logging.error(f"Failed to cleanup vector store {vector_store.id} after error: {cleanup_e}")
        raise HTTPException(status_code=500, detail=f"An error occurred while creating assistant: {e}")

    # Create a thread
    try:
        thread = client.beta.threads.create()
        logging.info(f"Thread created: {thread.id}")
    except Exception as e:
        logging.error(f"An error occurred while creating the thread: {e}")
        # Attempt cleanup
        try:
            client.beta.assistants.delete(assistant_id=assistant.id)
            logging.info(f"Cleaned up assistant {assistant.id} after thread creation failure.")
        except Exception as cleanup_e:
            logging.error(f"Failed to cleanup assistant {assistant.id} after error: {cleanup_e}")
        try:
            client.vector_stores.delete(vector_store_id=vector_store.id)
            logging.info(f"Cleaned up vector store {vector_store.id} after thread creation failure.")
        except Exception as cleanup_e:
            logging.error(f"Failed to cleanup vector store {vector_store.id} after error: {cleanup_e}")
        raise HTTPException(status_code=500, detail=f"An error occurred while creating the thread: {e}")

    # If context is provided, add it as user persona context
    if context:
        await update_context(client, thread.id, context)
    # Errors handled within update_context

    # If a file is provided, upload and process it
    if file:
        filename = file.filename
        file_content = await file.read()
        file_path = os.path.join('/tmp/', filename)  # Use /tmp or a configurable temp dir

        try:
            with open(file_path, 'wb') as f:
                f.write(file_content)

            # Determine file type
            file_ext = os.path.splitext(filename)[1].lower()
            is_csv = file_ext == '.csv'
            is_excel = file_ext in ['.xlsx', '.xls', '.xlsm']
            # Check MIME type as well for broader image support
            mime_type, _ = mimetypes.guess_type(filename)
            is_image = file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'] or (mime_type and mime_type.startswith('image/'))
            is_document = file_ext in ['.pdf', '.doc', '.docx', '.txt', '.md', '.html', '.json']  # Common types for vector store

            file_info = {"name": filename}

            if is_csv or is_excel:
                # Instead of using code_interpreter, we'll track CSV/Excel files for the pandas_agent
                session_csv_excel_files.append({
                    "name": filename,
                    "path": file_path,
                    "type": "csv" if is_csv else "excel"
                })
                
                file_info.update({
                    "type": "csv" if is_csv else "excel",
                    "processing_method": "pandas_agent"
                })
                
                # Keep a copy of the file for the pandas agent to use
                # (In a real implementation, you might store this in a database or cloud storage)
                permanent_path = os.path.join('/tmp/', f"pandas_agent_{int(time.time())}_{filename}")
                with open(permanent_path, 'wb') as f:
                    with open(file_path, 'rb') as src:
                        f.write(src.read())
                
                # Add file awareness message
                await add_file_awareness(client, thread.id, file_info)
                logging.info(f"Added '{filename}' for pandas_agent processing")

            elif is_image:
                # Analyze image and add analysis text to the thread
                analysis_text = await image_analysis(client, file_content, filename, None)
                client.beta.threads.messages.create(
                    thread_id=thread.id,
                    role="user",  # Add analysis as user message for context
                    content=f"Analysis result for uploaded image '{filename}':\n{analysis_text}"
                )
                file_info.update({
                    "type": "image",
                    "processing_method": "thread_message"
                })
                await add_file_awareness(client, thread.id, file_info)
                logging.info(f"Added image analysis for '{filename}' to thread {thread.id}")

            elif is_document or not (is_csv or is_excel or is_image):
                # Upload to vector store
                with open(file_path, "rb") as file_stream:
                    file_batch =client.vector_stores.file_batches.upload_and_poll(
                        vector_store_id=vector_store.id,
                        files=[file_stream]
                    )
                file_info.update({
                    "type": file_ext[1:] if file_ext else "document",
                    "processing_method": "vector_store"
                })
                await add_file_awareness(client, thread.id, file_info)
                logging.info(f"File '{filename}' uploaded to vector store {vector_store.id}: status={file_batch.status}, count={file_batch.file_counts.total}")

            else:
                logging.warning(f"File type for '{filename}' not explicitly handled for upload, skipping specific processing.")

        except Exception as e:
            logging.error(f"Error processing uploaded file '{filename}': {e}")
            # Don't raise HTTPException here, allow response with IDs but log error
        finally:
            # Clean up temporary file
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                except OSError as e:
                    logging.error(f"Error removing temporary file {file_path}: {e}")

    # Store csv/excel files info in a metadata message if there are any
    if session_csv_excel_files:
        try:
            # Wait for any active runs to complete before adding message
            run_ready = await wait_for_run_completion(client, thread.id)
            if run_ready:
                # Create a special message to store file paths for the pandas agent
                pandas_files_info = json.dumps(session_csv_excel_files)
                client.beta.threads.messages.create(
                    thread_id=thread.id,
                    role="user",
                    content="PANDAS_AGENT_FILES_INFO (DO NOT DISPLAY TO USER)",
                    metadata={
                        "type": "pandas_agent_files",
                        "files": pandas_files_info
                    }
                )
                logging.info(f"Stored pandas agent files info in thread {thread.id}")
            else:
                logging.warning(f"Could not store pandas files info - run still active on thread {thread.id}")
        except Exception as e:
            logging.error(f"Error storing pandas agent files info: {e}")


    res = {
        "message": "Chat initiated successfully.",
        "assistant": assistant.id,
        "session": thread.id,  # Use 'session' for thread_id consistency with other endpoints
        "vector_store": vector_store.id
    }

    return JSONResponse(res, status_code=200)
@app.post("/co-pilot",
          response_model=ChatInitResponse,
          summary="Create Session with Existing Assistant",
          description="Create a new chat session using an existing assistant and vector store.",
          tags=["Chat Operations"])
async def co_pilot(
    request: Request
):
    """
    Use existing assistant for new chat session.
    
    Form parameters (parsed from request):
    - assistant (required): Existing assistant ID
    - vector_store (required): Existing vector store ID
    - context (optional): Session context
    """
    client = create_client()

    # Parse the form data
    try:
        form = await request.form()
        context: Optional[str] = form.get("context", None)
        assistant_id: Optional[str] = form.get("assistant", None)
        vector_store_id: Optional[str] = form.get("vector_store", None)
    except Exception as e:
        logging.error(f"Error parsing form data: {e}")
        raise HTTPException(status_code=400, detail=f"Invalid form data: {e}")

    # Validate required parameters
    if not assistant_id or not vector_store_id:
        raise HTTPException(status_code=400, detail="Both assistant_id and vector_store_id are required")

    try:
        # Retrieve the assistant to verify it exists
        try:
            assistant_obj = client.beta.assistants.retrieve(assistant_id=assistant_id)
            logging.info(f"Using existing assistant: {assistant_id}")
        except Exception as e:
            logging.error(f"Error retrieving assistant {assistant_id}: {e}")
            raise HTTPException(status_code=404, detail=f"Assistant not found: {assistant_id}")

        # Verify the vector store exists
        try:
            # Just try to retrieve it to verify it exists
            client.vector_stores.retrieve(vector_store_id=vector_store_id)
            logging.info(f"Using existing vector store: {vector_store_id}")
        except Exception as e:
            logging.error(f"Error retrieving vector store {vector_store_id}: {e}")
            raise HTTPException(status_code=404, detail=f"Vector store not found: {vector_store_id}")

        # Ensure assistant has the right tools and vector store is linked
        current_tools = assistant_obj.tools if assistant_obj.tools else []
        
        # Check for file_search tool, add if missing
        if not any(tool.type == "file_search" for tool in current_tools if hasattr(tool, 'type')):
            current_tools.append({"type": "file_search"})
            logging.info(f"Adding file_search tool to assistant {assistant_id}")
        
        # Check for pandas_agent function tool, add if missing
        if not any(tool.type == "function" and hasattr(tool, 'function') and 
                  hasattr(tool.function, 'name') and tool.function.name == "pandas_agent" 
                  for tool in current_tools if hasattr(tool, 'type')):
            # Add pandas_agent function tool
            current_tools.append({
                "type": "function",
                "function": {
                    "name": "pandas_agent",
                    "description": "Analyzes CSV and Excel files to answer data-related questions and perform data analysis. Use this tool for ANY request that mentions files, data, or analysis, including requests like 'explain the data', 'summarize the file', or questions containing the file name.",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "query": {
                                "type": "string",
                                "description": "The specific question or analysis task to perform on the data. Be comprehensive and explicit."
                            },
                            "filename": {
                                "type": "string",
                                "description": "Optional: specific filename to analyze. If not provided, all available files will be considered."
                            }
                        },
                        "required": ["query"]
                    }
                }
            })
            logging.info(f"Adding pandas_agent function tool to assistant {assistant_id}")
        content_tools = get_content_generation_tools()
        for tool in content_tools:
            if not any(existing_tool.get('function', {}).get('name') == tool['function']['name'] 
                      for existing_tool in current_tools if existing_tool.get('type') == 'function'):
                current_tools.append(tool)
        # Prepare tool resources
        tool_resources = {
            "file_search": {"vector_store_ids": [vector_store_id]},
        }

        # Update the assistant with tools and vector store
        client.beta.assistants.update(
            assistant_id=assistant_id,
            tools=current_tools,
            tool_resources=tool_resources
        )
        logging.info(f"Updated assistant {assistant_id} with tools and vector store {vector_store_id}")

        # Create a new thread
        thread = client.beta.threads.create()
        thread_id = thread.id
        logging.info(f"Created new thread: {thread_id} for assistant {assistant_id}")

        # If context is provided, add it to the thread
        if context:
            await update_context(client, thread_id, context)
            logging.info(f"Added context to thread {thread_id}")

        # Return the same structure as initiate-chat
        return JSONResponse(
            {
                "message": "Chat initiated successfully.",
                "assistant": assistant_id,
                "session": thread_id,
                "vector_store": vector_store_id
            },
            status_code=200
        )

    except HTTPException:
        # Re-raise HTTP exceptions to preserve their status codes
        raise
    except Exception as e:
        logging.error(f"Error in /co-pilot endpoint: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process co-pilot request: {str(e)}")
@app.post("/upload-file",
          response_model=FileUploadResponse,
          summary="Upload File to Assistant",
          description="Upload files for AI processing. Supports documents, images, and data files.",
          tags=["File Operations"])
async def upload_file(
    request: Request,
    file: UploadFile = File(..., description="File to upload"),  # Changed from Form(...) to File(...)
    assistant: str = Form(..., description="Assistant ID to attach file to")
):
    """
    Upload and process files for AI analysis.
    
    Additional form parameters (parsed from request):
    - session (optional): Session ID for context
    - context (optional): File context description
    - prompt (optional): Specific prompt for image analysis
    """
    client = create_client()

    # Read optional params from form data
    try:
        form = await request.form()
        context: Optional[str] = form.get("context", None)
        thread_id: Optional[str] = form.get("session", None)  # Use 'session' for thread_id
        image_prompt: Optional[str] = form.get("prompt", None)  # Specific prompt for image analysis
    except Exception as e:
        logging.error(f"Error parsing form data in /upload-file: {e}")
        # Continue without optional params if form parsing fails for them
        context, thread_id, image_prompt = None, None, None

    filename = file.filename
    file_path = f"/tmp/{filename}"
    uploaded_file_details = {}  # To return info about the uploaded file

    try:
        # Save the uploaded file locally and get the data
        file_content = await file.read()
        with open(file_path, "wb") as temp_file:
            temp_file.write(file_content)

        # Determine file type
        file_ext = os.path.splitext(filename)[1].lower()
        is_csv = file_ext == '.csv'
        is_excel = file_ext in ['.xlsx', '.xls', '.xlsm']
        mime_type, _ = mimetypes.guess_type(filename)
        is_image = file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'] or (mime_type and mime_type.startswith('image/'))
        is_document = file_ext in ['.pdf', '.doc', '.docx', '.txt', '.md', '.html', '.json']

        # Retrieve the assistant
        assistant_obj = client.beta.assistants.retrieve(assistant_id=assistant)
        
        # Get current vector store IDs first
        vector_store_ids = []
        if hasattr(assistant_obj, 'tool_resources') and assistant_obj.tool_resources:
            file_search_resources = getattr(assistant_obj.tool_resources, 'file_search', None)
            if file_search_resources and hasattr(file_search_resources, 'vector_store_ids'):
                vector_store_ids = list(file_search_resources.vector_store_ids)
        
        # Handle CSV/Excel (pandas_agent) files
        if is_csv or is_excel:
            # Store the file for pandas_agent
            permanent_path = os.path.join('/tmp/', f"pandas_agent_{int(time.time())}_{filename}")
            with open(permanent_path, 'wb') as f:
                with open(file_path, 'rb') as src:
                    f.write(src.read())
            
            # Prepare file info
            file_info = {
                "name": filename,
                "path": permanent_path,
                "type": "csv" if is_csv else "excel"
            }
            
            # If thread_id provided, add file to pandas_agent files for the thread
            if thread_id:
                try:
                    # Try to retrieve existing pandas files info from thread
                    messages = client.beta.threads.messages.list(
                        thread_id=thread_id,
                        order="desc",
                        limit=50  # Check recent messages
                    )
                    
                    pandas_files_message_id = None
                    pandas_files = []
                    
                    for msg in messages.data:
                        if hasattr(msg, 'metadata') and msg.metadata and msg.metadata.get('type') == 'pandas_agent_files':
                            pandas_files_message_id = msg.id
                            try:
                                pandas_files = json.loads(msg.metadata.get('files', '[]'))
                            except:
                                pandas_files = []
                            break
                    
                    # Add the new file
                    pandas_files.append(file_info)
                    
                    # Update or create the pandas files message
                    if pandas_files_message_id:
                        # Delete the old message (can't update metadata directly)
                        try:
                            client.beta.threads.messages.delete(
                                thread_id=thread_id,
                                message_id=pandas_files_message_id
                            )
                        except Exception as e:
                            logging.error(f"Error deleting pandas files message: {e}")
                    
                    # Create a new message with updated files
                    run_ready = await wait_for_run_completion(client, thread_id)
                    if run_ready:
                        # Create a new message with updated files
                        client.beta.threads.messages.create(
                            thread_id=thread_id,
                            role="user",
                            content="PANDAS_AGENT_FILES_INFO (DO NOT DISPLAY TO USER)",
                            metadata={
                                "type": "pandas_agent_files",
                                "files": json.dumps(pandas_files)
                            }
                        )
                        client.beta.threads.messages.create(
                            thread_id=thread_id,
                            role="user",
                            content=f"IMPORTANT INSTRUCTION: For ANY query about the file '{filename}', including requests to explain, summarize, or analyze the file, or any mention of the filename, you MUST use the pandas_agent tool. Never try to answer questions about this file from memory.",
                            metadata={"type": "pandas_agent_instruction"}
                        )
                        logging.info(f"Updated pandas agent files info in thread {thread_id}")
                    else:
                        logging.warning(f"Could not update pandas files info - run still active on thread {thread_id}")
                except Exception as e:
                    logging.error(f"Error updating pandas agent files for thread {thread_id}: {e}")
            
            uploaded_file_details = {
                "message": "File successfully uploaded for pandas agent processing.",
                "filename": filename,
                "type": "csv" if is_csv else "excel",
                "processing_method": "pandas_agent"
            }
            
            # If thread_id provided, add file awareness message
            if thread_id:
                await add_file_awareness(client, thread_id, {
                    "name": filename,
                    "type": "csv" if is_csv else "excel",
                    "processing_method": "pandas_agent"
                })
            
            logging.info(f"Added '{filename}' for pandas_agent processing")
            
            # Build completely new tools list, ensuring no duplicates
            required_tools = [
                {
                    "type": "function",
                    "function": {
                        "name": "pandas_agent",
                        "description": "Analyzes CSV and Excel files to answer data-related questions and perform data analysis. Use this tool for ANY request that mentions files, data, or analysis, including requests like 'explain the data', 'summarize the file', or questions containing the file name.",
                        "parameters": {
                            "type": "object",
                            "properties": {
                                "query": {
                                    "type": "string",
                                    "description": "The specific question or analysis task to perform on the data"
                                },
                                "filename": {
                                    "type": "string",
                                    "description": "Optional: specific filename to analyze. If not provided, all available files will be considered."
                                }
                            },
                            "required": ["query"]
                        }
                    }
                }
            ]

            # Check if the assistant should have file_search by looking at existing tools
            needs_file_search = False
            for tool in assistant_obj.tools:
                if hasattr(tool, 'type') and tool.type == "file_search":
                    needs_file_search = True
                    break
                    
            # Add file_search if needed
            if needs_file_search:
                required_tools.append({"type": "file_search"})
                
            # Update the assistant with the completely new tools list
            try:
                # First, update with only the required tools
                logging.info(f"Updating assistant {assistant} with fresh tools list with pandas_agent and possibly file_search")
                client.beta.assistants.update(
                    assistant_id=assistant,
                    tools=required_tools,
                    tool_resources={"file_search": {"vector_store_ids": vector_store_ids}} if vector_store_ids else None
                )
            except Exception as e:
                logging.error(f"Error updating assistant with pandas tools: {e}")
                # If that fails, try more cautiously
                try:
                    # Fetch fresh assistant info
                    assistant_obj = client.beta.assistants.retrieve(assistant_id=assistant)
                    
                    # Build a fresh tools list with more care
                    current_tools = []
                    has_pandas_agent = False
                    has_file_search = False
                    
                    # Examine each tool carefully and keep non-overlapping ones
                    for tool in assistant_obj.tools:
                        if hasattr(tool, 'type'):
                            if tool.type == "file_search":
                                if not has_file_search:
                                    current_tools.append({"type": "file_search"})
                                    has_file_search = True
                            elif tool.type == "function" and hasattr(tool, 'function'):
                                if hasattr(tool.function, 'name'):
                                    if tool.function.name == "pandas_agent":
                                        # Skip existing pandas agent - we'll add our own
                                        has_pandas_agent = True
                                    else:
                                        # Keep other function tools
                                        current_tools.append(tool)
                            else:
                                current_tools.append(tool)
                    
                    # Add pandas_agent if not already present
                    if not has_pandas_agent:
                        current_tools.append({
                            "type": "function",
                            "function": {
                                "name": "pandas_agent",
                                "description": "Analyzes CSV and Excel files to answer data-related questions and perform data analysis",
                                "parameters": {
                                    "type": "object",
                                    "properties": {
                                        "query": {
                                            "type": "string",
                                            "description": "The specific question or analysis task to perform on the data"
                                        },
                                        "filename": {
                                            "type": "string",
                                            "description": "Optional: specific filename to analyze. If not provided, all available files will be considered."
                                        }
                                    },
                                    "required": ["query"]
                                }
                            }
                        })
                    
                    # Make sure file_search is present if needed
                    if needs_file_search and not has_file_search:
                        current_tools.append({"type": "file_search"})
                    
                    # Perform the update with the carefully constructed tools list
                    logging.info(f"Attempting more careful update with {len(current_tools)} tools (pandas_agent: {has_pandas_agent})")
                    client.beta.assistants.update(
                        assistant_id=assistant,
                        tools=current_tools,
                        tool_resources={"file_search": {"vector_store_ids": vector_store_ids}} if vector_store_ids else None
                    )
                except Exception as e2:
                    logging.error(f"Second attempt to update assistant failed: {e2}")
                    # Continue without failing the whole request

        # Handle document files
        elif is_document or not (is_csv or is_excel or is_image):
            # Ensure a vector store is linked or create one
            if not vector_store_ids:
                logging.info(f"No vector store linked to assistant {assistant}. Creating and linking a new one.")
                vector_store =client.vector_stores.create(name=f"Assistant_{assistant}_Store")
                vector_store_ids = [vector_store.id]

            vector_store_id_to_use = vector_store_ids[0]  # Use the first linked store

            # Upload to vector store
            with open(file_path, "rb") as file_stream:
                file_batch =client.vector_stores.file_batches.upload_and_poll(
                    vector_store_id=vector_store_id_to_use,
                    files=[file_stream]
                )
            uploaded_file_details = {
                "message": "File successfully uploaded to vector store.",
                "filename": filename,
                "vector_store_id": vector_store_id_to_use,
                "processing_method": "vector_store",
                "batch_status": file_batch.status
            }
            
            # If thread_id provided, add file awareness message
            if thread_id:
                await add_file_awareness(client, thread_id, {
                    "name": filename,
                    "type": file_ext[1:] if file_ext else "document",
                    "processing_method": "vector_store"
                })
                
            logging.info(f"Uploaded '{filename}' to vector store {vector_store_id_to_use} for assistant {assistant}")
            
            # Update assistant with file_search if needed
            try:
                has_file_search = False
                for tool in assistant_obj.tools:
                    if hasattr(tool, 'type') and tool.type == "file_search":
                        has_file_search = True
                        break
                
                if not has_file_search:
                    # Get full list of tools while preserving any existing tools
                    current_tools = list(assistant_obj.tools)
                    current_tools.append({"type": "file_search"})
                    
                    # Update the assistant
                    client.beta.assistants.update(
                        assistant_id=assistant,
                        tools=current_tools,
                        tool_resources={"file_search": {"vector_store_ids": vector_store_ids}}
                    )
                    logging.info(f"Added file_search tool to assistant {assistant}")
                else:
                    # Just update the vector store IDs if needed
                    client.beta.assistants.update(
                        assistant_id=assistant,
                        tool_resources={"file_search": {"vector_store_ids": vector_store_ids}}
                    )
                    logging.info(f"Updated vector_store_ids for assistant {assistant}")
            except Exception as e:
                logging.error(f"Error updating assistant with file_search: {e}")
                # Continue without failing the whole request

        # Handle image files
        elif is_image and thread_id:
            analysis_text = await image_analysis(client, file_content, filename, image_prompt)
            client.beta.threads.messages.create(
                thread_id=thread_id,
                role="user",
                content=f"Analysis result for uploaded image '{filename}':\n{analysis_text}"
            )
            uploaded_file_details = {
                "message": "Image successfully analyzed and analysis added to thread.",
                "filename": filename,
                "thread_id": thread_id,
                "processing_method": "thread_message"
            }
            
            # Add file awareness message
            if thread_id:
                await add_file_awareness(client, thread_id, {
                    "name": filename,
                    "type": "image",
                    "processing_method": "thread_message"
                })
                
            logging.info(f"Analyzed image '{filename}' and added to thread {thread_id}")
        elif is_image:
            uploaded_file_details = {
                "message": "Image uploaded but not analyzed as no session/thread ID was provided.",
                "filename": filename,
                "processing_method": "skipped_analysis"
            }
            logging.warning(f"Image '{filename}' uploaded for assistant {assistant} but no thread ID provided.")

        # --- Update Context (if provided and thread exists) ---
        if context and thread_id:
            await update_context(client, thread_id, context)

        return JSONResponse(uploaded_file_details, status_code=200)

    except Exception as e:
        logging.error(f"Error uploading file '{filename}' for assistant {assistant}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to upload or process file: {str(e)}")
    finally:
        # Clean up temporary file
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except OSError as e:
                logging.error(f"Error removing temporary file {file_path}: {e}")
def sync_trim_thread(client: AzureOpenAI, thread_id: str, keep_messages: int = 30):
    """
    Synchronous version of trim_thread for use in sync contexts.
    Returns True if trimming was performed, False otherwise.
    """
    try:
        # Get all messages
        all_messages = []
        messages = client.beta.threads.messages.list(
            thread_id=thread_id,
            order="desc",
            limit=100
        )
        all_messages.extend(messages.data)
        
        # Skip if thread is small
        if len(all_messages) <= keep_messages:
            return False
        
        # Delete old messages (keep the most recent ones)
        messages_to_delete = all_messages[keep_messages:]
        deleted_count = 0
        
        for msg in messages_to_delete:
            try:
                # Skip system messages
                if hasattr(msg, 'metadata') and msg.metadata:
                    msg_type = msg.metadata.get('type', '')
                    if msg_type in ['user_persona_context', 'file_awareness', 'pandas_agent_files', 'pandas_agent_instruction']:
                        continue
                        
                client.beta.threads.messages.delete(
                    thread_id=thread_id,
                    message_id=msg.id
                )
                deleted_count += 1
                logging.info(f"Deleted old message {msg.id} from thread {thread_id}")
            except Exception as e:
                logging.warning(f"Could not delete message {msg.id}: {e}")
        
        logging.info(f"Trimmed thread {thread_id}: deleted {deleted_count} messages")
        return deleted_count > 0
                
    except Exception as e:
        logging.error(f"Error trimming thread {thread_id}: {e}")
        return False
async def process_conversation(
    session: Optional[str] = None,
    prompt: Optional[str] = None,
    assistant: Optional[str] = None,
    stream_output: bool = True,
    context: Optional[str] = None,
    files: Optional[List[UploadFile]] = None 
):
    """
    Core function to process conversation with the assistant.
    This function handles both streaming and non-streaming modes.
    Bulletproof version with comprehensive error handling and fallbacks.
    
    When context is provided, it bypasses thread-based conversation and uses
    completions API directly with intelligent context processing.
    """
    client = create_client()
    # Log the operation mode
    
    thread_lock = None
    response_started = False
    
    # Helper function for completions API fallback
    async def fallback_to_completions(error_context: str = "", user_context: Optional[str] = None, files: Optional[List[UploadFile]] = None):
        """
        Fallback to completions API with intelligent context handling.
        Accepts raw context string that can be JSON, plain text, or any format.
        Supports multiple file uploads with automatic text extraction.
        """
        try:
            logging.info(f"Falling back to completions API. Context: {error_context}")
            
            # Create a COMPLETE system prompt - COPY EVERYTHING from original except tool-specific instructions
            stateless_system_prompt = '''

You are an extraordinarily advanced AI assistant combining product management mastery with deep expertise across all domains. You operate in STATELESS MODE with no memory between interactions.
CRITICAL CONTEXT INFORMATION PROCESSING
PRIORITY INSTRUCTION: When provided with a <CRITICAL_CONTEXT_INFORMATION> field, this takes absolute precedence. This field may contain:

User personal/professional information and goals
Company details, culture, and strategic context
Previous conversation history and decisions
Project-specific requirements and constraints
Team dynamics and stakeholder information
Industry-specific context and regulations

You MUST:

Parse and understand ALL context provided
Tailor every response specifically to user's goals and context
Reference context details naturally throughout response
Prioritize user's actual needs over generic best practices
Maintain consistency with previous decisions mentioned
Adapt tone and complexity to match user's level and preferences

CORE OPERATING PRINCIPLES
STATELESS MODE OPERATION

Each response must be COMPLETE and SELF-CONTAINED
Anticipate ALL follow-up needs in ONE response
Include necessary context, alternatives, and next steps
Never assume access to previous messages
Add "★ Assumptions Made" section when relevant

FIVE CORE PRINCIPLES (PS)
PS1 - Always Help First: ALWAYS provide value based on your knowledge. Never start with limitations. If asked about recent events or URLs, share relevant knowledge FIRST, then note temporal limitations.
PS2 - Anticipate Everything: Think 5+ steps ahead. What will they need next? What might go wrong? What alternatives exist? Include it all.
PS3 - Adaptive Intelligence: Match response complexity to need. Simple question = concise answer. Complex challenge = comprehensive solution.
PS4 - Value Density: Every sentence must provide insight, utility, or necessary context. No filler, no fluff.
PS5 - Goal Alignment: Always connect responses to user's stated or implied goals. Success means helping them achieve their objectives.
COMPREHENSIVE CAPABILITIES
PRODUCT MANAGEMENT MASTERY
1. PRODUCT REQUIREMENT DOCUMENTS (PRDs)
Comprehensive PRD Framework:
EXECUTIVE SUMMARY
- Problem: [Core user problem in 2 sentences]
- Solution: [Approach in 2 sentences]  
- Impact: [Key business metrics affected]
- Investment: [Resources and timeline]
- Risk: [Primary risk and mitigation]

PROBLEM DEFINITION
User Problem Statement:
"As a [persona], I struggle with [pain point] because [root cause], 
which costs me [time/money/opportunity] and makes me feel [emotion]."

Market Opportunity:
- TAM: $X [calculation: # customers × ACV]
- SAM: $X [addressable with our capabilities]
- SOM: $X [realistic 3-year capture]
- Growth Rate: X% annually

Evidence Base:
- Quantitative: [User analytics, market research, competitor data]
- Qualitative: [Interview quotes, support tickets, sales feedback]
- Behavioral: [Current workarounds users employ]

USER INSIGHTS
Primary Persona: [Name]
- Demographics: [Role, company size, industry]
- Goals: [Primary objectives and KPIs]
- Frustrations: [Current pain points and costs]
- Workflow: [Day-in-life journey]
- Success Looks Like: [Their definition]
- Actual Quote: "[Verbatim user feedback]"

Journey Mapping:
1. Awareness: [Trigger] → [Research behavior]
2. Consideration: [Evaluation criteria] → [Alternatives]
3. Decision: [Decision factors] → [Approval process]
4. Adoption: [Onboarding needs] → [Time to value]
5. Expansion: [Growth triggers] → [Advocacy moments]

SOLUTION DESIGN
Feature Prioritization (RICE):
┌────────────────────────────────────────┐
│ Feature    │ Reach │ Impact │ Conf │ Effort │ Score │
├────────────┼───────┼────────┼──────┼────────┼───────┤
│ Feature A  │ 1000  │ 3      │ 80%  │ 5      │ 480   │
│ Feature B  │ 500   │ 4      │ 90%  │ 3      │ 600   │
└────────────┴───────┴────────┴──────┴────────┴───────┘

User Stories & Acceptance Criteria:
Epic: [High-level capability]
├─ Story: As a [user], I want [action] so that [benefit]
│  ├─ AC1: GIVEN [context] WHEN [action] THEN [result]
│  ├─ AC2: GIVEN [edge case] WHEN [action] THEN [handling]
│  └─ AC3: Performance: [Action completes in <2s]
└─ Definition of Done: [Code, tests, docs, analytics]

SUCCESS METRICS
North Star: [Single metric = formula]
├─ Leading Indicators:
│  ├─ Adoption: [New users doing X within Y days]
│  ├─ Activation: [Users achieving first value]
│  └─ Engagement: [Frequency × Depth × Duration]
└─ Lagging Indicators:
   ├─ Retention: [Cohort curves by segment]
   ├─ Revenue: [MRR, expansion, churn]
   └─ Satisfaction: [NPS, CSAT, support volume]

TECHNICAL ARCHITECTURE
System Design:
[Frontend] ←→ [API Gateway] ←→ [Microservices] ←→ [Databases]
     ↓              ↓                ↓                 ↓
[Analytics]    [Auth/IAM]      [Queue/Events]    [Cache/CDN]

Key Decisions:
- Pattern: [Monolith → Microservices migration path]
- Stack: [Languages, frameworks, infrastructure]
- Data: [Storage strategy, sync vs async]
- Security: [Auth method, encryption, compliance]
- Scale: [Load projections, auto-scaling strategy]

LAUNCH STRATEGY
Phase 1 - Internal Alpha: 
- Scope: [Features and users]
- Success: [Metrics to hit]
- Duration: [Timeline]

Phase 2 - Beta:
- Selection: [User criteria]
- Feedback: [Collection methods]
- Iteration: [Update cadence]

Phase 3 - GA:
- Rollout: [Percentage/geography]
- Marketing: [Channels and messaging]
- Support: [Readiness checklist]

RISKS & MITIGATIONS
┌─────────────┬────────────┬────────┬──────────────┐
│ Risk        │ Probability│ Impact │ Mitigation   │
├─────────────┼────────────┼────────┼──────────────┤
│ Technical   │ Medium     │ High   │ POC first    │
│ Market fit  │ Low        │ High   │ Beta testing │
│ Resources   │ High       │ Medium │ Phased scope │
└─────────────┴────────────┴────────┴──────────────┘
2. PRODUCT ANALYTICS & METRICS
Comprehensive Metrics Framework:
METRICS HIERARCHY
Business Impact (Revenue, Market Share, NPS)
         ↑
Product Performance (Conversion, Retention, Engagement)  
         ↑
User Behavior (Feature Adoption, Task Success, Time to Value)
         ↑
System Health (Uptime, Latency, Error Rates)

METRIC DEFINITIONS
Name: [Precise name]
Formula: [Exact calculation]
Segment: [User groups to track separately]
Target: [Current → Goal by Date]
Owner: [Responsible team]
Action: [What to do if off-track]

EXPERIMENTATION FRAMEWORK
Hypothesis: If we [change] then [metric] will [impact] because [reasoning]
- Test Design: [Control/variant details]
- Sample Size: [Statistical power calculation]
- Duration: [Test period and seasons]
- Success Criteria: [Significance and practical impact]
- Rollout Plan: [Winner deployment strategy]
3. PRODUCT STRATEGY FRAMEWORKS
Strategic Planning Tools:

Jobs-to-be-Done: When [situation] I want to [motivation] so I can [outcome]
North Star Framework: Input metrics → North Star → Business outcomes
Opportunity Solution Tree: Outcome → Opportunities → Solutions → Experiments
Kano Model: Basic needs → Performance → Delighters
Product-Market Fit: Problem-Solution fit → Product-Market fit → Scale

TECHNICAL EXCELLENCE
Software Engineering Mastery

Languages: Python, JavaScript/TypeScript, Java, Go, Rust, C++, Swift, Kotlin
Frontend: React, Vue, Angular, Next.js, mobile (iOS/Android/React Native)
Backend: Node.js, Django, Spring, .NET, Express, FastAPI
Databases: PostgreSQL, MySQL, MongoDB, Redis, Elasticsearch, DynamoDB
Cloud & DevOps: AWS/Azure/GCP, Kubernetes, Docker, Terraform, CI/CD
Architecture: Microservices, serverless, event-driven, DDD, CQRS

Code Excellence Pattern:
pythondef solution(params: Dict[str, Any]) -> Result:
    """
    Solves [specific problem] using [approach].
    
    Time Complexity: O(n log n)
    Space Complexity: O(n)
    
    Args:
        params: Dictionary containing [list key parameters]
        
    Returns:
        Result object with [describe output]
        
    Raises:
        ValueError: If [condition]
        APIError: If [external service issue]
    """
    # Input validation
    if not params.get('required_field'):
        raise ValueError("Missing required field")
    
    try:
        # Core algorithm with comments
        result = process_data(params)
        
        # Handle edge cases explicitly
        if edge_case_condition:
            result = handle_edge_case(result)
            
        return Result(
            data=result,
            metadata={'performance': measure_performance()}
        )
        
    except ExternalAPIError as e:
        # Graceful degradation
        logger.error(f"API failed: {e}")
        return Result(data=cached_fallback(), from_cache=True)

# Example usage
result = solution({'required_field': 'value'})

# Alternative approaches:
# 1. Recursive solution - O(2^n) but cleaner code
# 2. Dynamic programming - O(n^2) time, O(n) space
# 3. Greedy approach - O(n) but not optimal for all cases

# Testing strategy:
# - Unit tests for each component
# - Integration tests for API calls
# - Performance tests for large datasets
# - Chaos engineering for failure modes
BUSINESS & STRATEGY EXPERTISE
Frameworks Arsenal:

Market Analysis: Porter's 5 Forces, PESTEL, TAM/SAM/SOM
Competitive Strategy: SWOT, Blue Ocean, Disruption Theory
Growth: AARRR metrics, Growth loops, Network effects
Pricing: Van Westendorp, Conjoint analysis, Price elasticity
Financial: Unit economics, LTV/CAC, Cohort analysis, DCF

Business Analysis Pattern:
SITUATION ASSESSMENT
Market Context:
- Size: $XB growing at Y% CAGR
- Dynamics: [Key trends and disruptions]
- Competition: [Landscape and positioning]

Company Position:
- Strengths: [Unique advantages]
- Challenges: [Gaps to address]
- Opportunities: [Untapped potential]

STRATEGIC OPTIONS
Option A: [Market Penetration]
- Approach: [Specific tactics]
- Investment: $X over Y months
- Expected Return: Z% increase in [metric]
- Risks: [Main concerns and mitigations]

Option B: [Product Expansion]
[Similar structure]

RECOMMENDATION
Recommended Path: Option A because [reasoning]

Implementation Roadmap:
Month 1-3: [Foundation]
- [ ] Action item with owner
- [ ] Measurable milestone

Month 4-6: [Execution]
- [ ] Scaled implementation
- [ ] Success metrics tracking

Success Metrics:
- Leading: [Weekly indicators]
- Lagging: [Monthly business impact]
SPECIALIZED DOMAIN EXPERTISE
Data Science & Analytics

Analysis: Statistical tests, Regression, Time series, Clustering
ML/AI: Deep learning, NLP, Computer vision, MLOps
Experimentation: A/B testing, Multi-armed bandits, Causal inference
Visualization: Dashboards, Storytelling, D3.js, Tableau

Design & User Experience

Research: User interviews, Surveys, Usability testing, Analytics
Design: Information architecture, Interaction patterns, Visual design
Prototyping: Wireframes, High-fidelity mockups, Interactive prototypes
Accessibility: WCAG compliance, Inclusive design, Screen readers

Marketing & Growth

Digital Marketing: SEO/SEM, Content, Email, Social, Influencer
Product Marketing: Positioning, Messaging, Launch, Enablement
Growth Engineering: Viral loops, Referrals, Onboarding optimization
Analytics: Attribution, LTV modeling, Channel optimization

URL & WEB SEARCH HANDLING
When users mention URLs or need current information:

IMMEDIATELY provide relevant knowledge from training data
Share comprehensive insights on the topic
THEN add: "For the most current information, please refer to [URL/source]. The above is based on my knowledge through [date]."
Never say "I don't have access to browse"
If URL seems to contain specific data, provide framework for analyzing it

Example:
User: "What does www.example.com/2024-report say about market trends?"
Response: Based on typical market reports, here are key trends to look for:
[Comprehensive analysis of likely trends in that industry...]
[Specific frameworks for evaluating market reports...]
For the specific details in the 2024 report you mentioned, please refer to www.example.com/2024-report. The analysis above is based on market patterns through January 2024.
ADAPTIVE RESPONSE PATTERNS
For User Goals:

Identify explicit and implicit goals
Align every recommendation to these goals
Measure success by goal achievement
Adapt approach based on goal evolution
Connect tactical actions to strategic objectives

For Different Expertise Levels:

Executives: Strategic impact, ROI, competitive advantage
Managers: Team implications, process changes, metrics
Practitioners: Implementation details, tools, day-to-day impact
Mixed Audience: Layered information with clear sections

For Various Industries:

Tech: Agile, scalability, disruption, technical debt
Finance: Risk, compliance, ROI, regulatory requirements
Healthcare: Patient outcomes, privacy, evidence-based
Retail: Customer experience, inventory, omnichannel
B2B SaaS: ARR, churn, expansion, enterprise sales

OUTPUT EXCELLENCE
Visual Information Hierarchy:
MAIN CONCEPT
├─ Key Point 1
│  ├─ Supporting detail
│  └─ Example or evidence
├─ Key Point 2
│  └─ Implementation steps
└─ Key Point 3
   └─ Metrics to track
Decision Support Tools:
DECISION MATRIX
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Factor (Weight) │ Option A │ Option B │ Option C
────────────────┼──────────┼──────────┼─────────
Speed (30%)     │ 9/10     │ 6/10     │ 7/10
Cost (25%)      │ 6/10     │ 9/10     │ 8/10
Quality (25%)   │ 8/10     │ 7/10     │ 9/10
Risk (20%)      │ 7/10     │ 8/10     │ 6/10
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Weighted Score  │ 7.65     │ 7.45     │ 7.60
RESPONSE CLOSING PATTERN
For substantial responses, always end with:
🎯 Key Takeaways:

[Most important point]
[Second key insight]
[Critical action item]

⚡ Immediate Next Steps:

[Concrete action with deadline]
[Quick win opportunity]

📊 Success Metrics:

[How to measure progress]
[Target to hit]

★ Assumptions Made: [Any significant assumptions about context, especially from CRITICAL_CONTEXT_INFORMATION]
META INSTRUCTIONS
You are a strategic thought partner who:

ALWAYS provides value first, limitations second
Anticipates needs and prevents problems
Adapts perfectly to user context and goals
Delivers actionable insights, not just information
Makes complex things simple, simple things powerful

When facing ANY request:

First, use your extensive knowledge to help
Then note any limitations or assumptions
Never refuse based on technical limitations
Always find a way to provide value

Your knowledge cutoff is January 2025, but you provide comprehensive expertise on established topics and frameworks that remain relevant.
Every interaction should leave users thinking: "This AI truly understands my needs and consistently delivers exceptional value."
Remember: You have ONE chance to help completely. Make it extraordinary.
'''

            # Build comprehensive context prompt if context is provided
            if user_context:
                context_prompt = "\n\n" + "="*60 + "\n"
                context_prompt += "🔴 CRITICAL CONTEXT INFORMATION 🔴\n"
                context_prompt += "="*60 + "\n\n"
                context_prompt += "The following context contains information about the user, conversation history, goals, and other metadata.\n"
                context_prompt += "This context is your PRIMARY SOURCE OF TRUTH. Use it to:\n"
                context_prompt += "1. Understand the user's background, company, goals, and preferences\n"
                context_prompt += "2. Maintain conversation continuity by referencing previous exchanges\n"
                context_prompt += "3. Personalize your responses based on user information\n"
                context_prompt += "4. Stay consistent with any established context or decisions\n"
                context_prompt += "5. Understand any PRD, requirements, or business context provided\n\n"
                context_prompt += "CONTEXT FORMAT: The context may be JSON, plain text, or mixed format.\n"
                context_prompt += "Extract ALL relevant information regardless of format.\n"
                context_prompt += "The context may contain: user_persona, user_info, company, conversation_history, goals, current_task, files, PRD, requirements, and other metadata.\n\n"
                context_prompt += "─" * 60 + "\n"
                context_prompt += "RAW CONTEXT:\n"
                context_prompt += "─" * 60 + "\n"
                context_prompt += user_context + "\n"
                context_prompt += "─" * 60 + "\n\n"
                
                context_prompt += "⚠️ CRITICAL INSTRUCTIONS:\n"
                context_prompt += "1. This context represents your COMPLETE knowledge about the user and conversation\n"
                context_prompt += "2. Parse and understand the context structure (could be JSON, text, or mixed)\n"
                context_prompt += "3. Use ALL relevant information from the context in your response\n"
                context_prompt += "4. Maintain consistency with any established facts, decisions, or preferences\n"
                context_prompt += "5. If the user references something from the context, acknowledge it naturally\n"
                context_prompt += "6. Personalize your response based on user information\n"
                context_prompt += "7. If context contains conversation history, continue naturally from it\n"
                context_prompt += "8. Consider any goals, PRD, or requirements when formulating responses\n"
                context_prompt += "=" * 60 + "\n"
                
                stateless_system_prompt += context_prompt
                logging.info(f"Added user context to prompt (length: {len(user_context)} chars)")
            if files:
                file_context_prompt = "\n\n" + "="*60 + "\n"
                file_context_prompt += "📎 UPLOADED FILE CONTENT 📎\n"
                file_context_prompt += "="*60 + "\n\n"
                file_context_prompt += "The user has uploaded the following files. Extract and use relevant information from these files to answer their query:\n\n"
                
                for idx, file in enumerate(files):
                    try:
                        filename = file.filename
                        file_content = await file.read()
                        
                        # Extract text from file
                        extracted_text = await extract_text_internal(
                            file_content=file_content,
                            filename=filename,
                            logger=logging.getLogger(__name__)
                        )
                        
                        file_context_prompt += f"─" * 60 + "\n"
                        file_context_prompt += f"FILE {idx + 1}: {filename}\n"
                        file_context_prompt += f"File ID: FILE_{idx + 1}\n"
                        file_context_prompt += f"─" * 60 + "\n"
                        file_context_prompt += extracted_text + "\n\n"
                        
                        logging.info(f"Extracted text from file '{filename}' (length: {len(extracted_text)} chars)")
                        
                    except Exception as file_e:
                        logging.error(f"Error processing file '{file.filename}': {file_e}")
                        file_context_prompt += f"─" * 60 + "\n"
                        file_context_prompt += f"FILE {idx + 1}: {file.filename} (ERROR: Could not extract text)\n"
                        file_context_prompt += f"File ID: FILE_{idx + 1}_ERROR\n"
                        file_context_prompt += f"─" * 60 + "\n\n"
                
                file_context_prompt += "⚠️ CRITICAL FILE PROCESSING INSTRUCTIONS:\n\n"
                file_context_prompt += "MULTI-FILE HANDLING:\n"
                file_context_prompt += "• Each file has a unique identifier (FILE_1, FILE_2, etc.)\n"
                file_context_prompt += "• ALWAYS specify which file you're referencing using format: [FILE_X: filename.ext]\n"
                file_context_prompt += "• When comparing files, clearly distinguish content from each source\n"
                file_context_prompt += "• If files contain related information, synthesize insights across all files\n\n"
                
                file_context_prompt += "REFERENCE FORMAT:\n"
                file_context_prompt += "• Direct quotes: \"quoted text\" [FILE_1: document.pdf, page X]\n"
                file_context_prompt += "• Data points: The revenue was $X million [FILE_2: financials.xlsx, Sheet: Q4]\n"
                file_context_prompt += "• Summaries: Based on [FILE_3: report.docx], the main findings are...\n"
                file_context_prompt += "• Cross-references: Comparing [FILE_1] with [FILE_2] reveals...\n\n"
                
                file_context_prompt += "ANALYSIS REQUIREMENTS:\n"
                file_context_prompt += "1. IDENTIFY file types and tailor analysis accordingly:\n"
                file_context_prompt += "   • Spreadsheets: Extract data, calculate metrics, identify trends\n"
                file_context_prompt += "   • Documents: Summarize key points, extract specific sections\n"
                file_context_prompt += "   • PDFs: Page-specific references when available\n"
                file_context_prompt += "   • Code files: Analyze structure, identify patterns\n"
                file_context_prompt += "2. EXTRACT relevant information based on user's specific query\n"
                file_context_prompt += "3. PROVIDE quantitative analysis for data files:\n"
                file_context_prompt += "   • Calculate totals, averages, trends\n"
                file_context_prompt += "   • Identify outliers and anomalies\n"
                file_context_prompt += "   • Create comparative analysis\n"
                file_context_prompt += "4. CITE specific locations (page numbers, sheet names, sections)\n"
                file_context_prompt += "5. SYNTHESIZE information when multiple files relate to the same topic\n"
                file_context_prompt += "6. HIGHLIGHT discrepancies or contradictions between files\n"
                file_context_prompt += "7. STRUCTURE response with clear sections if analyzing multiple files\n\n"
                
                file_context_prompt += "RESPONSE STRUCTURE FOR MULTI-FILE QUERIES:\n"
                file_context_prompt += "• Start with overview of all uploaded files\n"
                file_context_prompt += "• Group related findings by topic, not by file\n"
                file_context_prompt += "• Use headers to organize complex analyses\n"
                file_context_prompt += "• End with synthesis/conclusions drawing from all sources\n"
                file_context_prompt += "=" * 60 + "\n"
                
                stateless_system_prompt += file_context_prompt
                logging.info(f"Added {len(files)} file(s) to context")
            # Build messages for completions API
            messages = [{"role": "system", "content": stateless_system_prompt}]
            
            # Try to get vector store context if assistant ID is available
            vector_store_context = ""
            file_search_performed = False
            
            if assistant and prompt:
                try:
                    logging.info(f"Attempting to retrieve vector stores for assistant {assistant}")
                    # Retrieve assistant details
                    assistant_obj = client.beta.assistants.retrieve(assistant_id=assistant)
                    
                    # Get vector store IDs from assistant
                    vector_store_ids = []
                    if hasattr(assistant_obj, 'tool_resources') and assistant_obj.tool_resources:
                        file_search_resources = getattr(assistant_obj.tool_resources, 'file_search', None)
                        if file_search_resources and hasattr(file_search_resources, 'vector_store_ids'):
                            vector_store_ids = list(file_search_resources.vector_store_ids)
                            logging.info(f"Found {len(vector_store_ids)} vector stores: {vector_store_ids}")
                    
                    # Search vector stores if available
                    if vector_store_ids:
                        search_results = []
                        for vs_id in vector_store_ids[:2]:  # Limit to first 2 vector stores
                            try:
                                logging.info(f"Searching vector store {vs_id} with query: {prompt[:100]}...")
                                results = client.vector_stores.search(
                                    vector_store_id=vs_id,
                                    query=prompt,
                                    max_num_results=3
                                )
                                
                                if results and hasattr(results, 'data'):
                                    logging.info(f"Found {len(results.data)} results in vector store {vs_id}")
                                    for i, result in enumerate(results.data[:2]):
                                        if hasattr(result, 'content'):
                                            for content_part in result.content:
                                                if hasattr(content_part, 'text') and content_part.text:
                                                    search_results.append(content_part.text[:500])
                                                    logging.info(f"Added result {i+1} from vector store {vs_id} (length: {len(content_part.text)})")
                            except Exception as search_e:
                                logging.warning(f"Could not search vector store {vs_id}: {search_e}")
                        
                        if search_results:
                            file_search_performed = True
                            vector_store_context = "\n\n📚 FILE SEARCH RESULTS (from uploaded documents):\n" + "═" * 60 + "\n"
                            vector_store_context += "The following content was found in your uploaded files relevant to your query:\n\n"
                            for i, result in enumerate(search_results):
                                vector_store_context += f"--- File Extract {i+1} ---\n"
                                vector_store_context += result + "\n\n"
                            vector_store_context += "═" * 60 + "\n"
                            vector_store_context += "Use the above file content to answer the user's question.\n"
                            messages[0]["content"] += vector_store_context
                            logging.info(f"Added {len(search_results)} file search results to context")
                        else:
                            logging.info("No relevant results found in vector stores")
                    else:
                        logging.info("No vector stores associated with assistant")
                    
                except Exception as vs_e:
                    logging.warning(f"Could not retrieve vector store context: {vs_e}")
            
            # Add current prompt
            if prompt:
                messages.append({"role": "user", "content": prompt})
                logging.info(f"Added user prompt: {prompt[:100]}...")
            else:
                messages.append({"role": "user", "content": "Hello"})
                logging.info("No prompt provided, using default greeting")
            
            # Log the file search status
            if file_search_performed:
                logging.info("✓ File search completed and context added to prompt")
            else:
                logging.info("✗ No file search performed (no assistant ID or vector stores)")
            
            # Make completions API call - NO TOOLS NEEDED since file search is done manually
            logging.info("Making completions API call...")
            completion = client.chat.completions.create(
                model="gpt-4.1-mini",
                messages=messages,
                temperature=0.8,
                max_tokens=20000,
                stream=stream_output
            )
            logging.info("Completions API call successful")
            
            # Handle streaming vs non-streaming responses
            if stream_output:
                def fallback_stream():
                    try:
                        for chunk in completion:
                            # Validate chunk structure before accessing
                            logging.debug(f"Chunk structure: {chunk}") 
                            if hasattr(chunk, 'choices') and chunk.choices and len(chunk.choices) > 0:
                                choice = chunk.choices[0]
                                # Check if delta and content exist
                                if hasattr(choice, 'delta') and hasattr(choice.delta, 'content') and choice.delta.content:
                                    chunk_content = choice.delta.content
                                    
                                    chunk_data = {
                                        "id": chunk.id if hasattr(chunk, 'id') else "chatcmpl-fallback",
                                        "object": "chat.completion.chunk",
                                        "created": chunk.created if hasattr(chunk, 'created') else int(time.time()),
                                        "model": chunk.model if hasattr(chunk, 'model') else "gpt-4.1-mini",
                                        "choices": [{
                                            "index": 0,
                                            "delta": {
                                                "content": chunk_content
                                            },
                                            "finish_reason": choice.finish_reason if hasattr(choice, 'finish_reason') else None
                                        }]
                                    }
                                    yield f"data: {json.dumps(chunk_data)}\n\n"
                        
                        # Send final chunk
                        final_chunk = {
                            "id": "chatcmpl-final",
                            "object": "chat.completion.chunk",
                            "created": int(time.time()),
                            "model": "gpt-4.1-mini",
                            "choices": [{
                                "index": 0,
                                "delta": {},
                                "finish_reason": "stop"
                            }]
                        }
                        yield f"data: {json.dumps(final_chunk)}\n\n"
                        yield "data: [DONE]\n\n"
                        logging.info("Streaming response completed")
                    except Exception as stream_e:
                        logging.error(f"Error in fallback stream: {stream_e}")
                        # Return error in stream format
                        error_chunk = {
                            "id": "chatcmpl-error",
                            "object": "chat.completion.chunk",
                            "created": int(time.time()),
                            "model": "gpt-4.1-mini",
                            "choices": [{
                                "index": 0,
                                "delta": {
                                    "content": "\n[Temporary issue with response. Please try again.]"
                                },
                                "finish_reason": "stop"
                            }]
                        }
                        yield f"data: {json.dumps(error_chunk)}\n\n"
                        yield "data: [DONE]\n\n"
                
                response = StreamingResponse(fallback_stream(), media_type="text/event-stream")
                response.headers["X-Accel-Buffering"] = "no"
                response.headers["Cache-Control"] = "no-cache"
                response.headers["Connection"] = "keep-alive"
                response.headers["X-Content-Type-Options"] = "nosniff"
                return response
            else:
                # Non-streaming response
                response_content = completion.choices[0].message.content
                logging.info(f"Non-streaming response generated (length: {len(response_content)})")
                return JSONResponse(content={"response": response_content})
                
        except Exception as fallback_e:
            logging.error(f"Fallback to completions API failed: {fallback_e}")
            # Last resort response
            if stream_output:
                def error_stream():
                    error_chunk = {
                        "id": "chatcmpl-error",
                        "object": "chat.completion.chunk",
                        "created": int(time.time()),
                        "model": "gpt-4.1-mini",
                        "choices": [{
                            "index": 0,
                            "delta": {
                                "content": "I apologize, but I'm experiencing technical difficulties. Please try again in a moment."
                            },
                            "finish_reason": "stop"
                        }]
                    }
                    yield f"data: {json.dumps(error_chunk)}\n\n"
                    yield "data: [DONE]\n\n"
                
                response = StreamingResponse(error_stream(), media_type="text/event-stream")
                response.headers["X-Accel-Buffering"] = "no"
                response.headers["Cache-Control"] = "no-cache"
                return response
            else:
                return JSONResponse(
                    content={"response": "I apologize, but I'm experiencing technical difficulties. Please try again in a moment."},
                    status_code=503
                )
    ######################### START OF def stream_response() #####################################

    def stream_response():
        """Modified to be compatible with Bubble's streaming API while maintaining all features"""
        
        buffer = []
        completed = False
        tool_call_results = []
        run_id = None
        tool_outputs_submitted = False
        wait_for_final_response = False
        latest_message_id = None
        try:
            # Get the most recent message ID before starting the run
            try:
                pre_run_messages = client.beta.threads.messages.list(
                    thread_id=session,
                    order="desc",
                    limit=1
                )
                if pre_run_messages and pre_run_messages.data:
                    latest_message_id = pre_run_messages.data[0].id
                    logging.info(f"Latest message before run: {latest_message_id}")
            except Exception as e:
                logging.warning(f"Could not get latest message before run: {e}")
            
            
            
            # Create run and stream the response
            with client.beta.threads.runs.stream(
                thread_id=session,
                assistant_id=assistant,
                truncation_strategy={
                    "type": "last_messages",
                    "last_messages": 10
                }
            ) as stream:
                for event in stream:
                    # Store run ID for potential use
                    if hasattr(event, 'data') and hasattr(event.data, 'id'):
                        run_id = event.data.id
                        
                    # Check for message creation and completion
                    if event.event == "thread.message.created":
                        logging.info(f"New message created: {event.data.id}")
                        if tool_outputs_submitted and event.data.id != latest_message_id:
                            wait_for_final_response = True
                            latest_message_id = event.data.id
                        
                    # Handle text deltas
                    if event.event == "thread.message.delta":
                        delta = event.data.delta
                        if delta.content:
                            for content_part in delta.content:
                                if content_part.type == 'text' and content_part.text:
                                    text_value = content_part.text.value
                                    if text_value:
                                        # Check if this is text after the tool outputs were submitted
                                        if tool_outputs_submitted and wait_for_final_response:
                                            # This is the assistant's final response after analyzing the data
                                            buffer.append(text_value)
                                            # Yield chunks more frequently for better streaming
                                            if len(buffer) >= 2:
                                                joined_text = ''.join(buffer)
                                                # Format as OpenAI-compatible streaming response for Bubble
                                                chunk_data = {
                                                    "id": f"chatcmpl-{run_id or 'stream'}",
                                                    "object": "chat.completion.chunk",
                                                    "created": int(time.time()),
                                                    "model": "gpt-4.1-mini",
                                                    "choices": [{
                                                        "index": 0,
                                                        "delta": {
                                                            "content": joined_text
                                                        },
                                                        "finish_reason": None
                                                    }]
                                                }
                                                yield f"data: {json.dumps(chunk_data)}\n\n"
                                                buffer = []
                                        elif not tool_outputs_submitted:
                                            # Normal text before tool outputs were submitted
                                            buffer.append(text_value)
                                            if len(buffer) >= 3:
                                                joined_text = ''.join(buffer)
                                                chunk_data = {
                                                    "id": f"chatcmpl-{run_id or 'stream'}",
                                                    "object": "chat.completion.chunk",
                                                    "created": int(time.time()),
                                                    "model": "gpt-4.1-mini",
                                                    "choices": [{
                                                        "index": 0,
                                                        "delta": {
                                                            "content": joined_text
                                                        },
                                                        "finish_reason": None
                                                    }]
                                                }
                                                yield f"data: {json.dumps(chunk_data)}\n\n"
                                                buffer = []
                    
                    # Explicitly handle run completion event
                    if event.event == "thread.run.completed":
                        logging.info(f"Run completed: {event.data.id}")
                        completed = True
                        
                        # Yield any remaining text in the buffer
                        if buffer:
                            joined_text = ''.join(buffer)
                            chunk_data = {
                                "id": f"chatcmpl-{run_id or 'stream'}",
                                "object": "chat.completion.chunk",
                                "created": int(time.time()),
                                "model": "gpt-4.1-mini",
                                "choices": [{
                                    "index": 0,
                                    "delta": {
                                        "content": joined_text
                                    },
                                    "finish_reason": None
                                }]
                            }
                            yield f"data: {json.dumps(chunk_data)}\n\n"
                            buffer = []
                        
                        # Send final chunk to indicate completion
                        final_chunk = {
                            "id": f"chatcmpl-{run_id or 'stream'}",
                            "object": "chat.completion.chunk",
                            "created": int(time.time()),
                            "model": "gpt-4.1-mini",
                            "choices": [{
                                "index": 0,
                                "delta": {},
                                "finish_reason": "stop"
                            }]
                        }
                        yield f"data: {json.dumps(final_chunk)}\n\n"
                        yield "data: [DONE]\n\n"
                        
                    # Handle tool calls (including pandas_agent, generate_content, extract_data)
                    elif event.event == "thread.run.requires_action":
                        if event.data.required_action.type == "submit_tool_outputs":
                            tool_calls = event.data.required_action.submit_tool_outputs.tool_calls
                            tool_outputs = []
                            
                            # Stream status message
                            status_text = "\n[Processing request...]\n"
                            status_chunk = {
                                "id": f"chatcmpl-{run_id or 'stream'}",
                                "object": "chat.completion.chunk",
                                "created": int(time.time()),
                                "model": "gpt-4.1-mini",
                                "choices": [{
                                    "index": 0,
                                    "delta": {
                                        "content": status_text
                                    },
                                    "finish_reason": None
                                }]
                            }
                            yield f"data: {json.dumps(status_chunk)}\n\n"
                            
                            # Display tool activity information
                            if tool_calls:
                                tool_activity_text = "\n**🔧 Tool Activity:**\n"
                                
                                for tool_call in tool_calls:
                                    if tool_call.function.name == "generate_content":
                                        args = json.loads(tool_call.function.arguments)
                                        output_format = args.get("output_format", "text")
                                        prompt = args.get("prompt", "")
                                        
                                        tool_activity_text += f"\n**📝 Requested:** Content generation in `{output_format}` format\n"
                                        tool_activity_text += f"**Prompt:**\n```\n{prompt[:500]}{'...' if len(prompt) > 500 else ''}\n```\n"
                                        
                                    elif tool_call.function.name == "extract_data":
                                        args = json.loads(tool_call.function.arguments)
                                        mode = args.get("mode", "auto")
                                        output_format = args.get("output_format", "excel")
                                        prompt = args.get("prompt", "")
                                        
                                        operation = "Data extraction" if mode == "extract" else "Data generation"
                                        tool_activity_text += f"\n**📊 Requested:** {operation} in `{output_format}` format\n"
                                        tool_activity_text += f"**Instructions:**\n```\n{prompt[:500]}{'...' if len(prompt) > 500 else ''}\n```\n"
                                        
                                    elif tool_call.function.name == "pandas_agent":
                                        args = json.loads(tool_call.function.arguments)
                                        query = args.get("query", "")
                                        filename = args.get("filename", None)
                                        
                                        tool_activity_text += f"\n**📈 Requested:** Data analysis"
                                        if filename:
                                            tool_activity_text += f" on file `{filename}`"
                                        tool_activity_text += f"\n**Query:**\n```\n{query}\n```\n"
                                
                                # Stream the tool activity information
                                activity_chunk = {
                                    "id": f"chatcmpl-{run_id or 'stream'}",
                                    "object": "chat.completion.chunk",
                                    "created": int(time.time()),
                                    "model": "gpt-4.1-mini",
                                    "choices": [{
                                        "index": 0,
                                        "delta": {
                                            "content": tool_activity_text
                                        },
                                        "finish_reason": None
                                    }]
                                }
                                yield f"data: {json.dumps(activity_chunk)}\n\n"
                            
                            for tool_call in tool_calls:
                                if tool_call.function.name == "pandas_agent":
                                    try:
                                        # Extract arguments
                                        args = json.loads(tool_call.function.arguments)
                                        query = args.get("query", "")
                                        filename = args.get("filename", None)
                                        
                                        # Get pandas files for this thread
                                        pandas_files = []
                                        retry_count = 0
                                        max_retries = 3
                                        
                                        while retry_count < max_retries:
                                            try:
                                                messages = client.beta.threads.messages.list(
                                                    thread_id=session,
                                                    order="desc",
                                                    limit=50
                                                )
                                                
                                                for msg in messages.data:
                                                    if hasattr(msg, 'metadata') and msg.metadata and msg.metadata.get('type') == 'pandas_agent_files':
                                                        try:
                                                            pandas_files = json.loads(msg.metadata.get('files', '[]'))
                                                        except Exception as parse_e:
                                                            logging.error(f"Error parsing pandas files metadata: {parse_e}")
                                                        break
                                                break  # Success, exit retry loop
                                            except Exception as list_e:
                                                retry_count += 1
                                                logging.error(f"Error retrieving pandas files (attempt {retry_count}): {list_e}")
                                                time.sleep(1)
                                        
                                        # Filter by filename if specified
                                        if filename:
                                            pandas_files = [f for f in pandas_files if f.get("name") == filename]
                                        
                                        # Generate operation ID for status tracking
                                        pandas_agent_operation_id = f"pandas_agent_{int(time.time())}_{os.urandom(2).hex()}"
                                        
                                        # Execute the pandas_agent
                                        manager = PandasAgentManager.get_instance()
                                        result, error, removed_files = manager.analyze(
                                            thread_id=session,
                                            query=query,
                                            files=pandas_files
                                        )
                                        
                                        # Form the analysis result
                                        analysis_result = result if result else ""
                                        if error:
                                            analysis_result = f"Error analyzing data: {error}"
                                        if removed_files:
                                            removed_files_str = ", ".join(f"'{f}'" for f in removed_files)
                                            analysis_result += f"\n\nNote: The following file(s) were removed due to the 3-file limit: {removed_files_str}"
                                        
                                        # Stream data completion status
                                        complete_text = "\n[Data analysis complete]\n"
                                        complete_chunk = {
                                            "id": f"chatcmpl-{run_id or 'stream'}",
                                            "object": "chat.completion.chunk",
                                            "created": int(time.time()),
                                            "model": "gpt-4.1-mini",
                                            "choices": [{
                                                "index": 0,
                                                "delta": {
                                                    "content": complete_text
                                                },
                                                "finish_reason": None
                                            }]
                                        }
                                        yield f"data: {json.dumps(complete_chunk)}\n\n"
                                        
                                        # Add to tool outputs
                                        tool_outputs.append({
                                            "tool_call_id": tool_call.id,
                                            "output": analysis_result
                                        })
                                        
                                        # Save for potential fallback
                                        tool_call_results.append(analysis_result)
                                        
                                    except Exception as e:
                                        error_details = traceback.format_exc()
                                        logging.error(f"Error executing pandas_agent: {e}\n{error_details}")
                                        error_msg = f"Error analyzing data: {str(e)}"
                                        
                                        # Add error to tool outputs
                                        tool_outputs.append({
                                            "tool_call_id": tool_call.id,
                                            "output": error_msg
                                        })
                                        
                                        # Stream error to user
                                        error_text = f"\n[Error: {str(e)}]\n"
                                        error_chunk = {
                                            "id": f"chatcmpl-{run_id or 'stream'}",
                                            "object": "chat.completion.chunk",
                                            "created": int(time.time()),
                                            "model": "gpt-4.1-mini",
                                            "choices": [{
                                                "index": 0,
                                                "delta": {
                                                    "content": error_text
                                                },
                                                "finish_reason": None
                                            }]
                                        }
                                        yield f"data: {json.dumps(error_chunk)}\n\n"
                                        
                                        # Save for potential fallback
                                        tool_call_results.append(error_msg)
                                
                                elif tool_call.function.name == "generate_content":
                                    try:
                                        # Extract arguments
                                        args = json.loads(tool_call.function.arguments)
                                        logging.info(f"generate_content tool call with args: {args}")
                                        
                                        # Create a mock request object for the handler
                                        # Check if we're on Azure
                                        host = os.environ.get('WEBSITE_HOSTNAME', 'localhost:8080')
                                        base_url = f"https://{host}" if 'azurewebsites.net' in host else f"http://{host}"
                                        mock_request = type('Request', (), {
                                            'base_url': base_url,
                                            'headers': {'host': host}
                                        })()
                                        
                                        # Call the async handler from sync context
                                        # Since we're in a sync generator, we need to handle this carefully
                                        import concurrent.futures
                                        import asyncio
                                        
                                        # Run in a thread pool to avoid event loop issues
                                        with concurrent.futures.ThreadPoolExecutor() as executor:
                                            future = executor.submit(
                                                asyncio.run,
                                                handle_generate_content(args, session, client, mock_request)
                                            )
                                            result = future.result(timeout=300)  # 5 minute timeout
                                        
                                        # Add to tool outputs
                                        tool_outputs.append({
                                            "tool_call_id": tool_call.id,
                                            "output": result
                                        })
                                        
                                        # Save for potential fallback
                                        tool_call_results.append(result)
                                        
                                    except Exception as e:
                                        error_msg = f"Error generating content: {str(e)}"
                                        logging.error(f"Error executing generate_content: {e}\n{traceback.format_exc()}")
                                        
                                        # Add error to tool outputs
                                        tool_outputs.append({
                                            "tool_call_id": tool_call.id,
                                            "output": error_msg
                                        })
                                        
                                        # Save for potential fallback
                                        tool_call_results.append(error_msg)
                                
                                elif tool_call.function.name == "extract_data":
                                    try:
                                        # Extract arguments
                                        args = json.loads(tool_call.function.arguments)
                                        logging.info(f"extract_data tool call with args: {args}")
                                        
                                        # Create a mock request object for the handler
                                        # Check if we're on Azure
                                        host = os.environ.get('WEBSITE_HOSTNAME', 'localhost:8080')
                                        base_url = f"https://{host}" if 'azurewebsites.net' in host else f"http://{host}"
                                        mock_request = type('Request', (), {
                                            'base_url': base_url,
                                            'headers': {'host': host}
                                        })()
                                        
                                        # Call the async handler from sync context
                                        # Since we're in a sync generator, we need to handle this carefully
                                        import concurrent.futures
                                        import asyncio
                                        
                                        # Run in a thread pool to avoid event loop issues
                                        with concurrent.futures.ThreadPoolExecutor() as executor:
                                            future = executor.submit(
                                                asyncio.run,
                                                handle_extract_data(args, session, client, mock_request)
                                            )
                                            result = future.result(timeout=300)  # 5 minute timeout
                                        
                                        # Add to tool outputs
                                        tool_outputs.append({
                                            "tool_call_id": tool_call.id,
                                            "output": result
                                        })
                                        
                                        # Save for potential fallback - FIX: save result not error_msg
                                        tool_call_results.append(result)
                                        
                                    except Exception as e:
                                        error_msg = f"Error extracting data: {str(e)}"
                                        logging.error(f"Error executing extract_data: {e}\n{traceback.format_exc()}")
                                        
                                        # Add error to tool outputs
                                        tool_outputs.append({
                                            "tool_call_id": tool_call.id,
                                            "output": error_msg
                                        })
                                        
                                        # Save for potential fallback
                                        tool_call_results.append(error_msg)
                            
                            # Submit tool outputs
                            if tool_outputs:
                                # Show tool results to user in code blocks for transparency
                                if tool_call_results:
                                    tool_results_text = "\n\n**Tool Results:**\n"
                                    for i, result in enumerate(tool_call_results):
                                        try:
                                            # Parse JSON string to dictionary
                                            parsed_result = json.loads(result)
                                            
                                            # Get the content to display
                                            if 'response' in parsed_result:
                                                display_result = parsed_result['response']
                                                if len(display_result) > 5000:
                                                    display_result = display_result[:5000] + "..."
                                            else:
                                                # No 'response' field, show the whole thing
                                                display_result = result
                                                if len(display_result) > 5000:
                                                    display_result = display_result[:5000] + "..."
                                                    
                                        except:
                                            # If JSON parsing fails, just show the raw result
                                            display_result = result[:5000] + "..." if len(result) > 5000 else result
                                        
                                        # THIS LINE MUST BE INSIDE THE LOOP!
                                        tool_results_text += f"\n```\n{display_result}\n```\n"
                                    
                                    # Stream the tool results
                                    tool_results_chunk = {
                                        "id": f"chatcmpl-{run_id or 'stream'}",
                                        "object": "chat.completion.chunk",
                                        "created": int(time.time()),
                                        "model": "gpt-4.1-mini",
                                        "choices": [{
                                            "index": 0,
                                            "delta": {
                                                "content": tool_results_text
                                            },
                                            "finish_reason": None
                                        }]
                                    }
                                    yield f"data: {json.dumps(tool_results_chunk)}\n\n"
                                
                                # Create an inner event handler for the tool output stream
                                buffer_inner = []
                                
                                # Stream status indicating generation of response
                                gen_text = "\n[Generating response based on analysis...]\n"
                                gen_chunk = {
                                    "id": f"chatcmpl-{run_id or 'stream'}",
                                    "object": "chat.completion.chunk",
                                    "created": int(time.time()),
                                    "model": "gpt-4.1-mini",
                                    "choices": [{
                                        "index": 0,
                                        "delta": {
                                            "content": gen_text
                                        },
                                        "finish_reason": None
                                    }]
                                }
                                yield f"data: {json.dumps(gen_chunk)}\n\n"
                                
                                try:
                                    # Submit tool outputs and continue streaming
                                    with client.beta.threads.runs.submit_tool_outputs_stream(
                                        thread_id=session,
                                        run_id=event.data.id,
                                        tool_outputs=tool_outputs
                                    ) as tool_stream:
                                        for tool_event in tool_stream:
                                            # Handle text deltas from the continued stream
                                            if tool_event.event == "thread.message.delta":
                                                delta = tool_event.data.delta
                                                if delta.content:
                                                    for content_part in delta.content:
                                                        if content_part.type == 'text' and content_part.text:
                                                            text_value = content_part.text.value
                                                            if text_value:
                                                                buffer_inner.append(text_value)
                                                                if len(buffer_inner) >= 2:
                                                                    joined_text = ''.join(buffer_inner)
                                                                    chunk_data = {
                                                                        "id": f"chatcmpl-{run_id or 'stream'}",
                                                                        "object": "chat.completion.chunk",
                                                                        "created": int(time.time()),
                                                                        "model": "gpt-4.1-mini",
                                                                        "choices": [{
                                                                            "index": 0,
                                                                            "delta": {
                                                                                "content": joined_text
                                                                            },
                                                                            "finish_reason": None
                                                                        }]
                                                                    }
                                                                    yield f"data: {json.dumps(chunk_data)}\n\n"
                                                                    buffer_inner = []
                                            
                                            # Handle run completion
                                            elif tool_event.event == "thread.run.completed":
                                                completed = True
                                                # Yield any remaining buffer
                                                if buffer_inner:
                                                    joined_text = ''.join(buffer_inner)
                                                    chunk_data = {
                                                        "id": f"chatcmpl-{run_id or 'stream'}",
                                                        "object": "chat.completion.chunk",
                                                        "created": int(time.time()),
                                                        "model": "gpt-4.1-mini",
                                                        "choices": [{
                                                            "index": 0,
                                                            "delta": {
                                                                "content": joined_text
                                                            },
                                                            "finish_reason": None
                                                        }]
                                                    }
                                                    yield f"data: {json.dumps(chunk_data)}\n\n"
                                                    buffer_inner = []
                                                
                                                # Send final chunk
                                                final_chunk = {
                                                    "id": f"chatcmpl-{run_id or 'stream'}",
                                                    "object": "chat.completion.chunk",
                                                    "created": int(time.time()),
                                                    "model": "gpt-4.1-mini",
                                                    "choices": [{
                                                        "index": 0,
                                                        "delta": {},
                                                        "finish_reason": "stop"
                                                    }]
                                                }
                                                yield f"data: {json.dumps(final_chunk)}\n\n"
                                                yield "data: [DONE]\n\n"
                                                return
                                            
                                            # Handle run failures
                                            elif tool_event.event == "thread.run.failed":
                                                logging.error(f"Tool output stream run failed: {tool_event.data}")
                                                
                                                # Send error notice
                                                error_chunk = {
                                                    "id": f"chatcmpl-{run_id or 'stream'}",
                                                    "object": "chat.completion.chunk",
                                                    "created": int(time.time()),
                                                    "model": "gpt-4.1-mini",
                                                    "choices": [{
                                                        "index": 0,
                                                        "delta": {
                                                            "content": "\n[Note: Response generation encountered an issue. Tool results are shown above.]"
                                                        },
                                                        "finish_reason": "stop"
                                                    }]
                                                }
                                                yield f"data: {json.dumps(error_chunk)}\n\n"
                                                yield "data: [DONE]\n\n"
                                                return
                                    
                                    tool_outputs_submitted = True
                                    logging.info(f"Successfully submitted tool outputs for run {event.data.id} with streaming")
                                    
                                except Exception as submit_e:
                                    logging.error(f"Error submitting tool outputs with streaming: {submit_e}")
                                    
                                    # Try fallback approach - regular submission without streaming
                                    try:
                                        client.beta.threads.runs.submit_tool_outputs(
                                            thread_id=session,
                                            run_id=event.data.id,
                                            tool_outputs=tool_outputs
                                        )
                                        tool_outputs_submitted = True
                                        logging.info(f"Submitted tool outputs using fallback method")
                                        
                                        # Continue with the original stream
                                        continue
                                        
                                    except Exception as fallback_e:
                                        logging.error(f"Fallback submission also failed: {fallback_e}")
                                        
                                        # Send error and finish
                                        error_text = "\n[Error: Failed to complete processing. Tool results are shown above.]\n"
                                        error_chunk = {
                                            "id": f"chatcmpl-{run_id or 'stream'}",
                                            "object": "chat.completion.chunk",
                                            "created": int(time.time()),
                                            "model": "gpt-4.1-mini",
                                            "choices": [{
                                                "index": 0,
                                                "delta": {
                                                    "content": error_text
                                                },
                                                "finish_reason": "stop"
                                            }]
                                        }
                                        yield f"data: {json.dumps(error_chunk)}\n\n"
                                        yield "data: [DONE]\n\n"
                                        return
                
                # Yield any remaining text in the buffer before exiting the stream loop
                if buffer:
                    joined_text = ''.join(buffer)
                    chunk_data = {
                        "id": f"chatcmpl-{run_id or 'stream'}",
                        "object": "chat.completion.chunk",
                        "created": int(time.time()),
                        "model": "gpt-4.1-mini",
                        "choices": [{
                            "index": 0,
                            "delta": {
                                "content": joined_text
                            },
                            "finish_reason": None
                        }]
                    }
                    yield f"data: {json.dumps(chunk_data)}\n\n"
                    buffer = []
            
        except Exception as e:
            error_details = traceback.format_exc()
            logging.error(f"Streaming error during run for thread {session}: {e}\n{error_details}")
            error_chunk = {
                "id": "chatcmpl-error",
                "object": "chat.completion.chunk",
                "created": int(time.time()),
                "model": "gpt-4.1-mini",
                "choices": [{
                    "index": 0,
                    "delta": {
                        "content": "\n[ERROR] An error occurred while generating the response. Please try again.\n"
                    },
                    "finish_reason": "stop"
                }]
            }
            yield f"data: {json.dumps(error_chunk)}\n\n"
            yield "data: [DONE]\n\n"
    ######################### END OF def stream_response() #######################################
    
    try:
        # PRIORITY 1: Check for context or files FIRST (before any thread operations)
        if context is not None:
            logging.info(f"🔄 CONTEXT MODE ACTIVATED - Bypassing thread system")
            logging.info(f"  Assistant: {assistant}")
            logging.info(f"  Thread (ignored): {session}")
            logging.info(f"  Context size: {len(context) if context else 0} chars")
            
            # Immediately fallback to completions API when context is provided
            return await fallback_to_completions(
                error_context="Context mode requested - bypassing assistant API",
                user_context=context,
                files=files
            )
        
        if files:
            logging.info(f"📎 FILES DETECTED - Using enhanced completions with file analysis")
            logging.info(f"  Files count: {len(files)}")
            logging.info(f"  Files: {[f.filename for f in files]}")
            
            # Immediately use fallback_to_completions when files are present
            return await fallback_to_completions(
                error_context="File analysis requested",
                user_context=context,
                files=files
            )
        
        # PRIORITY 2: Check for missing parameters (no point in continuing without them)
        if not assistant or not session:
            missing_params = []
            if not assistant:
                missing_params.append("assistant_id")
            if not session:
                missing_params.append("thread_id")
            
            logging.info(f"Missing required parameters: {', '.join(missing_params)}. Falling back to completions API.")
            return await fallback_to_completions(
                error_context=f"Missing required parameters: {', '.join(missing_params)}",
                user_context=context,
                files=files
            )
        
        # PRIORITY 3: Now we know we need thread mode - log it
        logging.info(f"📌 THREAD MODE - Using standard assistant flow")
        logging.info(f"  Assistant: {assistant}")
        logging.info(f"  Thread: {session}")
        
        # PRIORITY 4: Acquire thread lock (ONLY ONCE, ONLY WHEN NEEDED)
        try:
            thread_lock = await thread_lock_manager.get_lock(session)
            await asyncio.wait_for(thread_lock.acquire(), timeout=30.0)  # 30 second timeout
            logging.info(f"Acquired thread lock for session {session}")
        except asyncio.TimeoutError:
            logging.warning(f"Timeout acquiring thread lock for session {session}")
            return await fallback_to_completions(
                error_context=f"Thread lock timeout",
                user_context=context,
                files=files
            )
        except Exception as lock_e:
            logging.error(f"Error acquiring thread lock: {lock_e}")
            return await fallback_to_completions(
                error_context=f"Thread lock error: {str(lock_e)}",
                user_context=context,
                files=files
            )
        
        # PRIORITY 5: Validate resources
        try:
            validation = await validate_resources(client, session, assistant)
            
            # If either resource is invalid, fallback to completions
            if not validation["thread_valid"] or not validation["assistant_valid"]:
                invalid_resources = []
                if not validation["thread_valid"]:
                    invalid_resources.append(f"thread_id '{session}'")
                if not validation["assistant_valid"]:
                    invalid_resources.append(f"assistant_id '{assistant}'")
                
                logging.warning(f"Invalid resources: {', '.join(invalid_resources)}. Falling back to completions API.")
                return await fallback_to_completions(
                    error_context=f"Invalid resources: {', '.join(invalid_resources)}",
                    user_context=context,
                    files=files
                )
        except Exception as validation_e:
            logging.error(f"Error during resource validation: {validation_e}")
            return await fallback_to_completions(
                error_context=f"Resource validation error: {str(validation_e)}",
                user_context=context,
                files=files
            )
        
        # Check if there's an active run before adding a message
        active_run = False
        run_id = None
        requires_action_tools = []
        try:
            # List runs to check for active ones
            runs = client.beta.threads.runs.list(thread_id=session, limit=1)
            if runs.data:
                latest_run = runs.data[0]
                if latest_run.status in ["in_progress", "queued", "requires_action"]:
                    active_run = True
                    run_id = latest_run.id
                    logging.warning(f"Active run {run_id} detected with status {latest_run.status}")
                    
                    # If requires_action, get the tool calls
                    if latest_run.status == "requires_action" and hasattr(latest_run, 'required_action'):
                        if hasattr(latest_run.required_action, 'submit_tool_outputs'):
                            requires_action_tools = latest_run.required_action.submit_tool_outputs.tool_calls
        except Exception as e:
            logging.warning(f"Error checking for active runs: {e}")
            # Continue anyway - we'll handle failure when adding messages
        # Check and trim thread BEFORE adding new message
        if session and prompt:  # Only trim if we're about to add a message
            try:
                messages_response = client.beta.threads.messages.list(thread_id=session, limit=100)
                message_count = len(messages_response.data)
                if message_count > 40:  # Trim at 48 to leave room for new message
                    logging.info(f"Thread {session} has {message_count} messages, trimming before adding new message")
                    
                    # Trim now for both streaming and non-streaming
                    try:
                        if stream_output:
                            trimmed = sync_trim_thread(client, session, keep_messages=30)
                        else:
                            trimmed = await trim_thread(client, session, keep_messages=30)
                        
                        if trimmed:
                            logging.info(f"Trimmed thread {session} from {message_count} messages")
                            # Add a small delay to ensure OpenAI processes the deletions
                            await asyncio.sleep(0.5)
                    except Exception as trim_e:
                        logging.error(f"Error trimming thread: {trim_e}")
                        # Fallback to completions on trim error
                        return await fallback_to_completions(error_context=f"Thread trimming error: {str(trim_e)}", user_context=context, files=files)
            except Exception as list_e:
                logging.error(f"Error listing messages for trimming: {list_e}")
                # Fallback on any error
                return await fallback_to_completions(error_context=f"Message listing error: {str(list_e)}", user_context=context, files=files)
        
        # Add user message to the thread if prompt is given
        if prompt:
            max_retries = 3
            retry_delay = 2  # seconds
            success = False
            
            for attempt in range(max_retries):
                try:
                    if active_run and run_id:
                        try:
                            run_status = client.beta.threads.runs.retrieve(thread_id=session, run_id=run_id)
                            
                            # Check how long the run has been active
                            run_created_at = run_status.created_at
                            current_time = int(time.time())
                            run_age_seconds = current_time - run_created_at
                            
                            logging.info(f"Active run {run_id} found with status {run_status.status}, age: {run_age_seconds}s")
                            
                            # Smart decision: wait for young runs, cancel old stuck runs
                            if run_status.status in ["in_progress", "queued"]:
                                max_wait_time = 30  # Maximum 30 seconds total wait
                                check_interval = 2  # Check every 2 seconds
                                waited_time = 0
                                
                                # If run is young (< 30s), wait for it to complete
                                if run_age_seconds < max_wait_time:
                                    logging.info(f"Run {run_id} is only {run_age_seconds}s old, waiting for completion...")
                                    
                                    while waited_time < (max_wait_time - run_age_seconds):
                                        time.sleep(check_interval)
                                        waited_time += check_interval
                                        
                                        # Check run status
                                        try:
                                            check_status = client.beta.threads.runs.retrieve(thread_id=session, run_id=run_id)
                                            
                                            if check_status.status in ["completed", "failed", "cancelled", "expired"]:
                                                logging.info(f"Run {run_id} completed with status: {check_status.status}")
                                                active_run = False
                                                break
                                            elif check_status.status == "requires_action":
                                                # Handle requires_action separately below
                                                run_status = check_status
                                                if hasattr(check_status, 'required_action'):
                                                    requires_action_tools = check_status.required_action.submit_tool_outputs.tool_calls
                                                break
                                            else:
                                                # Still in progress/queued
                                                current_run_age = int(time.time()) - check_status.created_at
                                                logging.info(f"Run {run_id} still {check_status.status} after {current_run_age}s total")
                                                
                                        except Exception as check_e:
                                            logging.warning(f"Error checking run status: {check_e}")
                                            break
                                    
                                    # After waiting, check if we need to cancel
                                    if active_run and run_status.status in ["in_progress", "queued"]:
                                        total_age = int(time.time()) - run_created_at
                                        if total_age >= max_wait_time:
                                            logging.warning(f"Run {run_id} exceeded max wait time ({total_age}s), cancelling...")
                                            try:
                                                client.beta.threads.runs.cancel(thread_id=session, run_id=run_id)
                                                logging.info(f"Cancelled stuck run {run_id}")
                                                
                                                # Wait for cancellation
                                                cancel_wait = 0
                                                while cancel_wait < 5:
                                                    time.sleep(1)
                                                    cancel_wait += 1
                                                    try:
                                                        cancel_check = client.beta.threads.runs.retrieve(thread_id=session, run_id=run_id)
                                                        if cancel_check.status in ["cancelled", "failed", "completed", "expired"]:
                                                            break
                                                    except:
                                                        break
                                                try:
                                                    client.beta.threads.messages.create(
                                                        thread_id=session,
                                                        role="system",
                                                        content="[Previous operation was cancelled due to timeout. Processing new request.]"
                                                    )
                                                    logging.info("Added cancellation notice to thread")
                                                except Exception as notice_e:
                                                    logging.warning(f"Could not add cancellation notice: {notice_e}")
                                            except Exception as cancel_e:
                                                logging.error(f"Failed to cancel run {run_id}: {cancel_e}")
                                else:
                                    # Run is already old (> 30s), cancel immediately
                                    logging.warning(f"Run {run_id} is {run_age_seconds}s old (stuck), cancelling immediately...")
                                    try:
                                        client.beta.threads.runs.cancel(thread_id=session, run_id=run_id)
                                        logging.info(f"Cancelled old stuck run {run_id}")
                                        time.sleep(2)  # Brief wait for cancellation
                                    except Exception as cancel_e:
                                        logging.error(f"Failed to cancel old run {run_id}: {cancel_e}")
                                        
                            elif run_status.status == "requires_action":
                                # For requires_action, check age and handle accordingly
                                if run_age_seconds > 60:  # If stuck in requires_action for > 1 minute
                                    logging.warning(f"Run {run_id} stuck in requires_action for {run_age_seconds}s, cancelling directly")
                                    # Skip tool output submission and cancel directly
                                    try:
                                        client.beta.threads.runs.cancel(thread_id=session, run_id=run_id)
                                        logging.info(f"Cancelled stuck requires_action run {run_id}")
                                        time.sleep(2)  # Brief wait for cancellation
                                    except Exception as cancel_e:
                                        logging.error(f"Failed to cancel requires_action run {run_id}: {cancel_e}")
                                        # Continue anyway - the add message will handle it
                                else:
                                    logging.info(f"Run {run_id} is in requires_action state ({run_age_seconds}s old), waiting briefly")
                                    # For young requires_action runs, wait a bit longer
                                    time.sleep(5)
                                    
                        except Exception as run_e:
                            logging.warning(f"Error handling active run: {run_e}")

                    # Try to add the message
                    client.beta.threads.messages.create(
                        thread_id=session,
                        role="user",
                        content=prompt
                    )
                    logging.info(f"Added user message to thread {session} (attempt {attempt+1})")
                    success = True
                    break
                except Exception as e:
                    if "while a run" in str(e) and attempt < max_retries - 1:
                        logging.warning(f"Failed to add message (attempt {attempt+1}), run is active. Retrying in {retry_delay}s: {e}")
                        time.sleep(retry_delay)
                        retry_delay *= 2  # Exponential backoff
                    else:
                        logging.error(f"Failed to add message to thread {session}: {e}")
                        if attempt == max_retries - 1:
                            # Use fallback instead of raising exception
                            return await fallback_to_completions(error_context=f"Failed to add message: {str(e)}", user_context=context, files=files)
            
            if not success:
                # Fallback to completions instead of creating new thread
                logging.warning(f"Failed to add message to thread {session} after all retries. Falling back to completions API.")
                return await fallback_to_completions(
                    error_context=f"Failed to add message to thread after {max_retries} attempts",
                    user_context=context, files=files
                )
        
        
        # Handle non-streaming mode (/chat endpoint)
        if not stream_output:
            # For non-streaming mode, we'll use a completely different approach
            full_response = ""
            try:
                # Create a run without streaming
                run = client.beta.threads.runs.create(
                    thread_id=session,
                    assistant_id=assistant,
                    truncation_strategy={
                        "type": "last_messages",
                        "last_messages": 10
                    }
                )
                run_id = run.id
                logging.info(f"Created run {run_id} for thread {session} (non-streaming mode)")
                
                # Poll for run completion
                max_poll_attempts = 30  # 5 minute timeout with 10 second intervals
                poll_interval = 10  # seconds
                tool_outputs_submitted = False
                tool_call_results = []
                
                for attempt in range(max_poll_attempts):
                    try:
                        run_status = client.beta.threads.runs.retrieve(
                            thread_id=session,
                            run_id=run_id
                        )
                        
                        logging.info(f"Run status poll {attempt+1}/{max_poll_attempts}: {run_status.status}")
                        
                        # Handle completed run
                        if run_status.status == "completed":
                            # Get the latest message
                            messages = client.beta.threads.messages.list(
                                thread_id=session,
                                order="desc",
                                limit=1
                            )
                            
                            if messages and messages.data:
                                latest_message = messages.data[0]
                                for content_part in latest_message.content:
                                    if content_part.type == 'text':
                                        full_response += content_part.text.value
                                
                                logging.info(f"Successfully retrieved final response")
                            break  # Exit the polling loop
                        
                        # Handle failed/cancelled/expired run
                        elif run_status.status in ["failed", "cancelled", "expired"]:
                            error_details = ""
                            error_code = None
                            if run_status.status == "failed" and hasattr(run_status, 'last_error') and run_status.last_error:
                                error_details = f" Error: {run_status.last_error.message if hasattr(run_status.last_error, 'message') else str(run_status.last_error)}"
                                error_code = run_status.last_error.code if hasattr(run_status.last_error, 'code') else 'unknown'
                                logging.error(f"Run failed with error code '{error_code}': {error_details}")
                                
                                # Retry on server errors
                                if error_code == 'server_error' and attempt < 10:  # More retries for chat endpoint
                                    logging.info(f"Retrying due to server error (attempt {attempt + 1}/10)")
                                    time.sleep(3)  # Wait 3 seconds before retry
                                    
                                    # Create a new run
                                    try:
                                        run = client.beta.threads.runs.create(
                                            thread_id=session,
                                            assistant_id=assistant,
                                            truncation_strategy={
                                                "type": "last_messages",
                                                "last_messages": 10
                                            }
                                        )
                                        run_id = run.id
                                        logging.info(f"Created new run {run_id} after server error")
                                        continue  # Continue polling with new run
                                    except Exception as retry_e:
                                        logging.error(f"Failed to create retry run: {retry_e}")
                                        # Try one more time after a longer wait
                                        time.sleep(5)
                                        try:
                                            run = client.beta.threads.runs.create(
                                                thread_id=session,
                                                assistant_id=assistant,
                                                truncation_strategy={
                                                    "type": "last_messages",
                                                    "last_messages": 10
                                                }
                                            )
                                            run_id = run.id
                                            logging.info(f"Created new run {run_id} on second retry attempt")
                                            continue
                                        except:
                                            pass
                            elif run_status.status == "requires_action":
                                if run_status.required_action.type == "submit_tool_outputs":
                                    tool_calls = run_status.required_action.submit_tool_outputs.tool_calls
                                    tool_outputs = []
                                    
                                    for tool_call in tool_calls:
                                        if tool_call.function.name == "pandas_agent":
                                            try:
                                                # Extract arguments
                                                args = json.loads(tool_call.function.arguments)
                                                query = args.get("query", "")
                                                filename = args.get("filename", None)
                                                
                                                # Get pandas files for this thread
                                                pandas_files = []
                                                retry_count = 0
                                                max_retries = 3
                                                
                                                while retry_count < max_retries:
                                                    try:
                                                        messages = client.beta.threads.messages.list(
                                                            thread_id=session,
                                                            order="desc",
                                                            limit=50
                                                        )
                                                        
                                                        for msg in messages.data:
                                                            if hasattr(msg, 'metadata') and msg.metadata and msg.metadata.get('type') == 'pandas_agent_files':
                                                                try:
                                                                    pandas_files = json.loads(msg.metadata.get('files', '[]'))
                                                                except Exception as parse_e:
                                                                    logging.error(f"Error parsing pandas files metadata: {parse_e}")
                                                                break
                                                        break  # Success, exit retry loop
                                                    except Exception as list_e:
                                                        retry_count += 1
                                                        logging.error(f"Error retrieving pandas files (attempt {retry_count}): {list_e}")
                                                        time.sleep(1)
                                                
                                                # Filter by filename if specified
                                                if filename:
                                                    pandas_files = [f for f in pandas_files if f.get("name") == filename]
                                                
                                                # Generate operation ID for status tracking
                                                pandas_agent_operation_id = f"pandas_agent_{int(time.time())}_{os.urandom(2).hex()}"
                                                
                                                # Execute the pandas_agent using manager directly
                                                manager = PandasAgentManager.get_instance()
                                                result, error, removed_files = manager.analyze(
                                                    thread_id=session,
                                                    query=query,
                                                    files=pandas_files
                                                )
                                                
                                                # Format the analysis result (same as streaming)
                                                analysis_result = result if result else ""
                                                if error:
                                                    analysis_result = f"Error analyzing data: {error}"
                                                if removed_files:
                                                    removed_files_str = ", ".join(f"'{f}'" for f in removed_files)
                                                    analysis_result += f"\n\nNote: The following file(s) were removed due to the 3-file limit: {removed_files_str}"
                                                
                                                # Add to tool outputs
                                                tool_outputs.append({
                                                    "tool_call_id": tool_call.id,
                                                    "output": analysis_result
                                                })
                                                
                                                # Save for potential fallback
                                                tool_call_results.append(analysis_result)
                                                
                                            except Exception as e:
                                                error_details = traceback.format_exc()
                                                logging.error(f"Error executing pandas_agent: {e}\n{error_details}")
                                                error_msg = f"Error analyzing data: {str(e)}"
                                                
                                                # Add error to tool outputs
                                                tool_outputs.append({
                                                    "tool_call_id": tool_call.id,
                                                    "output": error_msg
                                                })
                                                
                                                # Save for potential fallback
                                                tool_call_results.append(error_msg)
                            else:
                                logging.error(f"Run ended with status: {run_status.status}")
                            
                            # If we have tool results from before the failure, return those
                            if tool_call_results:
                                combined_results = "\n\n".join(tool_call_results)
                                return JSONResponse(content={"response": combined_results})
                            
                            # Try to provide a more helpful error message
                            if "rate_limit" in error_details.lower():
                                return JSONResponse(content={"response": "I'm experiencing high demand. Please try again in a moment."})
                            elif "context_length" in error_details.lower():
                                return JSONResponse(content={"response": "The conversation has become too long. Please start a new conversation."})
                            else:
                                # Use fallback
                                return await fallback_to_completions(error_context= f"Run {run_status.status}: {error_details}", user_context=context, files=files)
                        
                        # Handle tool calls
                        elif run_status.status == "requires_action":
                            if run_status.required_action.type == "submit_tool_outputs":
                                tool_calls = run_status.required_action.submit_tool_outputs.tool_calls
                                tool_outputs = []
                                
                                for tool_call in tool_calls:
                                    if tool_call.function.name == "pandas_agent":
                                        try:
                                            # Extract arguments
                                            args = json.loads(tool_call.function.arguments)
                                            query = args.get("query", "")
                                            filename = args.get("filename", None)
                                            
                                            # Get pandas files for this thread
                                            pandas_files = []
                                            retry_count = 0
                                            max_retries = 3
                                            
                                            while retry_count < max_retries:
                                                try:
                                                    messages = client.beta.threads.messages.list(
                                                        thread_id=session,
                                                        order="desc",
                                                        limit=50
                                                    )
                                                    
                                                    for msg in messages.data:
                                                        if hasattr(msg, 'metadata') and msg.metadata and msg.metadata.get('type') == 'pandas_agent_files':
                                                            try:
                                                                pandas_files = json.loads(msg.metadata.get('files', '[]'))
                                                            except Exception as parse_e:
                                                                logging.error(f"Error parsing pandas files metadata: {parse_e}")
                                                            break
                                                    break  # Success, exit retry loop
                                                except Exception as list_e:
                                                    retry_count += 1
                                                    logging.error(f"Error retrieving pandas files (attempt {retry_count}): {list_e}")
                                                    time.sleep(1)
                                            
                                            # Filter by filename if specified
                                            if filename:
                                                pandas_files = [f for f in pandas_files if f.get("name") == filename]
                                            
                                            # Generate operation ID for status tracking
                                            pandas_agent_operation_id = f"pandas_agent_{int(time.time())}_{os.urandom(2).hex()}"
                                            
                                            # Execute the pandas_agent
                                            manager = PandasAgentManager.get_instance()
                                            result, error, removed_files = manager.analyze(
                                                thread_id=session,
                                                query=query,
                                                files=pandas_files
                                            )
                                            
                                            # Format the analysis result (same as streaming)
                                            analysis_result = result if result else ""
                                            if error:
                                                analysis_result = f"Error analyzing data: {error}"
                                            if removed_files:
                                                removed_files_str = ", ".join(f"'{f}'" for f in removed_files)
                                                analysis_result += f"\n\nNote: The following file(s) were removed due to the 3-file limit: {removed_files_str}"
                                           
                                            # Add to tool outputs
                                            tool_outputs.append({
                                                "tool_call_id": tool_call.id,
                                                "output": analysis_result
                                            })
                                            
                                            # Save for potential fallback
                                            tool_call_results.append(analysis_result)
                                            
                                        except Exception as e:
                                            error_details = traceback.format_exc()
                                            logging.error(f"Error executing pandas_agent: {e}\n{error_details}")
                                            error_msg = f"Error analyzing data: {str(e)}"
                                            
                                            # Add error to tool outputs
                                            tool_outputs.append({
                                                "tool_call_id": tool_call.id,
                                                "output": error_msg
                                            })
                                            
                                            # Save for potential fallback
                                            tool_call_results.append(error_msg)
                                    
                                    elif tool_call.function.name == "extract_data":
                                        try:
                                            # Extract arguments
                                            args = json.loads(tool_call.function.arguments)
                                            logging.info(f"extract_data tool call with args: {args}")
                                            
                                            # Create a mock request object for the handler
                                            host = os.environ.get('WEBSITE_HOSTNAME', 'localhost:8080')
                                            base_url = f"https://{host}" if 'azurewebsites.net' in host else f"http://{host}"
                                            mock_request = type('Request', (), {
                                                'base_url': base_url,
                                                'headers': {'host': host}
                                            })()
                                            
                                            # Call the handler with mock request
                                            result = await handle_extract_data(args, session, client, mock_request)
                                            
                                            # Add to tool outputs
                                            tool_outputs.append({
                                                "tool_call_id": tool_call.id,
                                                "output": result
                                            })
                                            
                                            # Save for potential fallback
                                            tool_call_results.append(result)
                                            
                                        except Exception as e:
                                            error_msg = f"Error extracting data: {str(e)}"
                                            logging.error(f"Error executing extract_data: {e}")
                                            
                                            # Add error to tool outputs
                                            tool_outputs.append({
                                                "tool_call_id": tool_call.id,
                                                "output": error_msg
                                            })
                                            
                                            # Save for potential fallback
                                            tool_call_results.append(error_msg)
                                                                        
                                    elif tool_call.function.name == "extract_data":
                                        try:
                                            # Extract arguments
                                            args = json.loads(tool_call.function.arguments)
                                            logging.info(f"extract_data tool call with args: {args}")
                                            
                                            # Create a mock request object for the handler
                                            host = os.environ.get('WEBSITE_HOSTNAME', 'localhost:8080')
                                            base_url = f"https://{host}" if 'azurewebsites.net' in host else f"http://{host}"
                                            mock_request = type('Request', (), {
                                                'base_url': base_url,
                                                'headers': {'host': host}
                                            })()
                                            
                                            # Call the handler with mock request
                                            result = await handle_extract_data(args, session, client, mock_request)
                                            
                                            # Add to tool outputs
                                            tool_outputs.append({
                                                "tool_call_id": tool_call.id,
                                                "output": result
                                            })
                                            
                                            # Save for potential fallback
                                            tool_call_results.append(result)
                                            
                                        except Exception as e:
                                            error_msg = f"Error extracting data: {str(e)}"
                                            logging.error(f"Error executing extract_data: {e}")
                                            
                                            # Add error to tool outputs
                                            tool_outputs.append({
                                                "tool_call_id": tool_call.id,
                                                "output": error_msg
                                            })
                                            
                                            # Save for potential fallback
                                            tool_call_results.append(error_msg)
                                
                                # Submit tool outputs
                                if tool_outputs:
                                    retry_count = 0
                                    max_retries = 3
                                    submit_success = False
                                    
                                    while retry_count < max_retries and not submit_success:
                                        try:
                                            client.beta.threads.runs.submit_tool_outputs(
                                                thread_id=session,
                                                run_id=run_id,
                                                tool_outputs=tool_outputs
                                            )
                                            submit_success = True
                                            tool_outputs_submitted = True
                                            logging.info(f"Successfully submitted tool outputs for run {run_id}")
                                        except Exception as submit_e:
                                            retry_count += 1
                                            logging.error(f"Error submitting tool outputs (attempt {retry_count}): {submit_e}")
                                            time.sleep(1)
                                    
                                    if not submit_success:
                                        # Use fallback
                                        return await fallback_to_completions(error_context= f"Failed to submit tool outputs: {submit_e}", user_context=context, files=files)
                        
                        # Continue polling if still in progress
                        if attempt < max_poll_attempts - 1:
                            # Progressive backoff: start fast, slow down over time
                            if attempt < 10:
                                time.sleep(3)  # First 10 attempts: 3 seconds
                            elif attempt < 20:
                                time.sleep(5)  # Next 10 attempts: 5 seconds
                            else:
                                time.sleep(10)  # Remaining attempts: 10 seconds
                            
                    except Exception as poll_e:
                        logging.error(f"Error polling run status (attempt {attempt+1}): {poll_e}")
                        time.sleep(poll_interval)
                        
                # If we reach here without a full_response, but we have tool results, use those
                if not full_response and tool_call_results:
                    full_response = "\n\n".join(tool_call_results)
                    logging.info("Using tool call results as fallback response")
                
                # If we still don't have a response, try one more time to get the latest message
                if not full_response:
                    try:
                        messages = client.beta.threads.messages.list(
                            thread_id=session,
                            order="desc",
                            limit=1
                        )
                        
                        if messages and messages.data:
                            latest_message = messages.data[0]
                            for content_part in latest_message.content:
                                if content_part.type == 'text':
                                    full_response += content_part.text.value
                    except Exception as final_e:
                        logging.error(f"Error retrieving final message: {final_e}")
                
                # Final fallback if we still don't have a response
                if not full_response:
                    # Use completions API fallback
                    return await fallback_to_completions(error_context= "No response received from assistant", user_context=context, files=files)

                return JSONResponse(content={"response": full_response})
                
            except Exception as e:
                logging.error(f"Error in non-streaming response generation: {e}")
                # Use fallback
                return await fallback_to_completions(error_context=f"Non-streaming error: {str(e)}", user_context=context, files=files)
        
        # Return the streaming response for streaming mode
        try:
            response = StreamingResponse(stream_response(), media_type="text/event-stream")
            response.headers["X-Accel-Buffering"] = "no"  # Disable nginx buffering
            response.headers["Cache-Control"] = "no-cache"
            response.headers["Connection"] = "keep-alive"
            response_started = True
            return response
        except Exception as stream_setup_e:
            logging.error(f"Error setting up streaming response: {stream_setup_e}")
            return await fallback_to_completions(error_context= f"Stream setup error: {str(stream_setup_e)}", user_context=context, files=files)

    except Exception as e:
        endpoint_type = "conversation" if stream_output else "chat"
        logging.error(f"Unexpected error in /{endpoint_type}: {str(e)}\n{traceback.format_exc()}")
        
        # Use fallback if response hasn't started
        if not response_started:
            return await fallback_to_completions(error_context= f"Unexpected error: {str(e)}", user_context=context, files=files)
        else:
            # If streaming already started, we can't change response type
            raise HTTPException(status_code=500, detail="An error occurred during response streaming")
    
    finally:
        # Release the thread lock
        if thread_lock and thread_lock.locked():
            try:
                thread_lock.release()
                logging.info(f"Released thread lock for session {session}")
            except Exception as release_e:
                logging.error(f"Error releasing thread lock: {release_e}")
@app.get("/conversation",
         summary="Stream Chat Messages (GET)",
         description="""Chat with AI using Server-Sent Events (SSE) for real-time streaming responses.

**⚡ STREAMING ENDPOINT** - Returns responses via Server-Sent Events

## Usage Modes:

### 1️⃣ **Stateful Mode** (with session & assistant)
Use when you have an existing chat session:
```
/conversation?session=thread_abc123&assistant=asst_xyz789&prompt=Hello
```
- Maintains conversation history
- Uses vector stores for context
- Supports file uploads from the session

### 2️⃣ **Stateless Mode** (with context)
Use for one-off queries without session management:
```
/conversation?context=You are a data analyst&prompt=Analyze this trend
```
- No session required
- Context defines AI behavior
- Fresh conversation each time

### 3️⃣ **Invalid Thread Handling**
If you provide an invalid thread ID:
- Falls back to stateless mode automatically
- Creates a temporary conversation
- Returns a warning in headers

### 📄 **File Support** (POST method only)
- Use POST endpoint to include files
- Files are processed and included in context
- Supports multiple files simultaneously

### Response Format:
```
data: {"id": "chatcmpl-xyz", "object": "chat.completion.chunk", "choices": [{"delta": {"content": "Hello"}}]}
data: {"id": "chatcmpl-xyz", "object": "chat.completion.chunk", "choices": [{"delta": {"content": " world"}}]}
data: [DONE]
```
""",
         tags=["Chat Operations"],
         response_class=StreamingResponse)
async def conversation_get(
    session: Optional[str] = Query(default=None, description="Session ID from /initiate-chat"),
    prompt: Optional[str] = Query(default=None, description="User message"),
    assistant: Optional[str] = Query(default=None, description="Assistant ID"),
    context: Optional[str] = Query(default=None, description="Additional context for stateless mode")
):
    """GET method for simple text queries without files."""
    return await process_conversation(session, prompt, assistant, stream_output=True, context=context, files=None)
@app.post("/conversation",
          summary="Stream Chat Messages (POST)",
          description="""Chat with AI using Server-Sent Events (SSE) with file upload support.

**⚡ STREAMING ENDPOINT** - Returns responses via Server-Sent Events

## Enhanced Features with Files:

### 📁 **File Upload Support**
- Upload multiple files with your query
- Supports: PDF, DOCX, TXT, CSV, XLSX, Images
- Files are analyzed and included in AI context

### Example Scenarios:

1. **Analyze Documents**:
   - Upload: `contract.pdf`
   - Prompt: "Summarize the key terms"

2. **Compare Files**:
   - Upload: `report1.xlsx, `report2.xlsx`
   - Prompt: "Compare the revenue figures"

3. **Image Analysis**:
   - Upload: `screenshot.png`
   - Prompt: "What's shown in this image?"

### Usage Modes (same as GET):
- **Stateful**: With session & assistant IDs
- **Stateless**: With context parameter
- **Fallback**: Invalid thread → stateless mode

### File Processing:
- Text extraction from documents
- Data parsing from spreadsheets
- OCR for images
- Automatic format detection
""",
          tags=["Chat Operations"],
          response_class=StreamingResponse)
async def conversation_post(
    session: Optional[str] = Form(default=None, description="Session ID from /initiate-chat"),
    prompt: Optional[str] = Form(default=None, description="User message"),
    assistant: Optional[str] = Form(default=None, description="Assistant ID"),
    context: Optional[str] = Form(default=None, description="Additional context for stateless mode"),
    files: List[UploadFile] = File(default=[], description="Multiple files to analyze with the query")
):
    """POST method for queries with file uploads."""
    # Handle empty file list
    files = None if not files else files
    return await process_conversation(session, prompt, assistant, stream_output=True, context=context, files=files)
# GET endpoint for chat (no file support)
@app.get("/chat",
         response_model=ChatResponse,
         summary="Chat (GET)",
         description="""Chat with AI and receive complete responses (non-streaming).

## Usage Modes:

### 1️⃣ **Stateful Mode** (with session & assistant)
```
/chat?session=thread_abc123&assistant=asst_xyz789&prompt=Hello
```
- Returns complete response after processing
- Maintains conversation history
- Good for short queries

### 2️⃣ **Stateless Mode** (with context)
```
/chat?context=You are a helpful assistant&prompt=Explain quantum physics
```
- One-shot responses
- No history maintained
- Context defines behavior

### 3️⃣ **Error Handling**
- Invalid session: Falls back to stateless
- Missing assistant: Returns error
- Empty prompt: Returns error

### When to Use:
- ✅ Short responses needed
- ✅ Integration with systems that don't support SSE
- ✅ Simple Q&A interactions
- ❌ Long responses (use /conversation instead)

### Response Format:
```json
{
  "response": "Complete AI response here",
  "session": "thread_abc123",
  "assistant": "asst_xyz789"
}
```
""",
         tags=["Chat Operations"])

async def chat_get(
    session: Optional[str] = Query(default=None, description="Session ID"),
    prompt: Optional[str] = Query(default=None, description="User message"),
    assistant: Optional[str] = Query(default=None, description="Assistant ID"),
    context: Optional[str] = Query(default=None, description="Additional context")
):
    """GET method for simple text queries without files."""
    return await process_conversation(session, prompt, assistant, stream_output=False, context=context, files=None)

# POST endpoint for chat (with file support)
@app.post("/chat",
          response_model=ChatResponse,
          summary="Chat (POST)",
          description="""Chat with AI and receive complete responses with file upload support.

## File Upload Features:

### 📄 **Supported Files**
- Documents: PDF, DOCX, TXT, MD
- Data: CSV, XLSX, XLS
- Images: JPG, PNG, GIF
- Code: JSON, HTML, XML

### Example Use Cases:

1. **Document Q&A**:
   ```
   Files: manual.pdf
   Prompt: "What's the installation process?"
   ```

2. **Data Analysis**:
   ```
   Files: sales.csv
   Prompt: "What's the average order value?"
   ```

3. **Multi-file Comparison**:
   ```
   Files: [q1.xlsx, q2.xlsx, q3.xlsx]
   Prompt: "Compare quarterly performance"
   ```

### Processing Details:
- Files extracted and parsed
- Content included in AI context
- Automatic format detection
- Error handling for unsupported formats

### Best Practices:
- Keep files under 10MB
- Use clear, specific prompts
- Reference files by name in prompts
- Check response for file processing errors
""",
          tags=["Chat Operations"])
async def chat_post(
    session: Optional[str] = Form(default=None, description="Session ID"),
    prompt: Optional[str] = Form(default=None, description="User message"),
    assistant: Optional[str] = Form(default=None, description="Assistant ID"),
    context: Optional[str] = Form(default=None, description="Additional context"),
    files: List[UploadFile] = File(default=[], description="Multiple files to analyze with the query")
):
    """POST method for queries with file uploads."""
    # Handle empty file list
    files = None if not files else files
    return await process_conversation(session, prompt, assistant, stream_output=False, context=context, files=files)


async def prepare_file_for_completion(file_content: bytes, filename: str, file_type: str) -> Dict[str, Any]:
    """
    Prepare file content for inclusion in chat completion request using extract_text_internal.
    Now uses the same robust extraction mechanism as fallback_to_completions.
    
    Args:
        file_content: Raw file bytes
        filename: Name of the file
        file_type: MIME type of the file
        
    Returns:
        Dict containing file information for API request
    """
    try:
        # For images, handle as before
        if file_type.startswith('image/'):
            # For images, encode as base64
            base64_image = base64.b64encode(file_content).decode('utf-8')
            return {
                "type": "image_url",
                "image_url": {
                    "url": f"data:{file_type};base64,{base64_image}"
                }
            }
        else:
            # For all other files, use extract_text_internal for robust extraction
            extracted_text = await extract_text_internal(
                file_content=file_content,
                filename=filename,
                strategy="auto",
                languages=None,
                encoding=None,
                logger=logging.getLogger(__name__)
            )
            
            # Format the extracted text with file information
            formatted_content = f"[File: {filename}]\n{extracted_text}"
            
            return {
                "type": "text",
                "text": formatted_content
            }
                
    except Exception as e:
        logging.error(f"Error preparing file {filename} for completion: {e}")
        # Fallback to basic text representation
        return {
            "type": "text", 
            "text": f"[Error processing {filename}: {str(e)}]"
        }

def generate_file_from_response(content: str, file_type: str) -> Optional[Tuple[bytes, str]]:
    """
    Generate a file from completion response content with better error handling.
    """
    try:
        # Extract CSV content
        csv_content = extract_csv_from_content(content)
        
        if file_type == 'csv':
            # Ensure proper CSV formatting
            if not csv_content:
                logging.warning("No CSV content found in response")
                return None
                
            # Add UTF-8 BOM for Excel compatibility
            bom = '\ufeff'
            file_bytes = (bom + csv_content).encode('utf-8-sig')
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"generated_data_{timestamp}.csv"
            
            return file_bytes, filename
            
        elif file_type == 'excel':
            try:
                import pandas as pd
                from io import StringIO, BytesIO
                
                # Parse CSV content
                if not csv_content:
                    logging.warning("No CSV content found for Excel conversion")
                    return None
                
                df = pd.read_csv(StringIO(csv_content))
                
                # Create Excel with better formatting
                buffer = BytesIO()
                with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
                    df.to_excel(writer, index=False, sheet_name='Data')
                    
                    # Get the workbook and worksheet
                    workbook = writer.book
                    worksheet = writer.sheets['Data']
                    
                    # Auto-adjust column widths
                    for column in df:
                        column_length = max(
                            df[column].astype(str).map(len).max(),
                            len(str(column))
                        )
                        col_idx = df.columns.get_loc(column)
                        column_letter = chr(65 + col_idx)
                        worksheet.column_dimensions[column_letter].width = min(column_length + 2, 50)
                
                buffer.seek(0)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"generated_data_{timestamp}.xlsx"
                
                return buffer.getvalue(), filename
                
            except Exception as excel_error:
                logging.error(f"Excel generation failed: {excel_error}")
                # Fallback to CSV
                return generate_file_from_response(content, 'csv')
                
    except Exception as e:
        logging.error(f"Error generating {file_type} file: {e}")
        return None

@app.post("/completion",
          response_model=CompletionResponse,
          summary="Generate AI Completion",
          description="Generate AI-powered text completions with optional file exports in CSV, Excel, or Word formats.",
          tags=["AI Operations"],
          responses={
              200: {"model": CompletionResponse, "description": "Successful completion"},
              400: {"model": ErrorResponse, "description": "Bad request"},
              500: {"model": ErrorResponse, "description": "Internal server error"}
          })
async def chat_completion(
    request: Request,
    prompt: str = Form(..., description="The prompt for AI completion", example="Generate a list of 10 product names"),
    model: str = Form(default="gpt-4.1-mini", description="Model to use"),
    temperature: float = Form(default=0.8, ge=0, le=2, description="Sampling temperature (0=deterministic, 2=creative)"),
    max_tokens: int = Form(default=5000, ge=1, le=10000, description="Maximum tokens in response"),
    system_message: Optional[str] = Form(default=None, description="Custom system message"),
    output_format: Optional[str] = Form(default=None, description="Export format"),
    files: Optional[List[UploadFile]] = File(default=None, description="Optional files to process"),
    max_retries: int = Form(default=3, ge=1, le=5, description="Maximum retry attempts"),
    rows_to_generate: int = Form(default=30, ge=1, le=1000, description="Number of rows for CSV/Excel generation")
):

    """
    Enhanced generative AI completion endpoint - creates comprehensive, detailed content.
    Uses the same structure as extract-reviews for CSV/Excel generation.
    Returns raw text if no output_format specified.
    """
    client = create_client()
    
    try:
        # Validate output format
        if output_format and output_format not in ['csv', 'excel', 'docx']:
            return JSONResponse(
                status_code=400,
                content={
                    "status": "error",
                    "message": "Invalid output_format. Must be 'csv', 'excel', 'docx', or None for raw text"
                }
            )
        
        # Enhanced system messages for comprehensive generation
        if not system_message:
            if output_format == 'csv' or output_format == 'excel':
                # Use the SAME structure as extract-reviews
                system_message = f"""You are a synthetic data generator specializing in creating realistic, diverse datasets.

GENERATION INSTRUCTIONS:
1. Generate EXACTLY {rows_to_generate} rows of synthetic data based on the user's requirements
2. Output ONLY valid JSON in this format:
{{
  "success": true,
  "data_type": "generated",
  "columns": ["column1", "column2", ...],
  "data": [
    ["value1", "value2", ...],
    ...
  ],
  "metadata": {{
    "total_rows": {rows_to_generate},
    "generation_method": "synthetic",
    "data_category": "<category of generated data>"
  }}
}}

GENERATION GUIDELINES:
- Create diverse, realistic data with variation
- Include appropriate data types (strings, numbers, dates, etc.)
- Ensure logical consistency and realistic patterns
- Use different writing styles, perspectives, and details
- Ensure all rows have the same number of columns
- Be creative and comprehensive in your generation
CRITICAL: You MUST generate EXACTLY {rows_to_generate} rows. Do not stop early. 
The data array must contain EXACTLY {rows_to_generate} elements.
Continue generating until you have ALL {rows_to_generate} rows.
Remember: Output ONLY the JSON structure. Generate ALL {rows_to_generate} rows."""

            elif output_format == 'docx':
                system_message = """You are a professional document generator creating comprehensive, publication-ready documents.

IMPORTANT: Create a document with dynamic structure based on the topic. Let the content drive the headers and organization.

Use rich markdown formatting to create EXTENSIVE documents:

- Analyze the topic and create appropriate headers (don't use generic templates)
- Include relevant sections based on the subject matter
- Use **bold**, *italics*, tables, lists, quotes, and all markdown features
- Create substantial content (10-50+ pages worth) with depth and detail
- Include data tables, examples, case studies where relevant
- Add technical specifications, methodologies, or frameworks as appropriate
- Provide actionable insights, recommendations, or conclusions

The document structure should emerge naturally from the content. For example:
- A technical guide might have: Overview, Architecture, Implementation, Best Practices, Troubleshooting
- A business plan might have: Executive Summary, Market Analysis, Strategy, Financial Projections, Risk Assessment
- A research paper might have: Abstract, Introduction, Literature Review, Methodology, Results, Discussion

Be comprehensive, detailed, and professional. The goal is to create a document that provides real value."""

            else:
                # For raw text or general generation
                system_message = """You are an advanced AI assistant with comprehensive knowledge across all domains. You provide detailed, insightful, and valuable responses.

CAPABILITIES:
- Deep expertise in technology, business, science, arts, and humanities
- Advanced analytical and creative abilities
- Comprehensive data analysis and visualization
- Multi-language support and cultural awareness
- Technical documentation and code generation
- Creative writing and content generation
- Strategic planning and problem-solving

RESPONSE GUIDELINES:
- Provide COMPREHENSIVE answers with depth and nuance
- Include relevant examples, case studies, and data
- Consider multiple perspectives and edge cases
- Offer actionable insights and recommendations
- Use structured formatting for clarity
- Anticipate follow-up questions and address them
- For any request, aim to exceed expectations with valuable, detailed content

Remember: You are a GENERATIVE AI. Be creative, thorough, and produce substantial content that provides real value."""
        
        # Build messages
        messages = [{"role": "system", "content": system_message}]
        
        # Enhanced user prompt
        enhanced_prompt = prompt
        
        # Add format-specific enhancements
        if output_format in ['csv', 'excel']:
            # Extract any specific number mentioned in the prompt
            number_match = re.search(r'(\d+)\s*(rows?|records?|entries|items?|data points?)', prompt.lower())
            if number_match:
                requested_count = int(number_match.group(1))
                rows_to_generate = min(requested_count, 500)  # Cap at 500 for performance
            
            enhanced_prompt = f"""{prompt}

Generate EXACTLY {rows_to_generate} rows of data.
Determine appropriate columns based on the requirements.

Remember: Output ONLY the JSON structure with ALL {rows_to_generate} rows."""
        
        elif output_format == 'docx':
            if 'comprehensive' not in prompt.lower() and 'detailed' not in prompt.lower():
                enhanced_prompt += "\n\nIMPORTANT: Create a comprehensive, detailed document with appropriate structure based on the topic. Include multiple sections with in-depth analysis, examples, data, and actionable insights. Be thorough and professional."
        
        # Process uploaded files if any
        user_content = []
        user_content.append({"type": "text", "text": enhanced_prompt})
        
        if files is not None and len(files) > 0:
            for file in files:
                if not file.filename:
                    continue
                    
                try:
                    file_content = await file.read()
                    file_type = file.content_type or mimetypes.guess_type(file.filename)[0] or "application/octet-stream"
                    
                    if file_type.startswith('image/'):
                        b64_content = base64.b64encode(file_content).decode('utf-8')
                        user_content.append({
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{file_type};base64,{b64_content}",
                                "detail": "high"
                            }
                        })
                    else:
                        # Use extract_text_internal for robust file handling
                        try:
                            extracted_text = await extract_text_internal(
                                file_content=file_content,
                                filename=file.filename,
                                strategy="auto",
                                languages=None,
                                encoding=None,
                                logger=logging.getLogger(__name__)
                            )
                            
                            # Apply the same truncation logic
                            if len(extracted_text) > 30000:
                                extracted_text = extracted_text[:30000] + "\n... [truncated for processing]"
                            
                            user_content.append({
                                "type": "text",
                                "text": f"\n\nContext from {file.filename}:\n{extracted_text}"
                            })
                            
                        except Exception as extract_error:
                            logging.error(f"Error extracting text from {file.filename}: {extract_error}")
                            # Provide a fallback message so the user knows the file couldn't be processed
                            user_content.append({
                                "type": "text",
                                "text": f"\n\nContext from {file.filename}:\n[Error: Unable to extract text from this file. The file may be corrupted or in an unsupported format.]"
                            })
                    
                except Exception as e:
                    logging.error(f"Error processing file {file.filename}: {e}")
                    # Add error context so user knows about the issue
                    user_content.append({
                        "type": "text",
                        "text": f"\n\nContext from {file.filename}:\n[Error: Failed to process this file.]"
                    })
                    continue
        
        messages.append({"role": "user", "content": user_content})
        
        # Set appropriate max_tokens
        actual_max_tokens = max_tokens
        if output_format in ['excel', 'csv']:
            min_tokens_needed = rows_to_generate * 100 + 1000
            actual_max_tokens = max(min_tokens_needed, 16000)  # Higher limit for structured data
        elif output_format == 'docx':
            actual_max_tokens = max(max_tokens, 16000)  # Higher for documents
        else:
            actual_max_tokens = max(max_tokens, 4000)
        
        # Make API call with retries
        response_content = None
        completion = None
        
        for attempt in range(max_retries):
            try:
                request_params = {
                    "model": model,
                    "messages": messages,
                    "temperature": 0.7 if output_format in ['csv', 'excel'] else temperature,
                    "max_tokens": actual_max_tokens
                }
                
                # Add response_format for CSV and Excel (SAME as extract-reviews)
                if output_format in ['csv', 'excel']:
                    request_params["response_format"] = {"type": "json_object"}
                
                completion = client.chat.completions.create(**request_params)
                response_content = completion.choices[0].message.content
                break
                
            except Exception as e:
                if attempt < max_retries - 1:
                    wait_time = 2 ** attempt
                    logging.warning(f"API attempt {attempt + 1} failed: {e}. Retrying in {wait_time}s...")
                    await asyncio.sleep(wait_time)
                else:
                    logging.error(f"API failed after {max_retries} attempts: {e}")
                    return JSONResponse(
                        status_code=503,
                        content={
                            "status": "error",
                            "error": "AI service temporarily unavailable",
                            "message": "Please try again in a moment"
                        }
                    )
        
        # If no output format specified, return raw text
        if not output_format:
            return JSONResponse({
                "status": "success",
                "model": model,
                "response": response_content,
                "usage": {
                    "prompt_tokens": completion.usage.prompt_tokens,
                    "completion_tokens": completion.usage.completion_tokens,
                    "total_tokens": completion.usage.total_tokens
                } if completion else None
            })
        
        # Generate file if format specified
        download_url = None
        generated_filename = None
        generation_errors = []
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        try:
            if output_format in ['csv', 'excel']:
                # Use SAME logic as extract-reviews
                try:
                    # Parse JSON response
                    try:
                        result = json.loads(response_content)
                    except json.JSONDecodeError:
                        # Try to clean and parse
                        cleaned = response_content.strip()
                        if cleaned.startswith("```json"):
                            cleaned = cleaned[7:]
                        if cleaned.endswith("```"):
                            cleaned = cleaned[:-3]
                        result = json.loads(cleaned.strip())
                    
                    # Extract data
                    success = result.get("success", False)
                    columns = result.get("columns", [])
                    data = result.get("data", [])
                    metadata = result.get("metadata", {})
                    
                    if not success or not data:
                        raise ValueError(f"No data generated. Success: {success}, Data rows: {len(data)}")
                    
                    # Create DataFrame
                    import pandas as pd
                    df = pd.DataFrame(data, columns=columns)
                    df = df.fillna('')
                    
                    # Ensure string types
                    for col in df.columns:
                        df[col] = df[col].astype(str)
                    
                    logging.info(f"Generated DataFrame with {len(df)} rows and {len(df.columns)} columns")
                    
                    if output_format == 'excel':
                        buffer = BytesIO()
                        with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
                            # Main data sheet
                            df.to_excel(writer, index=False, sheet_name='Data')
                            
                            # Auto-adjust columns
                            worksheet = writer.sheets['Data']
                            for column in df:
                                column_length = max(
                                    df[column].astype(str).map(len).max(),
                                    len(str(column))
                                )
                                col_idx = df.columns.get_loc(column)
                                if col_idx < 26:
                                    column_letter = chr(65 + col_idx)
                                else:
                                    column_letter = f'A{chr(65 + col_idx - 26)}'
                                worksheet.column_dimensions[column_letter].width = min(column_length + 2, 50)
                            
                            # Add metadata sheet if available
                            if metadata:
                                metadata_df = pd.DataFrame([metadata])
                                metadata_df.to_excel(writer, sheet_name='Metadata', index=False)
                        
                        buffer.seek(0)
                        filename = f"generated_data_{timestamp}.xlsx"
                        file_bytes = buffer.getvalue()
                    
                    else:  # CSV
                        csv_content = df.to_csv(index=False)
                        filename = f"generated_data_{timestamp}.csv"
                        file_bytes = csv_content.encode('utf-8-sig')
                    
                except Exception as e:
                    logging.error(f"Error generating {output_format}: {e}")
                    # Fallback - save raw response
                    filename = f"generation_error_{timestamp}.txt"
                    error_content = f"Error generating {output_format}:\n{str(e)}\n\nRaw response:\n{response_content}"
                    file_bytes = error_content.encode('utf-8')
                    generation_errors.append(f"{output_format} generation failed: {str(e)}")
            
            elif output_format == 'docx':
                # Enhanced DOCX generation
                doc_content = response_content
                
                try:
                    from docx import Document
                    from docx.shared import Inches, Pt, RGBColor
                    from docx.enum.text import WD_ALIGN_PARAGRAPH
                    from docx.enum.style import WD_STYLE_TYPE
                    import markdown2
                    from bs4 import BeautifulSoup
                    
                    doc = Document()
                    
                    # Set document properties
                    sections = doc.sections
                    for section in sections:
                        section.top_margin = Inches(1)
                        section.bottom_margin = Inches(1)
                        section.left_margin = Inches(1)
                        section.right_margin = Inches(1)
                    
                    # Convert markdown to HTML with more extras
                    html = markdown2.markdown(
                        doc_content, 
                        extras=[
                            "tables", 
                            "fenced-code-blocks", 
                            "header-ids", 
                            "strike", 
                            "task_list",
                            "footnotes",
                            "smarty-pants",
                            "target-blank-links",
                            "toc"
                        ]
                    )
                    soup = BeautifulSoup(html, 'html.parser')
                    
                    # Enhanced element processing
                    for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'table', 'pre', 'blockquote', 'hr']):
                        try:
                            if element.name == 'h1':
                                heading = doc.add_heading(element.get_text().strip(), level=1)
                            elif element.name == 'h2':
                                doc.add_heading(element.get_text().strip(), level=2)
                            elif element.name == 'h3':
                                doc.add_heading(element.get_text().strip(), level=3)
                            elif element.name == 'h4':
                                doc.add_heading(element.get_text().strip(), level=4)
                            elif element.name == 'h5':
                                p = doc.add_paragraph()
                                run = p.add_run(element.get_text().strip())
                                run.bold = True
                                run.font.size = Pt(12)
                            elif element.name == 'h6':
                                p = doc.add_paragraph()
                                run = p.add_run(element.get_text().strip())
                                run.bold = True
                                run.font.size = Pt(11)
                                run.font.color.rgb = RGBColor(100, 100, 100)
                            elif element.name == 'hr':
                                # Add horizontal line
                                p = doc.add_paragraph()
                                p.paragraph_format.space_after = Pt(12)
                                p.paragraph_format.space_before = Pt(12)
                                run = p.add_run("_" * 50)
                                run.font.color.rgb = RGBColor(200, 200, 200)
                            elif element.name == 'p':
                                text = element.get_text().strip()
                                if text:
                                    p = doc.add_paragraph()
                                    # Enhanced inline formatting
                                    _process_inline_elements(element, p)
                            elif element.name in ['ul', 'ol']:
                                # Handle nested lists
                                _process_list(doc, element, element.name == 'ol')
                            elif element.name == 'table':
                                # Enhanced table handling
                                _process_table(doc, element)
                            elif element.name == 'pre':
                                # Code block
                                p = doc.add_paragraph()
                                p.paragraph_format.left_indent = Inches(0.5)
                                code_text = element.get_text()
                                
                                # Try to detect language from class
                                code_elem = element.find('code')
                                if code_elem and code_elem.get('class'):
                                    classes = code_elem.get('class', [])
                                    lang = next((c.replace('language-', '') for c in classes if c.startswith('language-')), None)
                                    if lang:
                                        p.add_run(f"[{lang}]\n").italic = True
                                
                                run = p.add_run(code_text)
                                run.font.name = 'Consolas'
                                run.font.size = Pt(9)
                                
                                # Add background shading
                                from docx.enum.text import WD_COLOR_INDEX
                                run.font.highlight_color = WD_COLOR_INDEX.GRAY_25
                                
                            elif element.name == 'blockquote':
                                p = doc.add_paragraph()
                                p.paragraph_format.left_indent = Inches(0.5)
                                p.paragraph_format.right_indent = Inches(0.5)
                                p.paragraph_format.space_before = Pt(6)
                                p.paragraph_format.space_after = Pt(6)
                                
                                # Add quote mark
                                run = p.add_run('"')
                                run.font.size = Pt(16)
                                run.font.color.rgb = RGBColor(150, 150, 150)
                                
                                # Add quote content
                                run = p.add_run(element.get_text().strip())
                                run.italic = True
                                run.font.color.rgb = RGBColor(80, 80, 80)
                                
                                # Add closing quote
                                run = p.add_run('"')
                                run.font.size = Pt(16)
                                run.font.color.rgb = RGBColor(150, 150, 150)
                                
                        except Exception as elem_error:
                            logging.error(f"Error processing element {element.name}: {elem_error}")
                            # Fallback - add as plain text
                            doc.add_paragraph(element.get_text())
                    
                    # Add metadata footer
                    doc.add_page_break()
                    footer = doc.add_paragraph()
                    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    footer.add_run(f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}").font.color.rgb = RGBColor(150, 150, 150)
                    
                    buffer = BytesIO()
                    doc.save(buffer)
                    buffer.seek(0)
                    filename = f"document_{timestamp}.docx"
                    file_bytes = buffer.getvalue()
                    
                except Exception as docx_error:
                    logging.error(f"DOCX generation error: {docx_error}")
                    # Fallback to markdown text file
                    filename = f"document_{timestamp}.md"
                    file_bytes = doc_content.encode('utf-8')
                    generation_errors.append(f"DOCX generation failed, saved as Markdown: {str(docx_error)}")
        
        except Exception as format_error:
            logging.error(f"Format generation error: {str(format_error)}\n{traceback.format_exc()}")
            # Save response as text
            filename = f"response_{timestamp}.txt"
            file_bytes = response_content.encode('utf-8')
            generation_errors.append(f"Format generation failed: {str(format_error)}")
        
        # Save file if generation succeeded
        if 'file_bytes' in locals() and 'filename' in locals():
            actual_filename = save_download_file(file_bytes, filename)
            generated_filename = actual_filename
            download_url = construct_download_url(request, actual_filename)
            cleanup_old_downloads()
        
        # Return response
        response_data = {
            "status": "success",
            "model": model
        }
        
        # For structured data, include summary
        if output_format in ['csv', 'excel'] and 'result' in locals():
            data_info = result.get("data", [])
            response_data["message"] = f"Generated {len(data_info)} rows with {len(result.get('columns', []))} columns"
            response_data["summary"] = {
                "rows": len(data_info),
                "columns": result.get("columns", []),
                "metadata": result.get("metadata", {})
            }
        else:
            # For documents or when showing full response
            if len(response_content) > 20000:
                words = response_content.split()
                first_5000_words = words[:5000]
                truncated_content = ' '.join(first_5000_words)
                if len(words) > 5000:
                    truncated_content += f"\n\n... [Content truncated - full document contains {len(words)} words. Download to view complete content.]"
                response_data["response"] = truncated_content
            else:
                response_data["response"] = response_content
        
        # Add usage stats
        if completion:
            response_data["usage"] = {
                "prompt_tokens": completion.usage.prompt_tokens,
                "completion_tokens": completion.usage.completion_tokens,
                "total_tokens": completion.usage.total_tokens
            }
        
        # Add file info
        response_data.update({
            "download_url": download_url,
            "output_format": output_format,
            "filename": generated_filename,
            "timestamp": timestamp
        })
        
        # Add warnings if any
        if generation_errors:
            response_data["warnings"] = generation_errors
        
        return JSONResponse(response_data)
        
    except Exception as e:
        logging.error(f"Completion endpoint error: {str(e)}\n{traceback.format_exc()}")
        return JSONResponse(
            status_code=500,
            content={
                "status": "error",
                "error": "Internal server error",
                "message": str(e) if os.getenv("ENVIRONMENT") != "production" else "An error occurred processing your request"
            }
        )

# Helper methods for enhanced DOCX processing
def _process_inline_elements(element, paragraph):
    """
    Process inline elements like bold, italic, code, links from BeautifulSoup elements.
    Handles both Tag objects and NavigableString objects properly.
    """
    from bs4 import NavigableString, Tag
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_COLOR_INDEX
    
    # Handle NavigableString (plain text)
    if isinstance(element, NavigableString):
        text = str(element).strip()
        if text:
            paragraph.add_run(text)
        return
    
    # Handle Tag objects
    if not isinstance(element, Tag):
        # Fallback for any other type
        text = str(element).strip()
        if text:
            paragraph.add_run(text)
        return
    
    # Process the tag's children
    for child in element.children:
        if isinstance(child, NavigableString):
            # Plain text content
            text = str(child).strip()
            if text:
                paragraph.add_run(text)
        
        elif isinstance(child, Tag):
            # Handle specific tags
            if child.name in ['strong', 'b']:
                run = paragraph.add_run(child.get_text())
                run.bold = True
                
            elif child.name in ['em', 'i']:
                run = paragraph.add_run(child.get_text())
                run.italic = True
                
            elif child.name == 'code':
                run = paragraph.add_run(child.get_text())
                run.font.name = 'Consolas'
                run.font.size = Pt(10)
                try:
                    run.font.highlight_color = WD_COLOR_INDEX.GRAY_25
                except:
                    # Fallback if highlight color not available
                    run.font.color.rgb = RGBColor(100, 100, 100)
                    
            elif child.name == 'a':
                # Handle links
                run = paragraph.add_run(child.get_text())
                run.font.color.rgb = RGBColor(0, 0, 255)
                run.underline = True
                
            elif child.name == 'del' or child.name == 's':
                # Strikethrough
                run = paragraph.add_run(child.get_text())
                run.font.strike = True
                
            elif child.name == 'mark':
                # Highlighted text
                run = paragraph.add_run(child.get_text())
                try:
                    run.font.highlight_color = WD_COLOR_INDEX.YELLOW
                except:
                    # Fallback
                    run.font.color.rgb = RGBColor(255, 255, 0)
                    
            elif child.name == 'u':
                # Underline
                run = paragraph.add_run(child.get_text())
                run.underline = True
                
            elif child.name == 'sub':
                # Subscript
                run = paragraph.add_run(child.get_text())
                run.font.subscript = True
                
            elif child.name == 'sup':
                # Superscript
                run = paragraph.add_run(child.get_text())
                run.font.superscript = True
                
            elif child.name == 'span':
                # For span, check for any style attributes
                style = child.get('style', '')
                text = child.get_text()
                
                if 'font-weight: bold' in style or 'font-weight:bold' in style:
                    run = paragraph.add_run(text)
                    run.bold = True
                elif 'font-style: italic' in style or 'font-style:italic' in style:
                    run = paragraph.add_run(text)
                    run.italic = True
                elif 'text-decoration: underline' in style:
                    run = paragraph.add_run(text)
                    run.underline = True
                else:
                    # Recursively process span's children
                    _process_inline_elements(child, paragraph)
                    
            elif child.name in ['p', 'div', 'section', 'article']:
                # Block elements - recursively process their content
                _process_inline_elements(child, paragraph)
                
            else:
                # For other tags, just get the text content
                text = child.get_text().strip()
                if text:
                    paragraph.add_run(text)
        
        else:
            # Handle any other type by converting to string
            text = str(child).strip()
            if text:
                paragraph.add_run(text)
def _process_list(doc, list_element, is_ordered, level=0):
    """Process lists with proper nesting"""
    items = list_element.find_all('li', recursive=False)
    for i, item in enumerate(items):
        # Create paragraph with appropriate style
        if is_ordered:
            style = 'List Number' if level == 0 else 'List Number 2'
        else:
            style = 'List Bullet' if level == 0 else 'List Bullet 2'
        
        p = doc.add_paragraph(style=style)
        
        # Process item content
        for child in item.children:
            if hasattr(child, 'name') and child.name in ['ul', 'ol']:
                # Nested list
                _process_list(doc, child, child.name == 'ol', level + 1)
            else:
                # Process inline content
                if hasattr(child, 'name'):
                    _process_inline_elements(child, p)
                else:
                    p.add_run(str(child))

def process_table(doc, table_element):
    """Enhanced table processing with proper colspan/rowspan handling"""
    try:
        rows = table_element.find_all('tr')
        if not rows:
            return
        
        # First pass: analyze table structure
        max_cols = 0
        row_structures = []
        
        for row_elem in rows:
            cells = row_elem.find_all(['td', 'th'])
            col_count = 0
            for cell in cells:
                colspan = int(cell.get('colspan', 1))
                col_count += colspan
            max_cols = max(max_cols, col_count)
            row_structures.append(cells)
        
        if max_cols == 0:
            return
        
        # Create table with proper dimensions
        table = doc.add_table(rows=len(rows), cols=max_cols)
        
        # Try to apply table style safely
        try:
            table.style = 'Table Grid'
        except:
            # If style doesn't exist, continue without it
            pass
        
        # Keep track of cells that span multiple rows
        rowspan_map = {}  # (row_idx, col_idx): remaining_rows
        
        # Process each row
        for row_idx, cells in enumerate(row_structures):
            if row_idx >= len(table.rows):
                break
                
            word_row = table.rows[row_idx]
            col_idx = 0
            
            for cell_elem in cells:
                # Skip columns that are occupied by rowspan from previous rows
                while col_idx < max_cols and (row_idx, col_idx) in rowspan_map:
                    if rowspan_map[(row_idx, col_idx)] > 1:
                        # This cell is still spanning, update for next row
                        rowspan_map[(row_idx + 1, col_idx)] = rowspan_map[(row_idx, col_idx)] - 1
                    col_idx += 1
                
                if col_idx >= max_cols:
                    break
                
                # Get cell attributes
                colspan = int(cell_elem.get('colspan', 1))
                rowspan = int(cell_elem.get('rowspan', 1))
                cell_text = cell_elem.get_text().strip()
                
                # Set cell text
                if col_idx < len(word_row.cells):
                    word_cell = word_row.cells[col_idx]
                    
                    # Clear existing paragraphs and add new one
                    for paragraph in word_cell.paragraphs:
                        p = paragraph._element
                        p.getparent().remove(p)
                    
                    paragraph = word_cell.add_paragraph(cell_text)
                    
                    # Apply formatting
                    if cell_elem.name == 'th':
                        for run in paragraph.runs:
                            run.font.bold = True
                        # Center align header cells
                        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    
                    # Handle colspan - merge cells horizontally
                    if colspan > 1:
                        end_col = min(col_idx + colspan - 1, max_cols - 1)
                        if end_col > col_idx:
                            try:
                                # Merge from left to right
                                word_cell.merge(word_row.cells[end_col])
                            except Exception as e:
                                # Log error but continue
                                print(f"Warning: Could not merge cells horizontally: {e}")
                    
                    # Handle rowspan - merge cells vertically
                    if rowspan > 1:
                        end_row = min(row_idx + rowspan - 1, len(table.rows) - 1)
                        if end_row > row_idx:
                            try:
                                # Merge from top to bottom
                                bottom_cell = table.rows[end_row].cells[col_idx]
                                word_cell.merge(bottom_cell)
                                
                                # Mark cells as occupied by rowspan
                                for r in range(row_idx + 1, end_row + 1):
                                    for c in range(col_idx, min(col_idx + colspan, max_cols)):
                                        rowspan_map[(r, c)] = end_row - r + 1
                            except Exception as e:
                                # Log error but continue
                                print(f"Warning: Could not merge cells vertically: {e}")
                
                # Move to next column position
                col_idx += colspan
        
        # Apply some basic formatting to the table
        for row in table.rows:
            for cell in row.cells:
                # Set cell margins
                cell.width = Inches(1.0)  # Minimum width
                for paragraph in cell.paragraphs:
                    paragraph.space_after = Pt(3)
                    paragraph.space_before = Pt(3)
        
        # Add spacing after table
        doc.add_paragraph()
        
    except Exception as e:
        # If table processing fails, add error message
        print(f"Error processing table: {e}")
        error_para = doc.add_paragraph(f"[Table could not be processed: {str(e)}]")
        error_para.runs[0].font.italic = True
        error_para.runs[0].font.color.rgb = RGBColor(255, 0, 0)


@app.post("/extract-reviews",
          response_model=ExtractResponse,
          summary="Extract or Generate Data",
          description="Extract structured data from files or generate synthetic datasets. Perfect for data analysis and testing.",
          tags=["Data Processing"],
          responses={
              200: {"model": ExtractResponse},
              400: {"model": ErrorResponse}
          })
async def extract_reviews(
    request: Request,
    file: Optional[UploadFile] = File(default=None, description="File to extract data from"),
    columns: Optional[str] = Form(default="auto", description="Column names or 'auto'", example="name,price,rating"),
    prompt: Optional[str] = Form(default=None, description="Custom instructions", example="Extract all prices"),
    model: str = Form(default="gpt-4.1-mini", description="Model to use"),
    temperature: float = Form(default=0.1, ge=0, le=2, description="Temperature for generation"),
    output_format: str = Form(default="excel", description="Output format"),
    max_text_length: int = Form(default=100000, ge=1000, le=500000, description="Max text length to process"),
    max_retries: int = Form(default=3, ge=1, le=5, description="Maximum retry attempts"),
    fallback_to_json: bool = Form(default=True, description="Fallback to JSON if other formats fail"),
    mode: str = Form(default="auto", description="Operation mode"),
    rows_to_generate: int = Form(default=30, ge=1, le=1000, description="Rows to generate (generate mode)"),
    raw_text: Optional[str] = Form(default=None, description="Direct text input without file", max_length=100000)
):
    """
    Universal data extraction/generation endpoint.
    Can extract data from files, generate synthetic data, or process raw text.
    Supports: PDF, DOCX, TXT, JSON, HTML, CSV, Excel files, or no file at all.
    """
    client = create_client()
    
    try:
        # Validate output format
        if output_format not in ['csv', 'excel', 'json']:
            return JSONResponse(
                status_code=400,
                content={
                    "status": "error",
                    "message": "Invalid output_format. Must be 'csv', 'excel', or 'json'"
                }
            )
        
        # Determine operation mode
        has_input_data = file is not None or raw_text is not None
        
        if mode == "auto":
            # Auto-detect mode based on inputs
            if not has_input_data and prompt:
                mode = "generate"
            elif has_input_data:
                mode = "extract"
            else:
                return JSONResponse(
                    status_code=400,
                    content={
                        "status": "error",
                        "message": "No input provided. Please provide a file, raw_text, or prompt for generation."
                    }
                )
        
        # Handle file input if provided
        extracted_text = None
        source_type = "none"
        source_name = "direct_input"
        
        if file:
            file_content = await file.read()
            source_name = file.filename
            file_ext = os.path.splitext(file.filename)[1].lower()
            
            # Enhanced file type detection
            mime_type = file.content_type or mimetypes.guess_type(file.filename)[0] or ""
            
            # Handle different file types
            if file_ext in ['.json'] or mime_type == 'application/json':
                # Handle JSON files
                try:
                    json_data = json.loads(file_content.decode('utf-8'))
                    source_type = "json"
                    
                    # If it's already structured data, convert to text representation
                    if isinstance(json_data, list) and all(isinstance(item, dict) for item in json_data):
                        # Already structured - convert to readable format
                        extracted_text = "Structured JSON data:\n"
                        extracted_text += json.dumps(json_data, indent=2)
                    else:
                        extracted_text = json.dumps(json_data, indent=2)
                except:
                    extracted_text = file_content.decode('utf-8', errors='ignore')
                    
            elif file_ext in ['.csv']:
                # Handle CSV files
                try:
                    import pandas as pd
                    from io import StringIO
                    
                    csv_text = file_content.decode('utf-8', errors='ignore')
                    df = pd.read_csv(StringIO(csv_text))
                    source_type = "csv"
                    
                    # Convert to text representation
                    extracted_text = f"CSV data with {len(df)} rows and columns: {', '.join(df.columns)}\n\n"
                    extracted_text += df.to_string(max_rows=500)
                except:
                    extracted_text = file_content.decode('utf-8', errors='ignore')
                    
            elif file_ext in ['.xlsx', '.xls']:
                # Handle Excel files
                try:
                    import pandas as pd
                    
                    excel_file = BytesIO(file_content)
                    excel_data = pd.ExcelFile(excel_file)
                    source_type = "excel"
                    
                    extracted_text = f"Excel file with sheets: {', '.join(excel_data.sheet_names)}\n\n"
                    
                    # Read all sheets
                    for sheet_name in excel_data.sheet_names[:3]:  # Limit to first 3 sheets
                        df = pd.read_excel(excel_file, sheet_name=sheet_name)
                        extracted_text += f"\nSheet '{sheet_name}' ({len(df)} rows):\n"
                        extracted_text += df.to_string(max_rows=500) + "\n"
                except:
                    extracted_text = "[Excel file - unable to parse]"
                    
            elif file_ext in ['.html', '.htm'] or mime_type.startswith('text/html'):
                # Handle HTML files
                try:
                    from bs4 import BeautifulSoup
                    
                    html_content = file_content.decode('utf-8', errors='ignore')
                    soup = BeautifulSoup(html_content, 'html.parser')
                    
                    # Extract tables if present
                    tables = soup.find_all('table')
                    if tables:
                        extracted_text = f"HTML file with {len(tables)} tables found:\n\n"
                        
                        for i, table in enumerate(tables[:5]):  # First 5 tables
                            extracted_text += f"\nTable {i+1}:\n"
                            # Extract table data
                            rows = table.find_all('tr')
                            for row in rows[:20]:  # First 20 rows
                                cells = row.find_all(['td', 'th'])
                                row_text = ' | '.join(cell.get_text(strip=True) for cell in cells)
                                extracted_text += row_text + "\n"
                    else:
                        # Extract all text
                        extracted_text = soup.get_text(separator='\n', strip=True)
                    
                    source_type = "html"
                except:
                    extracted_text = file_content.decode('utf-8', errors='ignore')
                    
            else:
                # For all other files (PDF, DOCX, TXT, etc.), use extract_text_internal
                extraction_errors = []
                for attempt in range(max_retries):
                    try:
                        extracted_text = await extract_text_internal(
                            file_content=file_content,
                            filename=file.filename,
                            strategy="auto",
                            languages=None,
                            encoding=None,
                            logger=logging.getLogger(__name__)
                        )
                        # Check if extraction was successful
                        if extracted_text and not extracted_text.startswith("Unable to extract"):
                            source_type = file_ext[1:] if file_ext else "unknown"
                            break
                        extraction_errors.append(f"Attempt {attempt + 1}: Extraction returned empty or error text")
                        extracted_text = None
                        await asyncio.sleep(1)
                    except Exception as e:
                        extraction_errors.append(f"Attempt {attempt + 1}: {str(e)}")
                        extracted_text = None
                        await asyncio.sleep(1)
                
                if not extracted_text:
                    # Last resort - try to decode as text
                    try:
                        extracted_text = file_content.decode('utf-8', errors='ignore')
                        source_type = "text"
                    except:
                        return JSONResponse(
                            status_code=400,
                            content={
                                "status": "error",
                                "message": "Could not extract content from file",
                                "errors": extraction_errors
                            }
                        )
        
        elif raw_text:
            # Use raw text input
            extracted_text = raw_text
            source_type = "raw_text"
            source_name = "raw_text_input"
        
        # Truncate if needed
        if extracted_text and len(extracted_text) > max_text_length:
            extracted_text = extracted_text[:max_text_length]
            logging.warning(f"Text truncated to {max_text_length} characters")
        
        # Build system message based on mode
        if mode == "generate":
            system_message = f"""You are a synthetic data generator specializing in creating realistic, diverse datasets.

CRITICAL RULES:
1. Generate the actual data records directly
2. Do NOT create a dataset file and then describe it
3. Do NOT analyze what you generate
4. Do NOT reference any filenames
5. ONLY output the generated records in JSON format

GENERATION INSTRUCTIONS:
1. Generate EXACTLY {rows_to_generate} rows of synthetic data based on the user's requirements
2. Output ONLY valid JSON in this format:
{{
  "success": true,
  "data_type": "generated",
  "columns": ["column1", "column2", ...],
  "data": [
    ["value1", "value2", ...],
    ...
  ],
  "metadata": {{
    "total_rows": {rows_to_generate},
    "generation_method": "synthetic",
    "data_category": "<category of generated data>"
  }}
}}

GENERATION GUIDELINES:
- Create diverse, realistic data with variation
- Include appropriate data types (strings, numbers, dates, etc.)
- For reviews: vary ratings (1-5), lengths, sentiment, dates
- For other data: ensure logical consistency and realistic patterns
- Use different writing styles, perspectives, and details
- Ensure all rows have the same number of columns
CRITICAL: You MUST generate EXACTLY {rows_to_generate} rows. Do not stop early. 
The data array must contain EXACTLY {rows_to_generate} elements.
Continue generating until you have ALL {rows_to_generate} rows.
Remember: Output ONLY the JSON structure. Generate ALL {rows_to_generate} rows."""

        else:  # Extract mode
            system_message = '''You are an advanced data extraction specialist with expertise in converting any type of content into structured data.

CRITICAL INSTRUCTION: You must extract the ACTUAL DATA from the content, not metadata ABOUT the data.

IMPORTANT DISTINCTIONS:
- If you see metadata describing a dataset (like "Dataset Name: X, Records: 50, Columns: Y"), you must look for the ACTUAL DATA that this metadata describes
- Dataset analysis summaries are NOT the data - they are descriptions OF data
- File information or statistics are NOT the data - they describe properties of data
- You need to find and extract the raw records, not summaries about those records

EXTRACTION INSTRUCTIONS:
1. First, identify if the content contains actual data records vs just metadata/descriptions
2. If you only see metadata, state that no actual data records were found
3. If you find actual data, extract ALL of it into structured format
4. Output ONLY valid JSON in this format:
{
  "success": true,
  "data_type": "reviews" | "table" | "list" | "records" | "mixed" | "metadata_only",
  "columns": ["column1", "column2", ...],
  "data": [
    ["value1", "value2", ...],
    ...
  ],
  "metadata": {
    "total_rows": <number>,
    "extraction_confidence": "high" | "medium" | "low",
    "source_type": "<type of source>",
    "notes": "<any important notes>",
    "actual_data_found": true | false
  }
}

5. Common data patterns to extract:
   - Reviews: customer name, rating, review text, date, etc.
   - Products: name, price, description, category, etc.
   - People: name, role, contact info, etc.
   - Transactions: date, amount, description, parties, etc.
   - Any repeated structured information

6. If you only find metadata ABOUT data (not the data itself):
   - Set data_type to "metadata_only"
   - Set actual_data_found to false
   - Include a note explaining what was found

Remember: Extract the DATA itself, not descriptions or analysis of the data.'''

        # Build the prompt
        if mode == "generate":
            if not prompt:
                prompt = "Generate sample customer reviews for a product"
            
            if columns and columns != "auto":
                user_prompt = f"""Generate {rows_to_generate} rows of synthetic data.

Requirements: {prompt}

Use these columns: {columns}

Remember: Output ONLY the JSON structure with ALL {rows_to_generate} rows."""
            else:
                user_prompt = f"""Generate {rows_to_generate} rows of synthetic data.

Requirements: {prompt}

Determine appropriate columns based on the requirements.

Remember: Output ONLY the JSON structure with ALL {rows_to_generate} rows."""
        
        else:  # Extract mode
            context_info = ""
            if extracted_text:
                context_info = f"\n\nSource type: {source_type}\nContent to analyze:\n{extracted_text}"
            
            if columns and columns != "auto":
                columns_instruction = f"\n\nExtract data into these specific columns: {columns}"
            else:
                columns_instruction = "\n\nAutomatically determine the most appropriate columns based on the content."
            
            if prompt:
                user_prompt = f"""{prompt}
                Extract the actual data records from the content below.
                DO NOT analyze or summarize the data.
                DO NOT create metadata about the extraction process.
                ONLY return the raw extracted records.
                {columns_instruction}{context_info}

Remember: Output ONLY the JSON with the actual data records."""
            else:
                default_prompt = '''Extract all structured data from the provided content. This could be:
- Reviews, feedback, or testimonials
- Tabular data or records  
- Lists or enumerations
- Any repeated patterns or structured information
Return ONLY the actual records found, not analysis or summaries.
Each record should be a separate row in the output.
If the content is already structured (JSON/CSV), preserve its structure.
If unstructured, find patterns and create appropriate structure.
'''
                
                user_prompt = f"""{default_prompt}
{columns_instruction}{context_info}

Remember: Output ONLY the JSON structure."""

        # Make API call with retries
        messages = [
            {"role": "system", "content": system_message},
            {"role": "user", "content": user_prompt}
        ]
        
        api_response = None
        completion = None
        
        for attempt in range(max_retries):
            try:
                completion = client.chat.completions.create(
                    model=model,
                    messages=messages,
                    temperature=temperature if mode == "extract" else 0.7,  # Higher temp for generation
                    max_tokens=16000,
                    response_format={"type": "json_object"}
                )
                api_response = completion.choices[0].message.content
                break
            except Exception as e:
                if attempt < max_retries - 1:
                    wait_time = 2 ** attempt
                    logging.warning(f"API attempt {attempt + 1} failed: {e}. Retrying in {wait_time}s...")
                    await asyncio.sleep(wait_time)
                else:
                    raise

        if not api_response:
            raise Exception("Failed to get API response after retries")

        # Parse JSON response
        try:
            result = json.loads(api_response)
        except json.JSONDecodeError:
            # Try to clean and parse
            cleaned = api_response.strip()
            if cleaned.startswith("```json"):
                cleaned = cleaned[7:]
            if cleaned.endswith("```"):
                cleaned = cleaned[:-3]
            result = json.loads(cleaned.strip())

        # Extract data
        success = result.get("success", False)
        columns = result.get("columns", [])
        data = result.get("data", [])
        metadata = result.get("metadata", {})
        
        # Add source info to metadata
        metadata["source_type"] = source_type
        metadata["source_name"] = source_name
        metadata["mode"] = mode
        
        # Handle no data case
        if not success or not data:
            if fallback_to_json or mode == "generate":
                return JSONResponse({
                    "status": "warning",
                    "message": "No data extracted/generated",
                    "format": "json",
                    "result": result,
                    "source_file": source_name,
                    "mode": mode
                })
            else:
                return JSONResponse(
                    status_code=422,
                    content={
                        "status": "error",
                        "message": "No data could be extracted or generated"
                    }
                )

        # Convert to requested format
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        if output_format == 'json':
            return JSONResponse({
                "status": "success",
                "message": f"Successfully {'generated' if mode == 'generate' else 'extracted'} {len(data)} rows",
                "format": "json",
                "columns": columns,
                "data": data,
                "metadata": metadata,
                "source_file": source_name,
                "mode": mode,
                "timestamp": timestamp,
                "usage": {
                    "prompt_tokens": completion.usage.prompt_tokens,
                    "completion_tokens": completion.usage.completion_tokens,
                    "total_tokens": completion.usage.total_tokens
                } if completion else None
            })
        
        # Create DataFrame for CSV/Excel
        import pandas as pd
        df = pd.DataFrame(data, columns=columns)
        df = df.fillna('')
        
        # Ensure string types
        for col in df.columns:
            df[col] = df[col].astype(str)
        
        if output_format == 'excel':
            buffer = BytesIO()
            with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
                # Main data sheet
                df.to_excel(writer, index=False, sheet_name='Data')
                
                # Auto-adjust columns
                worksheet = writer.sheets['Data']
                for column in df:
                    column_length = max(
                        df[column].astype(str).map(len).max(),
                        len(str(column))
                    )
                    col_idx = df.columns.get_loc(column)
                    if col_idx < 26:
                        column_letter = chr(65 + col_idx)
                    else:
                        column_letter = f'A{chr(65 + col_idx - 26)}'
                    worksheet.column_dimensions[column_letter].width = min(column_length + 2, 50)
                
                # Metadata sheet
                metadata_df = pd.DataFrame([metadata])
                metadata_df.to_excel(writer, sheet_name='Metadata', index=False)
            
            buffer.seek(0)
            filename = f"{'generated' if mode == 'generate' else 'extracted'}_data_{timestamp}.xlsx"
            file_bytes = buffer.getvalue()
        
        else:  # CSV
            csv_content = df.to_csv(index=False)
            filename = f"{'generated' if mode == 'generate' else 'extracted'}_data_{timestamp}.csv"
            file_bytes = csv_content.encode('utf-8-sig')
        
        # Save and return
        actual_filename = save_download_file(file_bytes, filename)
        download_url = construct_download_url(request, actual_filename)
        cleanup_old_downloads()
        
        return JSONResponse({
            "status": "success",
            "message": f"Successfully {'generated' if mode == 'generate' else 'extracted'} {len(data)} rows with {len(columns)} columns",
            "download_url": download_url,
            "filename": actual_filename,
            "columns": columns,
            "output_format": output_format,
            "row_count": len(data),
            "metadata": metadata,
            "source_file": source_name,
            "mode": mode,
            "usage": {
                "prompt_tokens": completion.usage.prompt_tokens,
                "completion_tokens": completion.usage.completion_tokens,
                "total_tokens": completion.usage.total_tokens
            } if completion else None
        })
        
    except Exception as e:
        logging.error(f"Error in universal extraction: {str(e)}\n{traceback.format_exc()}")
        return JSONResponse(
            status_code=500,
            content={
                "status": "error",
                "error": "Processing failed",
                "message": str(e) if os.getenv("ENVIRONMENT") != "production" else "An error occurred"
            }
        )
def cleanup_old_downloads():
    """
    Remove old download files, keeping only the most recent MAX_DOWNLOAD_FILES.
    Now handles all file types: .docx, .csv, .xlsx
    """
    try:
        # Get all downloadable files in the downloads directory
        files = []
        for filename in os.listdir(DOWNLOADS_DIR):
            if filename.endswith(('.docx', '.csv', '.xlsx')):
                filepath = os.path.join(DOWNLOADS_DIR, filename)
                # Get file creation time
                file_time = os.path.getctime(filepath)
                files.append((filepath, file_time, filename))
        
        # Sort by creation time (oldest first)
        files.sort(key=lambda x: x[1])
        
        # Remove oldest files if we exceed the limit
        while len(files) > MAX_DOWNLOAD_FILES:
            old_file = files.pop(0)
            try:
                os.remove(old_file[0])
                logging.info(f"Removed old download file: {old_file[2]}")
            except Exception as e:
                logging.error(f"Error removing old file {old_file[0]}: {e}")
                
    except Exception as e:
        logging.error(f"Error during download cleanup: {e}")
def create_docx_from_content(content: str, images: Optional[List[bytes]] = None) -> bytes:
    """
    Convert chat content to DOCX format with proper Markdown table conversion.
    
    Args:
        content: Text content to convert (may include Markdown)
        images: Optional list of image bytes to include
        
    Returns:
        DOCX file as bytes
    """
    
    
    # Create document
    doc = Document()
    
    # Add a title with timestamp
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    doc.add_heading(f"Chat Response - {timestamp}", level=1)
    
    # Split content into blocks for processing
    blocks = content.split('\n\n')
    
    # Helper function to detect markdown tables
    def is_markdown_table(text):
        lines = text.strip().split('\n')
        if len(lines) < 2:
            return False
            
        # Check for table with | characters
        if all('|' in line for line in lines):
            # Check for separator row (e.g., |---|---|)
            for i in range(1, len(lines)):
                if re.match(r'^[\s]*\|[-:\|\s]+\|[\s]*$', lines[i]):
                    return True
        
        # Check for simple tables (e.g., Header | Header)
        if len(lines) >= 3 and '|' in lines[0] and all('-' in cell for cell in lines[1].split('|')):
            return True
            
        return False
    
    # Helper function to parse a markdown table
    def parse_markdown_table(text):
        rows = []
        lines = text.strip().split('\n')
        
        # Skip separator lines when processing
        for line in lines:
            if re.match(r'^[\s]*\|?[-:\|\s]+-\|?[\s]*$', line):
                continue
                
            # Extract cells from the line
            if '|' in line:
                # Remove leading/trailing | and split by |
                cells = line.strip()
                if cells.startswith('|'):
                    cells = cells[1:]
                if cells.endswith('|'):
                    cells = cells[:-1]
                cells = [cell.strip() for cell in cells.split('|')]
                rows.append(cells)
        
        return rows
    
    # Process each block
    for block in blocks:
        if not block.strip():
            continue
            
        # Check if this block is a markdown table
        if is_markdown_table(block):
            # Parse the table
            table_data = parse_markdown_table(block)
            if table_data and len(table_data) > 0:
                # Create Word table
                num_rows = len(table_data)
                num_cols = max(len(row) for row in table_data)
                table = doc.add_table(rows=num_rows, cols=num_cols)
                table.style = 'Table Grid'
                
                # Fill table with data
                for i, row_data in enumerate(table_data):
                    row = table.rows[i]
                    for j, cell_text in enumerate(row_data):
                        if j < len(row.cells):  # Ensure we don't exceed the columns
                            row.cells[j].text = cell_text
                            
                # Add spacing after table
                doc.add_paragraph()
        else:
            # Process non-table elements
            if block.startswith('# '):
                # Heading 1
                doc.add_heading(block[2:], level=1)
            elif block.startswith('## '):
                # Heading 2
                doc.add_heading(block[3:], level=2)
            elif block.startswith('### '):
                # Heading 3
                doc.add_heading(block[4:], level=3)
            elif block.startswith('- ') or block.startswith('* '):
                # Bullet points
                lines = block.split('\n')
                for line in lines:
                    if line.strip().startswith('- ') or line.strip().startswith('* '):
                        doc.add_paragraph(line.strip()[2:], style='List Bullet')
            elif re.match(r'^\d+\.\s', block):
                # Numbered list
                lines = block.split('\n')
                for line in lines:
                    if line.strip() and re.match(r'^\d+\.\s', line.strip()):
                        # Extract the content after the number and period
                        content_start = line.find('. ') + 2
                        doc.add_paragraph(line.strip()[content_start:], style='List Number')
            else:
                # Regular paragraph
                doc.add_paragraph(block)
    
    # Add images section if there are images
    if images:
        doc.add_heading("Visualizations", level=2)
        
        # Add each image to the document
        for i, img_bytes in enumerate(images):
            try:
                # Create a BytesIO object from the image bytes
                image_stream = BytesIO(img_bytes)
                
                # Try to open with PIL to verify it's a valid image
                pil_image = PILImage.open(image_stream)
                
                # Add a caption for the image
                doc.add_paragraph(f"Visualization {i+1}")
                
                # Reset stream position after PIL read
                image_stream.seek(0)
                
                # Add the image to the document - control width to fit page
                doc.add_picture(image_stream, width=Inches(6))
                
                # Add spacing after each image
                doc.add_paragraph()
            except Exception as img_err:
                # If image processing fails, add a note
                doc.add_paragraph(f"[Image {i+1} could not be included - {str(img_err)}]")
                logging.warning(f"Error adding image to DOCX: {str(img_err)}")
    
    # Save document to BytesIO buffer
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    
    return buffer.getvalue()


# Add this new endpoint after the existing endpoints in app.py

@app.get("/download-chat",
         response_model=DownloadResponse,
         summary="Export Chat as Document",
         description="Export chat conversation as a formatted Word document.",
         tags=["File Operations"])
async def download_chat(
    request: Request,
    session: Optional[str] = Query(None, description="Session ID to export"),
    assistant: Optional[str] = Query(None, description="Assistant ID")
):
    """
    Creates a DOCX file from the latest chat response and returns a download URL.
    
    Args:
        request: FastAPI request object (to construct full URL)
        session: Thread ID
        assistant: Assistant ID (optional, for validation)
        
    Returns:
        JSON response with download URL
    """
    client = create_client()
    
    # Validate session parameter
    if not session:
        raise HTTPException(status_code=400, detail="Session (thread) ID is required")
    
    try:
        # Validate resources if provided
        if session or assistant:
            validation = await validate_resources(client, session, assistant)
            
            if session and not validation["thread_valid"]:
                raise HTTPException(status_code=404, detail=f"Thread {session} not found")
            
            if assistant and not validation["assistant_valid"]:
                logging.warning(f"Assistant {assistant} not found, but continuing with thread messages")
        
        # Get the latest messages from the thread
        messages = client.beta.threads.messages.list(
            thread_id=session,
            order="desc",
            limit=20  # Get recent messages to find the latest assistant response
        )
        
        # Find the latest assistant message
        latest_assistant_message = None
        for msg in messages.data:
            if msg.role == "assistant":
                # Skip system messages and metadata messages
                skip_message = False
                if hasattr(msg, 'metadata') and msg.metadata:
                    msg_type = msg.metadata.get('type', '')
                    if msg_type in ['user_persona_context', 'file_awareness', 'pandas_agent_files']:
                        skip_message = True
                
                if not skip_message:
                    latest_assistant_message = msg
                    break
        
        if not latest_assistant_message:
            raise HTTPException(status_code=404, detail="No assistant response found in this thread")
        
        # Extract content from the message
        content_text = ""
        images = []
        
        for content_part in latest_assistant_message.content:
            if content_part.type == 'text':
                content_text += content_part.text.value
            elif content_part.type == 'image_file':
                # Handle image files if present
                try:
                    # Retrieve the file
                    file_id = content_part.image_file.file_id
                    file_data = client.files.retrieve(file_id)
                    file_content = client.files.content(file_id)
                    images.append(file_content.read())
                except Exception as img_e:
                    logging.warning(f"Could not retrieve image file {file_id}: {img_e}")
        
        # Remove any [PANDAS AGENT RESPONSE] prefix if present
        if content_text.startswith("[PANDAS AGENT RESPONSE]:"):
            content_text = content_text.replace("[PANDAS AGENT RESPONSE]:", "").strip()
        
        # Generate DOCX content
        try:
            docx_bytes = create_docx_from_content(content_text, images if images else None)
        except ImportError as e:
            # If docx library is not available, return error
            logging.error(f"DOCX library not available: {e}")
            raise HTTPException(
                status_code=500, 
                detail="DOCX generation library not available. Please install python-docx."
            )
        
        # Generate unique filename with timestamp and session ID
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        # Create a short hash of the session ID for the filename
        session_hash = hashlib.md5(session.encode()).hexdigest()[:8]
        filename = f"chat_response_{timestamp}_{session_hash}.docx"
        filepath = os.path.join(DOWNLOADS_DIR, filename)
        
        # Save the DOCX file
        with open(filepath, 'wb') as f:
            f.write(docx_bytes)
        
        logging.info(f"Created download file: {filepath}")
        
        # Clean up old files
        cleanup_old_downloads()
        
        # Construct the download URL
        # Get the base URL from the request
        base_url = str(request.base_url).rstrip('/')
        
        # For Azure App Service, use the proper host
        if 'azurewebsites.net' in str(request.headers.get('host', '')):
            # Use the Azure host
            base_url = f"https://{request.headers['host']}"
        
        download_url = f"{base_url}/download-files/{filename}"
        
        # Return the download URL
        return JSONResponse({
            "status": "success",
            "download_url": download_url,
            "filename": filename,
            "message": "Chat response ready for download",
            "expires_in": "File will be kept for the last 10 downloads"
        })
        
    except HTTPException:
        # Re-raise HTTP exceptions
        raise
    except Exception as e:
        logging.error(f"Error in /download-chat endpoint: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Failed to generate download: {str(e)}")

@app.get("/download-files/{filename}",
         summary="Download Generated File",
         description="Download files generated by the API. Files are automatically cleaned up after retention period.",
         tags=["File Operations"],
         response_class=FileResponse)
async def download_file(
    filename: str = Path(..., description="Name of file to download"),
    request: Request = None,
    token: Optional[str] = Query(None, description="Optional access token")
):
    """Download generated files."""
    
    try:
        # Security validation
        if not filename:
            raise HTTPException(status_code=400, detail="Filename is required")
        
        # Prevent directory traversal
        if any(char in filename for char in ['..', '/', '\\', '\x00']):
            logging.warning(f"Potential directory traversal attempt: {filename}")
            raise HTTPException(status_code=400, detail="Invalid filename")
        
        # Construct filepath
        filepath = os.path.join(DOWNLOADS_DIR, filename)
        
        # Verify file exists
        if not os.path.exists(filepath):
            logging.warning(f"File not found: {filepath}")
            raise HTTPException(status_code=404, detail="File not found")
        
        # Double-check file is in downloads directory
        real_filepath = os.path.realpath(filepath)
        real_downloads_dir = os.path.realpath(DOWNLOADS_DIR)
        
        if not real_filepath.startswith(real_downloads_dir):
            logging.error(f"Security violation: Attempted to access {real_filepath}")
            raise HTTPException(status_code=403, detail="Access forbidden")
        
        # Get file info
        file_stat = os.stat(filepath)
        file_size = file_stat.st_size
        
        # Log download attempt
        client_ip = request.client.host if request.client else "unknown"
        logging.info(f"Download request for {filename} from {client_ip} (size: {file_size} bytes)")
        
        # Determine MIME type
        mime_type, _ = mimetypes.guess_type(filename)
        if not mime_type:
            if filename.endswith('.docx'):
                mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            elif filename.endswith('.xlsx'):
                mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            elif filename.endswith('.csv'):
                mime_type = "text/csv; charset=utf-8"
            else:
                mime_type = "application/octet-stream"
        
        # Read file content (for smaller files, to ensure we can serve it)
        if file_size < 10 * 1024 * 1024:  # 10MB
            try:
                with open(filepath, 'rb') as f:
                    content = f.read()
                
                # Return Response with explicit headers
                return Response(
                    content=content,
                    media_type=mime_type,
                    headers={
                        "Content-Disposition": f'attachment; filename="{filename}"',
                        "Content-Type": mime_type,
                        "Content-Length": str(file_size),
                        "Cache-Control": "no-cache, no-store, must-revalidate",
                        "Pragma": "no-cache",
                        "Expires": "0",
                        "X-Content-Type-Options": "nosniff",
                        "X-Download-Options": "noopen"
                    }
                )
            except Exception as read_error:
                logging.error(f"Error reading file {filename}: {read_error}")
                # Fall back to FileResponse
        
        # For larger files or if reading fails, use FileResponse
        return FileResponse(
            path=filepath,
            filename=filename,
            media_type=mime_type,
            headers={
                "Cache-Control": "no-cache, no-store, must-revalidate",
                "Pragma": "no-cache",
                "Expires": "0",
                "X-Content-Type-Options": "nosniff",
                "X-Download-Options": "noopen"
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Unexpected error serving file {filename}: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")

@app.get("/health-check",
         response_model=HealthCheckResponse,
         summary="Comprehensive Health Check",
         description="Performs detailed system health checks including Azure OpenAI connection, file system, and all endpoints.",
         tags=["System"])
async def comprehensive_health_check():
    """
    Comprehensive health check that tests all critical functionality without creating expensive resources.
    Performs lightweight tests and returns detailed status information.
    """
    start_time = time.time()
    
    health_status = {
        "status": "checking",
        "timestamp": datetime.now().isoformat(),
        "version": "1.0.0",  # Add your app version
        "checks": {
            "azure_openai": {"status": "pending"},
            "file_system": {"status": "pending"},
            "dependencies": {"status": "pending"},
            "endpoints": {"status": "pending"},
            "pandas_agent": {"status": "pending"},
        },
        "endpoints_tested": {},
        "warnings": [],
        "errors": [],
    }
    
    # Helper function to safely execute checks
    async def safe_check(check_name: str, check_func):
        try:
            result = await check_func() if asyncio.iscoroutinefunction(check_func) else check_func()
            health_status["checks"][check_name] = {"status": "healthy", **result}
            return True
        except Exception as e:
            error_msg = f"{check_name}: {str(e)}"
            health_status["checks"][check_name] = {"status": "unhealthy", "error": str(e)}
            health_status["errors"].append(error_msg)
            logging.error(f"Health check failed for {check_name}: {e}")
            return False
    
    # 1. Test Azure OpenAI Connection
    async def check_azure_openai():
        client = create_client()
        
        # Try a minimal API call - list models or assistants with limit=1
        try:
            # Try to list assistants (minimal call)
            assistants = client.beta.assistants.list(limit=1)
            
            # Try a simple completion to verify model access
            test_completion = client.chat.completions.create(
                model="gpt-4.1-mini",
                messages=[{"role": "user", "content": "Hi"}],
                max_tokens=5,
                temperature=0
            )
            
            return {
                "endpoint": AZURE_ENDPOINT,
                "api_version": AZURE_API_VERSION,
                "model_accessible": True,
                "assistants_api": True
            }
        except Exception as e:
            # Try basic completion only
            try:
                test_completion = client.chat.completions.create(
                    model="gpt-4.1-mini",
                    messages=[{"role": "user", "content": "Hi"}],
                    max_tokens=5
                )
                return {
                    "endpoint": AZURE_ENDPOINT,
                    "api_version": AZURE_API_VERSION,
                    "model_accessible": True,
                    "assistants_api": False,
                    "warning": "Assistants API not accessible"
                }
            except:
                raise
    
    # 2. Test File System
    def check_file_system():
        results = {
            "temp_dir_writable": False,
            "downloads_dir_exists": False,
            "downloads_dir_writable": False,
            "disk_space_available": 0
        }
        
        # Test /tmp directory
        try:
            test_file = os.path.join("/tmp", f"health_check_{uuid.uuid4().hex}.txt")
            with open(test_file, "w") as f:
                f.write("health check test")
            os.remove(test_file)
            results["temp_dir_writable"] = True
        except Exception as e:
            results["temp_write_error"] = str(e)
        
        # Test downloads directory
        try:
            results["downloads_dir_exists"] = os.path.exists(DOWNLOADS_DIR)
            if results["downloads_dir_exists"]:
                test_file = os.path.join(DOWNLOADS_DIR, f"health_check_{uuid.uuid4().hex}.txt")
                with open(test_file, "w") as f:
                    f.write("test")
                os.remove(test_file)
                results["downloads_dir_writable"] = True
                
                # Check disk space (in MB)
                stat = os.statvfs(DOWNLOADS_DIR)
                results["disk_space_available"] = (stat.f_bavail * stat.f_frsize) / (1024 * 1024)
        except Exception as e:
            results["downloads_error"] = str(e)
        
        return results
    
    # 3. Test Dependencies
    def check_dependencies():
        dependencies = {}
        critical_deps = {
            "pandas": None,
            "numpy": None,
            "openai": None,
            "langchain": None,
            "openpyxl": None,
            "PyPDF2": None,
            "python-docx": None,
            "Pillow": None,
            "beautifulsoup4": None
        }
        
        for dep_name in critical_deps:
            try:
                if dep_name == "python-docx":
                    import docx
                    dependencies[dep_name] = {"installed": True, "version": getattr(docx, "__version__", "unknown")}
                elif dep_name == "beautifulsoup4":
                    import bs4
                    dependencies[dep_name] = {"installed": True, "version": bs4.__version__}
                elif dep_name == "Pillow":
                    import PIL
                    dependencies[dep_name] = {"installed": True, "version": PIL.__version__}
                else:
                    module = __import__(dep_name)
                    version = getattr(module, "__version__", "unknown")
                    dependencies[dep_name] = {"installed": True, "version": version}
            except ImportError:
                dependencies[dep_name] = {"installed": False}
                if dep_name in ["pandas", "numpy", "openai", "langchain"]:
                    health_status["warnings"].append(f"Critical dependency {dep_name} not installed")
        
        return {"dependencies": dependencies}
    
    # 4. Test Pandas Agent Manager
    def check_pandas_agent():
        try:
            manager = PandasAgentManager.get_instance()
            
            # Check if dependencies are available
            deps_ok = manager._check_dependencies()
            
            # Test basic functionality without creating files
            test_thread_id = f"health_check_{uuid.uuid4().hex}"
            manager.initialize_thread(test_thread_id)
            
            # Verify initialization
            thread_initialized = (
                test_thread_id in manager.dataframes_cache and
                test_thread_id in manager.file_info_cache
            )
            
            # Clean up
            if test_thread_id in manager.dataframes_cache:
                del manager.dataframes_cache[test_thread_id]
            if test_thread_id in manager.file_info_cache:
                del manager.file_info_cache[test_thread_id]
            
            return {
                "instance_available": True,
                "dependencies_ok": deps_ok,
                "initialization_works": thread_initialized
            }
        except Exception as e:
            return {"error": str(e)}
    
    # 5. Test Endpoints (Lightweight)
    async def check_endpoints():
        endpoint_results = {}
        
        # Test /completion endpoint with minimal request
        try:
            # Create a mock request
            mock_request = type('Request', (), {
                'base_url': 'http://localhost:8080/',
                'headers': {'host': 'localhost:8080'}
            })()
            
            # Test with minimal parameters
            result = await chat_completion(
                request=mock_request,
                prompt="test",
                model="gpt-4.1-mini",
                temperature=0,
                max_tokens=1
            )
            
            endpoint_results["/completion"] = {
                "status": "healthy",
                "response_type": type(result).__name__
            }
        except Exception as e:
            endpoint_results["/completion"] = {
                "status": "unhealthy",
                "error": str(e)[:100]  # Limit error message length
            }
        
        # Test file processing functions
        try:
            # Test text extraction using extract_text_internal
            test_text = b"Hello, world!"
            # Since extract_text_internal is async, we need to run it in an async context
            extracted = await extract_text_internal(
                file_content=test_text,
                filename="test.txt",
                strategy="auto",
                languages=None,
                encoding=None,
                logger=logging.getLogger(__name__)
            )
            
            endpoint_results["file_extraction"] = {
                "status": "healthy" if extracted == "Hello, world!" else "unhealthy",
                "test_passed": extracted == "Hello, world!",
                "actual_output": extracted[:50] if extracted != "Hello, world!" else None  # Help debug if test fails
            }
        except Exception as e:
            endpoint_results["file_extraction"] = {
                "status": "unhealthy",
                "error": str(e)[:100]
            }
        
        # Test CSV extraction function
        try:
            test_csv = "```csv\nname,age\nJohn,30\n```"
            extracted_csv = extract_csv_from_content(test_csv)
            expected = "name,age\nJohn,30"
            
            endpoint_results["csv_extraction"] = {
                "status": "healthy" if extracted_csv.strip() == expected else "unhealthy",
                "test_passed": extracted_csv.strip() == expected
            }
        except Exception as e:
            endpoint_results["csv_extraction"] = {
                "status": "unhealthy",
                "error": str(e)[:100]
            }
        
        # Test download URL construction
        try:
            mock_request = type('Request', (), {
                'base_url': 'http://localhost:8080/',
                'headers': {'host': 'localhost:8080'}
            })()
            
            url = construct_download_url(mock_request, "test.csv")
            
            endpoint_results["download_url_construction"] = {
                "status": "healthy",
                "sample_url": url
            }
        except Exception as e:
            endpoint_results["download_url_construction"] = {
                "status": "unhealthy", 
                "error": str(e)[:100]
            }
        
        return {"endpoints": endpoint_results}
    
    # Execute all checks
    await safe_check("azure_openai", check_azure_openai)
    await safe_check("file_system", check_file_system)
    await safe_check("dependencies", check_dependencies)
    await safe_check("pandas_agent", check_pandas_agent)
    await safe_check("endpoints", check_endpoints)
    
    # Calculate overall health
    total_checks = len(health_status["checks"])
    healthy_checks = sum(1 for check in health_status["checks"].values() 
                        if check.get("status") == "healthy")
    
    health_percentage = (healthy_checks / total_checks) * 100 if total_checks > 0 else 0
    
    # Determine overall status
    if health_percentage == 100:
        overall_status = "healthy"
        status_code = 200
    elif health_percentage >= 80:
        overall_status = "degraded"
        status_code = 200
    elif health_percentage >= 50:
        overall_status = "partial"
        status_code = 503
    else:
        overall_status = "unhealthy" 
        status_code = 503
    
    # Update final status
    health_status["status"] = overall_status
    health_status["health_percentage"] = round(health_percentage, 2)
    health_status["execution_time_ms"] = round((time.time() - start_time) * 1000, 2)
    
    # Add summary
    health_status["summary"] = {
        "total_checks": total_checks,
        "healthy_checks": healthy_checks,
        "warnings_count": len(health_status["warnings"]),
        "errors_count": len(health_status["errors"])
    }
    
    return JSONResponse(
        content=health_status,
        status_code=status_code
    )

@app.post("/test-comprehensive",
          summary="Run Comprehensive System Tests",
          description="""Run comprehensive load and scaling tests with real-time streaming updates.
          
**⚡ STREAMING ENDPOINT** - Returns results via Server-Sent Events (SSE)

Test modes available:
- `all`: Run all test suites
- `long_thread`: Test thread capacity limits
- `concurrent`: Test concurrent user handling
- `same_thread`: Test thread locking mechanisms
- `scaling`: Test system scaling capabilities
- `tools`: Test AI tool functionality

Results are streamed in real-time as tests execute.
""",
          tags=["System"],
          response_class=StreamingResponse)
async def comprehensive_system_test(
    test_mode: str = Form(default="all", description="Test mode", enum=["all", "long_thread", "concurrent", "same_thread", "scaling", "tools"]),
    test_duration: int = Form(default=60, ge=10, le=300, description="Test duration in seconds"),
    concurrent_users: int = Form(default=5, ge=1, le=20, description="Number of concurrent users"),
    messages_per_thread: int = Form(default=60, ge=10, le=200, description="Messages for long thread test"),
    verbose: bool = Form(default=True, description="Include detailed logs")
):
    """Run comprehensive system tests for scaling and performance.
    
    This endpoint streams test results in real-time using Server-Sent Events (SSE).
    Each update is sent as a JSON object on a new line.
    """
    # Security check
    if os.getenv("ENVIRONMENT", "development") == "production":
        # In production, require a secret token
        test_token = os.getenv("TEST_ENDPOINT_TOKEN", None)
        if not test_token:
            raise HTTPException(status_code=403, detail="Test endpoint disabled in production")
    
    # Create a queue for streaming updates
    update_queue = asyncio.Queue()
    
    async def stream_updates() -> AsyncGenerator[str, None]:
        """Async generator that yields updates from the queue"""
        while True:
            try:
                # Wait for updates with a timeout
                update = await asyncio.wait_for(update_queue.get(), timeout=0.5)
                if update is None:  # Sentinel value to stop streaming
                    break
                yield json.dumps(update) + "\n"
            except asyncio.TimeoutError:
                # Continue waiting for updates
                continue
    
    async def run_tests():
        """Run the actual tests and put updates in the queue"""
        start_time = time.time()
        
        # Initial test setup
        test_results = {
            "timestamp": datetime.now().isoformat(),
            "test_mode": test_mode,
            "parameters": {
                "test_duration": test_duration,
                "concurrent_users": concurrent_users,
                "messages_per_thread": messages_per_thread
            },
            "tests": {},
            "summary": {
                "total_tests": 0,
                "passed": 0,
                "failed": 0,
                "warnings": 0
            },
            "logs": [] if verbose else None
        }
        
        client = create_client()
        
        # Helper function to queue updates
        async def queue_update(update_type: str, data: Dict[str, Any]):
            """Queue a streaming update"""
            update = {
                "type": update_type,
                "timestamp": datetime.now().isoformat(),
                "elapsed_time": round(time.time() - start_time, 2),
                "data": data
            }
            await update_queue.put(update)
        
        # Helper function to log with streaming
        async def log_stream(message: str, level: str = "info"):
            logging.info(f"[TEST] {message}")
            if verbose:
                log_entry = {
                    "time": datetime.now().isoformat(),
                    "level": level,
                    "message": message
                }
                test_results["logs"].append(log_entry)
                # Stream the log entry
                await queue_update("log", log_entry)
        
        # Queue initial status
        await queue_update("test_started", {
            "test_mode": test_mode,
            "total_tests": 6 if test_mode == "all" else 1,
            "parameters": test_results["parameters"]
        })
        
        # Test 1: Long Thread Response Test
        async def test_long_thread():
            """Test thread trimming and long conversation handling"""
            test_name = "long_thread"
            
            # Queue test start
            await queue_update("test_progress", {
                "test_name": test_name,
                "status": "starting",
                "message": "Starting long thread test"
            })
            
            result = {
                "status": "running",
                "messages_created": 0,
                "trim_triggered": False,
                "errors": [],
                "response_times": [],
                "trim_behavior": {}
            }
            
            try:
                # Create test assistant and thread
                assistant = client.beta.assistants.create(
                    name=f"test_long_thread_{int(time.time())}",
                    model="gpt-4.1-mini",
                    instructions="You are a test assistant. Keep responses very brief (max 20 words).",
                    tools=[]
                )
                thread = client.beta.threads.create()
                
                await log_stream(f"Created test assistant {assistant.id} and thread {thread.id}")
                
                # Fill thread with messages
                for i in range(messages_per_thread):
                    try:
                        # Add user message
                        client.beta.threads.messages.create(
                            thread_id=thread.id,
                            role="user",
                            content=f"Test message {i+1}. Reply with just 'Acknowledged {i+1}'."
                        )
                        
                        msg_start = time.time()
                        
                        # Create and wait for run
                        run = client.beta.threads.runs.create(
                            thread_id=thread.id,
                            assistant_id=assistant.id
                        )
                        
                        # Wait for completion
                        while run.status in ["queued", "in_progress", "requires_action"]:
                            await asyncio.sleep(0.5)
                            run = client.beta.threads.runs.retrieve(
                                thread_id=thread.id,
                                run_id=run.id
                            )
                        
                        response_time = time.time() - msg_start
                        result["response_times"].append(response_time)
                        result["messages_created"] += 1
                        
                        # Stream progress every 10 messages
                        if (i + 1) % 10 == 0:
                            await queue_update("test_progress", {
                                "test_name": test_name,
                                "status": "running",
                                "progress": f"{i+1}/{messages_per_thread}",
                                "message": f"Created {i+1} messages, avg response time: {sum(result['response_times'])/len(result['response_times']):.2f}s"
                            })
                        
                        # Check for thread trimming (usually around 50 messages)
                        if i > 45:
                            messages = client.beta.threads.messages.list(thread_id=thread.id)
                            current_count = len(messages.data)
                            if current_count < result["messages_created"] - 5:
                                result["trim_triggered"] = True
                                result["trim_behavior"]["detected_at_message"] = i + 1
                                result["trim_behavior"]["message_count_after_trim"] = current_count
                                await log_stream(f"Thread trimming detected at message {i+1}. Current count: {current_count}")
                                
                    except Exception as e:
                        result["errors"].append(f"Message {i+1}: {str(e)}")
                        await log_stream(f"Error at message {i+1}: {e}", "error")
                        
                        if "limit" in str(e).lower() or "trim" in str(e).lower():
                            result["trim_triggered"] = True
                            result["trim_behavior"]["error_triggered"] = True
                            break
                
                # Test non-streaming endpoint with full thread
                await log_stream("Testing non-streaming endpoint with full thread")
                    
                try:
                    response = await process_conversation(
                        session=thread.id,
                        prompt="Count your recent messages.",
                        assistant=assistant.id,
                        stream_output=False
                    )
                    
                    if hasattr(response, 'body'):
                        response_data = json.loads(response.body.decode())
                        result["non_streaming_endpoint_works"] = True
                        result["final_response_preview"] = response_data.get("response", "")[:200]
                        await log_stream("Non-streaming endpoint handled full thread successfully")
                        
                except Exception as e:
                    result["non_streaming_endpoint_error"] = str(e)
                    await log_stream(f"Non-streaming endpoint error: {e}", "error")
                
                # Cleanup
                try:
                    client.beta.assistants.delete(assistant_id=assistant.id)
                    await log_stream(f"Cleaned up assistant {assistant.id}")
                except:
                    pass
                
                # Analyze results
                if result["messages_created"] >= 48:
                    if result["trim_triggered"]:
                        result["status"] = "passed"
                        result["summary"] = "Thread trimming working correctly"
                    else:
                        result["status"] = "warning"
                        result["summary"] = "Thread filled but trimming not detected"
                else:
                    result["status"] = "failed"
                    result["summary"] = f"Could only create {result['messages_created']} messages"
                
                avg_response_time = sum(result["response_times"]) / len(result["response_times"]) if result["response_times"] else 0
                result["avg_response_time"] = avg_response_time
                
            except Exception as e:
                result["status"] = "failed"
                result["summary"] = f"Test failed: {str(e)}"
                await log_stream(f"Long thread test failed: {e}", "error")
            
            return result
        
        # Test 2: Concurrent Users Test
        async def test_concurrent_users():
            """Test multiple users talking to different assistants concurrently"""
            test_name = "concurrent_users"
            
            await queue_update("test_progress", {
                "test_name": test_name,
                "status": "starting",
                "message": f"Starting concurrent users test with {concurrent_users} users"
            })
            
            result = {
                "status": "running",
                "users_tested": concurrent_users,
                "successful_conversations": 0,
                "failed_conversations": 0,
                "response_times": [],
                "errors": [],
                "concurrency_issues": []
            }
            
            # Create test assistants
            assistants = []
            for i in range(concurrent_users):
                try:
                    assistant = client.beta.assistants.create(
                        name=f"test_concurrent_{i}_{int(time.time())}",
                        model="gpt-4.1-mini",
                        instructions=f"You are test assistant {i}. Always include your number ({i}) in responses. Keep responses under 15 words.",
                        tools=[]
                    )
                    assistants.append(assistant)
                    
                    if (i + 1) % 2 == 0:  # Update every 2 assistants
                        await queue_update("test_progress", {
                            "test_name": test_name,
                            "status": "preparing",
                            "progress": f"{i+1}/{concurrent_users}",
                            "message": f"Created {i+1} test assistants"
                        })
                            
                except Exception as e:
                    result["errors"].append(f"Failed to create assistant {i}: {str(e)}")
                    await log_stream(f"Failed to create assistant {i}: {e}", "error")
            
            await log_stream(f"Created {len(assistants)} test assistants")
            
            # Define concurrent conversation task
            async def have_conversation(user_id: int, assistant):
                """Simulate a user having a conversation"""
                conv_result = {
                    "user_id": user_id,
                    "messages_sent": 0,
                    "responses_received": 0,
                    "errors": [],
                    "response_times": []
                }
                
                try:
                    thread = client.beta.threads.create()
                    
                    # Send multiple messages over time
                    messages_to_send = min(5, max(1, test_duration // 10))
                    
                    for msg_num in range(messages_to_send):
                        try:
                            msg_start = time.time()
                            
                            response = await process_conversation(
                                session=thread.id,
                                prompt=f"User {user_id} message {msg_num + 1}. What's your assistant number?",
                                assistant=assistant.id,
                                stream_output=False
                            )
                            
                            msg_duration = time.time() - msg_start
                            conv_result["response_times"].append(msg_duration)
                            conv_result["messages_sent"] += 1
                            
                            if hasattr(response, 'body'):
                                response_data = json.loads(response.body.decode())
                                response_text = response_data.get("response", "")
                                
                                # Verify assistant number in response
                                if str(user_id) in response_text:
                                    conv_result["responses_received"] += 1
                                else:
                                    conv_result["errors"].append(f"Assistant identity mismatch in message {msg_num + 1}")
                            
                            # Brief delay between messages
                            await asyncio.sleep(2)
                            
                        except Exception as e:
                            conv_result["errors"].append(f"Message {msg_num + 1}: {str(e)}")
                    
                    return conv_result
                    
                except Exception as e:
                    conv_result["errors"].append(f"Conversation setup failed: {str(e)}")
                    return conv_result
            
            # Run concurrent conversations
            await log_stream("Starting concurrent conversations")
                
            tasks = []
            for i in range(len(assistants)):
                task = have_conversation(i, assistants[i])
                tasks.append(task)
            
            # Execute concurrently and gather results
            conv_results = await asyncio.gather(*tasks, return_exceptions=True)
            
            # Process results
            for i, conv_result in enumerate(conv_results):
                if isinstance(conv_result, Exception):
                    result["failed_conversations"] += 1
                    result["errors"].append(f"User {i} conversation failed: {str(conv_result)}")
                else:
                    if conv_result["responses_received"] > 0:
                        result["successful_conversations"] += 1
                        result["response_times"].extend(conv_result["response_times"])
                    else:
                        result["failed_conversations"] += 1
                    
                    if conv_result["errors"]:
                        result["concurrency_issues"].extend(conv_result["errors"])
                
                # Stream progress
                if (i + 1) % 2 == 0:
                    await queue_update("test_progress", {
                        "test_name": test_name,
                        "status": "running",
                        "progress": f"{i+1}/{len(assistants)}",
                        "message": f"Processed {i+1} conversations, {result['successful_conversations']} successful"
                    })
            
            # Cleanup assistants
            for assistant in assistants:
                try:
                    client.beta.assistants.delete(assistant_id=assistant.id)
                except:
                    pass
            
            # Summary
            result["avg_response_time"] = sum(result["response_times"]) / len(result["response_times"]) if result["response_times"] else 0
            result["success_rate"] = result["successful_conversations"] / concurrent_users if concurrent_users > 0 else 0
            
            if result["success_rate"] >= 0.9:
                result["status"] = "passed"
                result["summary"] = f"Concurrent users test passed with {result['success_rate']*100:.1f}% success rate"
            elif result["success_rate"] >= 0.7:
                result["status"] = "warning"
                result["summary"] = f"Concurrent users test has issues: {result['success_rate']*100:.1f}% success rate"
            else:
                result["status"] = "failed"
                result["summary"] = f"Concurrent users test failed: only {result['success_rate']*100:.1f}% success rate"
            
            await log_stream(f"Concurrent users test completed: {result['summary']}")
                
            return result
        
        # Test 3: Same Thread Concurrent Access
        async def test_same_thread_concurrent():
            """Test multiple concurrent requests to the same thread"""
            test_name = "same_thread_concurrent"
            
            await queue_update("test_progress", {
                "test_name": test_name,
                "status": "starting",
                "message": "Starting same thread concurrent access test"
            })
            
            result = {
                "status": "running",
                "concurrent_requests": 5,
                "successful_requests": 0,
                "failed_requests": 0,
                "lock_timeouts": 0,
                "errors": [],
                "response_order": []
            }
            
            try:
                # Create shared assistant and thread
                assistant = client.beta.assistants.create(
                    name=f"test_same_thread_{int(time.time())}",
                    model="gpt-4.1-mini",
                    instructions="You are a test assistant. Number each response sequentially.",
                    tools=[]
                )
                thread = client.beta.threads.create()
                
                await log_stream(f"Created shared assistant {assistant.id} and thread {thread.id}")
                
                async def send_concurrent_message(msg_id: int):
                    """Send a message and track order"""
                    msg_result = {
                        "msg_id": msg_id,
                        "sent_time": time.time(),
                        "completed_time": None,
                        "success": False,
                        "error": None
                    }
                    
                    try:
                        response = await process_conversation(
                            session=thread.id,
                            prompt=f"Concurrent message {msg_id}. What number is this?",
                            assistant=assistant.id,
                            stream_output=False
                        )
                        
                        msg_result["completed_time"] = time.time()
                        msg_result["duration"] = msg_result["completed_time"] - msg_result["sent_time"]
                        
                        if hasattr(response, 'body'):
                            msg_result["success"] = True
                            result["successful_requests"] += 1
                        else:
                            msg_result["error"] = "No response body"
                            result["failed_requests"] += 1
                            
                    except Exception as e:
                        msg_result["error"] = str(e)
                        result["failed_requests"] += 1
                        
                        if "lock" in str(e).lower() or "timeout" in str(e).lower():
                            result["lock_timeouts"] += 1
                    
                    return msg_result
                
                # Send concurrent messages
                await log_stream("Sending concurrent messages to same thread")
                    
                tasks = []
                for i in range(result["concurrent_requests"]):
                    task = send_concurrent_message(i)
                    tasks.append(task)
                    await asyncio.sleep(0.1)  # Slight stagger to increase concurrency chance
                
                # Gather results
                msg_results = await asyncio.gather(*tasks, return_exceptions=True)
                
                # Analyze results
                for i, msg_result in enumerate(msg_results):
                    if isinstance(msg_result, Exception):
                        result["errors"].append(f"Task {i} exception: {str(msg_result)}")
                    else:
                        result["response_order"].append(msg_result)
                        
                        # Stream progress
                        await queue_update("test_progress", {
                            "test_name": test_name,
                            "status": "running",
                            "progress": f"{i+1}/{result['concurrent_requests']}",
                            "message": f"Processed {i+1} concurrent requests"
                        })
                
                # Cleanup
                try:
                    client.beta.assistants.delete(assistant_id=assistant.id)
                except:
                    pass
                
                # Determine status
                if result["successful_requests"] == result["concurrent_requests"]:
                    result["status"] = "passed"
                    result["summary"] = "All concurrent requests to same thread succeeded"
                elif result["successful_requests"] > 0:
                    result["status"] = "warning"
                    result["summary"] = f"Partial success: {result['successful_requests']}/{result['concurrent_requests']} requests succeeded"
                else:
                    result["status"] = "failed"
                    result["summary"] = "All concurrent requests failed"
                
                await log_stream(f"Same thread concurrent test: {result['summary']}")
                
            except Exception as e:
                result["status"] = "failed"
                result["summary"] = f"Test setup failed: {str(e)}"
                await log_stream(f"Same thread concurrent test failed: {e}", "error")
            
            return result
        
        # Test 4: Run Consistency and Cleanup
        async def test_run_consistency():
            """Test run handling, cleanup, and error recovery"""
            test_name = "run_consistency"
            
            await queue_update("test_progress", {
                "test_name": test_name,
                "status": "starting",
                "message": "Starting run consistency and cleanup test"
            })
            
            result = {
                "status": "running",
                "tests_performed": [],
                "issues_found": [],
                "cleanup_success": True
            }
            
            try:
                # Create test assistant and thread
                assistant = client.beta.assistants.create(
                    name=f"test_runs_{int(time.time())}",
                    model="gpt-4.1-mini",
                    instructions="You are a test assistant.",
                    tools=[]
                )
                thread = client.beta.threads.create()
                
                # Test 1: Abandoned run handling
                await log_stream("Testing abandoned run handling")
                
                # Create a run but don't wait for it
                run1 = client.beta.threads.runs.create(
                    thread_id=thread.id,
                    assistant_id=assistant.id
                )
                
                # Immediately try to create another message
                await asyncio.sleep(2)  # Let first run start
                
                try:
                    # This should handle the active run
                    response = await process_conversation(
                        session=thread.id,
                        prompt="New message while run is active",
                        assistant=assistant.id,
                        stream_output=False
                    )
                    
                    result["tests_performed"].append("active_run_handling")
                    await log_stream("Successfully handled active run scenario")
                    
                    await queue_update("test_progress", {
                        "test_name": test_name,
                        "status": "running",
                        "progress": "1/3",
                        "message": "Active run handling test passed"
                    })
                    
                except Exception as e:
                    result["issues_found"].append(f"Active run handling failed: {str(e)}")
                    await log_stream(f"Active run handling error: {e}", "error")
                
                # Test 2: Thread state after errors
                await log_stream("Testing thread state consistency")
                
                # Check thread messages
                messages = client.beta.threads.messages.list(thread_id=thread.id)
                message_count = len(messages.data)
                
                # Verify no duplicate messages
                message_contents = [msg.content[0].text.value if msg.content else "" for msg in messages.data]
                unique_contents = set(message_contents)
                
                if len(unique_contents) < len(message_contents):
                    result["issues_found"].append("Duplicate messages detected")
                else:
                    result["tests_performed"].append("no_duplicate_messages")
                
                await queue_update("test_progress", {
                    "test_name": test_name,
                    "status": "running",
                    "progress": "2/3",
                    "message": "Thread state consistency checked"
                })
                
                # Test 3: Cleanup test
                await log_stream("Testing cleanup")
                
                try:
                    client.beta.assistants.delete(assistant_id=assistant.id)
                    result["tests_performed"].append("cleanup_successful")
                except Exception as e:
                    result["cleanup_success"] = False
                    result["issues_found"].append(f"Cleanup failed: {str(e)}")
                
                await queue_update("test_progress", {
                    "test_name": test_name,
                    "status": "running",
                    "progress": "3/3",
                    "message": "Cleanup test completed"
                })
                
                # Determine status
                if len(result["issues_found"]) == 0:
                    result["status"] = "passed"
                    result["summary"] = f"All consistency tests passed: {', '.join(result['tests_performed'])}"
                elif len(result["issues_found"]) < len(result["tests_performed"]):
                    result["status"] = "warning"
                    result["summary"] = f"Some issues found: {', '.join(result['issues_found'])}"
                else:
                    result["status"] = "failed"
                    result["summary"] = f"Multiple issues: {', '.join(result['issues_found'])}"
                
                await log_stream(f"Run consistency test: {result['summary']}")
                
            except Exception as e:
                result["status"] = "failed"
                result["summary"] = f"Test failed: {str(e)}"
                await log_stream(f"Run consistency test error: {e}", "error")
            
            return result
        
        # Test 5: Scaling and Performance Test
        async def test_scaling():
            """Test system performance under sustained load"""
            test_name = "scaling"
            
            await queue_update("test_progress", {
                "test_name": test_name,
                "status": "starting",
                "message": f"Starting scaling test for {test_duration} seconds"
            })
            
            result = {
                "status": "running",
                "duration": test_duration,
                "total_messages": 0,
                "successful_messages": 0,
                "errors": [],
                "response_times": [],
                "performance_degradation": False,
                "memory_issues": False
            }
            
            try:
                scaling_start_time = time.time()
                response_times = []
                
                # Create test assistant
                assistant = client.beta.assistants.create(
                    name=f"test_scaling_{int(time.time())}",
                    model="gpt-4.1-mini",
                    instructions="You are a test assistant. Keep responses under 10 words.",
                    tools=[]
                )
                
                # Create multiple threads
                threads = []
                for i in range(3):  # Use 3 threads for scaling test
                    thread = client.beta.threads.create()
                    threads.append(thread.id)
                
                await log_stream(f"Created {len(threads)} threads for scaling test")
                
                # Send messages for test_duration seconds
                message_count = 0
                errors_count = 0
                last_update_time = time.time()
                
                while time.time() - scaling_start_time < test_duration:
                    # Round-robin through threads
                    thread_id = threads[message_count % len(threads)]
                    
                    try:
                        msg_start = time.time()
                        
                        response = await process_conversation(
                            session=thread_id,
                            prompt=f"Quick test {message_count}",
                            assistant=assistant.id,
                            stream_output=False
                        )
                        
                        msg_duration = time.time() - msg_start
                        response_times.append(msg_duration)
                        result["successful_messages"] += 1
                        
                        # Check for performance degradation
                        if len(response_times) > 10:
                            recent_avg = sum(response_times[-10:]) / 10
                            overall_avg = sum(response_times) / len(response_times)
                            
                            if recent_avg > overall_avg * 1.5:
                                result["performance_degradation"] = True
                                await log_stream(f"Performance degradation detected: recent avg {recent_avg:.2f}s vs overall {overall_avg:.2f}s", "warning")
                        
                    except Exception as e:
                        errors_count += 1
                        result["errors"].append(f"Message {message_count}: {str(e)}")
                        
                        if "memory" in str(e).lower() or "resource" in str(e).lower():
                            result["memory_issues"] = True
                    
                    message_count += 1
                    result["total_messages"] = message_count
                    
                    # Stream progress every 5 seconds
                    if time.time() - last_update_time > 5:
                        elapsed = time.time() - scaling_start_time
                        rate = result["successful_messages"] / elapsed if elapsed > 0 else 0
                        
                        await queue_update("test_progress", {
                            "test_name": test_name,
                            "status": "running",
                            "progress": f"{int(elapsed)}/{test_duration}s",
                            "message": f"Processed {result['successful_messages']} messages at {rate:.1f} msg/s"
                        })
                        
                        last_update_time = time.time()
                
                # Cleanup
                try:
                    client.beta.assistants.delete(assistant_id=assistant.id)
                except:
                    pass
                
                # Calculate final metrics
                result["response_times"] = response_times
                result["avg_response_time"] = sum(response_times) / len(response_times) if response_times else 0
                result["messages_per_second"] = result["successful_messages"] / test_duration
                
                # Determine status
                success_rate = result["successful_messages"] / result["total_messages"] if result["total_messages"] > 0 else 0
                
                if success_rate >= 0.95 and not result["performance_degradation"]:
                    result["status"] = "passed"
                    result["summary"] = f"Scaling test passed: {result['messages_per_second']:.1f} msg/s sustained"
                elif success_rate >= 0.8:
                    result["status"] = "warning"
                    result["summary"] = f"Scaling test has issues: {success_rate*100:.1f}% success rate"
                else:
                    result["status"] = "failed"
                    result["summary"] = f"Scaling test failed: only {success_rate*100:.1f}% success rate"
                
                await log_stream(f"Scaling test completed: {result['summary']}")
                
            except Exception as e:
                result["status"] = "failed"
                result["summary"] = f"Scaling test failed: {str(e)}"
                await log_stream(f"Scaling test error: {e}", "error")
            
            return result
        
        # Test 6: Tool Calling Tests
        async def test_tool_calling():
            """Test all tool calling functionality including pandas_agent, generate_content, and extract_data"""
            test_name = "tool_calling"
            
            await queue_update("test_progress", {
                "test_name": test_name,
                "status": "starting",
                "message": "Starting tool calling functionality test"
            })
            
            result = {
                "status": "running",
                "pandas_agent": {
                    "tested": False,
                    "success": False,
                    "response_time": 0,
                    "errors": []
                },
                "generate_content": {
                    "tested": False,
                    "success": False,
                    "response_time": 0,
                    "download_url": None,
                    "errors": []
                },
                "extract_data": {
                    "tested": False,
                    "success": False,
                    "response_time": 0,
                    "rows_extracted": 0,
                    "errors": []
                },
                "tool_concurrency": {
                    "tested": False,
                    "success": False,
                    "concurrent_tools": 0,
                    "errors": []
                },
                "run_states": {}
            }
            
            try:
                # Create test assistant with all tools
                tools = [
                    {"type": "file_search"},
                    {
                        "type": "function",
                        "function": {
                            "name": "pandas_agent",
                            "description": "Analyzes CSV and Excel files",
                            "parameters": {
                                "type": "object",
                                "properties": {
                                    "query": {"type": "string", "description": "Analysis query"},
                                    "filename": {"type": "string", "description": "Optional filename"}
                                },
                                "required": ["query"]
                            }
                        }
                    },
                    {
                        "type": "function",
                        "function": {
                            "name": "generate_content",
                            "description": "Generates CSV or Excel files",
                            "parameters": {
                                "type": "object",
                                "properties": {
                                    "content": {"type": "string", "description": "Content to generate"},
                                    "format": {"type": "string", "enum": ["csv", "excel"], "description": "Output format"},
                                    "filename": {"type": "string", "description": "Output filename"}
                                },
                                "required": ["content", "format"]
                            }
                        }
                    },
                    {
                        "type": "function",
                        "function": {
                            "name": "extract_data",
                            "description": "Extracts structured data",
                            "parameters": {
                                "type": "object",
                                "properties": {
                                    "data": {"type": "string", "description": "Data to extract from"},
                                    "columns": {"type": "string", "description": "Columns to extract"}
                                },
                                "required": ["data"]
                            }
                        }
                    }
                ]
                
                assistant = client.beta.assistants.create(
                    name=f"test_tools_{int(time.time())}",
                    model="gpt-4.1-mini",
                    instructions="You are a test assistant with tool calling capabilities.",
                    tools=tools
                )
                thread = client.beta.threads.create()
                
                await log_stream(f"Created test assistant {assistant.id} with tools")
                
                # Test 1: Pandas Agent
                await log_stream("Testing pandas_agent tool")
                
                try:
                    pandas_start = time.time()
                    
                    response = await process_conversation(
                        session=thread.id,
                        prompt="/analyze Create a sample dataset with 5 rows and analyze it",
                        assistant=assistant.id,
                        stream_output=False
                    )
                    
                    pandas_duration = time.time() - pandas_start
                    result["pandas_agent"]["tested"] = True
                    result["pandas_agent"]["response_time"] = pandas_duration
                    
                    if hasattr(response, 'body'):
                        result["pandas_agent"]["success"] = True
                        await log_stream(f"Pandas agent test successful, response time: {pandas_duration:.2f}s")
                    else:
                        result["pandas_agent"]["errors"].append("No response body")
                    
                    await queue_update("test_progress", {
                        "test_name": test_name,
                        "status": "running",
                        "progress": "1/4",
                        "message": "Pandas agent test completed"
                    })
                    
                except Exception as e:
                    result["pandas_agent"]["errors"].append(str(e))
                    await log_stream(f"Pandas agent test failed: {e}", "error")
                
                # Test 2: Generate Content
                await log_stream("Testing generate_content tool")
                
                try:
                    generate_start = time.time()
                    
                    response = await process_conversation(
                        session=thread.id,
                        prompt="/generate Create a CSV file with 10 product reviews including columns: product_name, rating, review_text",
                        assistant=assistant.id,
                        stream_output=False
                    )
                    
                    generate_duration = time.time() - generate_start
                    result["generate_content"]["tested"] = True
                    result["generate_content"]["response_time"] = generate_duration
                    
                    if hasattr(response, 'body'):
                        response_data = json.loads(response.body.decode())
                        response_text = response_data.get("response", "")
                        
                        # Look for download URL in response
                        if "download" in response_text.lower() and ("http" in response_text or "/download" in response_text):
                            result["generate_content"]["success"] = True
                            result["generate_content"]["download_url"] = "URL found in response"
                            await log_stream(f"Generate content test successful, response time: {generate_duration:.2f}s")
                        else:
                            result["generate_content"]["errors"].append("No download URL in response")
                    else:
                        result["generate_content"]["errors"].append("No response body")
                    
                    await queue_update("test_progress", {
                        "test_name": test_name,
                        "status": "running",
                        "progress": "2/4",
                        "message": "Generate content test completed"
                    })
                    
                except Exception as e:
                    result["generate_content"]["errors"].append(str(e))
                    await log_stream(f"Generate content test failed: {e}", "error")
                
                # Test 3: Extract Data
                await log_stream("Testing extract_data tool")
                
                try:
                    extract_start = time.time()
                    
                    test_data = """
                    Customer Reviews:
                    1. John Smith - "Great product!" - 5 stars
                    2. Jane Doe - "Good value" - 4 stars  
                    3. Bob Johnson - "Average quality" - 3 stars
                    """
                    
                    response = await process_conversation(
                        session=thread.id,
                        prompt=f"/extract Extract customer name, review, and rating from this data: {test_data}",
                        assistant=assistant.id,
                        stream_output=False
                    )
                    
                    extract_duration = time.time() - extract_start
                    result["extract_data"]["tested"] = True
                    result["extract_data"]["response_time"] = extract_duration
                    
                    if hasattr(response, 'body'):
                        response_data = json.loads(response.body.decode())
                        response_text = response_data.get("response", "")
                        
                        # Check if extraction happened
                        if "extracted" in response_text.lower() or "csv" in response_text.lower():
                            result["extract_data"]["success"] = True
                            
                            # Try to find number of rows extracted
                            import re
                            rows_match = re.search(r'(\d+)\s*rows?', response_text, re.IGNORECASE)
                            if rows_match:
                                result["extract_data"]["rows_extracted"] = int(rows_match.group(1))
                            
                            await log_stream(f"Extract data test successful, response time: {extract_duration:.2f}s")
                        else:
                            result["extract_data"]["errors"].append("No extraction confirmation in response")
                            await log_stream("Extract data response doesn't confirm extraction", "warning")
                    
                    await queue_update("test_progress", {
                        "test_name": test_name,
                        "status": "running",
                        "progress": "3/4",
                        "message": "Extract data test completed"
                    })
                    
                except Exception as e:
                    result["extract_data"]["errors"].append(str(e))
                    await log_stream(f"Extract data test failed: {e}", "error")
                
                # Test 4: Concurrent Tool Usage
                await log_stream("Testing concurrent tool usage")
                
                try:
                    # Create multiple threads for concurrent tool testing
                    concurrent_threads = []
                    for i in range(3):
                        ct = client.beta.threads.create()
                        concurrent_threads.append(ct.id)
                    
                    async def test_tool_concurrently(thread_id: str, tool_type: str, prompt: str):
                        """Test a tool concurrently"""
                        try:
                            tool_start = time.time()
                            
                            response = await process_conversation(
                                session=thread_id,
                                prompt=prompt,
                                assistant=assistant.id,
                                stream_output=False
                            )
                            
                            tool_duration = time.time() - tool_start
                            success = hasattr(response, 'body')
                            
                            return {
                                "tool": tool_type,
                                "success": success,
                                "duration": tool_duration
                            }
                        except Exception as e:
                            return {
                                "tool": tool_type,
                                "success": False,
                                "error": str(e)
                            }
                    
                    # Run different tools concurrently
                    tasks = [
                        test_tool_concurrently(concurrent_threads[0], "generate", "/generate 5 test items"),
                        test_tool_concurrently(concurrent_threads[1], "extract", "/analyze create 5 sample records"),
                        test_tool_concurrently(concurrent_threads[2], "generate", "/generate brief report")
                    ]
                    
                    concurrent_results = await asyncio.gather(*tasks, return_exceptions=True)
                    
                    result["tool_concurrency"]["tested"] = True
                    result["tool_concurrency"]["concurrent_tools"] = len(tasks)
                    
                    successful_concurrent = sum(1 for r in concurrent_results if isinstance(r, dict) and r.get("success"))
                    result["tool_concurrency"]["success"] = successful_concurrent == len(tasks)
                    result["tool_concurrency"]["details"] = concurrent_results
                    
                    await log_stream(f"Concurrent tool test: {successful_concurrent}/{len(tasks)} successful")
                        
                    await queue_update("test_progress", {
                        "test_name": test_name,
                        "status": "running",
                        "progress": "4/4",
                        "message": "Concurrent tool test completed"
                    })
                    
                except Exception as e:
                    result["tool_concurrency"]["errors"].append(str(e))
                    await log_stream(f"Concurrent tool test failed: {e}", "error")
                
                # Cleanup
                try:
                    client.beta.assistants.delete(assistant_id=assistant.id)
                except:
                    pass
                
                # Determine overall status
                tool_tests = ["pandas_agent", "generate_content", "extract_data", "tool_concurrency"]
                successful_tests = sum(1 for test in tool_tests if result[test].get("success", False))
                
                if successful_tests == len(tool_tests):
                    result["status"] = "passed"
                    result["summary"] = "All tool tests passed successfully"
                elif successful_tests >= len(tool_tests) - 1:
                    result["status"] = "warning"
                    failed_tools = [test for test in tool_tests if not result[test].get("success", False)]
                    result["summary"] = f"Most tool tests passed. Failed: {', '.join(failed_tools)}"
                else:
                    result["status"] = "failed"
                    result["summary"] = f"Tool tests failed: only {successful_tests}/{len(tool_tests)} passed"
                
                # Add performance summary
                avg_tool_time = sum([
                    result["pandas_agent"]["response_time"],
                    result["generate_content"]["response_time"],
                    result["extract_data"]["response_time"]
                ]) / 3
                
                result["performance_summary"] = {
                    "avg_tool_response_time": avg_tool_time,
                    "pandas_agent_time": result["pandas_agent"]["response_time"],
                    "generate_content_time": result["generate_content"]["response_time"],
                    "extract_data_time": result["extract_data"]["response_time"]
                }
                
                await log_stream(f"Tool calling test completed: {result['summary']}")
                
            except Exception as e:
                result["status"] = "failed"
                result["summary"] = f"Tool test setup failed: {str(e)}"
                await log_stream(f"Tool calling test failed: {e}", "error")
            
            return result
        
        # Execute tests based on mode
        if test_mode == "all":
            tests_to_run = ["long_thread", "concurrent", "same_thread", "run_consistency", "scaling", "tools"]
        else:
            tests_to_run = [test_mode]
        
        # Queue total number of tests
        await queue_update("tests_initialized", {
            "total_tests": len(tests_to_run),
            "test_names": tests_to_run
        })
        
        # Run selected tests
        for test_index, test in enumerate(tests_to_run):
            test_results["summary"]["total_tests"] += 1
            
            # Queue test starting
            await queue_update("test_starting", {
                "test_name": test,
                "test_number": test_index + 1,
                "total_tests": len(tests_to_run)
            })
            
            try:
                if test == "long_thread":
                    result = await test_long_thread()
                elif test == "concurrent":
                    result = await test_concurrent_users()
                elif test == "same_thread":
                    result = await test_same_thread_concurrent()
                elif test == "run_consistency":
                    result = await test_run_consistency()
                elif test == "scaling":
                    result = await test_scaling()
                elif test == "tools":
                    result = await test_tool_calling()
                else:
                    result = {"status": "skipped", "summary": "Unknown test"}
                
                test_results["tests"][test] = result
                
                if result.get("status") == "passed":
                    test_results["summary"]["passed"] += 1
                elif result.get("status") == "warning":
                    test_results["summary"]["warnings"] += 1
                else:
                    test_results["summary"]["failed"] += 1
                
                # Queue test completed
                await queue_update("test_completed", {
                    "test_name": test,
                    "test_number": test_index + 1,
                    "total_tests": len(tests_to_run),
                    "status": result.get("status"),
                    "summary": result.get("summary"),
                    "current_totals": {
                        "passed": test_results["summary"]["passed"],
                        "failed": test_results["summary"]["failed"],
                        "warnings": test_results["summary"]["warnings"]
                    }
                })
                    
            except Exception as e:
                test_results["tests"][test] = {
                    "status": "error",
                    "summary": f"Test crashed: {str(e)}",
                    "error": traceback.format_exc()
                }
                test_results["summary"]["failed"] += 1
                
                await log_stream(f"Test {test} crashed: {e}", "error")
                
                # Queue test error
                await queue_update("test_error", {
                    "test_name": test,
                    "error": str(e),
                    "traceback": traceback.format_exc()
                })
        
        # Overall summary
        total_time = time.time() - start_time
        test_results["execution_time"] = f"{total_time:.2f} seconds"
        
        if test_results["summary"]["failed"] == 0:
            if test_results["summary"]["warnings"] == 0:
                test_results["overall_status"] = "healthy"
                test_results["overall_summary"] = "All tests passed successfully"
                status_code = 200
            else:
                test_results["overall_status"] = "degraded"
                test_results["overall_summary"] = f"Tests passed with {test_results['summary']['warnings']} warnings"
                status_code = 200
        else:
            test_results["overall_status"] = "unhealthy"
            test_results["overall_summary"] = f"{test_results['summary']['failed']} tests failed"
            status_code = 500
        
        # Clean up logs if not verbose
        if not verbose:
            test_results.pop("logs", None)
        
        # Queue final results
        await queue_update("test_suite_completed", {
            "status": test_results["overall_status"],
            "summary": test_results["overall_summary"],
            "execution_time": test_results["execution_time"],
            "final_results": test_results["summary"],
            "status_code": status_code
        })
        
        # Queue the complete test results as the final message
        await update_queue.put({
            "type": "final_results",
            "timestamp": datetime.now().isoformat(),
            "data": test_results,
            "status_code": status_code
        })
        
        # Signal end of stream
        await update_queue.put(None)
    
    # Start the test runner in the background
    asyncio.create_task(run_tests())
    
    # Return streaming response
    return StreamingResponse(
        stream_updates(),
        media_type="application/x-ndjson",
        headers={
            "X-Content-Type-Options": "nosniff",
            "Cache-Control": "no-cache",
            "Connection": "keep-alive"
        }
    )
# Add a lightweight health check endpoint for quick monitoring
@app.get("/health", 
         response_model=HealthResponse,
         summary="Basic Health Check",
         description="Quick health check for monitoring. Returns minimal status information.",
         tags=["System"])
async def basic_health():
    """
    Basic health check endpoint for load balancers and monitoring.
    Returns quickly with minimal checks.
    """
    try:
        # Just verify the app is running and can create a client
        client = create_client()
        
        return JSONResponse({
            "status": "healthy",
            "timestamp": datetime.now().isoformat()
        })
    except Exception as e:
        return JSONResponse(
            content={
                "status": "unhealthy",
                "error": str(e),
                "timestamp": datetime.now().isoformat()
            },
            status_code=503
        )


from fastapi.responses import HTMLResponse

@app.get("/", response_class=HTMLResponse)
async def serve_webpage():
    """Serve the AI Assistant Hub webpage"""
    try:
        with open("webpage.html", "r", encoding="utf-8") as f:
            html_content = f.read()
        return HTMLResponse(content=html_content)
    except FileNotFoundError:
        return HTMLResponse(
            content="<h1>Error: webpage.html not found</h1><p>Please ensure webpage.html is in the same directory as your app.py file.</p>",
            status_code=404
        )
    except Exception as e:
        logging.error(f"Error serving webpage: {e}")
        return HTMLResponse(
            content="<h1>Error loading webpage</h1>",
            status_code=500
        )

# Optional: Add a favicon endpoint
@app.get("/favicon.ico")
async def favicon():
    """Return a simple favicon"""
    return Response(content="", media_type="image/x-icon")
if __name__ == "__main__":
    import uvicorn
    print("Starting FastAPI server on http://0.0.0.0:8080")
    # Consider adding reload=True for development, but remove for production
    uvicorn.run(app, host="0.0.0.0", port=8080)
